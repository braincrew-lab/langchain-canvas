"""``DeckPptxExporter`` — canonical ``.slides.html`` decks compiled to pptx."""

from __future__ import annotations

import base64
import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
from langchain_canvas.deck.export import DeckPptxExporter, _verify_reopen
from langchain_canvas.exporters import PPTX_MIME

# A minimal 1x1 transparent PNG, base64-encoded — enough for python-pptx to
# read a native picture size from.
_TINY_PNG = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+A8AAQUBAScY42YAAAAASUVORK5CYII="
)
_TINY_PNG_URI = f"data:image/png;base64,{_TINY_PNG}"


def _slide_html(*nodes: str) -> str:
    return '<section class="slide">\n' + "\n".join(nodes) + "\n</section>"


def _text_node(
    node_id: str, text: str, *, shape_id: str | None = None, extra_style: str = ""
) -> str:
    shape_attr = f' data-pptx-shape-id="{shape_id}"' if shape_id else ""
    style = (
        "position: absolute; left: 100.00px; top: 50.00px; "
        f"width: 400.00px; height: 80.00px; {extra_style}"
    )
    return (
        f'<div class="lcx-block" data-node-id="{node_id}"{shape_attr} '
        f'style="{style}">{text}</div>'
    )


def _image_node(node_id: str) -> str:
    return (
        f'<img class="lcx-block" data-node-id="{node_id}" src="{_TINY_PNG_URI}" '
        'style="position: absolute; left: 500.00px; top: 200.00px; '
        'width: 100.00px; height: 100.00px;" alt="">'
    )


def _shape_node(node_id: str, kind: str) -> str:
    return (
        f'<div class="lcx-block {kind}" data-node-id="{node_id}" '
        'style="position: absolute; left: 50.00px; top: 400.00px; '
        'width: 200.00px; height: 100.00px; background: #ff0000;"></div>'
    )


def _raster_node(node_id: str) -> str:
    return (
        f'<img class="lcx-block" data-node-id="{node_id}" '
        'data-lcx-fallback="raster" alt="" '
        'style="position: absolute; left: 10.00px; top: 10.00px; '
        'width: 60.00px; height: 60.00px;">'
    )


def _deck_content(*slides: SlideTemplate, source: str | None = None) -> str:
    return serialize_deck(Deck(title="Deck", ratio="16:9", source=source, slides=list(slides)))


def test_export_maps_text_image_shape_to_native():
    slide = SlideTemplate(
        slide_id="slide-001",
        title=None,
        style_css="",
        body_html=_slide_html(
            _text_node("node-1", "Hello"),
            _image_node("node-2"),
            _shape_node("node-3", "ellipse"),
        ),
    )
    content = _deck_content(slide)

    result = DeckPptxExporter().export(content, path="deck.slides.html")

    assert result.media_type == PPTX_MIME
    reopened = Presentation(io.BytesIO(result.data))
    assert len(reopened.slides) == 1
    shapes = list(reopened.slides[0].shapes)
    assert len(shapes) == 3

    texts = [s for s in shapes if s.has_text_frame and s.text_frame.text]
    assert any(s.text_frame.text == "Hello" for s in texts)
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes)
    assert any(getattr(s, "adjustments", None) is not None for s in shapes)


def _skin_pptx_bytes() -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "Original"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _skin_uri() -> str:
    return f"data:{PPTX_MIME};base64,{base64.b64encode(_skin_pptx_bytes()).decode()}"


def test_export_patches_provenance_shape_on_skin():
    slide = SlideTemplate(
        slide_id="slide-001",
        title=None,
        style_css="",
        body_html=_slide_html(_text_node("node-1", "Patched", shape_id="e0")),
    )
    content = _deck_content(slide, source=_skin_uri())

    result = DeckPptxExporter().export(content, path="deck.slides.html")

    reopened = Presentation(io.BytesIO(result.data))
    assert len(reopened.slides) == 1
    shapes = list(reopened.slides[0].shapes)
    assert len(shapes) == 1  # patched in place, not added as a duplicate
    assert shapes[0].text_frame.text == "Patched"


def test_export_inlines_lcx_source_before_export():
    from langchain_canvas.deck import parse_deck
    from langchain_canvas.store import InMemoryCanvasStore
    from langchain_canvas.tools import inline_deck_skin

    store = InMemoryCanvasStore()
    skin_bytes = _skin_pptx_bytes()
    store.write_bytes("t1", "sources/deck.pptx", skin_bytes, "Upload", actor="human")
    slide = SlideTemplate(slide_id="slide-001", title=None, style_css="", body_html=_slide_html())
    content = _deck_content(slide, source="sources/deck.pptx")

    inlined = inline_deck_skin(content, store, "t1")

    deck = parse_deck(inlined)
    assert deck.source is not None
    assert deck.source.startswith(f"data:{PPTX_MIME};base64,")
    decoded = base64.b64decode(deck.source.split(",", 1)[1])
    assert decoded == skin_bytes


def test_export_without_render_adapter_reports_degraded_nodes():
    slide = SlideTemplate(
        slide_id="slide-001",
        title=None,
        style_css="",
        body_html=_slide_html(_raster_node("node-1")),
    )
    content = _deck_content(slide)

    exporter = DeckPptxExporter()
    result = exporter.export(content, path="deck.slides.html")

    assert exporter.degraded_nodes == ["slide-001:node-1"]
    reopened = Presentation(io.BytesIO(result.data))
    assert len(reopened.slides) == 1
    assert list(reopened.slides[0].shapes) == []


def test_export_with_render_adapter_rasterizes_degraded_nodes():
    slide = SlideTemplate(
        slide_id="slide-001",
        title=None,
        style_css="",
        body_html=_slide_html(_raster_node("node-1")),
    )
    content = _deck_content(slide)
    png_bytes = base64.b64decode(_TINY_PNG)
    calls: list[tuple[str, str]] = []

    def fake_adapter(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        calls.append((html, ratio))
        return {}, png_bytes

    exporter = DeckPptxExporter(render_adapter=fake_adapter)
    result = exporter.export(content, path="deck.slides.html")

    assert exporter.degraded_nodes == []
    assert calls and calls[0][1] == "16:9"
    reopened = Presentation(io.BytesIO(result.data))
    shapes = list(reopened.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_export_reopen_validation_fails_on_corrupt_output():
    with pytest.raises(ValueError):
        _verify_reopen(b"not a pptx file", expected_slide_count=1)


def test_export_reopen_validation_fails_on_slide_count_mismatch():
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    buffer = io.BytesIO()
    presentation.save(buffer)

    with pytest.raises(ValueError):
        _verify_reopen(buffer.getvalue(), expected_slide_count=2)
