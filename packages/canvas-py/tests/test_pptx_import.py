"""Reading a ``.pptx`` back into the deck model.

The exporter's tests prove a deck leaves as real shapes; these prove one
comes back. They are written as bounds rather than transcripts — a shape
count, a coordinate range, a field that must or must not appear — because a
reader pinned to one deck's exact numbers goes red on every unrelated change
and says nothing about the next deck.

Decks are built here with python-pptx rather than checked in, so what each
test covers is visible in the test itself.
"""

from __future__ import annotations

import base64
import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from langchain_canvas.pptx_import import PptxImportError, pptx_to_slides
from langchain_canvas.protocol.artifacts import SlidesData

# A real 1x1 red PNG — small enough to inline, real enough for pptx to embed.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _deck(build: Any = None, *, width_in: float = 13.333, height_in: float = 7.5) -> bytes:
    """Presentation bytes with one blank slide, shaped by ``build``."""
    deck = Presentation()
    deck.slide_width = Inches(width_in)
    deck.slide_height = Inches(height_in)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    if build is not None:
        build(slide)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _textbox(slide: Any, text: str = "Hello", **kwargs: Any) -> Any:
    box = slide.shapes.add_textbox(
        kwargs.get("left", Inches(1)),
        kwargs.get("top", Inches(1)),
        kwargs.get("width", Inches(4)),
        kwargs.get("height", Inches(1)),
    )
    box.text_frame.text = text
    return box


def _elements(data: dict[str, Any], index: int = 0) -> list[dict[str, Any]]:
    return data["slides"][index]["elements"]


# --- what comes across ----------------------------------------------------------


def test_the_deck_model_accepts_what_the_reader_produces() -> None:
    """The reader's output is a deck, not a dict that resembles one."""
    parsed = pptx_to_slides(_deck(lambda s: _textbox(s, "Title")))
    model = SlidesData.model_validate(parsed)
    assert model.page is not None
    assert model.slides


def test_the_page_is_the_deck_s_own_size() -> None:
    """Percent geometry means nothing without the page it is a percent of."""
    parsed = pptx_to_slides(_deck(width_in=10, height_in=5.625))
    assert parsed["page"] == {"widthIn": 10.0, "heightIn": 5.625}


def test_geometry_lands_inside_the_page_as_percent() -> None:
    """A box an inch in on a 10-inch page is at 10 percent, not at 914400."""
    parsed = pptx_to_slides(
        _deck(
            lambda s: _textbox(s, "Hi", left=Inches(1), top=Inches(1), width=Inches(2)),
            width_in=10,
            height_in=5,
        )
    )
    element = _elements(parsed)[0]
    assert element["x"] == pytest.approx(10, abs=0.01)
    assert element["y"] == pytest.approx(20, abs=0.01)
    assert element["w"] == pytest.approx(20, abs=0.01)
    assert 0 <= element["x"] <= 100 and 0 <= element["y"] <= 100


def test_a_picture_arrives_as_a_data_uri_the_exporter_can_place() -> None:
    """The exporter reads png/jpeg/gif data URIs and skips everything else."""

    def build(slide: Any) -> None:
        slide.shapes.add_picture(io.BytesIO(_PNG), Inches(1), Inches(1), Inches(1), Inches(1))

    element = _elements(pptx_to_slides(_deck(build)))[0]
    assert element["type"] == "image"
    assert element["src"].startswith("data:image/png;base64,")


def test_text_keeps_its_size_and_weight() -> None:
    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Bold and big"
        run.font.size = Pt(28)
        run.font.bold = True

    element = _elements(pptx_to_slides(_deck(build)))[0]
    assert element["type"] == "text"
    assert element["fontSize"] == pytest.approx(28 * 4 / 3, abs=0.1)  # px
    assert element["bold"] is True


def test_speaker_notes_survive() -> None:
    def build(slide: Any) -> None:
        _textbox(slide)
        slide.notes_slide.notes_text_frame.text = "Say this out loud"

    parsed = pptx_to_slides(_deck(build))
    assert parsed["slides"][0]["notes"] == "Say this out loud"


# --- what does not, and says so -------------------------------------------------


def test_a_table_keeps_its_place_among_the_other_shapes() -> None:
    def build(slide: Any) -> None:
        slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        _textbox(slide, "Kept")

    elements = _elements(pptx_to_slides(_deck(build)))
    assert [e["type"] for e in elements] == ["table", "text"]
    assert elements[1]["text"].strip() == "Kept"


def test_a_group_is_dropped_rather_than_flattened() -> None:
    """Flattening loses the grouping and keeps neither shape's real position."""

    def build(slide: Any) -> None:
        _textbox(slide, "Kept")

    deck = _deck(build)
    # python-pptx cannot author a group, so assert the guard directly.
    from langchain_canvas.pptx_import import _is_group

    class _Grouped:
        shape_type = "GROUP (6)"

    assert _is_group(_Grouped())
    assert len(_elements(pptx_to_slides(deck))) == 1


def test_a_box_whose_runs_disagree_takes_the_first_one() -> None:
    """One element carries one set of formatting. The words all survive; the
    variation does not, and the original keeps it."""

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        paragraph = box.text_frame.paragraphs[0]
        first = paragraph.add_run()
        first.text = "big "
        first.font.size = Pt(40)
        second = paragraph.add_run()
        second.text = "small"
        second.font.size = Pt(8)

    element = _elements(pptx_to_slides(_deck(build)))[0]
    assert element["fontSize"] == pytest.approx(40 * 4 / 3, abs=0.1)  # px
    assert "big " in element["text"] and "small" in element["text"]


def test_bytes_that_are_not_a_deck_raise_rather_than_return_an_empty_one() -> None:
    """An empty deck and an unreadable file are different answers, and the
    upload path shows a file card only for the second."""
    with pytest.raises(PptxImportError):
        pptx_to_slides(b"not a presentation")


# --- the round trip -------------------------------------------------------------


def test_a_deck_read_and_written_keeps_its_shapes_and_page() -> None:
    """Not byte equality — that was withdrawn as a requirement. What has to
    hold is that nothing silently disappears across one trip."""
    import json

    from langchain_canvas.exporters import SlidesPptxExporter

    def build(slide: Any) -> None:
        _textbox(slide, "One")
        _textbox(slide, "Two", top=Inches(3))
        slide.shapes.add_picture(io.BytesIO(_PNG), Inches(6), Inches(1), Inches(1), Inches(1))

    original = _deck(build)
    parsed = pptx_to_slides(original)
    payload = json.dumps({"type": "slides", "data": parsed}).encode()
    exported = SlidesPptxExporter().export(payload, path="deck.slides.json", title="Deck")

    again = pptx_to_slides(exported.data)
    assert sum(len(s["elements"]) for s in again["slides"]) == sum(
        len(s["elements"]) for s in parsed["slides"]
    )
    assert again["page"] == parsed["page"]


def test_an_upload_stays_a_file_until_it_is_copied_out() -> None:
    """Opening the upload itself as slides looked helpful and was not: nothing
    under sources/ is editable, and no exporter matches a .pptx name, so the
    deck rendered in an editor that could neither change nor export it. It
    stays a file card; open_deck_for_editing is the way in."""
    from langchain_canvas.replay import source_preview_events
    from langchain_canvas.store import InMemoryCanvasStore

    store = InMemoryCanvasStore()
    commit = store.write_bytes(
        "t1", "sources/deck.pptx", _deck(lambda s: _textbox(s, "Hi")), "Upload", actor="human"
    )
    events = source_preview_events(
        store,
        "t1",
        "sources/deck.pptx",
        is_new=True,
        revision=commit.revision,
        description="Upload",
    )
    created = next(e for e in events if e["type"] == "canvas.create")
    assert created["artifact"]["type"] == "file"


def test_an_unreadable_deck_still_gets_a_file_card() -> None:
    """Degrading to the ordinary preview beats showing nothing."""
    from langchain_canvas.replay import source_preview_events
    from langchain_canvas.store import InMemoryCanvasStore

    store = InMemoryCanvasStore()
    commit = store.write_bytes("t1", "sources/broken.pptx", b"PK-nope", "Upload", actor="human")
    events = source_preview_events(
        store,
        "t1",
        "sources/broken.pptx",
        is_new=True,
        revision=commit.revision,
        description="Upload",
    )
    created = next(e for e in events if e["type"] == "canvas.create")
    assert created["artifact"]["type"] == "file"


# --- copying an upload out so it can be edited and exported ---------------------


def _store_with_deck(deck: bytes = b"") -> Any:
    from langchain_canvas.store import InMemoryCanvasStore

    store = InMemoryCanvasStore()
    store.write_bytes(
        "t1",
        "sources/deck.pptx",
        deck or _deck(lambda s: _textbox(s, "Hi")),
        "Upload",
        actor="human",
    )
    return store


def _copy_tool(store: Any) -> Any:
    from langchain_canvas.tools import create_deck_tools

    return create_deck_tools(store)[0]


def _run(tool: Any, **kwargs: Any) -> str:
    """Through the same runtime stub the other tool tests use."""
    from test_tools import _runtime

    return tool.func(runtime=_runtime(thread_id="t1"), **kwargs)


def test_the_copy_lands_where_editing_and_export_both_work() -> None:
    """Under sources/ nothing is editable, and no exporter matches a .pptx
    name — the copy has to leave both behind to be worth making."""
    from langchain_canvas.exporters import default_exporters, exporter_for

    store = _store_with_deck()
    reply = _run(_copy_tool(store), source="sources/deck.pptx")
    assert reply.startswith("Copied "), reply

    paths = [f.path for f in store.list_files("t1")]
    assert "deck.slides.json" in paths
    assert "sources/deck.pptx" in paths  # the original is untouched
    assert exporter_for("deck.slides.json", "pptx", default_exporters()) is not None


def test_the_copy_names_the_original_as_its_template() -> None:
    """That is the whole reason the masters survive a round trip."""
    import json

    store = _store_with_deck()
    _run(_copy_tool(store), source="sources/deck.pptx")
    envelope = json.loads(store.read("t1", "deck.slides.json").content)
    assert envelope["data"]["template"] == "sources/deck.pptx"


def test_copying_the_same_deck_twice_is_refused_rather_than_overwriting() -> None:
    store = _store_with_deck()
    tool = _copy_tool(store)
    _run(tool, source="sources/deck.pptx")
    again = _run(tool, source="sources/deck.pptx")
    assert "already on the canvas" in again


def test_the_copy_cannot_be_written_back_into_sources() -> None:
    """The uploads folder stays the user's, whatever destination is passed."""
    store = _store_with_deck()
    reply = _run(
        _copy_tool(store),
        source="sources/deck.pptx",
        destination="sources/copy.slides.json",
    )
    assert reply.startswith("Error:")
    assert "sources/" in reply


def test_a_file_that_is_not_a_deck_is_pointed_somewhere_useful() -> None:
    store = _store_with_deck()
    reply = _run(_copy_tool(store), source="notes.md")
    assert reply.startswith("Error: this opens .pptx files")


# --- theme colours ---------------------------------------------------------------


def test_luminance_matches_what_powerpoint_shows() -> None:
    """White at -0.25 is #BFBFBF — the value the deck itself renders, and the
    one an independent reader recorded for the same slide."""
    from langchain_canvas.pptx_import import _shade

    assert _shade("FFFFFF", -0.25) == "BFBFBF"
    assert _shade("000000", 0.5) == "808080"
    assert _shade("FD7F00", 0) == "FD7F00"


def test_a_run_naming_a_theme_slot_still_gets_a_colour() -> None:
    """Half a real deck's runs name a slot rather than a value. Dropping them
    loses half the deck's colour, so the slot is resolved through the master's
    scheme and colour map."""
    from pptx.enum.dml import MSO_THEME_COLOR

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Themed"
        run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    element = _elements(pptx_to_slides(_deck(build)))[0]
    assert element.get("color", "").startswith("#"), element
    assert len(element["color"]) == 7


def test_a_slot_the_scheme_does_not_carry_stays_absent() -> None:
    """Better no colour than an invented one — the skin still knows."""
    from langchain_canvas.pptx_import import _colour

    class _Colour:
        rgb = None
        brightness = 0

        @property
        def theme_color(self) -> Any:
            return type("T", (), {"name": "ACCENT_1"})()

    class _Font:
        color = _Colour()

    assert _colour(_Font(), {}) is None


def test_the_copy_is_written_through_the_envelope_encoder() -> None:
    """A .slides.json without {"type","title","data"} parses as no artifact,
    and the canvas shows the raw JSON as a markdown document instead of a
    deck. The file has to go through encode_slides, not json.dumps."""
    import json

    from langchain_canvas.replay import events_for_commit

    store = _store_with_deck()
    _run(_copy_tool(store), source="sources/deck.pptx")
    content = store.read("t1", "deck.slides.json").content

    envelope = json.loads(content)
    assert envelope["type"] == "slides"
    assert "slides" in envelope["data"]
    assert envelope["data"]["template"] == "sources/deck.pptx"

    events = events_for_commit(
        "deck.slides.json", content, is_new=True, revision="v1", description="Copy"
    )
    created = next(e for e in events if e["type"] == "canvas.create")
    assert created["artifact"]["type"] == "slides"


def test_the_copy_is_broadcast_as_a_deck_not_as_a_text_preview() -> None:
    """A .json path is in the source-preview suffix list, so broadcasting the
    copy the way an upload is broadcast drew the deck as its own JSON."""
    from test_tools import _Runtime

    store = _store_with_deck()
    sent: list[dict] = []
    runtime = _Runtime(config={"configurable": {"thread_id": "t1"}}, )
    object.__setattr__(runtime, "stream_writer", sent.append)
    _copy_tool(store).func(runtime=runtime, source="sources/deck.pptx")

    created = [e for e in sent if e.get("type") == "canvas.create"]
    assert created, sent
    assert created[0]["artifact"]["type"] == "slides"


# --- connectors ------------------------------------------------------------------


def test_a_flat_connector_gets_its_stroke_as_thickness() -> None:
    """A horizontal connector is zero tall in the file — PowerPoint draws it
    from the line weight. The deck model paints shapes as boxes, so a zero
    box is an invisible line."""

    def build(slide: Any) -> None:
        slide.shapes.add_connector(1, Inches(1), Inches(2), Inches(5), Inches(2))

    element = next(e for e in _elements(pptx_to_slides(_deck(build))) if e["type"] == "shape")
    assert element["shape"] == "line"
    assert element["h"] > 0, element
    assert element["w"] > 0


# --- what the box alone cannot say -----------------------------------------------


def test_a_box_drawn_by_its_outline_alone_survives() -> None:
    """The common annotation in a real deck is an empty rectangle around
    content: no fill, a coloured border. With only a `fill` field it rendered
    as nothing, which is why red boxes went missing from an imported deck."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt as _Pt

    def build(slide: Any) -> None:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        box.fill.background()  # no fill at all
        box.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        box.line.width = _Pt(3)

    element = next(e for e in _elements(pptx_to_slides(_deck(build))) if e["type"] == "shape")
    assert element["stroke"] == "#FF0000"
    assert element["strokeWidth"] == pytest.approx(3 * 4 / 3, abs=0.1)  # px
    assert "fill" not in element  # an outline-only box has none


def test_the_type_face_comes_across() -> None:
    """Without the face, line breaks land in different places than the source
    file — the single largest source of 'the layout looks different'."""

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Hello"
        run.font.name = "Pretendard"

    assert _elements(pptx_to_slides(_deck(build)))[0]["fontFamily"] == "Pretendard"


def test_outline_and_face_survive_the_round_trip() -> None:
    """The fields are only real if the export writes them back."""
    import json

    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt as _Pt

    from langchain_canvas.exporters import SlidesPptxExporter

    def build(slide: Any) -> None:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        box.fill.background()
        box.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        box.line.width = _Pt(3)
        text = _textbox(slide, "", top=Inches(4))
        run = text.text_frame.paragraphs[0].add_run()
        run.text = "Hello"
        run.font.name = "Pretendard"

    parsed = pptx_to_slides(_deck(build))
    exported = SlidesPptxExporter().export(
        json.dumps({"type": "slides", "data": parsed}).encode(),
        path="deck.slides.json",
        title="Deck",
    )
    back = pptx_to_slides(exported.data)
    elements = back["slides"][0]["elements"]
    assert any(e.get("stroke") == "#FF0000" for e in elements), elements
    assert any(e.get("fontFamily") == "Pretendard" for e in elements), elements


def test_a_highlighted_heading_keeps_its_band() -> None:
    """`a:highlight` is a coloured bar behind the words. python-pptx exposes no
    accessor for it, so it is easy to miss — and missing it turns three marked
    headings into plain text."""
    from pptx.oxml.ns import qn

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "1. Marked"
        properties = run.font._rPr
        mark = properties.makeelement(qn("a:highlight"), {})
        mark.append(mark.makeelement(qn("a:srgbClr"), {"val": "FF0000"}))
        properties.append(mark)

    assert _elements(pptx_to_slides(_deck(build)))[0]["highlight"] == "#FF0000"


def test_a_highlight_survives_the_round_trip() -> None:
    import json

    from pptx.oxml.ns import qn

    from langchain_canvas.exporters import SlidesPptxExporter

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "1. Marked"
        properties = run.font._rPr
        mark = properties.makeelement(qn("a:highlight"), {})
        mark.append(mark.makeelement(qn("a:srgbClr"), {"val": "FF0000"}))
        properties.append(mark)

    parsed = pptx_to_slides(_deck(build))
    exported = SlidesPptxExporter().export(
        json.dumps({"type": "slides", "data": parsed}).encode(),
        path="deck.slides.json",
        title="Deck",
    )
    back = pptx_to_slides(exported.data)
    assert any(e.get("highlight") == "#FF0000" for e in back["slides"][0]["elements"])


# --- units ------------------------------------------------------------------------


def test_type_size_is_stored_in_the_unit_the_model_uses() -> None:
    """The deck model measures type in px on a 1280px slide; PowerPoint states
    it in points. Storing the point value renders every size a quarter too
    small, which moves every line break — the layout stops matching the file."""

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Sized"
        run.font.size = Pt(30)

    assert _elements(pptx_to_slides(_deck(build)))[0]["fontSize"] == pytest.approx(40.0, abs=0.1)


def test_type_size_comes_back_the_same_after_a_round_trip() -> None:
    """Import converts pt to px and the export converts back; a mismatch in
    either direction shrinks or grows the deck on every trip."""
    import json

    from langchain_canvas.exporters import SlidesPptxExporter

    def build(slide: Any) -> None:
        box = _textbox(slide, "")
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Sized"
        run.font.size = Pt(30)

    parsed = pptx_to_slides(_deck(build))
    exported = SlidesPptxExporter().export(
        json.dumps({"type": "slides", "data": parsed}).encode(),
        path="deck.slides.json",
        title="Deck",
    )
    back = pptx_to_slides(exported.data)
    assert back["slides"][0]["elements"][0]["fontSize"] == pytest.approx(40.0, abs=0.5)


def test_background_is_inherited_from_the_layout_and_the_master():
    """A slide with no background of its own wears the one above it.

    Reading only the slide's own element turns an inherited dark deck white.
    """
    pytest.importorskip("pptx")
    from lxml import etree
    from pptx import Presentation
    from pptx.oxml.ns import qn

    presentation = Presentation()
    master = presentation.slide_masters[0]
    common = master.element.find(qn("p:cSld"))
    common.insert(
        0,
        etree.fromstring(
            '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:bgPr><a:solidFill><a:srgbClr val="151515"/></a:solidFill></p:bgPr></p:bg>'
        ),
    )
    presentation.slides.add_slide(presentation.slide_layouts[6])
    buffer = io.BytesIO()
    presentation.save(buffer)

    deck = pptx_to_slides(buffer.getvalue())
    assert deck["slides"][0]["background"] == "#151515"


def test_a_theme_coloured_background_resolves_through_the_scheme():
    """``schemeClr`` names a colour the theme holds; the name is not a colour."""
    pytest.importorskip("pptx")
    from lxml import etree
    from pptx import Presentation
    from pptx.oxml.ns import qn

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    common = slide.element.find(qn("p:cSld"))
    common.insert(
        0,
        etree.fromstring(
            '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:bgPr><a:solidFill><a:schemeClr val="accent1"/></a:solidFill></p:bgPr></p:bg>'
        ),
    )
    buffer = io.BytesIO()
    presentation.save(buffer)

    background = pptx_to_slides(buffer.getvalue())["slides"][0]["background"]
    assert background is not None and background.startswith("#")
    assert background.lower() != "#accent1"


# --- pictures leave the copy for assets/ ------------------------------------------


def _deck_with_pictures(count: int = 2, *, same: bool = False) -> bytes:
    """A one-slide deck carrying ``count`` pictures (identical when ``same``)."""
    from PIL import Image

    def build(slide: Any) -> None:
        for index in range(count):
            buffer = io.BytesIO()
            colour = (200, 30, 30) if same else (200, 30 + index * 60, 30)
            Image.new("RGB", (4, 4), colour).save(buffer, format="PNG")
            buffer.seek(0)
            slide.shapes.add_picture(buffer, Inches(1 + index), Inches(1), Inches(1), Inches(1))

    return _deck(build)


def test_the_copy_stores_its_pictures_under_assets_and_references_them() -> None:
    """A deck with photos stored as data URIs is a megabyte of base64 the
    model reads the first page of and nothing after. The copy has to be the
    text it is, with the pictures where every canvas file finds them."""
    import json

    store = _store_with_deck(_deck_with_pictures(2))
    reply = _run(_copy_tool(store), source="sources/deck.pptx")
    assert "2 picture(s) under assets/deck/" in reply, reply

    content = store.read("t1", "deck.slides.json").content
    assert "data:image" not in content
    sources = [el["src"] for el in _elements(json.loads(content)["data"]) if el["type"] == "image"]
    assert len(sources) == 2
    assert all(src.startswith("assets/deck/") and src.endswith(".png") for src in sources)
    on_canvas = {f.path for f in store.list_files("t1")}
    assert set(sources) <= on_canvas
    # The bytes stored are the picture's own bytes, not a re-encoding.
    from PIL import Image

    stored = store.read_bytes("t1", sources[0]).data
    assert Image.open(io.BytesIO(stored)).size == (4, 4)


def test_a_picture_repeated_across_the_deck_is_stored_once() -> None:
    """A logo on every slide is one file, referenced from each slide."""
    import json

    store = _store_with_deck(_deck_with_pictures(3, same=True))
    reply = _run(_copy_tool(store), source="sources/deck.pptx")
    assert "1 picture(s)" in reply, reply
    sources = {
        el["src"]
        for el in _elements(json.loads(store.read("t1", "deck.slides.json").content)["data"])
        if el["type"] == "image"
    }
    assert len(sources) == 1


def test_the_copy_tells_the_model_to_edit_it_rather_than_start_over() -> None:
    """The reply is the one place the model learns the way to revise a deck."""
    store = _store_with_deck()
    reply = _run(_copy_tool(store), source="sources/deck.pptx")
    assert "edit_canvas" in reply
    assert "keeps the look" in reply


# --- the third way a file names a colour, and WordArt's size --------------------------


def _with_run_xml(deck_bytes: bytes, mutate: Any) -> bytes:
    """Re-open a deck and let ``mutate(run)`` edit the first run's XML."""
    from pptx.oxml.ns import qn

    deck = Presentation(io.BytesIO(deck_bytes))
    run = deck.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    mutate(run, qn)
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


def test_a_preset_colour_name_comes_across_as_hex() -> None:
    """Seen in a run: a 66pt title in prstClr white on a dark slide, read as
    no colour, drawn dark — the slide looked empty."""
    from lxml import etree

    def paint_white(run: Any, qn: Any) -> None:
        rpr = run._r.get_or_add_rPr()
        fill = etree.SubElement(rpr, qn("a:solidFill"))
        etree.SubElement(fill, qn("a:prstClr")).set("val", "white")

    data = pptx_to_slides(_with_run_xml(_deck(lambda s: _textbox(s, "Title")), paint_white))
    assert _elements(data)[0]["color"] == "#FFFFFF"


def test_preset_names_use_the_standard_table() -> None:
    from langchain_canvas.preset_colours import preset_colour

    assert preset_colour("white") == "FFFFFF"
    assert preset_colour("Black") == "000000"
    assert preset_colour("dkBlue") == "00008B"
    assert preset_colour("ltGray") == "D3D3D3"
    assert preset_colour("notacolour") is None


def test_word_art_takes_its_size_from_the_box() -> None:
    """WordArt warps text to fill its shape; a run with no size is not body
    text, it is a headline the height of the box."""
    from lxml import etree
    from pptx.util import Inches

    def warp(run: Any, qn: Any) -> None:
        body = run._r.getparent().getparent()  # a:r -> a:p -> p:txBody
        body_pr = body.find(qn("a:bodyPr"))
        etree.SubElement(body_pr, qn("a:prstTxWarp")).set("prst", "textArchUp")

    plain = pptx_to_slides(_deck(lambda s: _textbox(s, "BIG", height=Inches(1.5))))
    assert "fontSize" not in _elements(plain)[0]
    tall = _deck(lambda s: _textbox(s, "BIG", height=Inches(1.5)))
    warped = pptx_to_slides(_with_run_xml(tall, warp))
    assert _elements(warped)[0]["fontSize"] == round(1.5 * 72 * 4 / 3 * 0.8, 1)


# --- tables as cell boxes, charts as pictures ---------------------------------------


def _deck_with_table() -> bytes:
    from pptx.util import Inches

    def build(slide: Any) -> None:
        frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = frame.table
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "Header"
        table.cell(1, 0).text = "a"
        table.cell(1, 1).text = "b"

    return _deck(build)


def test_a_table_arrives_as_one_table_element_on_its_own_grid() -> None:
    """Dropping a table left a slide saying nothing where the original said
    the most; sixteen boxes said it but could not be widened by the column.
    The words go in rows, the grid in colWidths / rowHeights, a merge in cells."""
    data = pptx_to_slides(_deck_with_table())
    (table,) = [e for e in _elements(data) if e["type"] == "table"]
    assert table["id"] == "t0"
    assert table["rows"] == [["Header", ""], ["a", "b"]]
    assert table["colWidths"] == [50.0, 50.0] and table["rowHeights"] == [50.0, 50.0]
    assert table["header"] is True
    assert table["cells"] == [{"r": 0, "c": 0, "colSpan": 2}]
    assert abs(table["w"] / table["x"] - 4) < 0.01  # 1in in, 4in wide
    assert "stroke" not in table  # a built-in style names no line the sheet defines
    assert SlidesData.model_validate(data)


def test_a_chart_leaves_its_box_for_the_tool_and_the_tool_pictures_it() -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    from langchain_canvas.converters import ConvertedSource
    from langchain_canvas.store import InMemoryCanvasStore
    from langchain_canvas.tools import create_deck_tools

    PIL = pytest.importorskip("PIL")

    def build(slide: Any) -> None:
        chart_data = CategoryChartData()
        chart_data.categories = ["a", "b"]
        chart_data.add_series("s", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), chart_data
        )

    deck = _deck(build)
    data = pptx_to_slides(deck)
    assert data["slides"][0]["elements"] == []
    assert len(data["slides"][0]["charts"]) == 1

    class _Page:
        suffixes = (".pptx",)

        def convert(self, data: bytes, *, path: str) -> Any:
            return ConvertedSource(blocks=[{"type": "text", "text": "x"}], metadata={})

        def render_pages(self, data: bytes, *, path: str, pages: list[int]) -> Any:
            image = PIL.Image.new("RGB", (400, 225), (255, 255, 255))
            buffer = io.BytesIO()
            image.save(buffer, format="PNG")
            encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
            block = {
                "type": "image", "source_type": "base64", "data": encoded, "mime_type": "image/png",
            }
            return ConvertedSource(blocks=[block], metadata={})

        def render_grid(self, data: bytes, *, path: str) -> Any:
            return self.render_pages(data, path=path, pages=[1])

    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/deck.pptx", deck, "Upload", actor="human")
    tool = create_deck_tools(store, converters=[_Page()])[0]
    reply = _run(tool, source="sources/deck.pptx")
    assert "1 chart(s) as pictures" in reply, reply
    import json

    saved = json.loads(store.read("t1", "deck.slides.json").content)
    elements = saved["data"]["slides"][0]["elements"]
    assert "charts" not in saved["data"]["slides"][0]
    assert elements[-1]["type"] == "image"
    assert elements[-1]["src"].startswith("assets/deck/chart-s1-")
    assert any(f.path == elements[-1]["src"] for f in store.list_files("t1"))

    plain = InMemoryCanvasStore()
    plain.write_bytes("t1", "sources/deck.pptx", deck, "Upload", actor="human")
    reply = _run(create_deck_tools(plain)[0], source="sources/deck.pptx")
    assert "1 chart(s) dropped" in reply


def test_a_slide_carries_the_themes_text_colour_as_its_default() -> None:
    """Table cells and WordArt with no run colour were drawn in the screen's
    own theme colour — light grey on a white slide. The file does say what
    colourless text is: the theme's text slot through the colour map."""
    data = pptx_to_slides(_deck(lambda s: _textbox(s, "Hi")))
    slide = data["slides"][0]
    assert slide["textColor"].startswith("#") and len(slide["textColor"]) == 7
    assert slide["textColor"].upper() == "#000000"  # python-pptx's default theme: dk1 black
    assert SlidesData.model_validate(data)


# --- the table style sheet, and a text outline ---------------------------------------


_STYLE_SHEET = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'def="{838324F5-92EC-4F16-B9E8-4A18B4BBC7FD}">'
    '<a:tblStyle styleId="{838324F5-92EC-4F16-B9E8-4A18B4BBC7FD}" styleName="Table_0">'
    '<a:wholeTbl><a:tcTxStyle><a:font><a:latin typeface="Arial"/></a:font>'
    '<a:srgbClr val="000000"/></a:tcTxStyle><a:tcStyle><a:tcBdr><a:insideH>'
    '<a:ln w="9525"><a:solidFill><a:srgbClr val="9E9E9E"/></a:solidFill></a:ln>'
    "</a:insideH></a:tcBdr></a:tcStyle></a:wholeTbl></a:tblStyle></a:tblStyleLst>"
)


def _with_style_sheet(deck_bytes: bytes, style_id: str) -> bytes:
    from pptx.oxml.ns import qn

    deck = Presentation(io.BytesIO(deck_bytes))
    parts = deck.part.package.iter_parts()
    part = next(p for p in parts if str(p.partname).endswith("tableStyles.xml"))
    part._blob = _STYLE_SHEET.encode("utf-8")
    table = next(s for s in deck.slides[0].shapes if s.has_table).table
    ref = table._tbl.tblPr.find(qn("a:tableStyleId"))
    if ref is None:
        ref = table._tbl.tblPr.makeelement(qn("a:tableStyleId"), {})
        table._tbl.tblPr.append(ref)
    ref.text = style_id
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


def test_a_table_is_as_big_as_its_grid_not_its_frame() -> None:
    """PowerPoint draws a table from its columns and rows; the frame around
    it is advisory, and one file declared a square frame around a wide strip."""
    from pptx.util import Inches

    def build(slide: Any) -> None:
        frame = slide.shapes.add_table(2, 4, Inches(1), Inches(1), Inches(3), Inches(3))
        for column in frame.table.columns:
            column.width = Inches(2)
        for row in frame.table.rows:
            row.height = Inches(0.5)

    (table,) = [e for e in _elements(pptx_to_slides(_deck(build))) if e["type"] == "table"]
    assert abs(table["w"] / table["x"] - 8) < 0.01  # 8in wide, 1in in
    assert abs(table["h"] / table["y"] - 1) < 0.01  # 1in tall, 1in down


def test_a_tables_style_sheet_gives_its_cells_borders_colour_and_face() -> None:
    """A styled table with no cell borders is not an unbordered table — the
    borders live in the deck's style sheet when the file carries it."""
    styled = _with_style_sheet(_deck_with_table(), "{838324F5-92EC-4F16-B9E8-4A18B4BBC7FD}")
    data = pptx_to_slides(styled)
    (table,) = [e for e in _elements(data) if e["type"] == "table"]
    assert table["stroke"] == "#9E9E9E" and table["strokeWidth"] == 1.0
    assert table["color"] == "#000000" and table["fontFamily"] == "Arial"


def test_a_style_the_sheet_does_not_define_leaves_the_table_bare() -> None:
    unknown = _with_style_sheet(_deck_with_table(), "{00000000-0000-0000-0000-000000000000}")
    (table,) = [e for e in _elements(pptx_to_slides(unknown)) if e["type"] == "table"]
    assert "stroke" not in table and "color" not in table


def test_what_one_cell_does_differently_is_its_own() -> None:
    """A header fill and a bold total are the cells' — the rest of the table
    stays a grid of plain strings."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    def build(slide: Any) -> None:
        frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = frame.table
        for r in range(2):
            for c in range(2):
                cell = table.cell(r, c)
                cell.text = f"{r}{c}"
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        table.cell(1, 1).text_frame.paragraphs[0].runs[0].font.bold = True

    (table,) = [e for e in _elements(pptx_to_slides(_deck(build))) if e["type"] == "table"]
    assert table["fontSize"] == 16.0 and table["color"] == "#112233"
    assert "fill" not in table and "bold" not in table
    assert table["cells"] == [
        {"r": 0, "c": 0, "fill": "#DDEEFF"},
        {"r": 1, "c": 1, "bold": True},
    ]


def test_a_text_outline_comes_across_as_stroke_and_goes_back_out() -> None:
    """WordArt: a light fill drawn legible by a dark outline. Only the fill
    came across and the word faded into the page."""
    from lxml import etree

    from langchain_canvas.exporters import SlidesPptxExporter
    from langchain_canvas.replay import encode_slides

    def outline(run: Any, qn: Any) -> None:
        rpr = run._r.get_or_add_rPr()
        line = etree.SubElement(rpr, qn("a:ln"))
        line.set("w", "19050")
        fill = etree.SubElement(line, qn("a:solidFill"))
        etree.SubElement(fill, qn("a:srgbClr")).set("val", "595959")
        rpr.insert(0, line)

    data = pptx_to_slides(_with_run_xml(_deck(lambda s: _textbox(s, "BIG")), outline))
    element = _elements(data)[0]
    assert element["stroke"] == "#595959" and element["strokeWidth"] == 2.0

    exported = SlidesPptxExporter().export(encode_slides("D", data), path="d.slides.json")
    back = Presentation(io.BytesIO(exported.data))
    box = next(s for s in back.slides[0].shapes if s.has_text_frame)
    run = box.text_frame.paragraphs[0].runs[0]
    line = run.font._rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ln")
    assert line is not None and line.get("w") == "19050"


# --- autofit comes across --------------------------------------------------------------


def test_a_box_that_grows_with_its_text_comes_across_as_such() -> None:
    """Forty-five of forty-five text boxes in one uploaded deck grew with
    their text; the copy froze every one at its placeholder's height, and a
    body written longer than the placeholder ran past the box."""
    from pptx.enum.text import MSO_AUTO_SIZE

    def build(slide: Any) -> None:
        grows = _textbox(slide, "grows", top=Inches(1))
        grows.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        shrinks = _textbox(slide, "shrinks", top=Inches(2.5))
        shrinks.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        fixed = _textbox(slide, "fixed", top=Inches(4))
        fixed.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        # python-pptx's fresh textbox says spAutoFit on its own; a box whose
        # file says nothing at all has no autofit element.
        unsaid = _textbox(slide, "unsaid", top=Inches(5.5))
        unsaid.text_frame.auto_size = None

    by_text = {e["text"]: e for e in _elements(pptx_to_slides(_deck(build)))}
    assert by_text["grows"]["autofit"] == "shape"
    assert by_text["shrinks"]["autofit"] == "text"
    assert "autofit" not in by_text["fixed"]
    assert "autofit" not in by_text["unsaid"]
