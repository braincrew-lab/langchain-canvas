"""Tests for the standard canvas tools.

Tools are invoked directly (LangChain tools accept an injected ``runtime``
kwarg), with a minimal stand-in runtime — no model, no graph. The important
contracts: canvas-id resolution, line-numbered reads exposing the revision,
and edit_canvas rejecting stale or ambiguous edits with a retry hint instead
of raising.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from typing import Any

import pytest

from langchain_canvas import InMemoryCanvasStore, create_canvas_tools
from langchain_canvas.converters import default_converters


@dataclass
class _Runtime:
    """The slice of ToolRuntime the canvas tools read."""

    context: Any = None
    config: dict[str, Any] = field(default_factory=dict)


def _runtime(canvas_id: str | None = None, thread_id: str | None = None) -> _Runtime:
    configurable: dict[str, Any] = {}
    if canvas_id:
        configurable["canvas_id"] = canvas_id
    if thread_id:
        configurable["thread_id"] = thread_id
    return _Runtime(config={"configurable": configurable})


def _tools(store: InMemoryCanvasStore) -> dict[str, Any]:
    return {t.name: t for t in create_canvas_tools(store)}


def _invoke(tool_obj: Any, runtime: _Runtime, **kwargs: Any) -> str:
    return tool_obj.func(runtime=runtime, **kwargs)


# --- canvas id resolution --------------------------------------------------------


def test_thread_id_is_the_default_canvas_scope() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _runtime(thread_id="t1")
    _invoke(tools["write_canvas"], runtime, path="a.html", content="x", description="create")
    assert store.read("t1", "a.html").content == "x"


def test_explicit_canvas_id_wins_over_thread_id() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    _invoke(
        tools["write_canvas"],
        _runtime(canvas_id="c9", thread_id="t1"),
        path="a.html",
        content="x",
        description="create",
    )
    assert store.read("c9", "a.html").content == "x"
    assert store.history("t1") == []


def test_context_dict_canvas_id_wins() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _Runtime(context={"canvas_id": "ctx"}, config={"configurable": {"thread_id": "t1"}})
    _invoke(tools["write_canvas"], runtime, path="a.html", content="x", description="create")
    assert store.read("ctx", "a.html").content == "x"


# --- read ------------------------------------------------------------------------


def test_read_returns_revision_and_line_numbers() -> None:
    store = InMemoryCanvasStore()
    commit = store.write("t1", "a.html", "<h1>Hi</h1>\n<p>Body</p>", "create")
    out = _invoke(_tools(store)["read_canvas"], _runtime(thread_id="t1"), path="a.html")
    assert out.startswith(f"revision: {commit.revision}\n")
    assert "   1\t<h1>Hi</h1>" in out
    assert "   2\t<p>Body</p>" in out


def test_read_missing_file_returns_error_text() -> None:
    out = _invoke(
        _tools(InMemoryCanvasStore())["read_canvas"], _runtime(thread_id="t1"), path="a.html"
    )
    assert out.startswith("Error:")
    assert "list_canvas_files" in out


def _wide_table(sheets: int, rows: int) -> str:
    return json.dumps({
        "type": "table",
        "title": "Ledger",
        "data": {
            "columns": [{"key": "n"}, {"key": "city"}],
            "rows": [{"n": i, "city": f"city{i}"} for i in range(rows)],
            "sheet": [
                {"name": f"S{i}", "row": 1000, "column": 30, "config": {},
                 "celldata": [{"r": 0, "c": 0, "v": {"v": f"sheet {i}"}}]}
                for i in range(sheets)
            ],
        },
    })


def test_reading_a_table_answers_with_its_map_not_its_contents() -> None:
    # The grid state is the person's, and it is where the size is: a real
    # five-sheet import in the examples is 29.4M characters against 48k of
    # rows. Handing the file over whole spends the context on borders.
    store = InMemoryCanvasStore()
    content = _wide_table(sheets=3, rows=500)
    store.write("t1", "ledger.table.json", content, "create")
    out = _invoke(_tools(store)["read_canvas"], _runtime(thread_id="t1"), path="ledger.table.json")
    assert "[rows] 500 x 2 — what agents read and write" in out
    assert "[s0] S0 — 1000 x 30 grid, 1 value" in out
    assert "[s2] S2" in out
    assert "rewrites all 3 grid sheets" in out
    assert len(out) < len(content) / 10


def test_a_table_address_reads_one_rectangle_through_the_usual_window() -> None:
    store = InMemoryCanvasStore()
    store.write("t1", "ledger.table.json", _wide_table(sheets=2, rows=500), "create")
    tools = _tools(store)
    out = _invoke(tools["read_canvas"], _runtime(thread_id="t1"),
                  path="ledger.table.json", sheet="rows", offset=3, limit=2)
    assert "### sheet: rows" in out
    assert "   4	2,city2" in out
    assert "   5	3,city3" in out
    assert "of 501 — call read_canvas again with offset=5" in out
    grid = _invoke(tools["read_canvas"], _runtime(thread_id="t1"),
                   path="ledger.table.json", sheet="s1")
    assert grid.endswith("### sheet: S1 (columns A-A)\n   1\tsheet 1")


def test_a_file_that_is_not_really_a_table_is_still_read_as_itself() -> None:
    # The route is the suffix, so an envelope of another type can arrive
    # here. Reporting it as an empty table would hide the deck inside it.
    store = InMemoryCanvasStore()
    deck = json.dumps({"type": "slides", "data": {"slides": [{"title": "T"}]}})
    store.write("t1", "oops.table.json", deck, "create")
    out = _invoke(_tools(store)["read_canvas"], _runtime(thread_id="t1"), path="oops.table.json")
    assert "[rows]" not in out
    assert '"type": "slides"' in out


def test_a_table_address_that_does_not_exist_names_the_ones_that_do() -> None:
    store = InMemoryCanvasStore()
    store.write("t1", "ledger.table.json", _wide_table(sheets=2, rows=2), "create")
    store.write("t1", "page.html", "<p>x</p>", "create")
    tools = _tools(store)
    missing = _invoke(tools["read_canvas"], _runtime(thread_id="t1"),
                      path="ledger.table.json", sheet="s5")
    assert missing == "Error: sheet must be one of: rows, s0, s1 (got 's5')."
    wrong = _invoke(tools["read_canvas"], _runtime(thread_id="t1"),
                    path="page.html", sheet="s0")
    assert "`sheet` applies to .table.json tables" in wrong


def _table_tools(store: InMemoryCanvasStore) -> dict[str, Any]:
    from langchain_canvas import create_table_tools

    tools = _tools(store)
    tools.update({t.name: t for t in create_table_tools(store)})
    return tools


def _ledger() -> str:
    return json.dumps({"type": "table", "title": "Q3", "data": {
        "columns": [{"key": "region", "label": "region"}, {"key": "amount", "label": "amount"}],
        "rows": [{"region": "Seoul", "amount": 10}],
        "sheet": [{"name": "Ledger", "id": "sheet_0", "order": 0, "status": 1,
                   "row": 20, "column": 6, "config": {"merge": {}},
                   "celldata": [{"r": 0, "c": 0, "v": {"v": "region", "m": "region"}},
                                {"r": 0, "c": 1, "v": {"v": "amount", "m": "amount"}},
                                {"r": 1, "c": 0, "v": {"v": "Seoul", "m": "Seoul"}},
                                {"r": 1, "c": 1, "v": {"v": 10, "m": "10"}}]}],
    }})


def test_a_table_is_written_at_the_addresses_it_was_read_at() -> None:
    # The read prints the column letters and numbers the rows, so B2 is
    # visible rather than counted. Rewriting the whole file would take every
    # sheet and all the person's formatting with it.
    store = InMemoryCanvasStore()
    commit = store.write("t1", "q3.table.json", _ledger(), "create")
    tools = _table_tools(store)
    shown = _invoke(tools["read_canvas"], _runtime(thread_id="t1"),
                    path="q3.table.json", sheet="s0")
    assert "### sheet: Ledger (columns A-B)" in shown
    assert "   2\tSeoul,10" in shown

    out = _invoke(tools["write_table_cells"], _runtime(thread_id="t1"), path="q3.table.json",
                  sheet="s0", cells={"B2": 99}, description="fix", revision=commit.revision)
    assert out.startswith("Wrote B2 on Ledger.")
    # rows is the projection of s0, and the xlsx export reads it, so it moves
    # with the write instead of putting the old value back.
    assert "rows now has 1 entries" in out
    data = json.loads(store.read("t1", "q3.table.json").content)["data"]
    assert data["rows"] == [{"region": "Seoul", "amount": 99}]


def test_a_table_write_refuses_a_stale_revision() -> None:
    store = InMemoryCanvasStore()
    store.write("t1", "q3.table.json", _ledger(), "create")
    tools = _table_tools(store)
    out = _invoke(tools["write_table_cells"], _runtime(thread_id="t1"), path="q3.table.json",
                  sheet="s0", cells={"B2": 1}, description="fix", revision="v0")
    assert out.startswith("Error:") and "read_canvas again" in out
    stale = _invoke(tools["write_table_cells"], _runtime(thread_id="t1"), path="page.html",
                    sheet="s0", cells={"B2": 1}, description="fix", revision="v1")
    assert ".table.json tables" in stale


def test_a_sheet_added_through_the_tool_shows_up_on_the_map() -> None:
    store = InMemoryCanvasStore()
    commit = store.write("t1", "q3.table.json", _ledger(), "create")
    tools = _table_tools(store)
    out = _invoke(tools["add_table_sheet"], _runtime(thread_id="t1"), path="q3.table.json",
                  name="Summary", description="new", revision=commit.revision)
    assert out.startswith('Added sheet "Summary" as s1.')
    card = _invoke(tools["read_canvas"], _runtime(thread_id="t1"), path="q3.table.json")
    assert "[s0] Ledger" in card
    assert "[s1] Summary — 24 x 10 grid, no values" in card


# --- edit: read-before-update ----------------------------------------------------


def test_edit_with_current_revision_succeeds() -> None:
    store = InMemoryCanvasStore()
    commit = store.write("t1", "a.html", "<h1>[X] Title</h1>", "create")
    out = _invoke(
        _tools(store)["edit_canvas"],
        _runtime(thread_id="t1"),
        path="a.html",
        old="[X] ",
        new="",
        description="Remove marker",
        revision=commit.revision,
    )
    assert out.startswith("Edited a.html")
    assert store.read("t1", "a.html").content == "<h1>Title</h1>"


def test_edit_with_stale_revision_is_rejected_with_retry_hint() -> None:
    store = InMemoryCanvasStore()
    stale = store.write("t1", "a.html", "one", "create")
    store.write("t1", "a.html", "two", "human edit")
    out = _invoke(
        _tools(store)["edit_canvas"],
        _runtime(thread_id="t1"),
        path="a.html",
        old="two",
        new="three",
        description="edit",
        revision=stale.revision,
    )
    assert out.startswith("Error:")
    assert "read_canvas" in out
    assert store.read("t1", "a.html").content == "two"


def test_edit_ambiguous_match_is_rejected() -> None:
    store = InMemoryCanvasStore()
    commit = store.write("t1", "a.html", "dup dup", "create")
    out = _invoke(
        _tools(store)["edit_canvas"],
        _runtime(thread_id="t1"),
        path="a.html",
        old="dup",
        new="x",
        description="edit",
        revision=commit.revision,
    )
    assert out.startswith("Error:")


# --- list / history --------------------------------------------------------------


def test_list_canvas_empty_and_populated() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    assert _invoke(tools["list_canvas_files"], _runtime(thread_id="t1")) == "The canvas is empty."
    store.write("t1", "a.html", "abc", "create")
    assert _invoke(tools["list_canvas_files"], _runtime(thread_id="t1")) == "a.html (3 bytes)"


def test_tool_writes_produce_described_history() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _runtime(thread_id="t1")
    _invoke(tools["write_canvas"], runtime, path="a.html", content="one", description="Create")
    revision = store.read("t1", "a.html").revision
    _invoke(
        tools["edit_canvas"],
        runtime,
        path="a.html",
        old="one",
        new="two",
        description="Fix wording",
        revision=revision,
    )
    assert [c.description for c in store.history("t1")] == ["Fix wording", "Create"]


# --- live broadcast --------------------------------------------------------------


@dataclass
class _StreamingRuntime(_Runtime):
    """Runtime with a collecting stream writer, like a live LangGraph run."""

    events: list[dict] = field(default_factory=list)

    @property
    def stream_writer(self):  # noqa: ANN201 — mirrors ToolRuntime's attribute
        return self.events.append


def _streaming_runtime(thread_id: str) -> _StreamingRuntime:
    return _StreamingRuntime(config={"configurable": {"thread_id": thread_id}})


def test_write_broadcasts_create_then_patch() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _streaming_runtime("t1")
    _invoke(tools["write_canvas"], runtime, path="a.html", content="<p>1</p>", description="create")
    revision = store.read("t1", "a.html").revision
    _invoke(
        tools["write_canvas"],
        runtime,
        path="a.html",
        content="<p>2</p>",
        description="rewrite",
        revision=revision,
    )
    kinds = [e["type"] for e in runtime.events]
    assert kinds == [
        "canvas.create",
        "canvas.status",
        "canvas.commit",
        "canvas.patch",
        "canvas.commit",
    ]
    assert runtime.events[0]["artifact"]["data"]["html"] == "<p>1</p>"
    assert runtime.events[3]["patch"]["html"] == "<p>2</p>"


def test_edit_broadcasts_patch_and_commit() -> None:
    store = InMemoryCanvasStore()
    commit = store.write("t1", "a.html", "<p>one</p>", "create")
    runtime = _streaming_runtime("t1")
    _invoke(
        _tools(store)["edit_canvas"],
        runtime,
        path="a.html",
        old="one",
        new="two",
        description="tweak",
        revision=commit.revision,
    )
    kinds = [e["type"] for e in runtime.events]
    assert kinds == ["canvas.patch", "canvas.commit"]
    assert runtime.events[0]["patch"]["html"] == "<p>two</p>"
    assert runtime.events[1]["description"] == "tweak"


def test_non_html_and_failed_writes_broadcast_nothing() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _streaming_runtime("t1")
    _invoke(tools["write_canvas"], runtime, path="notes.md", content="hi", description="notes")
    stale = store.write("t1", "a.html", "one", "create")
    store.write("t1", "a.html", "two", "human edit")
    events_before = list(runtime.events)
    _invoke(
        tools["write_canvas"],
        runtime,
        path="a.html",
        content="three",
        description="stale",
        revision=stale.revision,
    )
    assert runtime.events == events_before  # rejected write emits nothing


def test_no_stream_writer_is_a_silent_noop() -> None:
    store = InMemoryCanvasStore()
    out = _invoke(
        _tools(store)["write_canvas"],
        _runtime(thread_id="t1"),
        path="a.html",
        content="<p>hi</p>",
        description="create",
    )
    assert out.startswith("Wrote a.html")


def test_broadcast_applies_title_and_meta_conventions() -> None:
    store = InMemoryCanvasStore()
    tools = {
        t.name: t
        for t in create_canvas_tools(
            store,
            title_for=lambda path: "Intro" if path == "01-intro.html" else path,
            meta_for=lambda path: {"kind": "slide", "ratio": "16:9"},
        )
    }
    runtime = _streaming_runtime("t1")
    _invoke(
        tools["write_canvas"],
        runtime,
        path="01-intro.html",
        content="<p>s</p>",
        description="slide",
    )
    artifact = runtime.events[0]["artifact"]
    assert artifact["title"] == "Intro"
    assert artifact["meta"] == {"kind": "slide", "ratio": "16:9"}


def test_table_write_broadcasts_table_artifact() -> None:
    from langchain_canvas import encode_table

    store = InMemoryCanvasStore()
    tools = _tools(store)
    runtime = _streaming_runtime("t1")
    content = encode_table("Compare", {"columns": [{"key": "a", "label": "A"}], "rows": [{"a": 1}]})
    _invoke(
        tools["write_canvas"],
        runtime,
        path="compare.table.json",
        content=content,
        description="table",
    )
    kinds = [e["type"] for e in runtime.events]
    assert kinds == ["canvas.create", "canvas.status", "canvas.commit"]
    assert runtime.events[0]["artifact"]["type"] == "table"
    assert runtime.events[0]["artifact"]["data"]["rows"] == [{"a": 1}]

    # A malformed table file persists (it may be mid-repair) but broadcasts nothing.
    revision = store.read("t1", "compare.table.json").revision
    _invoke(
        tools["write_canvas"],
        runtime,
        path="compare.table.json",
        content="not json {",
        description="broken",
        revision=revision,
    )
    assert [e["type"] for e in runtime.events] == kinds


# --- sources: uploads, converters, windows ---------------------------------------


def test_sources_are_read_only_for_the_agent() -> None:
    store = InMemoryCanvasStore()
    store.write("t1", "sources/notes.md", "# hi", "Upload notes.md", actor="human")
    tools = _tools(store)
    runtime = _runtime(thread_id="t1")
    for call in (
        lambda: _invoke(
            tools["write_canvas"],
            runtime,
            path="sources/notes.md",
            content="x",
            description="try",
        ),
        lambda: _invoke(
            tools["edit_canvas"],
            runtime,
            path="sources/notes.md",
            old="hi",
            new="bye",
            description="try",
            revision="v1",
        ),
    ):
        out = call()
        assert "read-only" in out
    # Nothing was committed by the rejected calls.
    assert len(store.history("t1")) == 1


def test_binary_source_without_converter_gets_honest_error() -> None:
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/archive.zip", b"PK\x03\x04\x00\xff", "Upload archive.zip")
    out = _invoke(
        _tools(store)["read_canvas"], _runtime(thread_id="t1"), path="sources/archive.zip"
    )
    assert "no converter handles it" in out
    assert ".md" in out  # the installed set is named


def test_custom_converter_renders_binary_source() -> None:
    from langchain_canvas import ConvertedSource, create_canvas_tools

    class ShoutConverter:
        suffixes = (".png",)

        def convert(self, data: bytes, *, path: str) -> ConvertedSource:
            return ConvertedSource(
                blocks=[{"type": "text", "text": f"IMAGE {path} ({len(data)} bytes)"}],
                metadata={"kind": "shout"},
            )

    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/photo.png", b"\x89PNG\x00\xff", "Upload photo.png")
    tools = {t.name: t for t in create_canvas_tools(store, converters=[ShoutConverter()])}
    out = _invoke(tools["read_canvas"], _runtime(thread_id="t1"), path="sources/photo.png")
    assert "converted view of sources/photo.png" in out
    assert "kind: shout" in out
    assert "IMAGE sources/photo.png (6 bytes)" in out


def test_read_canvas_windows_long_files() -> None:
    store = InMemoryCanvasStore()
    content = "\n".join(f"line {i}" for i in range(1, 1001))
    store.write("t1", "big.html", content, "create")
    tools = _tools(store)
    runtime = _runtime(thread_id="t1")

    first = _invoke(tools["read_canvas"], runtime, path="big.html")
    assert "line 400" in first and "line 401" not in first
    assert "offset=400" in first

    second = _invoke(tools["read_canvas"], runtime, path="big.html", offset=400, limit=100)
    assert "line 401" in second and "line 500" in second and "line 501" not in second
    assert "offset=500" in second


def test_write_canvas_fills_page_from_the_template_skin():
    import io as _io
    import json as _json

    pptx = __import__("pytest").importorskip("pptx")
    from pptx.util import Inches

    skin = pptx.Presentation()
    skin.slide_width = Inches(10)
    skin.slide_height = Inches(7.5)
    buf = _io.BytesIO()
    skin.save(buf)
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/brand.pptx", buf.getvalue(), "skin")
    tools = _tools(store)
    content = _json.dumps(
        {"type": "slides", "data": {"template": "sources/brand.pptx", "slides": []}}
    )
    result = _invoke(
        tools["write_canvas"], _runtime(thread_id="t1"), path="deck.slides.json",
        content=content, description="deck",
    )
    assert result.startswith("Wrote")
    saved = _json.loads(store.read("t1", "deck.slides.json").content)
    # The tool learned the skin's real page — the agent never types numbers.
    assert saved["data"]["page"] == {"widthIn": 10.0, "heightIn": 7.5}

    # The skin decides the page on a swap too: a stale page is replaced
    # with the template's real size.
    content2 = _json.dumps({"type": "slides", "data": {
        "template": "sources/brand.pptx",
        "page": {"widthIn": 12.0, "heightIn": 9.0}, "slides": [],
    }})
    _invoke(
        tools["write_canvas"], _runtime(thread_id="t1"), path="deck2.slides.json",
        content=content2, description="deck",
    )
    saved2 = _json.loads(store.read("t1", "deck2.slides.json").content)
    assert saved2["data"]["page"] == {"widthIn": 10.0, "heightIn": 7.5}


def _skin_store(width_in: float, height_in: float) -> InMemoryCanvasStore:
    import io as _io

    pptx = __import__("pytest").importorskip("pptx")
    from pptx.util import Inches

    skin = pptx.Presentation()
    skin.slide_width = Inches(width_in)
    skin.slide_height = Inches(height_in)
    buf = _io.BytesIO()
    skin.save(buf)
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/brand.pptx", buf.getvalue(), "skin")
    return store


# A perfect circle on the classic 16:9 canvas plus one text element.
_LEGACY_DECK_DATA: dict[str, Any] = {
    "template": "sources/brand.pptx",
    "slides": [{"elements": [
        {"id": "c", "type": "shape", "shape": "ellipse", "x": 10, "y": 10,
         "w": 10, "h": 17.7778, "fill": "#ff0000"},
        {"id": "t", "type": "text", "x": 30, "y": 60, "w": 40, "h": 20,
         "text": "Hi", "fontSize": 40},
    ]}],
}


def test_write_canvas_refits_a_legacy_deck_for_a_new_ratio():
    # Attaching a skin LATER is the common flow: the deck was drawn on the
    # 16:9 canvas, then "use our corporate template" arrives. Filling `page`
    # changes the coordinate space, so the stored percents must be
    # re-projected — leaving them made a circle export as ratio 0.750.
    import copy as _copy
    import io as _io
    import json as _json

    pptx = __import__("pytest").importorskip("pytest").importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    store = _skin_store(10, 7.5)  # 4:3 — a different ratio than the canvas
    tools = _tools(store)
    content = _json.dumps({"type": "slides",
                           "data": _copy.deepcopy(_LEGACY_DECK_DATA)})
    result = _invoke(
        tools["write_canvas"], _runtime(thread_id="t1"), path="deck.slides.json",
        content=content, description="deck",
    )
    # The ratio change is announced, never silent.
    assert "page changed to 10.0 x 7.5 in" in result
    assert "scaled to fit" in result

    saved = _json.loads(store.read("t1", "deck.slides.json").content)
    assert saved["data"]["page"] == {"widthIn": 10.0, "heightIn": 7.5}

    # The exported file keeps the circle a circle, centered in the letterbox.
    from langchain_canvas.assets import inline_slides_assets
    from langchain_canvas.exporters import SlidesPptxExporter

    content_saved = store.read("t1", "deck.slides.json").content
    inlined = inline_slides_assets(content_saved, store, "t1")
    exported = SlidesPptxExporter().export(inlined, path="deck.slides.json")
    deck = pptx.Presentation(_io.BytesIO(exported.data))
    (slide,) = deck.slides
    (circle,) = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert abs(circle.width / circle.height - 1.0) < 0.001
    # Same on-page picture as before: 1in circle, dropped by the centering
    # offset ((7.5 - 5.625) / 2 = 0.9375in) on the taller page.
    assert abs(circle.width / 914400 - 1.0) < 0.001
    assert abs(circle.top / 914400 - (0.5625 + 0.9375)) < 0.001
    (text,) = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text]
    assert text.text_frame.paragraphs[0].runs[0].font.size.pt == 30  # unchanged


def test_write_canvas_stays_quiet_for_a_same_ratio_skin():
    import copy as _copy
    import json as _json

    __import__("pytest").importorskip("pptx")
    store = _skin_store(10, 5.625)  # the classic canvas ratio
    tools = _tools(store)
    content = _json.dumps({"type": "slides",
                           "data": _copy.deepcopy(_LEGACY_DECK_DATA)})
    result = _invoke(
        tools["write_canvas"], _runtime(thread_id="t1"), path="deck.slides.json",
        content=content, description="deck",
    )
    assert "scaled to fit" not in result
    saved = _json.loads(store.read("t1", "deck.slides.json").content)
    # Coordinates untouched — same ratio needs no re-fit.
    assert saved["data"]["slides"][0]["elements"][0]["x"] == 10


def test_skin_swaps_never_push_content_off_the_page():
    # The refit contract, corrected: containment and uniform scale are hard;
    # round-trip identity and path-independence are NOT promised (min-scale
    # letterboxing is not invertible — slide software asks the user how to
    # re-fit on a page change for the same reason). The earlier
    # path-independent projection routed through the classic canvas and
    # silently pushed content re-placed for the current page off the page.
    import copy as _copy
    import io as _io
    import json as _json

    pptx = __import__("pytest").importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    from langchain_canvas.assets import inline_slides_assets
    from langchain_canvas.exporters import SlidesPptxExporter

    def _add_skin(store, w, h):
        skin = pptx.Presentation()
        skin.slide_width = Inches(w)
        skin.slide_height = Inches(h)
        buf = _io.BytesIO()
        skin.save(buf)
        store.write_bytes("t1", "sources/brand.pptx", buf.getvalue(), "skin")

    def _check_export(store):
        """Every shape inside the page; the circle still a circle."""
        content = store.read("t1", "deck.slides.json").content
        inlined = inline_slides_assets(content, store, "t1")
        deck = pptx.Presentation(
            _io.BytesIO(SlidesPptxExporter().export(inlined, path="d.slides.json").data)
        )
        page_w, page_h = deck.slide_width, deck.slide_height
        (slide,) = deck.slides
        for shape in slide.shapes:
            assert shape.left >= -1 and shape.top >= -1
            assert shape.left + shape.width <= page_w + 1
            assert shape.top + shape.height <= page_h + 1
        (circle,) = [
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]
        assert abs(circle.width / circle.height - 1.0) < 0.001

    def _rewrite(store, tools, content=None):
        got = store.read("t1", "deck.slides.json")
        _invoke(tools["write_canvas"], _runtime(thread_id="t1"),
                path="deck.slides.json", content=content or got.content,
                description="swap", revision=got.revision)

    pages = {"16:9": (10, 5.625), "4:3": (10, 7.5), "wide": (13.333, 7.5)}
    for order in (("4:3", "wide", "16:9"), ("wide", "4:3", "wide"),
                  ("16:9", "4:3", "wide", "4:3")):
        store = InMemoryCanvasStore()
        _add_skin(store, *pages[order[0]])
        tools = _tools(store)
        content = _json.dumps({"type": "slides",
                               "data": _copy.deepcopy(_LEGACY_DECK_DATA)})
        _invoke(tools["write_canvas"], _runtime(thread_id="t1"),
                path="deck.slides.json", content=content, description="deck")
        _check_export(store)
        for name in order[1:]:
            _add_skin(store, *pages[name])
            _rewrite(store, tools)
            _check_export(store)

    # The advertised flow: attach 4:3, RE-PLACE the elements to fill the 4:3
    # page (as the note suggests), then swap to wide. The re-placed content
    # is no longer inside any classic-canvas band — it must still stay on
    # the page.
    store = InMemoryCanvasStore()
    _add_skin(store, 10, 7.5)
    tools = _tools(store)
    content = _json.dumps({"type": "slides",
                           "data": _copy.deepcopy(_LEGACY_DECK_DATA)})
    _invoke(tools["write_canvas"], _runtime(thread_id="t1"),
            path="deck.slides.json", content=content, description="deck")
    replaced = _json.loads(store.read("t1", "deck.slides.json").content)
    for el in replaced["data"]["slides"][0]["elements"]:
        if el["type"] == "shape":  # blow the circle up to the full 4:3 page
            el.update({"x": 0, "y": 0, "w": 75, "h": 100})
        else:
            el.update({"x": 0, "y": 80, "w": 40, "h": 20})
    _rewrite(store, tools, content=_json.dumps(replaced))
    _add_skin(store, 13.333, 7.5)
    _rewrite(store, tools)
    _check_export(store)

    # A rewrite with the matching skin stays quiet and untouched.
    got = store.read("t1", "deck.slides.json")
    result = _invoke(tools["write_canvas"], _runtime(thread_id="t1"),
                     path="deck.slides.json", content=got.content,
                     description="no-op", revision=got.revision)
    assert "scaled to fit" not in result


def test_the_read_tool_names_the_formats_its_eye_actually_reaches() -> None:
    """A format the description omits is one the model will decline to look at."""
    tools = {t.name: t for t in create_canvas_tools(InMemoryCanvasStore())}
    description = tools["read_canvas"].description
    assert "{page_formats}" not in description
    assert ".pdf" in description


def test_a_host_converter_gets_announced_too() -> None:
    class _DeckPages:
        suffixes = (".pptx",)

        def convert(self, data: bytes, *, path: str):  # pragma: no cover - unused
            raise NotImplementedError

        def render_pages(self, data: bytes, *, path: str, pages: list[int]):  # pragma: no cover
            raise NotImplementedError

        def render_grid(self, data: bytes, *, path: str):  # pragma: no cover
            raise NotImplementedError

    tools = {
        t.name: t
        for t in create_canvas_tools(
            InMemoryCanvasStore(), converters=[_DeckPages(), *default_converters()]
        )
    }
    assert ".pptx" in tools["read_canvas"].description


def test_a_store_with_no_page_renderable_converter_says_so() -> None:
    class _TextOnly:
        suffixes = (".md",)

        def convert(self, data: bytes, *, path: str):  # pragma: no cover - unused
            raise NotImplementedError

    tools = {
        t.name: t
        for t in create_canvas_tools(InMemoryCanvasStore(), converters=[_TextOnly()])
    }
    assert "none installed" in tools["read_canvas"].description


# --- documents: editing what the user uploaded -----------------------------------


def _document_tools(store: InMemoryCanvasStore, converters: Any = None) -> dict[str, Any]:
    from langchain_canvas import create_document_tools

    built = [
        *create_canvas_tools(store, converters=converters),
        *create_document_tools(store, converters=converters),
    ]
    return {t.name: t for t in built}


def _uploaded(store: InMemoryCanvasStore, path: str = "sources/plan.docx") -> str:
    from documents import sample_document

    commit = store.write_bytes("t1", path, sample_document(), "Upload", actor="human")
    return commit.revision


def _open_copy(tools: dict[str, Any], runtime: Any) -> tuple[str, str]:
    """Copy the upload out of sources/ and return (path, revision).

    The path is read back out of the reply rather than spelled here, so the
    tests exercise whatever name the tool actually chose.
    """
    reply = _invoke(tools["open_document_for_editing"], runtime, source="sources/plan.docx")
    assert reply.startswith("Copied "), reply
    path = reply.split(" to ", 1)[1].rsplit(" (", 1)[0]
    revision = reply.split("revision ")[1].split(")")[0]
    return path, revision


def test_reading_an_uploaded_document_gives_addresses(tmp_path: Any = None) -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    text = _invoke(tools["read_canvas"], _runtime(thread_id="t1"), path="sources/plan.docx")
    assert "[p0] (Heading 1) 2026 반영계획안" in text
    assert "[t0] 3x3 table" in text
    assert "[img0]" in text


def test_uploads_stay_read_only_for_every_document_operation() -> None:
    """I4 — the guard over the user's originals does not open for documents."""
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    before = store.read_bytes("t1", "sources/plan.docx").data

    edited = _invoke(
        tools["edit_canvas"],
        runtime,
        path="sources/plan.docx",
        old="즉시 조치",
        new="즉각 조치",
        description="x",
        revision="1",
    )
    assert "read-only" in edited
    # Not the generic "make an .html page instead": a Word upload has a way
    # forward, and naming the wrong one tells the agent it cannot do this.
    assert "open_document_for_editing" in edited
    assert ".html" not in edited
    for name, extra in (
        ("insert_document_paragraph", {"anchor": "사진 1. 점검 당일 현장", "text": "x"}),
        ("insert_document_image", {"image_path": "assets/x.png"}),
        ("remove_document_paragraph", {"anchor": "사진 1. 점검 당일 현장"}),
        ("replace_document_image", {"index": 0, "image_path": "assets/x.png"}),
    ):
        reply = _invoke(
            tools[name], runtime, path="sources/plan.docx", description="x", revision="1", **extra
        )
        assert "read-only" in reply, name
        assert "open_document_for_editing" in reply, name
    assert store.read_bytes("t1", "sources/plan.docx").data == before


def test_the_working_copy_leaves_the_upload_alone() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    original = store.read_bytes("t1", "sources/plan.docx").data

    path, revision = _open_copy(tools, runtime)
    assert store.read_bytes("t1", path).data == original

    reply = _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치가 필요한",
        new="즉시 조치가 반드시 필요한",
        description="Sharpen the finding",
        revision=revision,
    )
    assert reply.startswith(f"Edited {path}")
    assert "반드시" in _invoke(tools["read_canvas"], runtime, path=path)
    assert store.read_bytes("t1", "sources/plan.docx").data == original


def test_the_copy_is_named_apart_from_the_original() -> None:
    """I3 — two tabs, and the user can tell which one is theirs."""
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    path, _ = _open_copy(tools, _runtime(thread_id="t1"))
    original = "sources/plan.docx"
    assert path != original
    assert path.rsplit("/", 1)[-1] != original.rsplit("/", 1)[-1]
    # The mark is in front: a tab shows the start of a name and clips the end.
    assert path.startswith("Editing - ")
    assert "Editing - " in tools["open_document_for_editing"].description


def test_copying_a_copy_does_not_stack_the_mark() -> None:
    store = InMemoryCanvasStore()
    from documents import sample_document

    store.write_bytes("t1", "Editing - plan.docx", sample_document(), "x", actor="agent")
    tools = _document_tools(store)
    reply = _invoke(
        tools["open_document_for_editing"], _runtime(thread_id="t1"), source="Editing - plan.docx"
    )
    assert "already on the canvas" in reply


def test_copying_twice_names_the_conflict() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    _open_copy(tools, runtime)
    again = _invoke(tools["open_document_for_editing"], runtime, source="sources/plan.docx")
    assert "already on the canvas" in again
    assert "destination" in again


def test_the_copy_can_be_given_another_name() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    reply = _invoke(
        tools["open_document_for_editing"],
        runtime,
        source="sources/plan.docx",
        destination="draft/plan-v2.docx",
    )
    assert "draft/plan-v2.docx" in reply
    assert store.read_bytes("t1", "draft/plan-v2.docx").data


def test_the_copy_cannot_land_back_in_sources() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    reply = _invoke(
        tools["open_document_for_editing"],
        _runtime(thread_id="t1"),
        source="sources/plan.docx",
        destination="sources/copy.docx",
    )
    assert "sources/" in reply and reply.startswith("Error")


def test_a_stale_revision_is_refused_the_same_way_text_files_are() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치",
        new="즉각 조치",
        description="first",
        revision=revision,
    )
    late = _invoke(
        tools["insert_document_paragraph"],
        runtime,
        path=path,
        anchor="사진 1. 점검 당일 현장",
        text="사진 2.",
        description="second",
        revision=revision,
    )
    assert "Call read_canvas again" in late


def test_insert_remove_and_picture_swap_land_on_the_copy() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    from documents import png_bytes

    store.write_bytes("t1", "assets/new.png", png_bytes(200, 50, (3, 3, 3)), "Add", actor="human")
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)

    added = _invoke(
        tools["insert_document_paragraph"],
        runtime,
        path=path,
        anchor="사진 1. 점검 당일 현장",
        text="4. 후속 조치",
        style="Heading 2",
        description="Add a closing section",
        revision=revision,
    )
    assert added.startswith(f"Added a paragraph to {path}")
    revision = added.split("revision ")[1].split(")")[0]
    assert "[p10] (Heading 2) 4. 후속 조치" in _invoke(tools["read_canvas"], runtime, path=path)

    removed = _invoke(
        tools["remove_document_paragraph"],
        runtime,
        path=path,
        anchor="담당 부서와 일정은 별도 협의한다.",
        description="Drop a bullet",
        revision=revision,
    )
    assert removed.startswith(f"Removed a paragraph from {path}")
    revision = removed.split("revision ")[1].split(")")[0]
    assert "담당 부서와 일정" not in _invoke(tools["read_canvas"], runtime, path=path)

    swapped = _invoke(
        tools["replace_document_image"],
        runtime,
        path=path,
        index=0,
        image_path="assets/new.png",
        description="Use the new photo",
        revision=revision,
    )
    assert swapped.startswith(f"Replaced a picture in {path}")
    assert "width kept, height refitted" in swapped


def test_an_anchor_that_matches_nothing_reaches_the_agent_intact() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    reply = _invoke(
        tools["remove_document_paragraph"],
        runtime,
        path=path,
        anchor="담당 부서와 일정은 별도로 협의한다.",
        description="Drop a bullet",
        revision=revision,
    )
    assert "0 matches" in reply
    assert "Closest paragraph [p6]" in reply
    assert "First difference at character" in reply


def test_document_operations_refuse_a_file_that_is_not_a_document() -> None:
    store = InMemoryCanvasStore()
    store.write("t1", "notes.md", "hello", "Write")
    tools = _document_tools(store)
    reply = _invoke(
        tools["remove_document_paragraph"],
        _runtime(thread_id="t1"),
        path="notes.md",
        anchor="hello",
        description="x",
        revision="1",
    )
    assert ".docx" in reply and "edit_canvas" in reply


def test_a_picture_has_to_come_from_the_canvas() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    reply = _invoke(
        tools["replace_document_image"],
        runtime,
        path=path,
        index=0,
        image_path="https://example.com/photo.png",
        description="x",
        revision=revision,
    )
    assert "not a canvas image path" in reply


# --- documents: nothing is saved that cannot be opened and seen ------------------


class _FakePages:
    """A page renderer whose render can be told to fail."""

    suffixes = (".docx",)

    def __init__(self, working: bool = True) -> None:
        self.working = working
        self.renders = 0

    def convert(self, data: bytes, *, path: str):
        from langchain_canvas.converters import DocxSourceConverter

        return DocxSourceConverter().convert(data, path=path)

    def render_pages(self, data: bytes, *, path: str, pages: list[int]):
        self.renders += 1
        if not self.working:
            raise ValueError("the renderer is down")
        from langchain_canvas.converters import ConvertedSource

        return ConvertedSource(blocks=[{"type": "text", "text": "page"}])

    def render_grid(self, data: bytes, *, path: str):  # pragma: no cover - unused here
        raise NotImplementedError


def test_a_document_that_no_longer_renders_is_not_saved() -> None:
    """I5 — the check runs before the store, so a failure changes nothing."""
    store = InMemoryCanvasStore()
    _uploaded(store)
    renderer = _FakePages(working=True)
    converters = [renderer, *default_converters()]
    tools = _document_tools(store, converters=converters)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    before = store.read_bytes("t1", path).data

    renderer.working = False
    reply = _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치",
        new="즉각 조치",
        description="Sharpen",
        revision=revision,
    )
    assert "did not save" in reply and "no longer renders" in reply
    assert store.read_bytes("t1", path).data == before
    assert store.read_bytes("t1", path).revision == revision


def test_a_working_render_saves_and_says_to_look_at_it() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    renderer = _FakePages(working=True)
    tools = _document_tools(store, converters=[renderer, *default_converters()])
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    reply = _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치",
        new="즉각 조치",
        description="Sharpen",
        revision=revision,
    )
    assert f'read_canvas(path="{path}", pages="grid")' in reply
    assert renderer.renders >= 1


def test_without_a_renderer_the_reply_says_the_change_was_not_seen() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    reply = _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치",
        new="즉각 조치",
        description="Sharpen",
        revision=revision,
    )
    assert "No page renderer is installed" in reply


def test_a_document_edit_redraws_the_file_card() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _StreamingRuntime(config={"configurable": {"thread_id": "t1"}})
    path, revision = _open_copy(tools, runtime)
    runtime.events.clear()
    _invoke(
        tools["edit_canvas"],
        runtime,
        path=path,
        old="즉시 조치",
        new="즉각 조치",
        description="Sharpen",
        revision=revision,
    )
    kinds = [event["type"] for event in runtime.events]
    assert "canvas.patch" in kinds and "canvas.commit" in kinds


# --- the contract the descriptions state -----------------------------------------


def test_no_description_ships_an_unfilled_placeholder() -> None:
    """I6 — a constant bumped without editing the text leaves the model lied to."""
    import re

    from langchain_canvas import (
        create_asset_tool,
        create_check_table_tool,
        create_document_tools,
        create_export_tool,
    )

    store = InMemoryCanvasStore()
    built = [
        *create_canvas_tools(store),
        *create_document_tools(store),
        create_export_tool(store),
        create_asset_tool(store),
        create_check_table_tool(store),
    ]
    for tool_obj in built:
        assert not re.search(r"\{[a-z_]+\}", tool_obj.description), tool_obj.name


def test_descriptions_only_name_tools_that_exist() -> None:
    """A description that sends the model to a tool nobody built is a dead end."""
    import re

    from langchain_canvas import (
        create_asset_tool,
        create_check_table_tool,
        create_document_tools,
        create_export_tool,
    )
    from langchain_canvas.tools import create_deck_tools, create_table_tools

    store = InMemoryCanvasStore()
    built = [
        *create_canvas_tools(store),
        *create_document_tools(store),
        *create_deck_tools(store),
        *create_table_tools(store),
        create_export_tool(store),
        create_asset_tool(store),
        create_check_table_tool(store),
    ]
    names = {tool_obj.name for tool_obj in built}
    pattern = re.compile(
        r"\b(?:read|write|edit|list|export|check|open|insert|remove|replace)_[a-z_]+\b"
    )
    mentioned: set[str] = set()
    for tool_obj in built:
        mentioned |= set(pattern.findall(tool_obj.description))
    assert mentioned <= names, mentioned - names


def test_the_edit_tool_names_the_document_formats_it_handles() -> None:
    from langchain_canvas.document_ops import DOCUMENT_OP_SUFFIXES

    tools = _tools(InMemoryCanvasStore())
    for suffix in DOCUMENT_OP_SUFFIXES:
        assert suffix in tools["edit_canvas"].description


def test_document_tools_only_claim_formats_the_operations_accept() -> None:
    import re

    from langchain_canvas import create_document_tools
    from langchain_canvas.document_ops import DOCUMENT_OP_SUFFIXES

    allowed = set(DOCUMENT_OP_SUFFIXES)
    for tool_obj in create_document_tools(InMemoryCanvasStore()):
        for suffix in re.findall(r"\.[a-z]{2,5}\b", tool_obj.description):
            assert suffix in allowed, f"{tool_obj.name} names {suffix}"


def test_a_non_document_upload_still_gets_the_general_refusal() -> None:
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/photo.png", b"\x89PNG\r\n\x1a\n", "Upload", actor="human")
    tools = _tools(store)
    reply = _invoke(
        tools["write_canvas"],
        _runtime(thread_id="t1"),
        path="sources/photo.png",
        content="x",
        description="d",
    )
    assert "read-only" in reply
    assert "open_document_for_editing" not in reply


def test_a_picture_can_be_added_to_the_copy() -> None:
    """The gap this closes: a Word file could only have a picture swapped."""
    store = InMemoryCanvasStore()
    _uploaded(store)
    from documents import png_bytes

    store.write_bytes("t1", "assets/site.png", png_bytes(240, 120, (7, 7, 7)), "Add", actor="human")
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)

    added = _invoke(
        tools["insert_document_image"],
        runtime,
        path=path,
        image_path="assets/site.png",
        description="Add the site photo",
        revision=revision,
        alt_text="현장 전경",
    )
    assert added.startswith(f"Added a picture to {path}")
    assert "in" in added
    # Read back: the document now addresses a second picture.
    assert "[img1]" in _invoke(tools["read_canvas"], runtime, path=path)


def test_markdown_image_syntax_is_refused_and_says_where_to_go() -> None:
    """The agent's own mistake: markdown habits carried into a Word file."""
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    before = store.read_bytes("t1", path)

    reply = _invoke(
        tools["insert_document_paragraph"],
        runtime,
        path=path,
        anchor="사진 1. 점검 당일 현장",
        text="![이미지](sources/photo.png)",
        description="Add the picture",
        revision=revision,
    )
    assert "markdown image syntax" in reply
    assert "insert_document_image" in reply
    after = store.read_bytes("t1", path)
    assert after.data == before.data
    assert after.revision == before.revision


def test_the_paragraph_tool_says_it_writes_text_only() -> None:
    tools = _document_tools(InMemoryCanvasStore())
    description = tools["insert_document_paragraph"].description
    assert "insert_document_image" in description


def test_the_picture_tool_only_takes_images_from_the_canvas() -> None:
    store = InMemoryCanvasStore()
    _uploaded(store)
    tools = _document_tools(store)
    runtime = _runtime(thread_id="t1")
    path, revision = _open_copy(tools, runtime)
    reply = _invoke(
        tools["insert_document_image"],
        runtime,
        path=path,
        image_path="https://example.com/photo.png",
        description="Add it",
        revision=revision,
    )
    assert "not a canvas image path" in reply


# --- a file the document operations cannot open says what does work -------------


class _DeckPages(_FakePages):
    """A renderer that covers decks, the way a host's office converter does."""

    suffixes = (".pptx",)


def test_a_deck_is_told_the_reads_that_work() -> None:
    """There is no .pptx import yet; the reply has to name what is possible.

    Told only "this opens .docx", an agent invents a summary of the deck and
    calls that opening it. Naming the page read keeps the answer a next step.
    """
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/deck.pptx", b"PK-deck", "Upload", actor="human")
    tools = _document_tools(store, converters=[_DeckPages(), *default_converters()])

    reply = _invoke(
        tools["open_document_for_editing"], _runtime(thread_id="t1"), source="sources/deck.pptx"
    )
    assert reply.startswith("Error: this opens")
    assert "`pages`" in reply, reply
    assert "sources/deck.pptx" in reply


def test_a_format_no_renderer_covers_is_not_promised_pages() -> None:
    """The opposite guard: never point at a read the host cannot perform."""
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/archive.zip", b"PK", "Upload", actor="human")
    tools = _document_tools(store, converters=[_DeckPages(), *default_converters()])

    reply = _invoke(
        tools["open_document_for_editing"], _runtime(thread_id="t1"), source="sources/archive.zip"
    )
    assert reply.startswith("Error: this opens")
    assert "`pages`" not in reply, reply


def test_the_editing_operations_answer_a_deck_the_same_way() -> None:
    """Both doors report the same thing, so the agent gets one answer."""
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "deck.pptx", b"PK-deck", "Upload", actor="human")
    tools = _document_tools(store, converters=[_DeckPages(), *default_converters()])

    reply = _invoke(
        tools["insert_document_paragraph"],
        _runtime(thread_id="t1"),
        path="deck.pptx",
        anchor="p0",
        text="Hello",
        description="Add it",
        revision="r1",
    )
    assert "`pages`" in reply, reply


# --- nothing to change, nothing saved ------------------------------------------------


def test_an_edit_whose_old_equals_new_saves_nothing() -> None:
    store = InMemoryCanvasStore()
    runtime = _runtime(thread_id="t1")
    _invoke(_tools(store)["write_canvas"], runtime, path="a.md", content="hello", description="c")
    before = len(store.history("t1"))
    revision = store.read("t1", "a.md").revision
    reply = _invoke(
        _tools(store)["edit_canvas"], runtime,
        path="a.md", old="hello", new="hello", description="noop", revision=revision,
    )
    assert reply.startswith("Error:") and "nothing to change" in reply
    assert len(store.history("t1")) == before


# --- the eye opens on its own --------------------------------------------------------


class _PptxEye:
    """A page renderer for .pptx that answers with one image block per page."""

    suffixes = (".pptx",)

    def __init__(self) -> None:
        self.calls: list[tuple[str, list[int] | str]] = []

    def convert(self, data: bytes, *, path: str) -> Any:
        from langchain_canvas.converters import ConvertedSource

        return ConvertedSource(blocks=[{"type": "text", "text": "deck"}], metadata={})

    def render_pages(self, data: bytes, *, path: str, pages: list[int]) -> Any:
        from langchain_canvas.converters import ConvertedSource

        self.calls.append((path, pages))
        return ConvertedSource(
            blocks=[
                {"type": "image", "source_type": "base64", "data": "AA==", "mime_type": "image/png"}
                for _ in pages
            ],
            metadata={"pages": ",".join(map(str, pages))},
        )

    def render_grid(self, data: bytes, *, path: str) -> Any:
        from langchain_canvas.converters import ConvertedSource

        self.calls.append((path, "grid"))
        return ConvertedSource(
            blocks=[
                {"type": "image", "source_type": "base64", "data": "AA==", "mime_type": "image/png"}
            ],
            metadata={},
        )


def _deck_with_overflow() -> str:
    from langchain_canvas import encode_slides

    return encode_slides("Deck", {"slides": [
        {"elements": [
            {"id": "ok", "type": "text", "x": 5, "y": 5, "w": 50, "h": 10, "text": "fine"}
        ]},
        {"elements": [
            {"id": "e0", "type": "text", "x": 60, "y": 10, "w": 58, "h": 10, "text": "off"}
        ]},
    ]})


def test_a_save_with_a_finding_arrives_with_that_slide_rendered() -> None:
    """Measured: told to look 14 times, looked once. So the slide the check
    named comes with the finding when a renderer is mounted."""
    pytest.importorskip("pptx")
    eye = _PptxEye()
    store = InMemoryCanvasStore()
    tools = {t.name: t for t in create_canvas_tools(store, converters=[eye])}
    reply = tools["write_canvas"].func(
        path="deck.slides.json",
        content=_deck_with_overflow(),
        description="c",
        runtime=_runtime(thread_id="t1"),
    )
    assert isinstance(reply, list), reply
    assert reply[0]["text"].startswith("Wrote deck.slides.json")
    assert "Deck check" in reply[0]["text"]
    assert [b["type"] for b in reply[1:]] == ["image"]
    assert eye.calls == [("deck.pptx", [2])]  # only the slide the check named


def test_a_clean_save_stays_text_even_with_a_renderer() -> None:
    from langchain_canvas import encode_slides

    store = InMemoryCanvasStore()
    tools = {t.name: t for t in create_canvas_tools(store, converters=[_PptxEye()])}
    reply = tools["write_canvas"].func(
        path="deck.slides.json",
        content=encode_slides("Deck", {"slides": [{"elements": [
            {"id": "t", "type": "text", "x": 10, "y": 10, "w": 50, "h": 10, "text": "Hi"}
        ]}]}),
        description="c",
        runtime=_runtime(thread_id="t1"),
    )
    assert isinstance(reply, str) and reply.startswith("Wrote deck.slides.json")


def test_an_export_comes_back_with_its_pages_when_a_renderer_is_mounted() -> None:
    from langchain_canvas import create_export_tool, encode_slides

    pytest.importorskip("pptx")
    eye = _PptxEye()
    store = InMemoryCanvasStore()
    runtime = _runtime(thread_id="t1")
    _invoke(
        _tools(store)["write_canvas"], runtime, path="deck.slides.json",
        content=encode_slides("Deck", {"slides": [{"title": "Hi", "layout": "title"}]}),
        description="c",
    )
    exporter = create_export_tool(store, converters=[eye])
    reply = exporter.func(path="deck.slides.json", target="pptx", runtime=runtime)
    assert isinstance(reply, list), reply
    assert reply[0]["text"].startswith("Exported deck.slides.json to exports/")
    assert "check them" in reply[0]["text"]
    assert reply[1]["type"] == "image"
    assert eye.calls and eye.calls[-1][1] == "grid"
    plain = create_export_tool(store).func(path="deck.slides.json", target="pptx", runtime=runtime)
    assert isinstance(plain, str)  # no renderer mounted: text as before


def test_the_eye_shows_the_slide_this_save_changed_not_the_templates_noise() -> None:
    """Fourteen saves once arrived with slides 1 and 2 — the template's own
    findings — while the slide being written was never shown."""
    pytest.importorskip("pptx")
    from langchain_canvas import encode_slides

    off = {"id": "e0", "type": "text", "x": 60, "y": 10, "w": 58, "h": 10, "text": "off page"}
    fine = {"id": "e0", "type": "text", "x": 5, "y": 5, "w": 50, "h": 10, "text": "fine"}
    eye = _PptxEye()
    store = InMemoryCanvasStore()
    tools = {t.name: t for t in create_canvas_tools(store, converters=[eye])}
    first = tools["write_canvas"].func(
        path="deck.slides.json",
        content=encode_slides("Deck", {"slides": [{"elements": [off]}, {"elements": [fine]}]}),
        description="c",
        runtime=_runtime(thread_id="t1"),
    )
    import re as regex

    revision = regex.search(r"revision (v\d+)", first[0]["text"]).group(1)
    # This save touches only slide 2, making it overflow too. Slide 1's
    # finding is older noise; the eye follows the change.
    tools["write_canvas"].func(
        path="deck.slides.json",
        content=encode_slides("Deck", {"slides": [{"elements": [off]}, {"elements": [
            {**off, "id": "e0"}]}]}),
        description="c",
        revision=revision,
        runtime=_runtime(thread_id="t1"),
    )
    assert eye.calls[-1] == ("deck.pptx", [2])


def test_an_eye_image_is_resized_to_glance_size() -> None:
    pytest.importorskip("pptx")
    pytest.importorskip("PIL")
    import base64 as b64
    import io as iolib

    from PIL import Image

    big = iolib.BytesIO()
    Image.new("RGB", (1920, 1080), "white").save(big, format="PNG")
    payload = b64.b64encode(big.getvalue()).decode()

    class _BigEye(_PptxEye):
        def render_pages(self, data: bytes, *, path: str, pages: list[int]):
            from langchain_canvas.converters import ConvertedSource

            return ConvertedSource(
                blocks=[{"type": "image", "source_type": "base64", "data": payload,
                         "mime_type": "image/png"} for _ in pages],
                metadata={},
            )

    store = InMemoryCanvasStore()
    tools = {t.name: t for t in create_canvas_tools(store, converters=[_BigEye()])}
    reply = tools["write_canvas"].func(
        path="deck.slides.json",
        content=_deck_with_overflow(),
        description="c",
        runtime=_runtime(thread_id="t1"),
    )
    image = reply[1]
    assert image["mime_type"] == "image/jpeg"
    got = Image.open(iolib.BytesIO(b64.b64decode(image["data"])))
    assert got.width == 1024 and got.height == 576


# --- set_slide_texts: one slide, one save ----------------------------------------------


def _two_slide_deck() -> str:
    from langchain_canvas import encode_slides

    return encode_slides("Deck", {"slides": [
        {"elements": [
            {"id": "e0", "type": "text", "x": 5, "y": 5, "w": 80, "h": 10,
             "fontSize": 40, "text": "표지 제목"},
            {"id": "e1", "type": "shape", "shape": "rect", "x": 5, "y": 20, "w": 20, "h": 10},
        ]},
        {"elements": [
            {"id": "e0", "type": "text", "x": 5, "y": 5, "w": 80, "h": 10,
             "fontSize": 32, "text": "본문 제목"},
            {"id": "e2", "type": "text", "x": 5, "y": 20, "w": 60, "h": 8,
             "fontSize": 24, "text": "자리표시자"},
        ]},
    ]})


def test_set_slide_texts_changes_several_elements_in_one_save() -> None:
    pytest.importorskip("pptx")
    store = InMemoryCanvasStore()
    tools = _tools(store)
    write = tools["write_canvas"].func(
        path="d.slides.json", content=_two_slide_deck(), description="c",
        runtime=_runtime(thread_id="t1"),
    )
    import re as regex

    text = write if isinstance(write, str) else write[0]["text"]
    revision = regex.search(r"revision (v\d+)", text).group(1)
    reply = tools["set_slide_texts"].func(
        path="d.slides.json", slide=2,
        texts={"e0": "협업 구조와 역할", "e2": "본문 한 줄"},
        description="slide 2 body", revision=revision, runtime=_runtime(thread_id="t1"),
    )
    reply_text = reply if isinstance(reply, str) else reply[0]["text"]
    assert reply_text.startswith("Set 2 text(s) on slide 2")
    saved = json.loads(store.read("t1", "d.slides.json").content)
    slide2 = {e["id"]: e for e in saved["data"]["slides"][1]["elements"]}
    assert slide2["e0"]["text"] == "협업 구조와 역할"
    assert slide2["e2"]["text"] == "본문 한 줄"
    assert slide2["e0"]["fontSize"] == 32  # the look stays
    # slide 1 untouched
    assert json.loads(_two_slide_deck())["data"]["slides"][0] == saved["data"]["slides"][0]


def test_set_slide_texts_refuses_with_the_slides_own_ids() -> None:
    pytest.importorskip("pptx")
    store = InMemoryCanvasStore()
    tools = _tools(store)
    tools["write_canvas"].func(path="d.slides.json", content=_two_slide_deck(),
                               description="c", runtime=_runtime(thread_id="t1"))
    revision = store.read("t1", "d.slides.json").revision
    run = tools["set_slide_texts"].func
    wrong_id = run(path="d.slides.json", slide=1, texts={"nope": "x"},
                   description="d", revision=revision, runtime=_runtime(thread_id="t1"))
    assert "has no element 'nope'" in wrong_id and "e0 (text)" in wrong_id
    not_text = run(path="d.slides.json", slide=1, texts={"e1": "x"},
                   description="d", revision=revision, runtime=_runtime(thread_id="t1"))
    assert "is a shape, not text" in not_text
    off_range = run(path="d.slides.json", slide=9, texts={"e0": "x"},
                    description="d", revision=revision, runtime=_runtime(thread_id="t1"))
    assert "out of range" in off_range and "2 slide(s)" in off_range
    stale = run(path="d.slides.json", slide=1, texts={"e0": "새 제목"},
                description="d", revision="v0", runtime=_runtime(thread_id="t1"))
    assert stale.startswith("Error:")


def test_set_slide_texts_overflow_arrives_with_the_check_and_the_eye() -> None:
    pytest.importorskip("pptx")
    eye = _PptxEye()
    store = InMemoryCanvasStore()
    tools = {t.name: t for t in create_canvas_tools(store, converters=[eye])}
    tools["write_canvas"].func(path="d.slides.json", content=_two_slide_deck(),
                               description="c", runtime=_runtime(thread_id="t1"))
    revision = store.read("t1", "d.slides.json").revision
    reply = tools["set_slide_texts"].func(
        path="d.slides.json", slide=2, texts={"e2": "자리표시자보다 훨씬 긴 본문 " * 8},
        description="overflow", revision=revision, runtime=_runtime(thread_id="t1"),
    )
    assert isinstance(reply, list)
    assert "Deck check" in reply[0]["text"]
    assert eye.calls[-1] == ("d.pptx", [2])


# --- the export gate and the review ----------------------------------------------------


def _crowded_template() -> bytes:
    """One slide: a title placeholder and a fixed box whose own text overflows."""
    import io as iolib

    from pptx import Presentation
    from pptx.util import Inches, Pt

    deck = Presentation()
    deck.slide_width, deck.slide_height = Inches(13.333), Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])

    def box(text: str, top: float, height: float, autofit: object) -> None:
        shape = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(6), Inches(height))
        shape.text_frame.text = text
        shape.text_frame.auto_size = autofit
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(20)

    box("핵심 주제를 입력해 주세요", 0.5, 0.8, None)
    box("상세 내용을 작성해 주세요 " * 6, 2.0, 0.5, None)  # overflows already
    out = iolib.BytesIO()
    deck.save(out)
    return out.getvalue()


def _deck_workbench(converters: list | None = None) -> tuple[Any, dict[str, Any]]:
    """Store with the crowded upload opened for editing, plus every tool."""
    from langchain_canvas.tools import create_deck_tools, create_export_tool

    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/skin.pptx", _crowded_template(), "Upload", actor="human")
    tools = {t.name: t for t in create_canvas_tools(store, converters=converters)}
    for t in create_deck_tools(store, converters=converters):
        tools[t.name] = t
    tools["export_canvas"] = create_export_tool(store, converters=converters)
    tools["open_deck_for_editing"].func(
        source="sources/skin.pptx", runtime=_runtime(thread_id="t1")
    )
    return store, tools


def test_an_inherited_overflow_does_not_block_export_but_a_new_one_does() -> None:
    pytest.importorskip("pptx")
    store, tools = _deck_workbench()
    fresh = tools["export_canvas"].func(path="skin.slides.json", target="pptx",
                                        runtime=_runtime(thread_id="t1"))
    text = fresh if isinstance(fresh, str) else fresh[0]["text"]
    assert not text.startswith("Error"), text  # the template's own overflow is folded
    # Write far more text into the crowded box: now it is the agent's doing.
    revision = store.read("t1", "skin.slides.json").revision
    deck = json.loads(store.read("t1", "skin.slides.json").content)
    crowded = deck["data"]["slides"][0]["elements"][1]["id"]
    tools["set_slide_texts"].func(
        path="skin.slides.json", slide=1, texts={crowded: "훨씬 더 길게 다시 쓴 본문 " * 30},
        description="longer", revision=revision, runtime=_runtime(thread_id="t1"),
    )
    blocked = tools["export_canvas"].func(path="skin.slides.json", target="pptx",
                                          runtime=_runtime(thread_id="t1"))
    assert isinstance(blocked, str) and blocked.startswith("Error: not exported")
    assert "run past the box" in blocked and "review_deck" in blocked
    assert "accept_findings=True" in blocked
    # The user said to ship it anyway.
    shipped = tools["export_canvas"].func(path="skin.slides.json", target="pptx",
                                          accept_findings=True, runtime=_runtime(thread_id="t1"))
    text = shipped if isinstance(shipped, str) else shipped[0]["text"]
    assert "exports/" in text


def test_review_deck_names_what_still_reads_as_the_original() -> None:
    pytest.importorskip("pptx")
    store, tools = _deck_workbench()
    revision = store.read("t1", "skin.slides.json").revision
    deck = json.loads(store.read("t1", "skin.slides.json").content)
    title = deck["data"]["slides"][0]["elements"][0]["id"]
    untouched = deck["data"]["slides"][0]["elements"][1]["id"]
    tools["set_slide_texts"].func(
        path="skin.slides.json", slide=1, texts={title: "협업 제안"},
        description="title", revision=revision, runtime=_runtime(thread_id="t1"),
    )
    report = tools["review_deck"].func(path="skin.slides.json", runtime=_runtime(thread_id="t1"))
    text = report if isinstance(report, str) else report[0]["text"]
    assert text.startswith("skin.slides.json — 1 slide(s).")
    assert "changed 1 text(s)" in text
    assert f"1 still read as the original ({untouched})" in text


def test_review_deck_renders_now_and_original_side_by_side() -> None:
    pytest.importorskip("pptx")
    eye = _PptxEye()
    store, tools = _deck_workbench(converters=[eye])
    grid = tools["review_deck"].func(path="skin.slides.json", runtime=_runtime(thread_id="t1"))
    assert isinstance(grid, list)
    assert "Images: every page of the deck now, then the original." in grid[0]["text"]
    assert [c for c in eye.calls if c[1] == "grid"][-2:] == [("skin.pptx", "grid")] * 2
    one = tools["review_deck"].func(path="skin.slides.json", slide=1,
                                    runtime=_runtime(thread_id="t1"))
    assert isinstance(one, list)
    assert eye.calls[-2:] == [("skin.pptx", [1]), ("skin.pptx", [1])]


# --- the master rides behind the copy --------------------------------------------------


def _deck_with_master_logo() -> bytes:
    import io as iolib

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.slides.add_slide(deck.slide_layouts[6])
    donor = deck.slides[0].shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(7), Inches(1), Inches(0.3)
    )
    donor.fill.solid()
    element = donor._element
    element.getparent().remove(element)
    deck.slide_masters[0].shapes._spTree.append(element)
    out = iolib.BytesIO()
    deck.save(out)
    return out.getvalue()


def test_the_copy_gets_the_master_as_a_display_backdrop() -> None:
    pytest.importorskip("pptx")
    from langchain_canvas.tools import create_deck_tools

    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/logo.pptx", _deck_with_master_logo(), "Upload", actor="human")
    eye = _PptxEye()
    tools = {t.name: t for t in create_deck_tools(store, converters=[eye])}
    reply = tools["open_deck_for_editing"].func(
        source="sources/logo.pptx", runtime=_runtime(thread_id="t1")
    )
    text = reply if isinstance(reply, str) else reply[0]["text"]
    assert "drawn behind 2 slide(s)" in text
    deck = json.loads(store.read("t1", "logo.slides.json").content)["data"]
    backdrops = [slide.get("masterImage") for slide in deck["slides"]]
    assert all(b and b.startswith("assets/logo/master-") for b in backdrops)
    # identical pages share one stored file
    assert len(set(backdrops)) == 1
    stored = [info.path for info in store.list_files("t1") if "master-" in info.path]
    assert len(stored) == 1


def test_without_a_renderer_the_reply_says_the_master_is_safe() -> None:
    pytest.importorskip("pptx")
    from langchain_canvas.tools import create_deck_tools

    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/logo.pptx", _deck_with_master_logo(), "Upload", actor="human")
    tools = {t.name: t for t in create_deck_tools(store, converters=None)}
    reply = tools["open_deck_for_editing"].func(
        source="sources/logo.pptx", runtime=_runtime(thread_id="t1")
    )
    text = reply if isinstance(reply, str) else reply[0]["text"]
    assert "returns on export" in text
    deck = json.loads(store.read("t1", "logo.slides.json").content)["data"]
    assert all("masterImage" not in slide for slide in deck["slides"])


def test_read_canvas_fields_projects_a_deck_and_refuses_other_files() -> None:
    from langchain_canvas import encode_slides

    store = InMemoryCanvasStore()
    runtime = _runtime(thread_id="t1")
    tools = _tools(store)
    _invoke(tools["write_canvas"], runtime, path="d.slides.json", description="d",
            content=encode_slides("d", {"slides": [{"elements": [
                {"id": "e0", "type": "text", "x": 1, "y": 1, "w": 50, "h": 10,
                 "fontSize": 20, "color": "#111111", "text": "hello"}]}]}))
    out = _invoke(tools["read_canvas"], runtime, path="d.slides.json", fields="color,fontSize")
    assert "revision:" in out and "[s1] e0 color=#111111 fontSize=20" in out
    assert '"hello"' not in out  # only the asked keys
    _invoke(tools["write_canvas"], runtime, path="a.md", content="x", description="c")
    refused = _invoke(tools["read_canvas"], runtime, path="a.md", fields="color")
    assert refused.startswith("Error: `fields` applies to .slides.json")


# --- export with a workbook-engine recalc ----------------------------------------------


def _table_store() -> InMemoryCanvasStore:
    from langchain_canvas.replay import encode_table

    store = InMemoryCanvasStore()
    data = {"columns": [{"key": "a"}], "rows": [{"a": 1}],
            "sheet": [{"name": "S", "celldata": [
                {"r": 0, "c": 0, "v": {"v": 2}},
                {"r": 1, "c": 0, "v": {"f": "=A1*3"}},
            ]}]}
    store.write("t1", "book.table.json", encode_table("Book", data), "seed", actor="agent")
    return store


def test_an_xlsx_export_leaves_through_the_recalc_hook() -> None:
    pytest.importorskip("openpyxl")
    from langchain_canvas.tools import create_export_tool

    store = _table_store()
    tool = create_export_tool(store, xlsx_recalc=lambda data: b"FRESH-XLSX-BYTES")
    reply = tool.func(path="book.table.json", target="xlsx", runtime=_runtime(thread_id="t1"))
    text = reply if isinstance(reply, str) else reply[0]["text"]
    assert "exports/" in text
    assert store.read_bytes("t1", "exports/book.xlsx").data == b"FRESH-XLSX-BYTES"


def test_a_failing_recalc_hook_still_lands_the_export_with_a_note() -> None:
    pytest.importorskip("openpyxl")
    from langchain_canvas.tools import create_export_tool

    def broken(data: bytes) -> bytes:
        raise RuntimeError("endpoint down")

    store = _table_store()
    tool = create_export_tool(store, xlsx_recalc=broken)
    reply = tool.func(path="book.table.json", target="xlsx", runtime=_runtime(thread_id="t1"))
    text = reply if isinstance(reply, str) else reply[0]["text"]
    assert "Formula caches were not recalculated" in text and "endpoint down" in text
    assert store.read_bytes("t1", "exports/book.xlsx").data[:2] == b"PK"  # the real workbook
