"""Source converters — the pluggable bytes-to-blocks contract."""

from __future__ import annotations

import io

from langchain_canvas.converters import (
    ConvertedSource,
    TextSourceConverter,
    converter_for,
    default_converters,
)


def test_text_converter_decodes_utf8_bom() -> None:
    data = "﻿분석명,값\n표본,64\n".encode()
    got = TextSourceConverter().convert(data, path="sources/report.csv")
    assert got.blocks[0]["text"].startswith("분석명")
    assert got.metadata["encoding"] == "utf-8-sig"


def test_text_converter_falls_back_to_cp949() -> None:
    data = "한글 문서".encode("cp949")
    got = TextSourceConverter().convert(data, path="sources/note.txt")
    assert got.blocks[0]["text"] == "한글 문서"
    assert got.metadata["encoding"] == "cp949"


def test_text_converter_is_lossy_but_honest_on_unknown_bytes() -> None:
    got = TextSourceConverter().convert(b"\xff\xfe\x00abc", path="sources/x.txt")
    assert got.metadata["encoding"] == "unknown (lossy utf-8)"
    assert "abc" in got.blocks[0]["text"]


def test_converter_for_matches_suffix_case_insensitively() -> None:
    converters = default_converters()
    assert converter_for("sources/NOTES.MD", converters) is not None
    assert converter_for("sources/archive.zip", converters) is None


def test_custom_converter_plugs_in() -> None:
    class UpperConverter:
        suffixes = (".shout",)

        def convert(self, data: bytes, *, path: str) -> ConvertedSource:
            return ConvertedSource(blocks=[{"type": "text", "text": data.decode().upper()}])

    got = converter_for("sources/a.shout", [UpperConverter()])
    assert got is not None
    assert got.convert(b"hi", path="sources/a.shout").blocks[0]["text"] == "HI"


def _tiny_xlsx() -> bytes:
    import io

    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Scores"
    ws.append(["model", "score"])
    ws.append(["A", 91])
    ws.append(["B", 87])
    wb.create_sheet("Notes").append(["비고", "테스트"])
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def test_xlsx_converter_renders_each_sheet() -> None:
    from langchain_canvas.converters import XlsxSourceConverter

    got = XlsxSourceConverter().convert(_tiny_xlsx(), path="sources/scores.xlsx")
    text = got.blocks[0]["text"]
    assert "### sheet: Scores" in text
    assert "model,score" in text
    assert "A,91" in text
    assert "### sheet: Notes" in text
    assert "비고,테스트" in text
    assert got.metadata["sheets"] == "Scores, Notes"


def test_xlsx_converter_is_in_the_default_set() -> None:
    assert converter_for("sources/report.xlsx", default_converters()) is not None


# --- images ----------------------------------------------------------------------


def test_image_converter_emits_a_vision_block() -> None:
    from langchain_canvas.converters import ImageSourceConverter

    got = ImageSourceConverter().convert(b"\x89PNG fake", path="sources/logo.png")
    kinds = [b["type"] for b in got.blocks]
    assert kinds == ["text", "image"]
    image = got.blocks[1]
    assert image["mime_type"] == "image/png"
    import base64

    assert base64.b64decode(image["data"]) == b"\x89PNG fake"
    assert got.metadata["inlined"] is True


def test_oversized_image_degrades_honestly() -> None:
    from langchain_canvas.converters import ImageSourceConverter

    converter = ImageSourceConverter()
    converter.max_bytes = 10
    got = converter.convert(b"x" * 11, path="sources/big.jpg")
    assert [b["type"] for b in got.blocks] == ["text"]
    assert "too large to inline" in got.blocks[0]["text"]
    assert got.metadata["inlined"] is False


# --- pdf -------------------------------------------------------------------------


def _tiny_pdf(text: str = "Hello canvas") -> bytes:
    # A minimal one-page PDF with a real text layer, built by hand so the test
    # exercises extraction (pypdf's writer only makes blank pages).
    stream = f"BT /F1 24 Tf 72 700 Td ({text}) Tj ET".encode()
    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n",
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n",
        b"4 0 obj << /Length "
        + str(len(stream)).encode()
        + b" >> stream\n"
        + stream
        + b"\nendstream endobj\n",
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n",
    ]
    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = []
    for obj in objects:
        offsets.append(out.tell())
        out.write(obj)
    xref_at = out.tell()
    out.write(b"xref\n0 6\n0000000000 65535 f \n")
    for off in offsets:
        out.write(f"{off:010d} 00000 n \n".encode())
    out.write(b"trailer << /Size 6 /Root 1 0 R >>\nstartxref\n")
    out.write(str(xref_at).encode())
    out.write(b"\n%%EOF")
    return out.getvalue()


def test_pdf_converter_extracts_page_text() -> None:
    from langchain_canvas.converters import PdfSourceConverter

    got = PdfSourceConverter().convert(_tiny_pdf(), path="sources/doc.pdf")
    text = got.blocks[0]["text"]
    assert "### page 1" in text
    assert "Hello canvas" in text
    assert got.metadata["pages"] == 1


def test_image_and_pdf_are_in_the_default_set() -> None:
    assert converter_for("sources/a.png", default_converters()) is not None
    assert converter_for("sources/a.pdf", default_converters()) is not None


# --- office (docx / pptx) --------------------------------------------------------


def test_docx_converter_keeps_paragraphs_and_tables_in_order() -> None:
    from docx import Document

    from langchain_canvas.converters import DocxSourceConverter

    doc = Document()
    doc.add_paragraph("서문 문단")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "항목"
    table.rows[0].cells[1].text = "값"
    table.rows[1].cells[0].text = "표본"
    table.rows[1].cells[1].text = "64"
    doc.add_paragraph("맺음 문단")
    out = io.BytesIO()
    doc.save(out)

    got = DocxSourceConverter().convert(out.getvalue(), path="sources/doc.docx")
    text = got.blocks[0]["text"]
    assert text.index("서문 문단") < text.index("항목,값") < text.index("맺음 문단")
    assert "표본,64" in text
    assert got.metadata == {"paragraphs": 2, "tables": 1}


def test_pptx_converter_renders_each_slide() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    from langchain_canvas.converters import PptxSourceConverter

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = "발표 제목"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "본문 텍스트"
    deck.slides.add_slide(deck.slide_layouts[6])  # blank
    out = io.BytesIO()
    deck.save(out)

    got = PptxSourceConverter().convert(out.getvalue(), path="sources/deck.pptx")
    text = got.blocks[0]["text"]
    assert "### slide 1" in text and "발표 제목" in text and "본문 텍스트" in text
    assert "### slide 2\n(no text on this slide)" in text
    assert got.metadata == {"slides": 2}


def test_office_converters_are_in_the_default_set() -> None:
    assert converter_for("sources/a.docx", default_converters()) is not None
    assert converter_for("sources/a.pptx", default_converters()) is not None


def _bomb_zip(uncompressed_mb: int = 201) -> bytes:
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        z.writestr("ppt/presentation.xml", b"\x00" * (uncompressed_mb * 1024 * 1024))
    return buf.getvalue()


def test_archive_guard_refuses_decompression_bombs() -> None:
    import pytest

    from langchain_canvas.converters import (
        UnsafeArchiveError,
        ensure_archive_within_limits,
    )

    # ~200 KB compressed, >200 MB unpacked — trips size (and ratio) limits.
    with pytest.raises(UnsafeArchiveError, match="unpacks to"):
        ensure_archive_within_limits(_bomb_zip(), path="bomb.pptx")
    # Not a ZIP at all is a formatting problem, not an attack — the guard
    # leaves it to each caller's own failure path.
    ensure_archive_within_limits(b"not a zip", path="junk.pptx")  # no raise


def test_archive_guard_passes_a_real_office_file() -> None:
    from langchain_canvas.converters import ensure_archive_within_limits

    try:
        from pptx import Presentation
    except ImportError:
        import pytest

        pytest.skip("python-pptx not installed")
    buf = io.BytesIO()
    Presentation().save(buf)
    ensure_archive_within_limits(buf.getvalue(), path="ok.pptx")  # no raise


def test_pptx_converter_refuses_a_bomb_before_parsing() -> None:
    import pytest

    from langchain_canvas.converters import PptxSourceConverter, UnsafeArchiveError

    pytest.importorskip("pptx")
    with pytest.raises(UnsafeArchiveError):
        PptxSourceConverter().convert(_bomb_zip(), path="sources/bomb.pptx")
