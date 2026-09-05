"""Exporter contract: table -> xlsx, html -> docx, slides -> pptx, the tool.

Round-trips are asserted with the real readers (openpyxl / python-docx /
python-pptx), so "exported" means a mainstream library opens the file and
finds the content.
"""

from __future__ import annotations

import base64
import io
import json
import re
import zipfile
from dataclasses import dataclass, field
from typing import Any

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches

from langchain_canvas import InMemoryCanvasStore, create_export_tool
from langchain_canvas.exporters import (
    PPTX_MIME,
    HtmlDocxExporter,
    SlidesPptxExporter,
    TableXlsxExporter,
    default_exporters,
    exporter_for,
)

# A real 1x1 red PNG — small enough to inline, real enough for pptx to embed.
PNG_1PX = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4z8DwHwAFAAH/q842iQAAAABJRU5ErkJggg=="
)

# --- routing ---------------------------------------------------------------------


def test_exporter_routing_matches_suffix_and_target():
    exporters = default_exporters()
    assert isinstance(exporter_for("sales.table.json", "xlsx", exporters), TableXlsxExporter)
    assert isinstance(exporter_for("report/01-a.html", "docx", exporters), HtmlDocxExporter)
    assert isinstance(exporter_for("deck.slides.json", "pptx", exporters), SlidesPptxExporter)
    assert exporter_for("sales.table.json", "docx", exporters) is None
    assert exporter_for("deck.slides.json", "docx", exporters) is None
    assert exporter_for("photo.png", "docx", exporters) is None


# --- table -> xlsx ---------------------------------------------------------------


def test_table_columns_rows_export_round_trip():
    content = json.dumps(
        {
            "type": "table",
            "title": "Sales",
            "data": {
                "columns": [
                    {"key": "region", "label": "Region"},
                    {"key": "total", "label": "Total"},
                ],
                "rows": [
                    {"region": "East", "total": 10},
                    {"region": "West", "total": 20},
                ],
            },
        }
    )
    result = TableXlsxExporter().export(content, path="sales.table.json")
    assert result.filename == "sales.xlsx"

    sheet = load_workbook(io.BytesIO(result.data)).worksheets[0]
    assert [cell.value for cell in sheet[1]] == ["Region", "Total"]
    assert sheet[1][0].font.bold
    assert sheet.cell(row=2, column=1).value == "East"
    assert sheet.cell(row=3, column=2).value == 20


def test_table_fortune_sheet_export_values_and_merges():
    content = json.dumps(
        {
            "type": "table",
            "data": {
                "sheet": [
                    {
                        "name": "Q1",
                        "celldata": [
                            {"r": 0, "c": 0, "v": {"v": "Head", "bl": 1}},
                            {"r": 1, "c": 1, "v": 42},
                        ],
                        "config": {"merge": {"0_0": {"r": 0, "c": 0, "rs": 1, "cs": 2}}},
                    }
                ]
            },
        }
    )
    result = TableXlsxExporter().export(content, path="q1.table.json")
    sheet = load_workbook(io.BytesIO(result.data))["Q1"]
    assert sheet.cell(row=1, column=1).value == "Head"
    assert sheet.cell(row=2, column=2).value == 42
    assert [str(r) for r in sheet.merged_cells.ranges] == ["A1:B1"]


def test_table_formula_rows_stay_formulas():
    # An "="-prefixed row value must land as a live formula, not a frozen string.
    content = json.dumps(
        {
            "type": "table",
            "data": {
                "columns": [{"key": "a", "label": "A"}],
                "rows": [{"a": 10}, {"a": 20}, {"a": "=SUM(A2:A3)"}],
            },
        }
    )
    result = TableXlsxExporter().export(content, path="sums.table.json")
    sheet = load_workbook(io.BytesIO(result.data)).worksheets[0]
    cell = sheet.cell(row=4, column=1)
    assert cell.value == "=SUM(A2:A3)"
    assert cell.data_type == "f"  # stored as a formula, recalculated on open


def test_table_fortune_sheet_typed_formula_stays_a_formula():
    # A grid-typed formula carries `f` next to the cached `v` — export the
    # formula (openpyxl stores no cached value; apps recalculate on open).
    content = json.dumps(
        {
            "type": "table",
            "data": {
                "sheet": [
                    {
                        "name": "S",
                        "celldata": [
                            {"r": 0, "c": 0, "v": {"v": 1}},
                            {"r": 1, "c": 0, "v": {"v": 2}},
                            {"r": 2, "c": 0, "v": {"v": 3, "f": "=SUM(A1:A2)"}},
                        ],
                    }
                ]
            },
        }
    )
    result = TableXlsxExporter().export(content, path="s.table.json")
    sheet = load_workbook(io.BytesIO(result.data))["S"]
    cell = sheet.cell(row=3, column=1)
    assert cell.value == "=SUM(A1:A2)"
    assert cell.data_type == "f"


# --- html -> docx ----------------------------------------------------------------

_PNG_1PX = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
    "AAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)

_HTML = f"""<html><head><title>t</title><style>body {{ font: 16px serif; }}</style></head>
<body>
<div class="kicker">REPORT · SECTION 02</div>
<h1>핵심 API 기능</h1>
<p>The store keeps <strong>every</strong> file and its <em>history</em>.</p>
<ul><li>read</li><li>write</li></ul>
<table><tr><th>Tool</th><th>Role</th></tr><tr><td>read_canvas</td><td>read</td></tr></table>
<img src="data:image/png;base64,{_PNG_1PX}"/>
<hr/>
<p>After the break.</p>
</body></html>"""


def test_html_docx_export_structure():
    result = HtmlDocxExporter().export(_HTML, path="report/02-features.html")
    assert result.filename == "02-features.docx"

    document = Document(io.BytesIO(result.data))
    texts = [p.text for p in document.paragraphs if p.text.strip()]
    assert "REPORT · SECTION 02" in texts  # kicker div survives as a paragraph
    assert "After the break." in texts
    assert "font: 16px" not in " ".join(texts)  # style content dropped

    heading = next(p for p in document.paragraphs if p.text == "핵심 API 기능")
    assert heading.style.name == "Heading 1"

    prose = next(p for p in document.paragraphs if "every" in p.text)
    assert any(run.bold for run in prose.runs)
    assert any(run.italic for run in prose.runs)

    bullets = [p.text for p in document.paragraphs if p.style.name == "List Bullet"]
    assert bullets == ["read", "write"]

    table = document.tables[0]
    assert table.cell(0, 0).text == "Tool"
    assert table.cell(0, 0).paragraphs[0].runs[0].bold  # th row bolded
    assert table.cell(1, 1).text == "read"

    assert len(document.inline_shapes) == 1  # the data: URI image
    assert "<w:pBdr>" in document.element.xml  # hr became a rule, not a new page


def test_html_docx_title_names_the_file():
    result = HtmlDocxExporter().export("<p>hi</p>", path="report/", title="My Report")
    assert result.filename == "My-Report.docx"


# --- slides -> pptx --------------------------------------------------------------


def _deck(slides: list[dict[str, Any]]) -> str:
    return json.dumps({"type": "slides", "title": "Deck", "data": {"slides": slides}})


def _png_uri() -> str:
    return f"data:image/png;base64,{base64.b64encode(PNG_1PX).decode()}"


def test_slides_pptx_elements_land_as_real_shapes():
    content = _deck(
        [
            {
                "background": "#112233",
                "notes": "speaker notes here",
                "elements": [
                    {"id": "t", "type": "text", "x": 10, "y": 10, "w": 60, "h": 12,
                     "text": "Hello deck", "fontSize": 40, "bold": True,
                     "color": "#ff0000", "align": "center"},
                    {"id": "i", "type": "image", "x": 20, "y": 30, "w": 40, "h": 40,
                     "src": _png_uri()},
                    {"id": "r", "type": "shape", "shape": "rect", "x": 5, "y": 80,
                     "w": 30, "h": 10, "fill": "#00ff00"},
                    {"id": "l", "type": "shape", "shape": "line", "x": 40, "y": 85,
                     "w": 50, "h": 0, "fill": "#0000ff"},
                ],
            }
        ]
    )
    result = SlidesPptxExporter().export(content, path="deck.slides.json")
    assert result.filename == "Deck.pptx"

    deck = Presentation(io.BytesIO(result.data))
    assert deck.slide_width == Inches(10)
    assert deck.slide_height == Inches(5.625)
    (slide,) = deck.slides

    # The slide is shapes, not one baked bitmap: text keeps its runs and
    # font, the image is its own picture with the original bytes intact.
    texts = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text]
    ((run,),) = [p.runs for p in texts[0].text_frame.paragraphs]
    assert run.text == "Hello deck"
    assert run.font.size.pt == 30  # 40 px on the 1280px canvas -> 30 pt
    assert run.font.bold is True
    assert str(run.font.color.rgb) == "FF0000"
    assert texts[0].text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

    (picture,) = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert picture.image.blob == PNG_1PX

    fills = {
        str(s.fill.fore_color.rgb)
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    }
    assert "00FF00" in fills
    assert any(s.shape_type == MSO_SHAPE_TYPE.LINE for s in slide.shapes)

    assert str(slide.background.fill.fore_color.rgb) == "112233"
    assert slide.notes_slide.notes_text_frame.text == "speaker notes here"


def test_slides_pptx_derives_structured_slides_like_the_canvas():
    content = _deck(
        [
            {"layout": "title", "title": "Big Title", "subtitle": "The subtitle"},
            {"title": "Agenda", "bullets": ["One", "Two"]},
        ]
    )
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    first, second = deck.slides

    first_texts = [s.text_frame.text for s in first.shapes if s.has_text_frame]
    assert "Big Title" in first_texts and "The subtitle" in first_texts
    title_shape = next(s for s in first.shapes if s.text_frame.text == "Big Title")
    assert title_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 36.0  # 48 px

    # The bullet glyph is a paragraph property, not part of the run — see
    # test_slides_pptx_draws_bullets_as_a_real_list.
    second_texts = [s.text_frame.text for s in second.shapes if s.has_text_frame]
    assert second_texts == ["Agenda", "One", "Two"]


def test_slides_pptx_boxes_hold_their_own_text():
    """A wrapping bullet gets a taller box, and nothing shrinks to fit.

    Shrink-to-fit is what rendered two bullets of the same size at two
    different sizes on the same slide.
    """
    from pptx.enum.text import MSO_AUTO_SIZE

    long_line = "a bullet long enough that it has to wrap onto a second line inside its box"
    content = _deck([{"title": "T", "bullets": ["short", long_line]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    boxes = {s.text_frame.text: s for s in next(iter(deck.slides)).shapes if s.has_text_frame}
    assert boxes[long_line].height == pytest.approx(2 * boxes["short"].height, abs=2)
    assert all(b.text_frame.auto_size == MSO_AUTO_SIZE.NONE for b in boxes.values())
    sizes = {
        box.text_frame.paragraphs[0].runs[0].font.size
        for box in (boxes[long_line], boxes["short"])
    }
    assert len(sizes) == 1


def test_slides_pptx_draws_bullets_as_a_real_list():
    """A literal '•' in the run picks up the font of whatever follows it.

    Mixed scripts on one slide then get mixed bullet glyphs and ragged left
    edges, and a wrapped line starts under the bullet instead of under the
    text. A paragraph bullet is drawn once, by the list.
    """
    content = _deck([{"title": "T", "bullets": ["매출 성장", "Revenue growth"]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    shapes = [s for s in next(iter(deck.slides)).shapes if s.has_text_frame]
    bullets = [s for s in shapes if s.text_frame.text != "T"]
    assert [s.text_frame.text for s in bullets] == ["매출 성장", "Revenue growth"]
    for shape in bullets:
        properties = shape.text_frame.paragraphs[0]._p.find(qn("a:pPr"))
        assert properties is not None
        assert properties.find(qn("a:buChar")).get("char") == "•"
        # A hanging indent: the marker sits left of the text it belongs to.
        assert int(properties.get("indent")) == -int(properties.get("marL"))
        assert int(properties.get("marL")) > 0

    # The heading is not a list item.
    heading = next(s for s in shapes if s.text_frame.text == "T")
    heading_properties = heading.text_frame.paragraphs[0]._p.find(qn("a:pPr"))
    assert heading_properties is None or heading_properties.find(qn("a:buChar")) is None


def test_slides_pptx_padding_insets_geometry():
    element = {"id": "t", "type": "text", "x": 0, "y": 0, "w": 100, "h": 10, "text": "x"}
    def exported(slides: list[dict[str, Any]]) -> Presentation:
        return Presentation(io.BytesIO(
            SlidesPptxExporter().export(_deck(slides), path="d.slides.json").data
        ))

    plain = exported([{"elements": [element]}])
    padded = exported([{"padding": 10, "elements": [element]}])
    plain_box = next(iter(plain.slides)).shapes[0]
    padded_box = next(iter(padded.slides)).shapes[0]
    assert plain_box.left == 0 and plain_box.width == Inches(10)
    assert padded_box.left == Inches(1.0)  # 10% of the 10in page
    assert padded_box.width == Inches(8.0)  # spans the inset content area


def test_slides_pptx_skips_what_it_cannot_embed():
    content = _deck(
        [
            {
                "elements": [
                    # Not inlined (a bare reference) — skipped, never a crash.
                    {"id": "i", "type": "image", "x": 0, "y": 0, "w": 50, "h": 50,
                     "src": "assets/logo.png"},
                    # url() background is out of contract for the pptx door.
                ],
                "background": "url(assets/bg.png)",
            }
        ]
    )
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    (slide,) = deck.slides
    assert [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE] == []


def test_slides_pptx_rejects_non_deck_content():
    with pytest.raises(ValueError):
        SlidesPptxExporter().export("not json", path="d.slides.json")
    with pytest.raises(ValueError):
        SlidesPptxExporter().export(json.dumps({"data": {"slides": "nope"}}), path="d.slides.json")


def test_slides_pptx_names_the_envelope_when_data_is_missing():
    # A deck written without the envelope must fail with the expected shape
    # in the message — not export as one silent blank slide.
    bare = json.dumps({"slides": [{"title": "T"}], "theme": {}})
    with pytest.raises(ValueError, match='"data" envelope'):
        SlidesPptxExporter().export(bare, path="d.slides.json")


def _skin_pptx_bytes() -> bytes:
    """A 4:3 template with a branded master shape and one slide of its own."""
    from pptx import Presentation as _P
    from pptx.enum.shapes import MSO_SHAPE

    skin = _P()
    skin.slide_width = Inches(10)
    skin.slide_height = Inches(7.5)
    slide = skin.slides.add_slide(skin.slide_layouts[6])
    logo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(1), Inches(0.4)
    )
    logo.name = "BrandLogo"
    # python-pptx has no master add_shape; moving the element is the test's
    # stand-in for a real branded master.
    skin.slide_masters[0].shapes._spTree.append(logo._element)
    buf = io.BytesIO()
    skin.save(buf)
    return buf.getvalue()


def _skin_uri() -> str:
    return f"data:{PPTX_MIME};base64,{base64.b64encode(_skin_pptx_bytes()).decode()}"


def test_slides_pptx_template_skin_keeps_master_and_page_size():
    content = json.dumps({
        "type": "slides",
        "title": "Skinned",
        "data": {
            "template": _skin_uri(),
            "slides": [
                {"elements": [{"id": "t", "type": "text", "x": 10, "y": 10,
                               "w": 80, "h": 20, "text": "Hello", "fontSize": 54}]},
                {"title": "Second", "bullets": ["One"]},
            ],
        },
    })
    result = SlidesPptxExporter().export(content, path="deck.slides.json")
    deck = Presentation(io.BytesIO(result.data))

    # The skin's page (4:3) survives; its own slide is dropped — only the
    # canvas content ships.
    assert deck.slide_width == Inches(10)
    assert deck.slide_height == Inches(7.5)
    assert len(list(deck.slides)) == 2

    # The branded master styles the new slides: its shape is reachable from
    # the exported slide's own layout chain.
    first = list(deck.slides)[0]
    master_names = [shape.name for shape in first.slide_layout.slide_master.shapes]
    assert "BrandLogo" in master_names
    assert first.slide_layout.name == "Blank"

    # Content still lands as real shapes with percent geometry projected
    # onto the skin's page.
    texts = [s for s in first.shapes if s.has_text_frame and s.text_frame.text]
    ((run,),) = [p.runs for p in texts[0].text_frame.paragraphs]
    assert run.text == "Hello"
    assert run.font.size.pt == 40.5


def test_slides_pptx_unusable_skin_degrades_to_blank_export():
    for template in ("sources/missing.pptx",  # never inlined — reference miss
                     f"data:{PPTX_MIME};base64,{base64.b64encode(b'junk').decode()}"):
        content = json.dumps({
            "type": "slides",
            "data": {"template": template, "slides": [{"title": "T"}]},
        })
        result = SlidesPptxExporter().export(content, path="d.slides.json")
        deck = Presentation(io.BytesIO(result.data))
        assert deck.slide_width == Inches(10)
        assert deck.slide_height == Inches(5.625)  # blank 16:9 default
        assert len(list(deck.slides)) == 1


def _face_uri(faces: list[str], *, on_slide: bool) -> str:
    """A skin whose runs name ``faces``, on its slides or on a layout.

    ``+mj-lt`` rides along as a theme reference — a pointer at the theme,
    not a face, and the theme is where the east-asian entry goes missing.
    """
    from pptx import Presentation as _P

    skin = _P()
    if on_slide:
        holder = skin.slides.add_slide(skin.slide_layouts[6]).shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1)
        )
    else:
        holder = skin.slide_layouts[0].placeholders[0]
    paragraph = holder.text_frame.paragraphs[0]
    for face in [*faces, "+mj-lt"]:
        run = paragraph.add_run()
        run.text = "x"
        run.font._rPr.get_or_add_latin().set("typeface", face)
    buf = io.BytesIO()
    skin.save(buf)
    return f"data:{PPTX_MIME};base64,{base64.b64encode(buf.getvalue()).decode()}"


def _exported_faces(template: str | None) -> list[list[tuple[str, str]]]:
    """Every exported run's script-to-face pairs, in document order."""
    data = {"slides": [
        {"elements": [{"id": "t", "type": "text", "x": 10, "y": 10, "w": 80,
                       "h": 20, "text": "Hangul 한글", "fontSize": 40}]},
        {"title": "Second", "bullets": ["One"]},
    ]}
    if template:
        data["template"] = template
    result = SlidesPptxExporter().export(
        json.dumps({"type": "slides", "data": data}), path="d.slides.json"
    )
    runs = []
    with zipfile.ZipFile(io.BytesIO(result.data)) as archive:
        for name in sorted(archive.namelist()):
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            xml = archive.read(name).decode("utf-8")
            for properties in re.findall(r"<a:rPr\b[^>]*(?:/>|>.*?</a:rPr>)", xml, re.S):
                runs.append(re.findall(r'<a:(latin|ea|cs) typeface="([^"]+)"', properties))
    return runs


def _bullet_font_uri() -> str:
    """A skin in PowerPoint's ordinary shape: the body left to the theme,
    the bullets naming a dingbat face of their own."""
    from pptx import Presentation as _P

    skin = _P()
    box = skin.slides.add_slide(skin.slide_layouts[6]).shapes.add_textbox(
        Inches(1), Inches(1), Inches(6), Inches(3)
    )
    paragraph = box.text_frame.paragraphs[0]
    properties = paragraph._p.get_or_add_pPr()
    properties.append(properties.makeelement(qn("a:buFont"), {"typeface": "Wingdings"}))
    run = paragraph.add_run()
    run.text = "x"
    run.font._rPr.get_or_add_latin().set("typeface", "+mn-lt")
    buf = io.BytesIO()
    skin.save(buf)
    return f"data:{PPTX_MIME};base64,{base64.b64encode(buf.getvalue()).decode()}"


def test_slides_pptx_never_takes_a_bullet_font_for_the_deck_face():
    # A bullet font is a dingbat picked for one glyph. Counting it as a face
    # left this skin with no other literal to beat it, and every exported
    # character shipped as Wingdings.
    runs = _exported_faces(_bullet_font_uri())
    assert runs
    assert all(run == [] for run in runs)


def test_slides_pptx_names_the_skin_face_for_every_script():
    # Latin, east-asian, and complex scripts each read their own element;
    # naming only a:latin leaves Hangul to whatever the theme says, and a
    # theme's east-asian entry is routinely empty.
    runs = _exported_faces(_face_uri(["Pretendard"] * 3 + ["Arial"], on_slide=True))
    assert runs
    for run in runs:
        assert run == [("latin", "Pretendard"), ("ea", "Pretendard"), ("cs", "Pretendard")]


def test_slides_pptx_reads_a_slideless_template_from_its_layouts():
    # A real template file carries layouts and no slides. Its faces tie one
    # to one there, and the plainest name is the family rather than one of
    # its weights.
    runs = _exported_faces(_face_uri(["Pretendard Light", "Pretendard"], on_slide=False))
    assert runs
    assert all(face == "Pretendard" for run in runs for _, face in run)


def test_slides_pptx_without_a_skin_names_no_face():
    # No skin, no evidence — inventing a face would be making one up.
    assert _exported_faces(None)
    assert all(run == [] for run in _exported_faces(None))


def _bare_skin_uri(width_in: float, height_in: float) -> str:
    from pptx import Presentation as _P

    skin = _P()
    skin.slide_width = Inches(width_in)
    skin.slide_height = Inches(height_in)
    buf = io.BytesIO()
    skin.save(buf)
    return f"data:{PPTX_MIME};base64,{base64.b64encode(buf.getvalue()).decode()}"


_CIRCLE_DECK_SLIDES = [{
    "elements": [
        # A perfect circle on the 16:9 canvas: 10% of 10in == 17.7778% of 5.625in.
        {"id": "c", "type": "shape", "shape": "ellipse", "x": 10, "y": 10,
         "w": 10, "h": 17.7778, "fill": "#ff0000"},
        {"id": "t", "type": "text", "x": 30, "y": 60, "w": 40, "h": 20,
         "text": "Hi", "fontSize": 40},
    ],
}]


def _circle_and_font(template: str | None) -> tuple[float, float, float]:
    data: dict[str, Any] = {"slides": _CIRCLE_DECK_SLIDES}
    if template:
        data["template"] = template
    content = json.dumps({"type": "slides", "data": data})
    result = SlidesPptxExporter().export(content, path="d.slides.json")
    deck = Presentation(io.BytesIO(result.data))
    (slide,) = deck.slides
    (circle,) = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    (text,) = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text]
    font_pt = text.text_frame.paragraphs[0].runs[0].font.size.pt
    return circle.width / circle.height, circle.top / 914400, font_pt


def test_slides_pptx_uniform_scale_never_distorts_shapes():
    # A circle approved on the 16:9 preview stays a circle on every page.
    # Width and height once scaled separately: on a 4:3 skin the same deck
    # measured ratio 0.750 (an ellipse).
    for template in (None, _bare_skin_uri(10, 7.5), _bare_skin_uri(13.333, 7.5)):
        ratio, _, _ = _circle_and_font(template)
        assert abs(ratio - 1.0) < 0.01


def test_slides_pptx_content_centers_on_a_taller_page():
    # 16:9 content on a 4:3 page: scale stays 1, the leftover 1.875in of
    # height splits evenly — the circle drops by 0.9375in and the margins
    # show the skin's own background.
    _, top_unskinned, _ = _circle_and_font(None)
    _, top_43, _ = _circle_and_font(_bare_skin_uri(10, 7.5))
    assert abs(top_unskinned - 0.5625) < 0.01
    assert abs(top_43 - (0.5625 + 0.9375)) < 0.01


def test_slides_pptx_font_rides_the_page_scale():
    # Shapes and images grow with the page (percent geometry); type must
    # grow with them or it shrinks relative to everything else.
    _, _, base_pt = _circle_and_font(None)
    _, _, wide_pt = _circle_and_font(_bare_skin_uri(13.333, 7.5))
    assert base_pt == 30  # 40px * 0.75, the slice-8 value — regression pin
    assert abs(wide_pt - 30 * 13.333 / 10) < 0.1  # ~39.99pt


def test_slides_pptx_deck_page_sets_the_unskinned_page_size():
    content = json.dumps({"type": "slides", "data": {
        "page": {"widthIn": 10.0, "heightIn": 7.5},
        "slides": [{"title": "T"}],
    }})
    result = SlidesPptxExporter().export(content, path="d.slides.json")
    deck = Presentation(io.BytesIO(result.data))
    assert deck.slide_width == Inches(10)
    assert deck.slide_height == Inches(7.5)


def test_slides_pptx_refuses_a_zip_bomb_skin():
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        z.writestr("ppt/presentation.xml", b"\x00" * (201 * 1024 * 1024))
    bomb_uri = f"data:{PPTX_MIME};base64,{base64.b64encode(buf.getvalue()).decode()}"
    content = json.dumps({"type": "slides", "data": {
        "template": bomb_uri, "slides": [{"title": "T"}],
    }})
    # An attack is refused loudly — never absorbed into a blank-layout degrade.
    with pytest.raises(ValueError, match="unpacks to"):
        SlidesPptxExporter().export(content, path="d.slides.json")


def test_slides_pptx_text_after_an_image_keeps_its_font_size():
    # The picture contain factor once rebound the projection `scale`, so any
    # text AFTER an image inherited a factor in the millions and the export
    # died on pptx's font-size limit. A 1x1 image maximizes the corruption.
    content = _deck([{
        "elements": [
            {"id": "i", "type": "image", "x": 10, "y": 10, "w": 30, "h": 30,
             "src": _png_uri()},
            {"id": "t", "type": "text", "x": 10, "y": 60, "w": 60, "h": 20,
             "text": "After the image", "fontSize": 36},
        ],
    }])
    result = SlidesPptxExporter().export(content, path="d.slides.json")
    deck = Presentation(io.BytesIO(result.data))
    (slide,) = deck.slides
    (text,) = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text]
    assert text.text_frame.paragraphs[0].runs[0].font.size.pt == 27  # 36px * 0.75


def test_pptx_page_size_inches_reads_the_declared_size():
    from pptx import Presentation as _P

    from langchain_canvas.exporters import pptx_page_size_inches

    skin = _P()
    skin.slide_width = Inches(13.333)
    skin.slide_height = Inches(7.5)
    buf = io.BytesIO()
    skin.save(buf)
    size = pptx_page_size_inches(buf.getvalue())
    assert size is not None
    assert abs(size[0] - 13.333) < 0.01
    assert abs(size[1] - 7.5) < 0.01
    assert pptx_page_size_inches(b"not a zip") is None


# --- the export tool -------------------------------------------------------------


@dataclass
class _Runtime:
    context: Any = None
    config: dict[str, Any] = field(default_factory=dict)


def _runtime(thread_id: str = "t1") -> _Runtime:
    return _Runtime(config={"configurable": {"thread_id": thread_id}})


def test_export_tool_writes_under_exports():
    store = InMemoryCanvasStore()
    store.write("t1", "page.html", "<h1>One</h1><p>Body</p>", "seed", actor="agent")
    tool_obj = create_export_tool(store)

    message = tool_obj.func(path="page.html", target="docx", runtime=_runtime())
    assert "exports/page.docx" in message

    exported = store.read_bytes("t1", "exports/page.docx")
    document = Document(io.BytesIO(exported.data))
    assert any(p.text == "One" for p in document.paragraphs)


def test_export_tool_merges_a_directory_in_name_order():
    store = InMemoryCanvasStore()
    store.write("t1", "report/02-b.html", "<h1>Second</h1>", "seed", actor="agent")
    store.write("t1", "report/01-a.html", "<h1>First</h1>", "seed", actor="agent")
    store.write("t1", "report/notes.txt", "not html", "seed", actor="agent")
    tool_obj = create_export_tool(store)

    message = tool_obj.func(path="report/", target="docx", runtime=_runtime())
    assert "exports/report.docx" in message

    document = Document(io.BytesIO(store.read_bytes("t1", "exports/report.docx").data))
    texts = [p.text for p in document.paragraphs if p.text.strip()]
    assert texts.index("First") < texts.index("Second")
    assert "<w:pBdr>" in document.element.xml  # sections split by a rule


def test_export_tool_is_honest_about_misses():
    store = InMemoryCanvasStore()
    store.write("t1", "sales.table.json", json.dumps({"data": {}}), "seed", actor="agent")
    tool_obj = create_export_tool(store)

    missing = tool_obj.func(path="nope.html", target="docx", runtime=_runtime())
    assert missing.startswith("Error:")

    wrong_target = tool_obj.func(path="sales.table.json", target="docx", runtime=_runtime())
    assert wrong_target.startswith("Error:")
    assert "xlsx" in wrong_target  # names the formats that would work

    empty_dir = tool_obj.func(path="deck/", target="docx", runtime=_runtime())
    assert empty_dir.startswith("Error:")


def test_export_tool_inlines_slide_assets_before_pptx():
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "assets/logo.png", PNG_1PX, "logo", actor="agent")
    store.write(
        "t1",
        "deck.slides.json",
        _deck([
            {"elements": [
                {"id": "i", "type": "image", "x": 10, "y": 10, "w": 50, "h": 50,
                 "src": "assets/logo.png"},
            ]},
        ]),
        "seed",
        actor="agent",
    )
    tool_obj = create_export_tool(store)

    message = tool_obj.func(path="deck.slides.json", target="pptx", runtime=_runtime())
    assert "exports/Deck.pptx" in message  # the deck's own title names the file

    deck = Presentation(io.BytesIO(store.read_bytes("t1", "exports/Deck.pptx").data))
    (picture,) = [
        s for s in next(iter(deck.slides)).shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    # The stored reference became the stored bytes — the deck is self-contained.
    assert picture.image.blob == PNG_1PX


def test_slides_pptx_table_lands_as_a_real_table():
    """A table element goes out as a PowerPoint table — columns at their
    widths, the merged header merged, the grid line and fills the element
    states — not as sixteen boxes a column cannot be dragged across."""
    content = _deck(
        [
            {
                "elements": [
                    {
                        "id": "t", "type": "table", "x": 10, "y": 20, "w": 80, "h": 40,
                        "rows": [["Header", ""], ["a", "b"]],
                        "header": True,
                        "colWidths": [75, 25],
                        "stroke": "#9E9E9E", "strokeWidth": 2,
                        "fontSize": 16, "color": "#000000", "fontFamily": "Arial",
                        "cells": [
                            {"r": 0, "c": 0, "colSpan": 2, "fill": "#DDEEFF", "align": "center"},
                            {"r": 1, "c": 1, "bold": True, "color": "#FF0000"},
                        ],
                    }
                ]
            }
        ]
    )
    result = SlidesPptxExporter().export(content, path="deck.slides.json")
    deck = Presentation(io.BytesIO(result.data))
    (slide,) = deck.slides
    (frame,) = [s for s in slide.shapes if s.has_table]
    table = frame.table
    assert len(table.rows) == 2 and len(table.columns) == 2
    assert abs(table.columns[0].width / table.columns[1].width - 3.0) < 0.01
    assert table.cell(0, 0).is_merge_origin and table.cell(0, 1).is_spanned
    header = table.cell(0, 0)
    assert header.text_frame.text == "Header"
    assert header.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert str(header.fill.fore_color.rgb) == "DDEEFF"
    run = header.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True  # the header row
    assert run.font.size.pt == 12  # 16 px -> 12 pt
    assert run.font.name == "Arial"
    b = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
    assert b.font.bold is True and str(b.font.color.rgb) == "FF0000"
    a = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert a.font.bold is False and str(a.font.color.rgb) == "000000"
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    edge = table.cell(1, 0)._tc.tcPr.find(f"{ns}lnL")
    assert edge is not None and edge.get("w") == str(int(2 * 0.75 * 12700))
    assert edge.find(f"{ns}solidFill/{ns}srgbClr").get("val") == "9E9E9E"
    style = table._tbl.tblPr.find(f"{ns}tableStyleId")
    assert style is not None and style.text == "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def test_slides_pptx_table_without_a_grid_line_declares_no_lines():
    content = _deck([{"elements": [{"id": "t", "type": "table", "x": 0, "y": 0, "w": 50, "h": 50,
                                     "rows": [["x"]]}]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    (frame,) = [s for s in deck.slides[0].shapes if s.has_table]
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    edge = frame.table.cell(0, 0)._tc.tcPr.find(f"{ns}lnT")
    assert edge is not None and edge.find(f"{ns}noFill") is not None


def test_slides_pptx_a_growing_box_is_written_grown_and_keeps_growing():
    """A box that grows with its text lands at the height its words need,
    flagged to keep growing — so a viewer that does not re-fit on open still
    shows every line, and one that does agrees."""
    from pptx.enum.text import MSO_AUTO_SIZE

    long_text = "가나다라마바사아자차카타파하 " * 8
    content = _deck([{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 10, "w": 40, "h": 5, "fontSize": 24,
         "text": long_text, "autofit": "shape"},
        {"id": "b", "type": "text", "x": 5, "y": 60, "w": 40, "h": 5, "fontSize": 24,
         "text": long_text},
    ]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    shapes = [s for s in next(iter(deck.slides)).shapes if s.has_text_frame]
    grown, fixed = shapes[0], shapes[1]
    assert grown.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    assert fixed.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert grown.height > fixed.height * 3


def test_slides_pptx_shrinking_type_is_written_with_its_scale():
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.oxml.ns import qn

    content = _deck([{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 10, "w": 40, "h": 5, "fontSize": 24,
         "text": "가나다라마바사아자차카타파하 " * 8, "autofit": "text"},
    ]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    box = next(s for s in next(iter(deck.slides)).shapes if s.has_text_frame)
    assert box.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    fit = box.text_frame._bodyPr.find(qn("a:normAutofit"))
    assert 25000 <= int(fit.get("fontScale")) < 100000


def test_markdown_docx_keeps_headings_lists_tables_and_breaks():
    """The .md door carries the same deliberate subset as the .html one —
    before it existed, a canvas document could not leave as a file at all."""
    from docx import Document

    from langchain_canvas.exporters import MarkdownDocxExporter

    md = "\n".join([
        "# 협업 제안서",
        "",
        "본 문서는 **브레인크루**와 *신한은행*의 협업 범위를 정리한다.",
        "이어지는 줄은 같은 문단이다.",
        "",
        "## 일정",
        "| 단계 | 기간 |",
        "|---|---|",
        "| 기획 | 2주 |",
        "| 개발 | 6주 |",
        "",
        "1. 첫째",
        "2) 둘째",
        "- 항목 A",
        "* 항목 B",
        "",
        "```",
        "uv run python -m app",
        "```",
        "",
        "---",
        "",
        "끝.",
    ])
    exported = MarkdownDocxExporter().export(md, path="report.md", title="협업 제안서")
    assert exported.filename.endswith(".docx")
    doc = Document(io.BytesIO(exported.data))
    styles = [p.style.name for p in doc.paragraphs if p.text.strip()]
    assert "Heading 1" in styles and "Heading 2" in styles
    assert styles.count("List Number") == 2 and styles.count("List Bullet") == 2
    body = next(p for p in doc.paragraphs if "협업 범위" in p.text)
    assert "이어지는 줄은 같은 문단이다" in body.text  # soft-wrap joins
    assert any(r.bold for r in body.runs) and any(r.italic for r in body.runs)
    assert len(doc.tables) == 1 and doc.tables[0].rows[0].cells[0].text == "단계"
    assert any("uv run python -m app" == p.text for p in doc.paragraphs)
    assert "<w:pBdr>" in doc.element.xml  # --- is a line, not a new page


def test_markdown_docx_keeps_inline_marks_in_cells_code_quotes_and_nesting():
    """What a model writes most — bold in table cells, ~~strike~~, ***both***,
    fenced code, block quotes, nested bullets, hard line breaks — left the
    door as raw markdown text before. Measured on a user's sample, 2026-09-04."""
    from docx import Document

    from langchain_canvas.exporters import MarkdownDocxExporter

    md = "\n".join([
        "첫 줄  ",
        "둘째 줄은 줄바꿈 뒤에 온다.",
        "",
        "***굵게 기울임*** 그리고 ~~취소선~~ 그리고 `코드`.",
        "",
        "> **인용문**",
        "> 둘째 줄",
        "",
        "- 항목",
        "  - 하위 항목",
        "    - 더 깊은 항목",
        "",
        "| 구분 | 예시 |",
        "|---|---|",
        "| **Bold** | ~~옛 값~~ |",
        "",
        "```python",
        "def f():",
        "    return 1",
        "```",
    ])
    doc = Document(io.BytesIO(MarkdownDocxExporter().export(md, path="s.md").data))
    xml = doc.element.xml
    first = doc.paragraphs[0]
    assert "<w:br/>" in first._p.xml and "둘째 줄은" in first.text  # hard break kept
    marks = doc.paragraphs[1]
    both = next(r for r in marks.runs if r.text == "굵게 기울임")
    assert both.bold and both.italic
    assert next(r for r in marks.runs if r.text == "취소선").font.strike
    assert next(r for r in marks.runs if r.text == "코드").font.name == "Consolas"
    quote = next(p for p in doc.paragraphs if "인용문" in p.text)
    assert quote.style.name == "Intense Quote" and "둘째 줄" in quote.text
    styles = [p.style.name for p in doc.paragraphs]
    assert "List Bullet" in styles and "List Bullet 2" in styles and "List Bullet 3" in styles
    cell = doc.tables[0].cell(1, 0)
    assert cell.text == "Bold" and cell.paragraphs[0].runs[0].bold  # no raw asterisks
    assert doc.tables[0].cell(1, 1).paragraphs[0].runs[0].font.strike
    code = next(p for p in doc.paragraphs if "def f():" in p.text)
    assert "    return 1" in code.text
    assert all(r.font.name == "Consolas" for r in code.runs if r.text.strip())
    assert "**" not in xml and "~~" not in xml and "`" not in xml


def test_markdown_docx_registered_for_md_canvas_files():
    from langchain_canvas.exporters import default_exporters, exporter_for

    exporter = exporter_for("notes.md", "docx", default_exporters())
    assert exporter is not None and exporter.target == "docx"


def test_markdown_docx_filename_drops_the_md_suffix():
    """A .md canvas file exports as `<stem>.docx`, never `<stem>.md.docx` —
    the stem strips every canvas extension, markdown included."""
    from langchain_canvas.exporters import MarkdownDocxExporter, _stem

    assert _stem("report.md") == "report"
    assert _stem("notes/보고서.markdown") == "보고서"
    exported = MarkdownDocxExporter().export("# 제목\n\n본문", path="report.md")
    assert exported.filename == "report.docx"


def qn_docx(tag: str) -> str:
    from docx.oxml.ns import qn as _qn

    return _qn(tag)


def test_slides_pptx_an_explicitly_unfilled_shape_stays_unfilled():
    """fill "none" is the shape saying it is transparent — the exporter must
    not paint the default over what the border (or the slide) shows."""
    from pptx.enum.dml import MSO_FILL

    content = _deck([{"elements": [
        {"id": "a", "type": "shape", "shape": "rect", "x": 5, "y": 5, "w": 20, "h": 10,
         "fill": "none"},
        {"id": "b", "type": "shape", "shape": "rect", "x": 40, "y": 5, "w": 20, "h": 10},
    ]}])
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    shapes = [s for s in next(iter(deck.slides)).shapes if not s.has_text_frame or True]
    unfilled, defaulted = shapes[0], shapes[1]
    assert unfilled.fill.type == MSO_FILL.BACKGROUND
    assert defaulted.fill.type == MSO_FILL.SOLID  # an authored box still shows


def test_slides_pptx_writes_element_rotation():
    """A rotated element turns in the .pptx too, so the exported deck matches
    what the editor drew."""
    content = _deck(
        [{"elements": [
            {"id": "t", "type": "text", "x": 10, "y": 10, "w": 40, "h": 12,
             "text": "Tilted", "rotation": 30},
        ]}]
    )
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="d.slides.json").data))
    (slide,) = deck.slides
    box = next(s for s in slide.shapes if s.has_text_frame and s.text_frame.text == "Tilted")
    assert round(box.rotation) == 30


def test_slides_pptx_honors_a_portrait_page():
    """A deck whose page is taller than wide exports a portrait .pptx — the
    page size follows `data.page`, not a fixed 16:9."""
    content = json.dumps({
        "type": "slides", "title": "Tall",
        "data": {"page": {"widthIn": 7.5, "heightIn": 10}, "slides": [
            {"elements": [{"id": "t", "type": "text", "x": 10, "y": 10, "w": 60,
                           "h": 10, "text": "Portrait"}]}
        ]},
    })
    deck = Presentation(io.BytesIO(SlidesPptxExporter().export(content, path="tall.slides.json").data))
    assert deck.slide_height > deck.slide_width
    assert deck.slide_width == Inches(7.5)
    assert deck.slide_height == Inches(10)
