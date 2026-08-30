"""Layout lint — certain-only findings, silence on clean decks, no crashes.

The false-positive suite feeds real deck shapes from actual use (skinned
decks, letterboxed refits, padded slides, structured slides) and demands
zero warnings — one warning here means the check must be narrowed or cut.
"""

from __future__ import annotations

import json
from typing import Any

import pytest
from pydantic import ValidationError

from langchain_canvas import encode_slides
from langchain_canvas.layout_lint import lint_slides_data
from langchain_canvas.protocol.artifacts import SlidesData
from langchain_canvas.store import InMemoryCanvasStore
from langchain_canvas.tools import (
    _refit_slides_to_page,
    create_canvas_tools,
)


class _Runtime:
    """The slice of ToolRuntime the canvas tools read."""

    def __init__(self, config: dict[str, Any]) -> None:
        self.context = None
        self.config = config
        self.stream_writer = None


def _runtime() -> _Runtime:
    return _Runtime(config={"configurable": {"thread_id": "t1"}})


def _tools(store: InMemoryCanvasStore) -> dict[str, Any]:
    return {t.name: t for t in create_canvas_tools(store)}


def _write(store: InMemoryCanvasStore, data: dict[str, Any], path: str = "deck.slides.json") -> str:
    tools = _tools(store)
    return tools["write_canvas"].func(
        path=path,
        content=encode_slides("Deck", data),
        description="save",
        runtime=_runtime(),
    )


def _deck(*elements: dict[str, Any], **slide_extra: Any) -> dict[str, Any]:
    return {"slides": [{**slide_extra, "elements": list(elements)}]}


def _el(
    eid: str, etype: str, x: float, y: float, w: float, h: float, **extra: Any
) -> dict[str, Any]:
    return {"id": eid, "type": etype, "x": x, "y": y, "w": w, "h": h, **extra}


# --- each certain check fires ----------------------------------------------------


def test_an_off_page_element_is_named_with_its_numbers() -> None:
    warnings = lint_slides_data(
        _deck(_el("title", "text", 60, 10, 58, 10, text="Hi"))
    )
    assert warnings == [
        'slide 1, element "title": x + w = 118 (off the page — the page runs 0 to 100)'
    ]


def test_negative_origin_is_off_page() -> None:
    warnings = lint_slides_data(
        _deck(_el("a", "shape", -5, -2, 10, 10, shape="rect"))
    )
    assert len(warnings) == 1
    assert "x = -5" in warnings[0] and "y = -2" in warnings[0]


def test_a_zero_sized_element_is_flagged() -> None:
    warnings = lint_slides_data(
        _deck(_el("z", "shape", 10, 10, 0, 5, shape="rect"))
    )
    assert warnings == [
        'slide 1, element "z": w = 0, h = 5 (zero or negative size renders nothing)'
    ]


def test_an_empty_text_element_is_flagged() -> None:
    warnings = lint_slides_data(
        _deck(_el("t", "text", 10, 10, 30, 10, text="   "))
    )
    assert "a text element with no text" in warnings[0]


def test_a_broken_image_reference_is_flagged_only_with_store_knowledge() -> None:
    deck = _deck(_el("img", "image", 10, 10, 30, 30, src="sources/missing.png"))
    assert lint_slides_data(deck) == []  # no ref_exists — no claim
    warnings = lint_slides_data(deck, ref_exists=lambda ref: False)
    assert 'src "sources/missing.png" is not on the canvas' in warnings[0]
    assert lint_slides_data(deck, ref_exists=lambda ref: True) == []


def test_a_fully_covered_element_is_flagged() -> None:
    warnings = lint_slides_data(
        _deck(
            _el("hidden", "text", 20, 20, 20, 10, text="Hi"),
            _el("panel", "shape", 10, 10, 60, 40, shape="rect", fill="#0d1b3e"),
        )
    )
    assert 'element "hidden" is completely covered by element "panel"' in warnings[0]


# --- schema: what the export will refuse ------------------------------------------


def _schema_warning(warnings: list[str]) -> str:
    matched = [w for w in warnings if "does not match the slides schema" in w]
    assert len(matched) == 1, warnings
    return matched[0]


def test_a_missing_element_id_is_flagged_with_the_export_consequence() -> None:
    deck = {"slides": [{"elements": [
        {"type": "text", "x": 6, "y": 8, "w": 88, "h": 10, "text": "Q3"}
    ]}]}
    warning = _schema_warning(lint_slides_data(deck))
    assert "exporting it fails" in warning
    assert 'slide 1, element #1: "id" is required' in warning
    assert "short unique id string" in warning


def test_an_invalid_layout_value_names_the_allowed_values() -> None:
    deck = {"slides": [{"layout": "title+bullets", "title": "Q3"}]}
    warning = _schema_warning(lint_slides_data(deck))
    assert "slide 1: \"layout\" = 'title+bullets'" in warning
    assert "'title', 'content', 'section', 'image', 'two-column' or 'blank'" in warning


def test_the_same_mistake_across_many_elements_is_summarized() -> None:
    deck = {"slides": [{"elements": [
        {"type": "text", "x": 1, "y": i, "w": 10, "h": 5, "text": "x"}
        for i in range(10)
    ]}]}
    warning = _schema_warning(lint_slides_data(deck))
    assert warning.count('"id" is required') == 6
    assert "... and 4 more" in warning


def test_a_field_the_schema_has_no_place_for_is_flagged_as_ignored() -> None:
    deck = _deck(_el("box", "shape", 5, 5, 10, 10, shape="rect", rotation=45))
    warnings = lint_slides_data(deck)
    assert len(warnings) == 1
    assert "the canvas and the export both ignore them" in warnings[0]
    assert 'slide 1, element "box": "rotation"' in warnings[0]


def test_unknown_fields_are_caught_at_every_level() -> None:
    deck = {
        "theme": "dark",
        "slides": [{"transition": "fade", "elements": []}],
        "page": {"widthIn": 10, "heightIn": 5.625, "dpi": 96},
    }
    warnings = lint_slides_data(deck)
    assert len(warnings) == 1
    assert 'the deck: "theme"' in warnings[0]
    assert 'slide 1: "transition"' in warnings[0]
    assert 'the deck page: "dpi"' in warnings[0]


def test_out_of_range_padding_is_reported_once_not_twice() -> None:
    warnings = lint_slides_data({"slides": [{"padding": 60, "elements": []}]})
    assert warnings == [
        "slide 1: padding = 60 leaves no content area (must be 0 to below 50)"
    ]


def test_a_padding_of_the_wrong_type_still_reaches_the_schema_line() -> None:
    # The dedicated padding check ignores non-numbers, so the schema line is
    # the only thing standing between this deck and a refused export.
    warning = _schema_warning(lint_slides_data({"slides": [{"padding": "wide"}]}))
    assert 'slide 1: "padding"' in warning


# --- shadowed content: written but never drawn ------------------------------------


def test_structured_text_next_to_elements_is_flagged() -> None:
    deck = {"slides": [{
        "layout": "content",
        "title": "지표",
        "bullets": ["a", "b"],
        "elements": [_el("kpi", "text", 60, 20, 35, 10, text="128건")],
    }]}
    warnings = lint_slides_data(deck)
    assert len(warnings) == 1
    assert '"title", "bullets" are set next to "elements"' in warnings[0]
    assert "only \"elements\" is drawn" in warnings[0]
    assert "move that text into an element" in warnings[0]


def test_one_shadowed_field_reads_as_singular() -> None:
    deck = {"slides": [{
        "subtitle": "요약",
        "elements": [_el("kpi", "text", 60, 20, 35, 10, text="128건")],
    }]}
    assert '"subtitle" is set next to "elements"' in lint_slides_data(deck)[0]


def test_an_empty_elements_array_does_not_shadow_anything() -> None:
    # An empty array is what the structured shape looks like on the wire.
    assert lint_slides_data(
        {"slides": [{"layout": "content", "title": "지표", "elements": []}]}
    ) == []


# --- what deliberately does NOT fire ---------------------------------------------


def test_partial_overlap_is_not_a_finding() -> None:
    # Text over a background panel is normal design, not a defect.
    assert lint_slides_data(
        _deck(
            _el("text", "text", 20, 20, 40, 10, text="Hi"),
            _el("panel", "shape", 30, 15, 40, 40, shape="rect"),
        )
    ) == []


def test_an_ellipse_or_translucent_fill_never_counts_as_cover() -> None:
    base = _el("under", "text", 20, 20, 10, 10, text="Hi")
    ellipse = _el("e", "shape", 0, 0, 100, 100, shape="ellipse")
    alpha = _el("a", "shape", 0, 0, 100, 100, shape="rect", fill="#0d1b3e80")
    assert lint_slides_data(_deck(base, ellipse)) == []
    assert lint_slides_data(_deck(base, alpha)) == []


def test_data_uri_images_make_no_reference_claim() -> None:
    deck = _deck(_el("img", "image", 10, 10, 30, 30, src="data:image/png;base64,AAAA"))
    assert lint_slides_data(deck, ref_exists=lambda ref: False) == []


def test_edge_touching_elements_are_clean() -> None:
    # Exactly on the page edge — and the re-fit's 4-decimal rounding just
    # past it — is fine. Text painted AFTER a full-bleed rect is not covered
    # (paint order: array end is the front).
    assert lint_slides_data(
        _deck(
            _el("full", "shape", 0, 0, 100, 100, shape="rect"),
            _el("r", "text", 55.0001, 0, 45, 10, text="Hi"),
        )
    ) == []


def test_real_deck_shapes_from_use_stay_silent() -> None:
    """The false-positive corpus: decks that mirror real saved decks."""
    corpus: list[dict[str, Any]] = [
        # structured slides (derived layout)
        {"slides": [
            {"layout": "title", "title": "Q3 결과 보고", "subtitle": "요약"},
            {"title": "성과", "bullets": ["a", "b"]},
        ]},
        # skinned circle deck (post-refit letterbox numbers, 4-decimal rounding)
        {
            "template": "sources/brand-skin.pptx",
            "page": {"widthIn": 10.0, "heightIn": 7.5},
            "slides": [{"elements": [
                _el("c", "shape", 0.0, 12.4991, 56.25, 75.0019, shape="ellipse", fill="#ff0000"),
                _el("t", "text", 0.0, 80.0008, 40.0, 7.5002, text="Hi", fontSize=40.0),
            ]}],
        },
        # padded slide with full-width text and a background panel behind it
        {"slides": [{"padding": 5, "elements": [
            _el("panel", "shape", 0, 0, 100, 30, shape="rect", fill="#0d1b3e"),
            _el("title", "text", 5, 8, 90, 14, text="제목", fontSize=36, color="#ffffff"),
        ]}]},
        # edge-touching full-bleed image
        {"slides": [{"elements": [
            _el("bg", "image", 0, 0, 100, 100, src="sources/photo.png"),
            _el("cap", "text", 4, 88, 60, 8, text="caption"),
        ]}]},
    ]
    for deck in corpus:
        assert lint_slides_data(deck, ref_exists=lambda ref: True) == [], deck


def test_both_wire_spellings_of_every_field_stay_silent() -> None:
    """camelCase and snake_case are both the schema — neither is unknown."""
    camel = {
        "page": {"widthIn": 10, "heightIn": 5.625},
        "slides": [{"textColor": "#111", "elements": [
            _el("t", "text", 5, 5, 40, 10, text="Hi", fontSize=28, bold=True),
        ]}],
    }
    snake = {
        "page": {"width_in": 10, "height_in": 5.625},
        "slides": [{"text_color": "#111", "elements": [
            _el("t", "text", 5, 5, 40, 10, text="Hi", font_size=28, bold=True),
        ]}],
    }
    assert lint_slides_data(camel) == []
    assert lint_slides_data(snake) == []


_LAYOUTS = [None, "title", "content", "section", "image", "two-column", "blank"]
_ELEMENTS = [
    {"type": "text", "text": "Hi", "fontSize": 24, "color": "#111", "align": "center"},
    {"type": "image", "src": "sources/photo.png"},
    {"type": "shape", "shape": "rect", "fill": "#0d1b3e"},
    {"type": "shape", "shape": "ellipse", "fill": "#ff0000"},
    {"type": "shape", "shape": "line", "fill": "#888888"},
]


def test_the_valid_deck_sweep_stays_silent() -> None:
    """Every combination of the schema's own options warns about nothing.

    One warning here means a check claims a valid deck is broken, which is
    the one failure mode that would teach a model to ignore the channel.
    """
    checked = 0
    for layout in _LAYOUTS:
        for element in _ELEMENTS:
            for padding in (None, 0, 5, 49.9):
                for page in (None, {"widthIn": 10, "heightIn": 5.625},
                             {"widthIn": 13.333, "heightIn": 7.5}):
                    for background in (None, "#ffffff"):
                        slide: dict[str, Any] = {
                            "elements": [{"id": "e1", "x": 5, "y": 5, "w": 40,
                                          "h": 10, **element}]
                        }
                        if layout is not None:
                            slide["layout"] = layout
                        if padding is not None:
                            slide["padding"] = padding
                        if background is not None:
                            slide["background"] = background
                        deck: dict[str, Any] = {"slides": [slide]}
                        if page is not None:
                            deck["page"] = page
                        assert lint_slides_data(
                            deck, ref_exists=lambda ref: True
                        ) == [], deck
                        checked += 1
    assert checked == 840


def test_the_structured_deck_sweep_stays_silent() -> None:
    """The same options with the derived layout instead of free elements."""
    checked = 0
    for layout in _LAYOUTS:
        for filled in ({"title": "제목"}, {"title": "제목", "bullets": ["a", "b"]},
                       {"title": "제목", "subtitle": "요약"},
                       {"title": "제목", "bullets": ["a"], "bullets2": ["b"]},
                       {"title": "제목", "image": "sources/photo.png"}):
            for notes in (None, "speaker notes"):
                for elements in (None, []):
                    slide: dict[str, Any] = dict(filled)
                    if layout is not None:
                        slide["layout"] = layout
                    if notes is not None:
                        slide["notes"] = notes
                    if elements is not None:
                        slide["elements"] = elements
                    assert lint_slides_data(
                        {"slides": [slide]}, ref_exists=lambda ref: True
                    ) == [], slide
                    checked += 1
    assert checked == 140


# --- the write/edit channel -------------------------------------------------------


def test_write_canvas_appends_the_layout_check_block() -> None:
    store = InMemoryCanvasStore()
    result = _write(store, {"slides": [{"elements": [
        {"id": "title", "type": "text", "x": 60, "y": 10, "w": 58, "h": 10, "text": "Hi"}
    ]}]})
    assert result.startswith("Wrote deck.slides.json")
    assert "Deck check:" in result
    assert 'element "title": x + w = 118' in result
    assert 'pages="grid"' in result


def test_write_canvas_stays_silent_for_a_clean_deck() -> None:
    store = InMemoryCanvasStore()
    result = _write(store, {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 10, "y": 10, "w": 50, "h": 10, "text": "Hi"}
    ]}]})
    assert "Deck check" not in result


def test_write_canvas_refuses_the_schema_and_shadow_findings() -> None:
    """What an agent actually wrote in a run. A deck the export cannot run,
    with text that would never be drawn, does not land — the reply names
    every fix and the canvas keeps what it had."""
    store = InMemoryCanvasStore()
    result = _write(store, {"slides": [{
        "layout": "title+bullets",
        "title": "3분기 실적",
        "bullets": ["매출 128억", "신규 고객 24곳"],
        "elements": [
            {"type": "text", "x": 60, "y": 20, "w": 35, "h": 12,
             "text": "128억", "fontSize": 40, "rotation": 45}
        ],
    }]})
    assert result.startswith("Error: deck.slides.json was not saved")
    assert "exporting it fails" in result
    assert '"id" is required' in result
    assert "'title', 'content', 'section'" in result
    assert '"title", "bullets" are set next to "elements"' in result
    assert "deck.slides.json" not in {f.path for f in store.list_files("t1")}


def test_an_unknown_field_alone_still_lands_with_a_warning() -> None:
    """A field the schema ignores loses that field, nothing more — advice,
    not a refusal, the same as before."""
    store = InMemoryCanvasStore()
    result = _write(store, {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 10, "y": 10, "w": 50, "h": 10,
         "text": "Hi", "rotation": 45}
    ]}]})
    assert result.startswith("Wrote deck.slides.json")
    assert '"rotation"' in result


def test_a_deck_key_outside_data_is_refused_with_its_place() -> None:
    """The failure seen in a run: `template` at the top level, saved without
    complaint, and an export with no skin. The reply now shows the shape."""
    store = InMemoryCanvasStore()
    tools = _tools(store)
    content = json.dumps({
        "type": "slides", "title": "Deck", "template": "sources/brand.pptx",
        "data": {"slides": [{"title": "Hi", "layout": "title"}]},
    })
    result = tools["write_canvas"].func(
        path="deck.slides.json", content=content, description="save", runtime=_runtime()
    )
    assert result.startswith("Error: deck.slides.json was not saved")
    assert '"template" is written at the top level' in result
    assert '"data": {"template": "sources/brand.pptx", "slides": [...]}' in result
    assert store.list_files("t1") == []


def test_a_deck_that_is_not_json_is_refused() -> None:
    store = InMemoryCanvasStore()
    tools = _tools(store)
    result = tools["write_canvas"].func(
        path="deck.slides.json", content="# not a deck", description="save", runtime=_runtime()
    )
    assert result.startswith("Error: deck.slides.json was not saved")
    assert "not valid JSON" in result


def _store_with_upload(*names: str) -> InMemoryCanvasStore:
    store = InMemoryCanvasStore()
    for name in names:
        store.write_bytes("t1", f"sources/{name}", b"PK-stub", "Upload", actor="human")
    return store


def test_a_new_deck_takes_the_only_pptx_upload_as_its_template() -> None:
    """A model composing a deck in an upload's style forgets the pointer more
    often than it writes it; with one upload there is nothing to choose."""
    store = _store_with_upload("brand.pptx")
    result = _write(store, {"slides": [{"title": "Hi", "layout": "title"}]})
    assert result.startswith("Wrote deck.slides.json")
    assert "template set to sources/brand.pptx" in result
    saved = json.loads(store.read("t1", "deck.slides.json").content)
    assert saved["data"]["template"] == "sources/brand.pptx"


def test_two_uploads_leave_the_template_to_the_model() -> None:
    store = _store_with_upload("a.pptx", "b.pptx")
    result = _write(store, {"slides": [{"title": "Hi", "layout": "title"}]})
    assert "template set to" not in result
    assert "template" not in json.loads(store.read("t1", "deck.slides.json").content)["data"]


def test_a_null_template_opts_out_of_the_default() -> None:
    store = _store_with_upload("brand.pptx")
    result = _write(store, {"template": None, "slides": [{"title": "Hi", "layout": "title"}]})
    assert "template set to" not in result
    assert json.loads(store.read("t1", "deck.slides.json").content)["data"]["template"] is None


def test_a_rewrite_never_adds_a_template_behind_the_model() -> None:
    """The default is for a brand-new deck only; replacing one keeps its choice."""
    store = InMemoryCanvasStore()
    _write(store, {"slides": [{"title": "Hi", "layout": "title"}]})
    store.write_bytes("t1", "sources/brand.pptx", b"PK-stub", "Upload", actor="human")
    tools = _tools(store)
    revision = store.read("t1", "deck.slides.json").revision
    result = tools["write_canvas"].func(
        path="deck.slides.json",
        content=encode_slides("Deck", {"slides": [{"title": "Hello", "layout": "title"}]}),
        description="rewrite",
        revision=revision,
        runtime=_runtime(),
    )
    assert result.startswith("Wrote deck.slides.json")
    assert "template" not in json.loads(store.read("t1", "deck.slides.json").content)["data"]


def test_edit_canvas_lints_the_edited_deck() -> None:
    store = InMemoryCanvasStore()
    _write(store, {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 10, "y": 10, "w": 50, "h": 10, "text": "Hi"}
    ]}]})
    tools = _tools(store)
    revision = store.read("t1", "deck.slides.json").revision
    result = tools["edit_canvas"].func(
        path="deck.slides.json",
        old='"x": 10,',
        new='"x": 70,',
        description="move",
        revision=revision,
        runtime=_runtime(),
    )
    assert result.startswith("Edited deck.slides.json")
    assert "x + w = 120" in result


# --- padding: no crash, honest schema ---------------------------------------------


@pytest.mark.parametrize("padding", [50, 60, 100, -5])
def test_refit_survives_degenerate_padding(padding: float) -> None:
    data = {"slides": [{"padding": padding, "elements": [
        {"id": "t", "type": "text", "x": 10, "y": 10, "w": 50, "h": 10, "text": "Hi"}
    ]}]}
    # Must not raise; the degenerate slide is skipped, not mangled.
    _refit_slides_to_page(data, 10.0, 5.625, 10.0, 7.5)
    assert data["slides"][0]["elements"][0]["x"] == 10


def test_the_schema_refuses_degenerate_padding_honestly() -> None:
    with pytest.raises(ValidationError, match="less_than|less than"):
        SlidesData.model_validate({"slides": [{"padding": 60}]})
    assert SlidesData.model_validate({"slides": [{"padding": 5}]}).slides[0].padding == 5


def test_a_padded_deck_still_exports_and_validates() -> None:
    deck = SlidesData.model_validate(
        json.loads(encode_slides("D", {"slides": [{"padding": 6, "elements": []}]}))["data"]
    )
    assert deck.slides[0].padding == 6


# --- structured slides: the checks look at what will be drawn -----------------


def _structured(**fields: Any) -> dict[str, Any]:
    return {"slides": [fields]}


def test_a_structured_slide_is_checked_through_its_derived_boxes() -> None:
    """The gap this closes: a slide with no coordinates used to be skipped."""
    warnings = lint_slides_data(_structured(title="T", bullets=["real", "", "also real"]))
    assert len(warnings) == 1
    assert "no text" in warnings[0]


def test_a_derived_finding_is_named_the_way_the_author_wrote_it() -> None:
    """`bullets[2]`, not `bul_1` — the author never typed `bul_1`."""
    (warning,) = lint_slides_data(_structured(title="T", bullets=["real", "  "]))
    assert '"bullets[2]"' in warning


def test_a_derived_finding_in_the_second_column_names_that_column() -> None:
    (warning,) = lint_slides_data(
        _structured(layout="two-column", title="T", bullets=["a"], bullets2=["b", ""])
    )
    assert '"bullets2[2]"' in warning


def test_a_structured_image_reference_is_checked() -> None:
    (warning,) = lint_slides_data(
        _structured(layout="image", title="T", image="assets/missing.png"),
        ref_exists=lambda _: False,
    )
    assert '"image"' in warning and "not on the canvas" in warning


def test_a_clean_structured_slide_stays_silent() -> None:
    assert lint_slides_data(_structured(title="Quarter", bullets=["one", "two", "three"])) == []


def test_a_slide_the_schema_rejects_is_not_reported_twice() -> None:
    """The schema line already says it; deriving it would repeat the finding."""
    warnings = lint_slides_data(_structured(layout="title_bullets", title="T", bullets=["a"]))
    assert len(warnings) == 1
    assert "slides schema" in warnings[0]


# --- room: too much body, and text too small ---------------------------------


def _many(count: int) -> dict[str, Any]:
    return _structured(
        title="T", bullets=[f"운영 지표 {i}: 전분기 대비 개선 추세" for i in range(count)]
    )


@pytest.mark.parametrize("count", [3, 6, 10, 12])
def test_a_slide_the_layout_can_hold_stays_silent(count: int) -> None:
    assert lint_slides_data(_many(count)) == []


@pytest.mark.parametrize("count", [13, 16, 20])
def test_a_slide_with_more_body_than_fits_says_so(count: int) -> None:
    """13 is where the band starts being tiled and the lines touch."""
    (warning,) = lint_slides_data(_many(count))
    assert f"{count} bullets do not fit" in warning
    assert "Split this across two slides" in warning


def test_long_bullets_can_overfill_a_slide_of_few_of_them() -> None:
    """Rows, not bullets: five bullets that each wrap four times do not fit."""
    warnings = lint_slides_data(_structured(title="T", bullets=["아주 긴 불릿 " * 30] * 5))
    assert any("do not fit" in warning for warning in warnings)


def test_text_smaller_than_the_slide_can_show_is_reported_once_per_slide() -> None:
    data = {"slides": [{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 5, "w": 40, "h": 5,
         "text": "tiny", "fontSize": 9},
        {"id": "b", "type": "text", "x": 5, "y": 15, "w": 40, "h": 5,
         "text": "small", "fontSize": 13},
        {"id": "c", "type": "text", "x": 5, "y": 25, "w": 40, "h": 5,
         "text": "fine", "fontSize": 24},
    ]}]}
    (warning,) = lint_slides_data(data)
    assert "2 text element(s) below 14px" in warning
    assert "9px, 13px" in warning


@pytest.mark.parametrize("size", [14, 19, 24, 30, 38, 48])
def test_text_at_or_above_the_floor_stays_silent(size: int) -> None:
    data = {"slides": [{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 5, "w": 40, "h": 5, "text": "hi", "fontSize": size}
    ]}]}
    assert lint_slides_data(data) == []


def test_a_shape_with_no_font_size_is_not_small_text() -> None:
    data = {"slides": [{"elements": [
        {"id": "s", "type": "shape", "shape": "rect", "x": 5, "y": 5, "w": 10, "h": 10}
    ]}]}
    assert lint_slides_data(data) == []


def test_derived_sizes_are_never_reported_as_too_small() -> None:
    """The layout's own sizes are ours; warning about them teaches nothing."""
    for count in (3, 8, 13, 30):
        assert not any("below 14px" in warning for warning in lint_slides_data(_many(count)))


def test_the_pitch_this_replaced_would_be_caught_now() -> None:
    """Nine bullets at the old fixed pitch ran to y+h = 108 without a word."""
    data = {"slides": [{"elements": [
        {"id": f"bul_{i}", "type": "text", "x": 8, "y": 28 + i * 9, "w": 84, "h": 8,
         "text": f"• item {i}", "fontSize": 20}
        for i in range(9)
    ]}]}
    warnings = lint_slides_data(data)
    assert any("off the page" in warning and "108" in warning for warning in warnings)


# --- the floor comes from the deck's own skin ----------------------------------------


def test_the_floor_can_be_the_decks_own_smallest_size() -> None:
    """An uploaded deck's 12px footnotes are the author's; only text set
    smaller than anything in the original is called out."""
    deck = {"slides": [{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 5, "w": 50, "h": 8,
         "text": "footnote", "fontSize": 12},
        {"id": "b", "type": "text", "x": 5, "y": 20, "w": 50, "h": 8,
         "text": "tiny", "fontSize": 9},
    ]}]}
    default = lint_slides_data(deck)
    assert any("2 text element(s) below 14px" in w for w in default)
    own = lint_slides_data(deck, min_text_px=12)
    assert len([w for w in own if "below 12px" in w]) == 1
    assert "9px" in next(w for w in own if "below 12px" in w)
    assert "smallest" in next(w for w in own if "below 12px" in w)


def test_write_canvas_reads_the_floor_from_the_template_skin() -> None:
    """The skin's smallest run size is the floor for the copy made from it."""
    import io as _io

    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches, Pt

    skin = pptx.Presentation()
    skin.slide_width, skin.slide_height = Inches(10), Inches(5.625)
    slide = skin.slides.add_slide(skin.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "footnote"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(9)  # 12px
    buf = _io.BytesIO()
    skin.save(buf)
    store = InMemoryCanvasStore()
    store.write_bytes("t1", "sources/brand.pptx", buf.getvalue(), "skin")
    result = _write(store, {"template": "sources/brand.pptx", "slides": [{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 5, "w": 50, "h": 8,
         "text": "note", "fontSize": 12},
    ]}]})
    assert "below" not in result, result
    result = _write(store, {"template": "sources/brand.pptx", "slides": [{"elements": [
        {"id": "a", "type": "text", "x": 5, "y": 5, "w": 50, "h": 8, "text": "note", "fontSize": 8},
    ]}]}, path="deck2.slides.json")
    assert "below 12px" in result and "smallest" in result


# --- text that wraps past its box ------------------------------------------------------


def test_a_title_that_wraps_past_its_box_is_called_out() -> None:
    """Seen in a run: a 88px title rewritten twice as long in a box one line
    tall. The JSON looked fine; the slide showed a cut-off second line."""
    data = {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 5, "y": 10, "w": 54, "h": 12,
         "text": "왜 지금 브레인크루 X 신한은행인가", "fontSize": 88}
    ]}]}
    found = [w for w in lint_slides_data(data) if "run past the box" in w]
    assert len(found) == 1 and "line(s)" in found[0] and "86.4px tall" in found[0]


def test_one_line_in_a_snug_box_stays_silent() -> None:
    data = {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 5, "y": 10, "w": 54, "h": 5,
         "text": "Short", "fontSize": 48}
    ]}]}
    assert not [w for w in lint_slides_data(data) if "run past" in w]


def test_the_skins_own_overhang_is_allowed() -> None:
    """The original's bleed box at 101 and footer at 105 are the author's."""
    data = {"slides": [{"elements": [
        {"id": "a", "type": "shape", "shape": "rect", "x": 0, "y": 0, "w": 101.1, "h": 10},
        {"id": "b", "type": "text", "x": 0, "y": 95, "w": 50, "h": 10.2,
         "text": "footer", "fontSize": 20},
    ]}]}
    assert len([w for w in lint_slides_data(data) if "off the page" in w]) == 2
    assert [w for w in lint_slides_data(data, max_overhang=5.2) if "off the page" in w] == []
    over = {"slides": [{"elements": [
        {"id": "c", "type": "shape", "shape": "rect", "x": 0, "y": 0, "w": 110, "h": 10},
    ]}]}
    assert len([w for w in lint_slides_data(over, max_overhang=5.2) if "off the page" in w]) == 1


def test_a_table_whose_rows_outgrow_its_box_is_reported() -> None:
    from langchain_canvas.layout_lint import lint_slides_data

    rows = [[f"row {i}", "a long cell of text that wraps in a narrow column"] for i in range(8)]
    tall = {
        "id": "t", "type": "table", "x": 5, "y": 10, "w": 50, "h": 10, "rows": rows, "fontSize": 18,
    }
    (warning,) = [w for w in lint_slides_data({"slides": [{"elements": [tall]}]}) if "table" in w]
    assert warning.startswith('slide 1, element "t": the table\'s 8 row(s) need about')
    assert "the box is 72px tall" in warning

    roomy = {**tall, "h": 90}
    assert [w for w in lint_slides_data({"slides": [{"elements": [roomy]}]}) if "table" in w] == []


def test_a_tables_small_font_counts_as_small_text() -> None:
    from langchain_canvas.layout_lint import lint_slides_data

    tiny = {
        "id": "t", "type": "table", "x": 5, "y": 10, "w": 50, "h": 50,
        "rows": [["a"]], "fontSize": 8,
    }
    (warning,) = lint_slides_data({"slides": [{"elements": [tiny]}]})
    assert "below" in warning and "8px" in warning


def test_a_table_without_rows_is_refused_and_a_cells_unknown_key_is_named() -> None:
    from langchain_canvas.layout_lint import blocking_deck_findings, lint_slides_data

    bare = {"id": "t", "type": "table", "x": 5, "y": 10, "w": 50, "h": 50}
    envelope = {"type": "slides", "data": {"slides": [{"elements": [bare]}]}}
    (finding,) = blocking_deck_findings(envelope)
    assert '"rows"' in finding

    odd = {**bare, "rows": [["a"]], "cells": [{"r": 0, "c": 0, "colour": "#fff"}]}
    warnings = lint_slides_data({"slides": [{"elements": [odd]}]})
    assert any('cell: "colour"' in w for w in warnings)


# --- autofit: a box that grows, type that shrinks ---------------------------------------


def test_a_growing_box_is_not_an_overflow() -> None:
    """The same title as above, in a box that grows with its text: nothing
    runs past anything, so the finding that told the model to shorten the
    words — which it did, six times, on the wrong slide — is gone."""
    data = {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 5, "y": 10, "w": 54, "h": 12, "autofit": "shape",
         "text": "왜 지금 브레인크루 X 신한은행인가", "fontSize": 88}
    ]}]}
    assert not [w for w in lint_slides_data(data) if "run past" in w]


def test_a_growing_box_that_reaches_off_the_page_is_named_with_its_grown_bottom() -> None:
    data = {"slides": [{"elements": [
        {"id": "body", "type": "text", "x": 5, "y": 80, "w": 30, "h": 5, "autofit": "shape",
         "text": "가나다라마바사아자차카타파하 " * 12, "fontSize": 24}
    ]}]}
    found = [w for w in lint_slides_data(data) if "grows with its text" in w]
    assert len(found) == 1
    assert "line(s)" in found[0] and "y + h = " in found[0] and "off the page" in found[0]


def test_a_growing_box_that_stays_on_the_page_is_silent() -> None:
    data = {"slides": [{"elements": [
        {"id": "body", "type": "text", "x": 5, "y": 10, "w": 60, "h": 5, "autofit": "shape",
         "text": "가나다라마바사아자차카타파하 " * 4, "fontSize": 24}
    ]}]}
    assert not [w for w in lint_slides_data(data) if "grows" in w or "run past" in w]


def test_shrinking_type_is_not_an_overflow_until_it_shrinks_below_readable() -> None:
    fits = {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 5, "y": 10, "w": 54, "h": 12, "autofit": "text",
         "text": "왜 지금 브레인크루 X 신한은행인가", "fontSize": 88}
    ]}]}
    assert not [w for w in lint_slides_data(fits) if "run past" in w or "shrinks" in w]
    tiny = {"slides": [{"elements": [
        {"id": "t", "type": "text", "x": 5, "y": 10, "w": 20, "h": 4, "autofit": "text",
         "text": "가나다라마바사아자차카타파하 " * 10, "fontSize": 24}
    ]}]}
    found = [w for w in lint_slides_data(tiny) if "shrinks to fit its box" in w]
    assert len(found) == 1 and "below the 14px" in found[0]


def test_a_growing_box_that_did_not_grow_is_named_once_by_the_off_page_check() -> None:
    """Seen on a skin: a footer at y + h = 105 that grows with its one line
    was named twice — as off the page, and as grown off the page."""
    data = {"slides": [{"elements": [
        {"id": "foot", "type": "text", "x": 5, "y": 100, "w": 30, "h": 5.2, "autofit": "shape",
         "text": "footer", "fontSize": 12}
    ]}]}
    found = [w for w in lint_slides_data(data) if "off the page" in w]
    assert len(found) == 1 and "grows" not in found[0]
