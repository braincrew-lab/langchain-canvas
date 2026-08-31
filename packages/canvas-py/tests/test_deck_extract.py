"""PPTX structural extraction and deterministic baseline slide HTML.

Reuses `test_pptx_import.py`'s fixture builders — the same synthetic decks
prove both the old percent-geometry reader and this dialect's extraction
read the same underlying shapes.
"""

from __future__ import annotations

import base64
import io
import zipfile
from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Inches

from langchain_canvas.converters import UnsafeArchiveError
from langchain_canvas.deck.baseline import baseline_slide_html
from langchain_canvas.deck.extract import PptxImportError, extract_slides, extracted_text
from langchain_canvas.deck.validate import ensure_text_equality, validate_slide_html

# A real 1x1 red PNG — small enough to inline, real enough for pptx to embed.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _deck(build: Any = None, *, width_in: float = 13.333, height_in: float = 7.5) -> bytes:
    """Presentation bytes with one blank slide, shaped by ``build``."""
    deck = Presentation()
    deck.slide_width = Inches(width_in)
    deck.slide_height = Inches(height_in)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    if build is not None:
        build(slide)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _textbox(slide: Any, text: str = "Hello", **kwargs: Any) -> Any:
    box = slide.shapes.add_textbox(
        kwargs.get("left", Inches(1)),
        kwargs.get("top", Inches(1)),
        kwargs.get("width", Inches(4)),
        kwargs.get("height", Inches(1)),
    )
    box.text_frame.text = text
    return box


def _bomb_zip(uncompressed_mb: int = 201) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("ppt/presentation.xml", b"\x00" * (uncompressed_mb * 1024 * 1024))
    return buf.getvalue()


def test_extract_slides_returns_text_geometry_assets() -> None:
    """One slide with text, a shape, and a picture extracts all three."""

    def build(slide: Any) -> None:
        _textbox(slide, "Hello world", left=Inches(1), top=Inches(1))
        slide.shapes.add_shape(1, Inches(2), Inches(2), Inches(1), Inches(1))  # rectangle
        slide.shapes.add_picture(io.BytesIO(_PNG), Inches(3), Inches(3), Inches(1), Inches(1))

    extractions = extract_slides(_deck(build), path="sources/deck.pptx")

    assert len(extractions) == 1
    slide = extractions[0]
    assert slide.index == 0
    assert [run.text for run in slide.texts] == ["Hello world"]
    assert len(slide.shapes) == 1
    assert slide.shapes[0].kind == "rect"
    assert len(slide.images) == 1
    assert slide.images[0].data == _PNG
    assert slide.images[0].ext == "png"
    assert len(slide.images[0].sha) == 64  # sha256 hex digest


def test_selected_pages_preserve_original_index_and_default_all() -> None:
    """`pages` extracts only the selected 1-based slides, index stays 0-based."""

    def build_deck() -> bytes:
        deck = Presentation()
        deck.slide_width = Inches(13.333)
        deck.slide_height = Inches(7.5)
        for label in ("Slide 1", "Slide 2", "Slide 3"):
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            _textbox(slide, label)
        buffer = io.BytesIO()
        deck.save(buffer)
        return buffer.getvalue()

    data = build_deck()

    full = extract_slides(data, path="sources/deck.pptx")
    assert [run.text for slide in full for run in slide.texts] == [
        "Slide 1",
        "Slide 2",
        "Slide 3",
    ]

    selected = extract_slides(data, path="sources/deck.pptx", pages=[3, 1])
    assert [slide.index for slide in selected] == [2, 0]
    assert [run.text for slide in selected for run in slide.texts] == [
        "Slide 3",
        "Slide 1",
    ]


def test_selected_pages_reject_empty_duplicate_and_out_of_range() -> None:
    data = _deck(lambda s: _textbox(s, "Only slide"))

    with pytest.raises(PptxImportError, match="non-empty"):
        extract_slides(data, path="sources/deck.pptx", pages=[])

    with pytest.raises(PptxImportError, match="unique"):
        extract_slides(data, path="sources/deck.pptx", pages=[1, 1])

    with pytest.raises(PptxImportError, match="1..1"):
        extract_slides(data, path="sources/deck.pptx", pages=[2])


def test_extract_rejects_zip_bomb() -> None:
    """A crafted decompression bomb is refused before any parser runs."""
    with pytest.raises(UnsafeArchiveError, match="unpacks to"):
        extract_slides(_bomb_zip(), path="sources/bomb.pptx")


def test_baseline_html_preserves_extracted_text_exactly() -> None:
    """Text carried into the baseline HTML passes the text-integrity gate."""
    extraction = extract_slides(
        _deck(lambda s: _textbox(s, "Quarterly Revenue Report")), path="sources/deck.pptx"
    )[0]

    html = baseline_slide_html(extraction, slide_id="s1", ratio="16:9")

    ensure_text_equality(extracted_text(extraction), html)  # no raise


def test_baseline_html_is_valid_dialect() -> None:
    """Baseline markup is rooted at <section class="slide"> with unique node ids."""

    def build(slide: Any) -> None:
        _textbox(slide, "Title", left=Inches(1), top=Inches(1))
        slide.shapes.add_shape(1, Inches(2), Inches(2), Inches(1), Inches(1))
        slide.shapes.add_picture(io.BytesIO(_PNG), Inches(3), Inches(3), Inches(1), Inches(1))

    extraction = extract_slides(_deck(build), path="sources/deck.pptx")[0]

    html = baseline_slide_html(extraction, slide_id="s1", ratio="16:9")

    assert validate_slide_html(html, slide_id="s1") == []
