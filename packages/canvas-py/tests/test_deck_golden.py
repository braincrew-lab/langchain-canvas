"""Golden-corpus integration test: import -> edit -> export, end to end.

Exercises the same steps ``open_deck_for_editing`` (tools.py) chains together
for a real uploaded ``.pptx`` — :func:`extract_slides` ->
:func:`baseline_slide_html` per slide -> :func:`serialize_deck` — against the
Phase 1 golden fixtures under ``tests/fixtures/golden/*.pptx`` (see
``generate_fixtures.py`` for how each one isolates a single PPTX feature).
The edit and export steps then run the deterministic pipeline a real edit
session would: a slide-scoped :func:`patch_slide`, then
:class:`DeckPptxExporter` with ``inline_canvas_assets``/``inline_deck_skin``
standing in for the export tool's own asset- and skin-inlining steps. No LLM
call is involved anywhere — ``convert_slide``'s own model step is faked
separately in ``apps/server``'s tests.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langchain_canvas.assets import inline_canvas_assets
from langchain_canvas.deck.baseline import baseline_slide_html
from langchain_canvas.deck.export import DeckPptxExporter
from langchain_canvas.deck.extract import SlideExtraction, extract_slides, extracted_text
from langchain_canvas.deck.model import Deck, SlideTemplate, parse_deck, patch_slide, serialize_deck
from langchain_canvas.deck.validate import ensure_text_equality, validate_deck
from langchain_canvas.store import InMemoryCanvasStore
from langchain_canvas.tools import inline_deck_skin

FIXTURES_DIR = Path(__file__).parent / "fixtures" / "golden"

# Fixed slide order for the combined multi-slide deck the edit/export test
# builds — one slide per golden fixture, in the same order
# ``generate_fixtures.py`` declares them.
FIXTURE_NAMES = (
    "typography.pptx",
    "image.pptx",
    "group.pptx",
    "gradient.pptx",
    "table.pptx",
    "rotation.pptx",
    "unsupported.pptx",
)


def _fixture_bytes(name: str) -> bytes:
    return (FIXTURES_DIR / name).read_bytes()


def _slide_id(index: int) -> str:
    return f"slide-{index + 1:03d}"


def _import_one(name: str, index: int) -> tuple[SlideExtraction, SlideTemplate]:
    """One golden fixture, imported the way ``open_deck_for_editing`` would.

    Each golden fixture is a single-slide presentation (see
    ``generate_fixtures.py``), so this always reads exactly one
    :class:`SlideExtraction` out of it.
    """
    extractions = extract_slides(_fixture_bytes(name), path=name)
    assert len(extractions) == 1, f"{name} fixture is expected to be a single slide"
    extraction = extractions[0]
    slide_id = _slide_id(index)
    body_html = baseline_slide_html(extraction, slide_id=slide_id, ratio="16:9")
    slide = SlideTemplate(slide_id=slide_id, title=None, style_css="", body_html=body_html)
    return extraction, slide


def _combined_deck_html(*, source: str | None = None) -> str:
    slides = [_import_one(name, i)[1] for i, name in enumerate(FIXTURE_NAMES)]
    return serialize_deck(
        Deck(title="Golden Corpus", ratio="16:9", source=source, slides=slides)
    )


def _write_image_assets(store: InMemoryCanvasStore, canvas_id: str) -> None:
    """Persist every golden fixture's extracted picture under ``assets/``.

    Mirrors what ``open_deck_for_editing`` itself does at import time (see
    ``tools.py``'s ``open_deck_for_editing``): each :class:`ImageAsset` is
    written to ``assets/{sha}.{ext}`` so the baseline HTML's relative
    ``src="assets/..."`` reference resolves before export.
    """
    for i, name in enumerate(FIXTURE_NAMES):
        extraction, _ = _import_one(name, i)
        for image in extraction.images:
            store.write_bytes(
                canvas_id,
                f"assets/{image.sha}.{image.ext}",
                image.data,
                f"Asset for {name}",
                actor="agent",
            )


@pytest.mark.parametrize("fixture_name", FIXTURE_NAMES)
def test_golden_import_produces_a_valid_text_preserving_deck(fixture_name: str) -> None:
    """Every golden fixture imports into a structurally valid, text-preserving slide.

    Mirrors ``open_deck_for_editing``'s own pipeline: extract -> baseline
    render -> serialize -> (implicitly) parse back by the caller that reads
    the deck. ``ensure_text_equality`` — called on the fixture's extracted
    text against the rendered slide's own body — must not raise, matching
    the hard gate a real save enforces.
    """
    extraction, slide = _import_one(fixture_name, index=0)
    deck_html = serialize_deck(
        Deck(title="Golden", ratio="16:9", source=fixture_name, slides=[slide])
    )

    # Round-trips through the canonical dialect parser.
    reparsed = parse_deck(deck_html)
    assert [s.slide_id for s in reparsed.slides] == [slide.slide_id]
    assert reparsed.slides[0].body_html == slide.body_html

    # Structurally valid: <section class="slide"> root, unique node ids,
    # correct dialect version — zero issues for a fresh import.
    issues = validate_deck(deck_html)
    assert issues == []

    # No source text was dropped or reworded rendering the baseline HTML.
    # (Golden fixtures whose content lives in a table/group/chart are
    # intentionally skipped by the extractor — see extract.py's shape
    # loop — so their extracted-text list is legitimately empty here.)
    ensure_text_equality(extracted_text(extraction), slide.body_html)


def test_golden_edit_changes_only_the_target_slides_bytes() -> None:
    """A slide-scoped ``patch_slide`` on the combined golden deck touches one slide only."""
    original_html = _combined_deck_html()
    original_deck = parse_deck(original_html)
    assert [s.slide_id for s in original_deck.slides] == [
        _slide_id(i) for i in range(len(FIXTURE_NAMES))
    ]

    target_id = _slide_id(0)  # typography.pptx's slide — the one with real text
    target_slide = next(s for s in original_deck.slides if s.slide_id == target_id)
    assert "Quarterly Review" in target_slide.body_html

    edited_body = target_slide.body_html.replace("Quarterly Review", "Quarterly Review — Updated")
    assert edited_body != target_slide.body_html
    new_template = f'<template data-slide-id="{target_id}">{edited_body}</template>'

    edited_html = patch_slide(original_html, target_id, new_template)
    edited_deck = parse_deck(edited_html)

    by_id_before = {s.slide_id: s for s in original_deck.slides}
    by_id_after = {s.slide_id: s for s in edited_deck.slides}

    assert by_id_after[target_id].body_html == edited_body
    assert "Updated" in by_id_after[target_id].body_html
    for slide_id in by_id_before:
        if slide_id == target_id:
            continue
        assert by_id_after[slide_id].body_html == by_id_before[slide_id].body_html, (
            f"{slide_id} bytes changed even though only {target_id} was patched"
        )

    # The edit did not corrupt deck-level structure.
    assert validate_deck(edited_html) == []


def test_golden_export_round_trips_through_pptx() -> None:
    """The edited combined deck exports to a real, reopenable ``.pptx``.

    ``inline_canvas_assets`` and ``inline_deck_skin`` stand in for
    ``export_canvas``'s own inlining steps (tools.py's ``create_export_tool``,
    same order): the image fixture's relative ``assets/...`` reference and
    the combined deck's ``lcx:source`` (pointing at ``typography.pptx``
    stored under ``sources/`` in an in-memory store) both become ``data:``
    URIs before the exporter runs. No ``RenderSlideAdapter`` is injected, so
    any node the exporter cannot map natively lands in ``degraded_nodes``
    instead of raising.
    """
    store = InMemoryCanvasStore()
    canvas_id = "golden-corpus"
    skin_bytes = _fixture_bytes("typography.pptx")
    store.write_bytes(canvas_id, "sources/typography.pptx", skin_bytes, "Upload", actor="human")
    _write_image_assets(store, canvas_id)

    original_html = _combined_deck_html(source="sources/typography.pptx")
    target_id = _slide_id(0)
    original_deck = parse_deck(original_html)
    target_slide = next(s for s in original_deck.slides if s.slide_id == target_id)
    edited_body = target_slide.body_html.replace("Quarterly Review", "Quarterly Review — Updated")
    new_template = f'<template data-slide-id="{target_id}">{edited_body}</template>'
    edited_html = patch_slide(original_html, target_id, new_template)

    # Same order `create_export_tool`'s `export_canvas` uses for a
    # `.slides.html` deck: relative asset references first, then the
    # `lcx:source` skin reference.
    inlined = inline_canvas_assets(edited_html, store, canvas_id)
    inlined = inline_deck_skin(inlined, store, canvas_id)
    inlined_deck = parse_deck(inlined)
    assert inlined_deck.source is not None
    assert inlined_deck.source.startswith("data:")  # skin reference was inlined, not left as a path

    exporter = DeckPptxExporter()
    result = exporter.export(inlined, path="golden-corpus.slides.html")

    reopened = Presentation(io.BytesIO(result.data))
    assert len(reopened.slides) == len(FIXTURE_NAMES)

    # Representative text survived the round trip: the edited typography
    # slide's updated title, and its untouched second paragraph.
    first_slide_text = "\n".join(
        shape.text_frame.text for shape in reopened.slides[0].shapes if shape.has_text_frame
    )
    assert "Quarterly Review — Updated" in first_slide_text
    assert "Revenue grew 12% year over year." in first_slide_text

    # The image fixture's picture made it across as a native picture shape.
    image_slide_shapes = list(reopened.slides[1].shapes)
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in image_slide_shapes)

    # None of the golden fixtures currently produce a `data-lcx-fallback`
    # node: `pptx_import._drawing` classifies FREEFORM shapes (what
    # `unsupported.pptx` uses to stand in for SmartArt/chart content, per
    # generate_fixtures.py) as a plain "rect", so the raster-degrade path
    # this corpus was meant to exercise is not actually reached today. This
    # is the current, verified pipeline behavior — asserted here rather
    # than the fixture's aspirational docstring — and is reported as a
    # concern below rather than "fixed" (fixing it would mean editing
    # `pptx_import.py`, out of scope for this test-only task).
    assert exporter.degraded_nodes == []
