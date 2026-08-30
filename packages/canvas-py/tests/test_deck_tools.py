"""The SDK deck tools: `open_deck_for_editing`, `read/edit/list_deck_slide`.

`open_deck_for_editing` used to copy an uploaded `.pptx` into a `.slides.json`
envelope (see `test_pptx_import.py`'s transitioned cases). This module covers
its replacement — `deck.extract` + `deck.baseline` straight into a
`.slides.html` deck — plus the slide-granular read/edit/list tools that make
that deck editable one slide at a time.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from langchain_canvas.deck import parse_deck
from langchain_canvas.store import InMemoryCanvasStore
from langchain_canvas.tools import create_deck_tools


def _deck_bytes(build: Any = None, *, width_in: float = 13.333, height_in: float = 7.5) -> bytes:
    """Presentation bytes with one blank slide, shaped by ``build``."""
    deck = Presentation()
    deck.slide_width = Inches(width_in)
    deck.slide_height = Inches(height_in)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    if build is not None:
        build(slide)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _textbox(slide: Any, text: str = "Hello") -> Any:
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    return box


def _store_with_deck(deck: bytes = b"") -> InMemoryCanvasStore:
    store = InMemoryCanvasStore()
    store.write_bytes(
        "t1",
        "sources/deck.pptx",
        deck or _deck_bytes(lambda s: _textbox(s, "Hi")),
        "Upload",
        actor="human",
    )
    return store


def _tools(store: InMemoryCanvasStore) -> dict[str, Any]:
    return {t.name: t for t in create_deck_tools(store)}


def _run(tool: Any, **kwargs: Any) -> str:
    """Through the same runtime stub the other tool tests use."""
    from test_tools import _runtime

    return tool.func(runtime=_runtime(thread_id="t1"), **kwargs)


# --- open_deck_for_editing --------------------------------------------------------


def test_open_deck_for_editing_writes_slides_html_assets_and_source_meta() -> None:
    """The copy is a `.slides.html` deck, not a `.slides.json` envelope: real
    boxes per extracted shape, pictures under `assets/`, and the original
    named as the deck's `lcx:source` meta rather than a JSON `template` key."""
    png = bytes.fromhex(
        "89504e470d0a1a0a0000000d494844520000000100000001080600000"
        "01f15c4890000000a49444154789c6360000002000155ff8f350000000049454e44ae426082"
    )

    def build(slide: Any) -> None:
        _textbox(slide, "Hello world")
        slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(3), Inches(2), Inches(2))

    store = _store_with_deck(_deck_bytes(build))
    reply = _run(_tools(store)["open_deck_for_editing"], source="sources/deck.pptx")
    assert reply.startswith("Copied "), reply

    paths = [f.path for f in store.list_files("t1")]
    assert "deck.slides.html" in paths
    assert "sources/deck.pptx" in paths  # the original is untouched
    assert any(p.startswith("assets/") for p in paths), paths

    content = store.read("t1", "deck.slides.html").content
    deck = parse_deck(content)
    assert deck.source == "sources/deck.pptx"
    assert len(deck.slides) == 1
    assert "Hello world" in deck.slides[0].body_html
    assert any("assets/" in slide.body_html for slide in deck.slides)


def test_open_deck_for_editing_broadcasts_deck_create_and_slide_status() -> None:
    from test_tools import _Runtime

    store = _store_with_deck()
    sent: list[dict] = []
    runtime = _Runtime(config={"configurable": {"thread_id": "t1"}})
    object.__setattr__(runtime, "stream_writer", sent.append)
    _tools(store)["open_deck_for_editing"].func(runtime=runtime, source="sources/deck.pptx")

    created = [e for e in sent if e.get("type") == "canvas.create"]
    assert created, sent
    assert created[0]["artifact"]["type"] == "slides"
    assert created[0]["artifact"]["meta"]["kind"] == "deck"

    statuses = [e for e in sent if e.get("type") == "canvas.slide_status"]
    assert statuses, sent
    assert all(s["stage"] == "complete" for s in statuses)


def test_copying_the_same_deck_twice_is_refused_rather_than_overwriting() -> None:
    store = _store_with_deck()
    tool = _tools(store)["open_deck_for_editing"]
    _run(tool, source="sources/deck.pptx")
    again = _run(tool, source="sources/deck.pptx")
    assert "already on the canvas" in again


def test_a_file_that_is_not_a_deck_is_pointed_somewhere_useful() -> None:
    store = _store_with_deck()
    reply = _run(_tools(store)["open_deck_for_editing"], source="notes.md")
    assert reply.startswith("Error: this opens .pptx files")


# --- read/edit/list_deck_slide ------------------------------------------------


def _opened_deck(store: InMemoryCanvasStore) -> str:
    """`deck.slides.html`'s path after `open_deck_for_editing` copies it out."""
    _run(_tools(store)["open_deck_for_editing"], source="sources/deck.pptx")
    return "deck.slides.html"


def test_list_deck_slides_lists_every_slide_id() -> None:
    store = _store_with_deck()
    path = _opened_deck(store)
    out = _run(_tools(store)["list_deck_slides"], path=path)
    deck = parse_deck(store.read("t1", path).content)
    for slide in deck.slides:
        assert slide.slide_id in out


def test_read_deck_slide_returns_the_slide_template_and_revision() -> None:
    store = _store_with_deck()
    path = _opened_deck(store)
    deck = parse_deck(store.read("t1", path).content)
    slide_id = deck.slides[0].slide_id
    out = _run(_tools(store)["read_deck_slide"], path=path, slide_id=slide_id)
    assert "revision:" in out
    assert f'data-slide-id="{slide_id}"' in out
    assert "<template" in out


def test_edit_deck_slide_rejects_stale_revision() -> None:
    store = _store_with_deck()
    path = _opened_deck(store)
    deck = parse_deck(store.read("t1", path).content)
    slide_id = deck.slides[0].slide_id
    new_fragment = (
        f'<template data-slide-id="{slide_id}"><section class="slide"></section></template>'
    )

    # A stale revision must be rejected up front, the same way edit_canvas
    # rejects one — never silently overwritten.
    reply = _run(
        _tools(store)["edit_deck_slide"],
        path=path,
        slide_id=slide_id,
        template_html=new_fragment,
        revision="not-a-real-revision",
    )
    assert reply.startswith("Error:")
    assert "revision" in reply.lower()

    unchanged = parse_deck(store.read("t1", path).content).slides[0]
    assert unchanged.slide_id == slide_id
    assert unchanged.body_html == deck.slides[0].body_html


def test_edit_deck_slide_sanitizes_before_write() -> None:
    store = _store_with_deck()
    path = _opened_deck(store)
    got = store.read("t1", path)
    deck = parse_deck(got.content)
    slide_id = deck.slides[0].slide_id
    malicious_fragment = (
        f'<template data-slide-id="{slide_id}">'
        '<section class="slide"><div data-node-id="node-1">safe'
        "<script>alert(1)</script></div></section></template>"
    )

    reply = _run(
        _tools(store)["edit_deck_slide"],
        path=path,
        slide_id=slide_id,
        template_html=malicious_fragment,
        revision=got.revision,
    )
    assert reply.startswith("Edited "), reply

    saved = store.read("t1", path).content
    assert "<script>" not in saved
    assert "alert(1)" not in saved
    saved_slide = parse_deck(saved).slides[0]
    assert "safe" in saved_slide.body_html


def test_edit_deck_slide_safe_with_duplicate_templates() -> None:
    """Two slides sharing byte-identical markup — editing one must not touch
    the other, because `patch_slide` matches on `slide_id`, never content
    (this is exactly why `edit_canvas`'s `store.edit` is unsafe here)."""
    store = InMemoryCanvasStore()
    deck_html = (
        '<!DOCTYPE html><html data-lcx-dialect="1" data-ratio="16:9">'
        "<head><title>Deck</title></head><body>"
        '<template data-slide-id="slide-001"><section class="slide">'
        '<div data-node-id="node-1">Same text</div></section></template>'
        '<template data-slide-id="slide-002"><section class="slide">'
        '<div data-node-id="node-1">Same text</div></section></template>'
        "</body></html>\n"
    )
    commit = store.write("t1", "deck.slides.html", deck_html, "Seed")

    new_fragment = (
        '<template data-slide-id="slide-002"><section class="slide">'
        '<div data-node-id="node-1">Changed text</div></section></template>'
    )
    reply = _run(
        _tools(store)["edit_deck_slide"],
        path="deck.slides.html",
        slide_id="slide-002",
        template_html=new_fragment,
        revision=commit.revision,
    )
    assert reply.startswith("Edited "), reply

    deck = parse_deck(store.read("t1", "deck.slides.html").content)
    by_id = {slide.slide_id: slide for slide in deck.slides}
    assert "Changed text" in by_id["slide-002"].body_html
    assert "Same text" in by_id["slide-001"].body_html
    assert "Changed text" not in by_id["slide-001"].body_html


def test_edit_deck_slide_broadcasts_slide_patch() -> None:
    from test_tools import _Runtime

    store = _store_with_deck()
    path = _opened_deck(store)
    got = store.read("t1", path)
    deck = parse_deck(got.content)
    slide_id = deck.slides[0].slide_id
    new_fragment = (
        f'<template data-slide-id="{slide_id}">'
        '<section class="slide"><div data-node-id="node-1">Updated</div></section>'
        "</template>"
    )

    sent: list[dict] = []
    runtime = _Runtime(config={"configurable": {"thread_id": "t1"}})
    object.__setattr__(runtime, "stream_writer", sent.append)
    _tools(store)["edit_deck_slide"].func(
        runtime=runtime,
        path=path,
        slide_id=slide_id,
        template_html=new_fragment,
        revision=got.revision,
    )

    patches = [e for e in sent if e.get("type") == "canvas.slide_patch"]
    assert patches, sent
    assert patches[0]["slideId"] == slide_id
    assert "Updated" in patches[0]["templateHtml"]
    assert any(e.get("type") == "canvas.commit" for e in sent)
