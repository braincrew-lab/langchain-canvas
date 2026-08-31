"""Render-free census of original PDF/PPTX pages — never renders, decodes, or calls a model.

Fixture PDFs are built directly with ``pypdfium2``'s low-level object API
(the same technique ``app.agent.pdf_source`` uses to read positioned text
without rasterizing), so tests exercise real bounding boxes instead of a
hand-rolled PDF byte string.
"""

from __future__ import annotations

import ctypes
import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Inches

from langchain_canvas.deck.source_inventory import (
    OversizedPageError,
    SourceLimits,
    inspect_source_pages,
)

pdfium = pytest.importorskip("pypdfium2")
from pypdfium2 import raw  # noqa: E402


def _text_object(document: Any, text: str, x: float, y: float, size: float = 18.0) -> Any:
    obj = raw.FPDFPageObj_NewTextObj(document.raw, b"Helvetica", size)
    encoded = (text + "\x00").encode("utf-16-le")
    buffer = (ctypes.c_ushort * (len(encoded) // 2)).from_buffer_copy(encoded)
    raw.FPDFText_SetText(obj, buffer)
    raw.FPDFPageObj_Transform(obj, 1, 0, 0, 1, x, y)
    return obj


def _pdf(
    pages: list[list[tuple[str, float, float]]], size: tuple[float, float] = (612.0, 792.0)
) -> bytes:
    """Text-bearing PDF bytes; each page is a list of (text, x, y) placements."""
    document = pdfium.PdfDocument.new()
    for items in pages:
        page = document.new_page(*size)
        for text, x, y in items:
            raw.FPDFPage_InsertObject(page.raw, _text_object(document, text, x, y))
        page.gen_content()
    out = io.BytesIO()
    document.save(out)
    document.close()
    return out.getvalue()


def _blank_pdf(pages: int, size: tuple[float, float] = (612.0, 792.0)) -> bytes:
    document = pdfium.PdfDocument.new()
    for _ in range(pages):
        document.new_page(*size)
    out = io.BytesIO()
    document.save(out)
    document.close()
    return out.getvalue()


def _pptx(build) -> bytes:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    build(slide)
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


# --- never renders, decodes, or calls a model -------------------------------------


def test_inventory_never_renders_decodes_images_or_calls_model(monkeypatch):
    def _boom(*_args, **_kwargs):
        raise AssertionError("census must not render or decode")

    monkeypatch.setattr(pdfium.PdfPage, "render", _boom)
    monkeypatch.setattr(pdfium.PdfImage, "get_bitmap", _boom)

    data = _pdf([[("Title", 100, 700), ("Body text here", 100, 600)]])
    result = inspect_source_pages(data, path="sources/deck.pdf", start_page=1, limit=10)

    assert result.fingerprint.page_count == 1
    assert result.pages[0].has_text is True


# --- native census sees what extraction would skip --------------------------------


def test_native_census_reports_group_before_extraction():
    def build(slide):
        first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        first.text_frame.text = "a"
        second = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
        second.text_frame.text = "b"
        slide.shapes.add_group_shape([first, second])
        slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(3), Inches(2))

    data = _pptx(build)
    result = inspect_source_pages(data, path="sources/deck.pptx", start_page=1, limit=1)

    page = result.pages[0]
    assert "group" in page.capability_issues
    assert "native_table" in page.capability_issues
    assert page.object_kind_counts.get("group") == 1
    assert page.object_kind_counts.get("table") == 1


# --- page-limited census + hash-pinned pagination ----------------------------------


def test_scan_limits_and_hash_pinned_cursor():
    pages = [[(f"Slide {n}", 100, 700)] for n in range(1, 61)]
    data = _pdf(pages)

    first = inspect_source_pages(data, path="sources/deck.pdf", start_page=1, limit=50)
    assert len(first.pages) == 50
    assert first.scope_complete is False
    assert first.next_start_page == 51

    second = inspect_source_pages(
        data, path="sources/deck.pdf", start_page=first.next_start_page, limit=50
    )
    assert len(second.pages) == 10
    assert second.scope_complete is True
    assert second.next_start_page is None
    assert second.fingerprint.sha256 == first.fingerprint.sha256

    overwritten = _pdf([[("Different content", 100, 700)]])
    restarted = inspect_source_pages(overwritten, path="sources/deck.pdf", start_page=1, limit=50)
    assert restarted.fingerprint.sha256 != first.fingerprint.sha256


# --- scanned page reports unknown style, not a guess -------------------------------


def test_scanned_page_reports_unknown_style():
    data = _blank_pdf(1)
    result = inspect_source_pages(data, path="sources/deck.pdf", start_page=1, limit=1)

    page = result.pages[0]
    assert page.has_text is False
    assert page.needs_visual_inspection is True
    assert page.text_boxes == ()


# --- oversized single page rejected before render ----------------------------------


def test_oversized_single_page_rejected_before_render(monkeypatch):
    def _boom(*_args, **_kwargs):
        raise AssertionError("an oversized page must be rejected before any render")

    monkeypatch.setattr(pdfium.PdfPage, "render", _boom)

    data = _pdf([[("Title", 100, 2900)]], size=(612.0, 3000.0))
    limits = SourceLimits(max_pdf_page_side_pt=2000.0)

    with pytest.raises(OversizedPageError):
        inspect_source_pages(data, path="sources/deck.pdf", start_page=1, limit=1, limits=limits)
