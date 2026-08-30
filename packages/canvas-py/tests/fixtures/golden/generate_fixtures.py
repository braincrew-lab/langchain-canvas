"""One-off generator for the golden ``.pptx`` fixture corpus.

Not a test — run manually (``uv run --extra dev python
tests/fixtures/golden/generate_fixtures.py``) whenever a fixture needs to be
regenerated. The committed ``.pptx`` files are the actual fixtures Task 10's
``test_deck_golden.py`` (import -> edit -> export integration test) reads;
this script only documents how each one was built.

Each fixture isolates one PPTX feature the extractor/exporter pipeline must
handle: typography, images, grouped shapes, gradient fills, tables, rotated
shapes, and unsupported content that must degrade gracefully.
"""

from __future__ import annotations

import base64
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

FIXTURES_DIR = os.path.dirname(os.path.abspath(__file__))

# A real 1x1 red PNG, small enough to inline.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _blank_deck() -> Presentation:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    return deck


def _blank_slide(deck: Presentation):
    return deck.slides.add_slide(deck.slide_layouts[6])


def build_typography() -> Presentation:
    """A slide exercising heading/body text, bold/italic runs, and alignment."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Quarterly Review"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
    body_frame = body_box.text_frame
    body_frame.text = "Revenue grew 12% year over year."
    run = body_frame.paragraphs[0].add_run()
    run.text = " Highlights follow."
    run.font.italic = True
    return deck


def build_image() -> Presentation:
    """A slide with one inline picture."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    import io

    slide.shapes.add_picture(
        io.BytesIO(_PNG), Inches(1), Inches(1), width=Inches(4), height=Inches(3)
    )
    return deck


def build_group() -> Presentation:
    """A slide with two shapes combined into one group."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(1))
    slide.shapes.add_group_shape([rect, oval])
    return deck


def build_gradient() -> Presentation:
    """A slide with a gradient-filled rectangle."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(6), Inches(3))
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
    stops[1].color.rgb = RGBColor(0xFB, 0xBF, 0x24)
    return deck


def build_table() -> Presentation:
    """A slide with a 2x2 table."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    graphic_frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    table = graphic_frame.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "APAC"
    table.cell(1, 1).text = "$1.2M"
    return deck


def build_rotation() -> Presentation:
    """A slide with a shape rotated 30 degrees."""
    deck = _blank_deck()
    slide = _blank_slide(deck)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(3), Inches(1.5)
    )
    shape.rotation = 30.0
    return deck


def build_unsupported() -> Presentation:
    """A slide with content the v1 exporter cannot map natively (SmartArt/chart placeholder).

    python-pptx cannot author SmartArt directly, so this fixture stands in a
    freeform (custom geometry) autoshape at an odd EMU offset — a shape kind
    the extractor's known-shape mapping does not cover, exercising the
    ``data-lcx-fallback="raster"`` degrade path.
    """
    deck = _blank_deck()
    slide = _blank_slide(deck)
    freeform_builder = slide.shapes.build_freeform(Emu(914400), Emu(914400))
    freeform_builder.add_line_segments(
        [(Emu(914400 * 3), Emu(914400)), (Emu(914400 * 2), Emu(914400 * 3))],
        close=True,
    )
    freeform_builder.convert_to_shape()
    return deck


FIXTURES = {
    "typography.pptx": build_typography,
    "image.pptx": build_image,
    "group.pptx": build_group,
    "gradient.pptx": build_gradient,
    "table.pptx": build_table,
    "rotation.pptx": build_rotation,
    "unsupported.pptx": build_unsupported,
}


def main() -> None:
    for filename, builder in FIXTURES.items():
        deck = builder()
        deck.save(os.path.join(FIXTURES_DIR, filename))
        print(f"wrote {filename}")


if __name__ == "__main__":
    main()
