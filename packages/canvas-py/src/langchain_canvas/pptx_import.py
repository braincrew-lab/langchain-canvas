"""Read an uploaded ``.pptx`` into the deck model the canvas renders.

The twin of :class:`~langchain_canvas.exporters.SlidesPptxExporter`: that
writes a deck out as real shapes, this reads one back in. A deck opened this
way lands on the canvas as editable slides instead of a file card, and the
original stays under ``sources/`` untouched.

The original is also named as the deck's ``template``. That is what carries
everything this model has no field for — masters, layouts, themes, logos,
headers — because the pptx export builds on the skin it points at. So the
reader does not have to capture the whole file to give it back intact; it
captures what a person edits, and the skin restores the rest.

What comes across: slide size, per-shape position and size as percent
geometry, text with its size / bold / colour / alignment, pictures as data
URIs, and rectangles, ellipses and lines.

Honest limits — the reader drops these rather than guessing at them, and each
one stays visible in the original and in the exported file through the skin:

* tables and charts (the deck model has no element type for them)
* grouped shapes
* a text box whose runs disagree takes the first run's formatting, since one
  element carries one set of it
* pictures that are not png / jpeg / gif (the exporter reads no others)
* gradient and theme-referenced backgrounds (the field holds one colour;
  a picture background is laid in as the bottom element instead)
* animations, transitions, SmartArt
"""

from __future__ import annotations

import base64
import json
from collections import Counter
from io import BytesIO
from typing import Any

from .preset_colours import preset_colour

# Pixels per inch used to turn point sizes into the canvas' font scale.
_EMU_PER_INCH = 914400

# Picture formats the pptx exporter can write back out; anything else is
# dropped rather than embedded as bytes no exporter will read.
_IMAGE_TYPES = {"png": "png", "jpg": "jpeg", "jpeg": "jpeg", "gif": "gif"}

# Longest edge a picture background is resampled to before it is inlined.
_BACKGROUND_MAX_PX = 1920

# PowerPoint's default line weight, used when a connector declares none.
_DEFAULT_STROKE_EMU = 12700  # 1pt

# The deck model measures type and spacing in px on a 1280px-wide slide;
# PowerPoint states them in points. Storing a point value in a pixel field
# renders every size a quarter too small and moves every line break.
_PT_TO_PX = 4 / 3

# Autoshape names that map onto the three shapes the deck model draws.
_ELLIPSE_NAMES = ("OVAL", "ELLIPSE", "CIRCLE")
_LINE_NAMES = ("LINE", "STRAIGHT_CONNECTOR", "CONNECTOR")


class PptxImportError(ValueError):
    """The bytes are not a presentation this reader can open."""


class DeckBaseline:
    """What the original deck already does, as the floor the copy is judged by.

    ``smallest_text_px`` — the smallest run size the author printed; a check
    that calls the original's 9pt footnotes an error is one the model learns
    to ignore along with the findings that matter. ``max_overhang`` — how far
    past the page edge (in percent points) the original's own shapes reach;
    decks routinely carry a bleed box at 101 or a footer at 105, and flagging
    those on every save buried the one real overflow among them.
    ``overflow`` — the boxes whose own text already needs more height than
    they have, keyed by rounded percent geometry, valued by the need/have
    ratio. A copy inherits those overflows untouched, and a finding the
    editor cannot bring to zero is a finding the editor learns to ignore.
    """

    def __init__(
        self,
        smallest_text_px: float | None,
        max_overhang: float,
        overflow: dict[tuple[float, float, float, float], float] | None = None,
    ) -> None:
        self.smallest_text_px = smallest_text_px
        self.max_overhang = max_overhang
        self.overflow = overflow or {}


def deck_baseline(data: bytes) -> DeckBaseline | None:
    """The original's own floor, from a light pass over its shapes (no pictures)."""
    from pptx import Presentation  # type: ignore[import-untyped]

    from .slide_text import metrics_page_px, needed_height, widest_line_px

    try:
        deck = Presentation(BytesIO(data))
    except Exception:  # noqa: BLE001 - an unreadable skin has no floor to offer
        return None
    width, height = int(deck.slide_width or 0), int(deck.slide_height or 0)
    # The box metric runs on the deck's own page shape (its inches), so an
    # imported portrait deck's overflow baseline lines up with the deck check,
    # which now measures on the same page. Falls back to 16:9 for a page-less
    # or malformed deck, exactly as before.
    _EMU_PER_INCH = 914400.0
    page_dims = (width / _EMU_PER_INCH, height / _EMU_PER_INCH) if width and height else None
    page_w_px, page_h_px = metrics_page_px(page_dims)
    smallest: float | None = None
    overhang = 0.0
    overflow: dict[tuple[float, float, float, float], float] = {}
    for slide in deck.slides:
        for shape in slide.shapes:
            geometry = (shape.left, shape.top, shape.width, shape.height)
            if width and height and None not in geometry:
                left, top, w, h = (int(v) for v in geometry)
                right = (left + w) / width * 100 - 100
                bottom = (top + h) / height * 100 - 100
                overhang = max(overhang, right, bottom, -left / width * 100, -top / height * 100)
            if getattr(shape, "has_table", False):
                frames = [cell.text_frame for row in shape.table.rows for cell in row.cells]
            elif getattr(shape, "has_text_frame", False):
                frames = [shape.text_frame]
            else:
                continue
            for paragraph in (p for frame in frames for p in frame.paragraphs):
                for run in paragraph.runs:
                    size = getattr(run.font, "size", None)
                    if size is None or not (run.text or "").strip():
                        continue
                    px = round(size.pt * _PT_TO_PX, 1)
                    if px > 0 and (smallest is None or px < smallest):
                        smallest = px
            # The box's own overflow, keyed by geometry rather than id so the
            # record survives slide reordering. Only frames the check would
            # look at count: fixed boxes with an explicit size (autofit boxes
            # never produce an overflow finding, and tables size from their
            # grid). Geometry is the key on purpose — an agent that edits the
            # words keeps the key, one that moves the box gives it up.
            if (
                getattr(shape, "has_text_frame", False)
                and width
                and height
                and None not in geometry
                and _autofit(shape.text_frame) is None
                and (shape.text_frame.text or "").strip()
            ):
                runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
                size = runs[0].font.size if runs else None
                if size is not None:
                    left, top, w, h = (int(v) for v in geometry)
                    # The same inset shrink the reader applies: the key and the
                    # measured width must be the geometry the imported element
                    # carries, or an inherited overflow stops matching its box.
                    ins_l, ins_t, ins_r, ins_b = _text_insets_emu(shape.text_frame)
                    if w - ins_l - ins_r > 0 and h - ins_t - ins_b > 0:
                        left, top = left + ins_l, top + ins_t
                        w, h = w - ins_l - ins_r, h - ins_t - ins_b
                    x_pct, y_pct = round(100 * left / width, 3), round(100 * top / height, 3)
                    w_pct, h_pct = round(100 * w / width, 3), round(100 * h / height, 3)
                    size_px = round(size.pt * _PT_TO_PX, 1)
                    box_w = w_pct / 100.0 * page_w_px
                    box_h = h_pct / 100.0 * page_h_px
                    if shape.text_frame.word_wrap is False:
                        # A no-wrap box overflows sideways; record the width
                        # ratio under the same geometry key the check folds on.
                        needed = widest_line_px(shape.text_frame.text, size_px)
                        box = box_w
                    else:
                        needed = needed_height(
                            shape.text_frame.text,
                            size_px,
                            box_w,
                            _line_height(list(shape.text_frame.paragraphs)),
                        )
                        box = box_h
                    if box > 0 and needed > box:
                        key = overflow_key(x_pct, y_pct, w_pct, h_pct)
                        overflow[key] = max(overflow.get(key, 0.0), round(needed / box, 3))
    return DeckBaseline(smallest, round(overhang, 3), overflow)


def overflow_key(x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    """A box's identity for the overflow baseline: its geometry, to one decimal."""
    return (round(x, 1), round(y, 1), round(w, 1), round(h, 1))


def smallest_text_px(data: bytes) -> float | None:
    """The smallest explicit run size in the deck, in the model's px, or ``None``."""
    baseline = deck_baseline(data)
    return baseline.smallest_text_px if baseline else None


def _table_styles(deck: Any) -> dict[str, dict[str, Any]]:
    """The deck's table style sheet, by style id: what a cell looks like when
    its own XML says nothing.

    A table that names a style and no cell borders is not an unbordered
    table — its borders live in ``ppt/tableStyles.xml``. PowerPoint's
    built-in styles are not written there (only their id), so those still
    come across bare; a deck that carries its definitions (Google Slides
    exports do) gets its whole-table border, text colour and face back.
    """
    from lxml import etree  # type: ignore[import-untyped]
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    styles: dict[str, dict[str, Any]] = {}
    for part in deck.part.package.iter_parts():
        if not str(part.partname).endswith("tableStyles.xml"):
            continue
        try:
            root = etree.fromstring(part.blob)
        except etree.XMLSyntaxError:
            continue
        for style in root.iter(qn("a:tblStyle")):
            style_id = style.get("styleId")
            whole = style.find(qn("a:wholeTbl"))
            if not style_id or whole is None:
                continue
            entry: dict[str, Any] = {}
            text_style = whole.find(qn("a:tcTxStyle"))
            if text_style is not None:
                entry["text"] = text_style  # colour child resolved per slide (scheme)
                face = text_style.find(".//" + qn("a:latin"))
                if face is not None and face.get("typeface"):
                    entry["font"] = face.get("typeface")
            borders = whole.find(qn("a:tcStyle") + "/" + qn("a:tcBdr"))
            if borders is not None:
                for side in ("a:insideH", "a:insideV", "a:left", "a:top", "a:bottom", "a:right"):
                    line = borders.find(qn(side) + "/" + qn("a:ln"))
                    if line is not None and line.find(qn("a:solidFill")) is not None:
                        entry["border"] = line
                        break
            styles[style_id] = entry
    return styles


def pptx_to_slides(data: bytes) -> dict[str, Any]:
    """Parse presentation bytes into ``SlidesData``-shaped dict.

    Raises :class:`PptxImportError` when the bytes are not a readable
    presentation. Callers that want a file card on failure catch it there.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ModuleNotFoundError as exc:  # pragma: no cover - install-time path
        raise PptxImportError(
            "reading .pptx needs python-pptx — install langchain-canvas[office]"
        ) from exc

    try:
        deck = Presentation(BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — any malformed file is one failure
        raise PptxImportError(f"could not read the presentation: {exc}") from exc

    width = deck.slide_width or 0
    height = deck.slide_height or 0
    if not width or not height:
        raise PptxImportError("the presentation declares no slide size")

    styles = _table_styles(deck)
    slides = [_slide(slide, width, height, styles) for slide in deck.slides]
    return {
        "slides": slides,
        "page": {
            "widthIn": round(width / _EMU_PER_INCH, 4),
            "heightIn": round(height / _EMU_PER_INCH, 4),
        },
    }


def _slide(
    slide: Any, width: int, height: int, table_styles: dict[str, dict[str, Any]] | None = None
) -> dict[str, Any]:
    """One slide's elements, in the order they stack on the page."""
    scheme = _scheme(slide)
    elements: list[dict[str, Any]] = []
    charts: list[dict[str, Any]] = []
    for index, shape in enumerate(slide.shapes):
        if getattr(shape, "has_table", False):
            elements.extend(
                _table_elements(shape, index, width, height, scheme, table_styles or {})
            )
            continue
        if getattr(shape, "has_chart", False):
            # A chart is a picture of data the deck model cannot hold; the
            # caller with a page renderer turns this box into that picture.
            frame = _frame(shape, index, width, height)
            if frame is not None:
                charts.append(frame)
            continue
        element = _element(shape, index, width, height, scheme)
        if element is not None:
            elements.append(element)
    picture = _background_picture(slide, scheme)
    if picture is not None:
        # Behind everything else, filling the page — a background is what the
        # rest sits on, and the element model has no field for one.
        elements.insert(
            0, {"id": "bg", "type": "image", "x": 0, "y": 0, "w": 100, "h": 100, "src": picture}
        )
    out: dict[str, Any] = {"elements": elements}
    if charts:
        out["charts"] = charts  # not a deck field: the copying tool consumes it
    # Text that names no colour takes the theme's text slot through the colour
    # map — dark on a light master, light on a dark one. That is the slide's
    # default text colour, and without it a table's cells and a WordArt with
    # no run colour were drawn in whatever the screen's own theme used.
    default_text = scheme.get("tx1")
    if default_text:
        out["textColor"] = f"#{default_text}"
    background = _background(slide, scheme)
    if background:
        out["background"] = background
    notes = _notes(slide)
    if notes:
        out["notes"] = notes
    return out


def _background_picture(slide: Any, scheme: dict[str, str]) -> str | None:
    """A picture background as a data URI, scaled down to a screen size.

    ``Slide.background`` holds one colour, so a deck built on photographs
    would come back white. Laying the picture in as the bottom element is
    what the model can express, and it survives the export as a real image.

    Originals are print-sized — one deck's backgrounds measured 25.8 MB
    across four slides, which is not a canvas file. They are resampled to
    1920px JPEG (90 KB for that same deck). Without Pillow there is no
    resampler, and an unshrunk background is worse than none, so it is left
    out and the slide keeps whatever colour it declares.
    """
    blob = _background_blob(slide)
    if blob is None:
        return None
    try:
        from PIL import Image  # type: ignore[import-untyped]
    except ModuleNotFoundError:
        return None
    try:
        image = Image.open(BytesIO(blob))
        image.thumbnail((_BACKGROUND_MAX_PX, _BACKGROUND_MAX_PX), Image.Resampling.LANCZOS)
        buffer = BytesIO()
        image.convert("RGB").save(buffer, "JPEG", quality=82, optimize=True)
    except Exception:  # noqa: BLE001 — an unreadable background drops out
        return None
    return f"data:image/jpeg;base64,{base64.b64encode(buffer.getvalue()).decode('ascii')}"


def _background_source(slide: Any) -> tuple[Any, Any] | None:
    """The part and the ``p:bg`` element that paint this slide's background.

    A slide that declares no background is not a white slide — it wears its
    layout's, and failing that its master's. Reading only the slide's own
    element turns an inherited dark deck white, which is what a reader that
    stops at the slide sees.

    Returns the owning part alongside the element, because a picture fill
    names its image by a relationship id that only that part can resolve.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return None
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None)
    for part in (slide, layout, master):
        element = getattr(part, "element", None)
        if element is None:
            continue
        common = element.find(qn("p:cSld"))
        background = common.find(qn("p:bg")) if common is not None else None
        if background is not None:
            return part, background
    return None


def _background_blob(slide: Any) -> bytes | None:
    """The bytes of the picture fill painting this slide, if there is one."""
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return None
    source = _background_source(slide)
    if source is None:
        return None
    part, background = source
    blip = background.find(
        f"./{qn('p:bgPr')}/{qn('a:blipFill')}/{qn('a:blip')}"
    )
    if blip is None:
        return None
    rid = blip.get(qn("r:embed"))
    if not rid:
        return None
    try:
        return part.part.related_part(rid).blob
    except (KeyError, AttributeError):
        return None


def _background(slide: Any, scheme: dict[str, str]) -> str | None:
    """The colour painting this slide, walking slide -> layout -> master.

    Handles the three fills a background can carry: a flat colour, a
    gradient (written as the CSS gradient the field can hold), and a theme
    reference resolved through the deck's own scheme. A picture fill returns
    ``None`` here — it comes in as the bottom element instead, where the
    model can hold the image.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return None
    source = _background_source(slide)
    if source is None:
        return None
    _part, background = source
    properties = background.find(qn("p:bgPr"))
    if properties is None:
        # ``p:bgRef`` points at a theme fill style and carries the colour to
        # tint it with; the colour alone is much closer than white.
        reference = background.find(qn("p:bgRef"))
        return _fill_colour(reference, scheme, qn) if reference is not None else None
    solid = properties.find(qn("a:solidFill"))
    if solid is not None:
        return _fill_colour(solid, scheme, qn)
    gradient = properties.find(qn("a:gradFill"))
    if gradient is not None:
        return _gradient_css(gradient, scheme, qn)
    return None


def _fill_colour(holder: Any, scheme: dict[str, str], qn: Any) -> str | None:
    """The ``#rrggbb`` of a colour child, theme references resolved."""
    srgb = holder.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return f"#{_with_brightness(srgb, srgb.get('val', ''), qn)}"
    themed = holder.find(qn("a:schemeClr"))
    if themed is not None:
        # The scheme already carries the master's colour map, so bg1/tx1
        # resolve here the same way the deck resolves them.
        value = scheme.get(themed.get("val", ""))
        if value:
            return f"#{_with_brightness(themed, value, qn)}"
    preset = holder.find(qn("a:prstClr"))
    if preset is not None:
        value = preset_colour(preset.get("val"))
        if value:
            return f"#{_with_brightness(preset, value, qn)}"
    return None


def _with_brightness(holder: Any, hex_rgb: str, qn: Any) -> str:
    """A colour with its ``lumMod``/``lumOff`` applied, as the deck shows it."""
    brightness = 0.0
    modulate = holder.find(qn("a:lumMod"))
    offset = holder.find(qn("a:lumOff"))
    if modulate is not None and modulate.get("val"):
        brightness = int(modulate.get("val", "100000")) / 100000 - 1
    if offset is not None and offset.get("val"):
        brightness = int(offset.get("val", "0")) / 100000
    return _shade(hex_rgb, brightness) if brightness else hex_rgb


def _gradient_css(gradient: Any, scheme: dict[str, str], qn: Any) -> str | None:
    """A gradient background as the CSS the ``background`` field can hold."""
    stops: list[str] = []
    for stop in gradient.findall(f"./{qn('a:gsLst')}/{qn('a:gs')}"):
        colour = _fill_colour(stop, scheme, qn)
        if colour is None:
            continue
        position = int(stop.get("pos", "0")) / 1000
        stops.append(f"{colour} {position:.0f}%")
    if len(stops) < 2:
        return stops[0].split(" ")[0] if stops else None
    linear = gradient.find(qn("a:lin"))
    # PowerPoint measures the sweep clockwise from east in 1/60000 degree;
    # CSS measures it clockwise from north, so the two are 90 degrees apart.
    angle = (int(linear.get("ang", "0")) / 60000 + 90) % 360 if linear is not None else 180
    return f"linear-gradient({angle:.0f}deg, {', '.join(stops)})"


# Theme-colour names as the scheme spells them, keyed by the enum member name
# python-pptx reports. ``bg1``/``tx1`` are indirections the master's colour map
# resolves (usually bg1 -> lt1), so they are looked up through it.
_THEME_KEYS = {
    "BACKGROUND_1": "bg1",
    "BACKGROUND_2": "bg2",
    "TEXT_1": "tx1",
    "TEXT_2": "tx2",
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
    **{f"ACCENT_{n}": f"accent{n}" for n in range(1, 7)},
}


def _scheme(slide: Any) -> dict[str, str]:
    """The slide master's colour scheme, as ``name -> rrggbb``.

    Half the runs in a real deck name a theme colour instead of a value, and
    a reader that skips those drops half the deck's colour on the floor. The
    scheme is where the values live, and the master's colour map is what
    ``bg1`` and ``tx1`` mean for this deck.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return {}
    try:
        master = slide.slide_layout.slide_master
        theme = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
    except (AttributeError, KeyError):
        return {}
    try:
        from lxml import etree  # type: ignore[import-untyped]

        root = etree.fromstring(theme.blob)
    except Exception:  # noqa: BLE001 - an unreadable theme leaves colours alone
        return {}

    values: dict[str, str] = {}
    scheme = root.find(f".//{qn('a:clrScheme')}")
    for child in scheme if scheme is not None else []:
        name = child.tag.split("}")[-1]
        srgb = child.find(qn("a:srgbClr"))
        system = child.find(qn("a:sysClr"))
        if srgb is not None and srgb.get("val"):
            values[name] = srgb.get("val", "")
        elif system is not None and system.get("lastClr"):
            values[name] = system.get("lastClr", "")

    mapping = master._element.find(qn("p:clrMap"))
    if mapping is not None:
        for key, target in mapping.attrib.items():
            if target in values:
                values[key] = values[target]
    return values


def _shade(hex_rgb: str, brightness: float) -> str:
    """Apply PowerPoint's luminance change to a colour.

    Negative brightness is ``lumMod`` (toward black), positive is ``lumOff``
    (toward white) — white at -0.25 is #BFBFBF, which is what the deck shows.
    """
    try:
        red, green, blue = (int(hex_rgb[i : i + 2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return hex_rgb
    if brightness < 0:
        factor = 1 + brightness
        channels = (round(c * factor) for c in (red, green, blue))
    else:
        channels = (round(c + (255 - c) * brightness) for c in (red, green, blue))
    return "".join(f"{max(0, min(255, c)):02X}" for c in channels)


def _notes(slide: Any) -> str:
    """The speaker notes, which the exporter writes back out."""
    if not getattr(slide, "has_notes_slide", False):
        return ""
    frame = getattr(slide.notes_slide, "notes_text_frame", None)
    return (frame.text or "").strip() if frame is not None else ""


def _element(
    shape: Any, index: int, width: int, height: int, scheme: dict[str, str]
) -> dict[str, Any] | None:
    """One shape as a deck element, or ``None`` when it has no equivalent."""
    frame = _frame(shape, index, width, height)
    if frame is None:
        return None

    if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
        return None
    if _is_group(shape):
        return None

    picture = _picture(shape)
    if picture is not None:
        return {**frame, "type": "image", "src": picture}

    text = _text(shape, scheme)
    if text is not None:
        return {**_inset_frame(frame, shape, width, height), "type": "text", **text}

    drawing = _drawing(shape, scheme)
    if drawing is not None:
        if drawing.get("shape") == "line":
            frame = _with_stroke(frame, shape, width, height)
        return {**frame, "type": "shape", **drawing}
    return None


def _inset_frame(
    frame: dict[str, Any], shape: Any, width: int, height: int
) -> dict[str, Any]:
    """The frame shrunk to the box's *text* area — where the words really are.

    PowerPoint keeps 0.1in side and 0.05in top/bottom insets between a box
    and its text unless the file says otherwise. The canvas draws text at
    the box edge and the exporter writes zero-margin boxes, so carrying the
    raw frame made every line measure against ~0.2in more width than the
    original wrapped in — and break a word or two later. Shrinking the frame
    to the inset area is what puts the line breaks back where the file's
    reader saw them. A degenerate shrink (insets meeting) keeps the frame.
    """
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None or not width or not height:
        return frame
    left, top, right, bottom = _text_insets_emu(text_frame)
    new_w = frame["w"] - 100.0 * (left + right) / width
    new_h = frame["h"] - 100.0 * (top + bottom) / height
    if new_w <= 0 or new_h <= 0:
        return frame
    return {
        **frame,
        "x": round(frame["x"] + 100.0 * left / width, 3),
        "y": round(frame["y"] + 100.0 * top / height, 3),
        "w": round(new_w, 3),
        "h": round(new_h, 3),
    }


class _CellText:
    """A table cell dressed as the text shape :func:`_text` reads."""

    has_text_frame = True

    def __init__(self, cell: Any) -> None:
        self.text_frame = cell.text_frame
        self._element = cell._tc
        self.height = None


def _table_elements(
    shape: Any,
    index: int,
    width: int,
    height: int,
    scheme: dict[str, str],
    table_styles: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    """A table as one ``table`` element on the table's own grid.

    The words go in ``rows``; the column widths and row heights come from
    the file; a merged cell keeps its span. The look the cells share — the
    grid line, the face, the colour, the size — becomes the element's own,
    from the cells' XML where they say it and from the deck's table style
    sheet where they do not (PowerPoint's built-in styles are not written
    to the sheet, so a deck that only names one comes across with no
    line). What single cells do differently — a header fill, a bold total —
    goes in ``cells``. A cell's own edge colours beyond the first are not
    carried; the original keeps them.
    """
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    frame = _frame(shape, index, width, height)
    if frame is None:
        return []
    table = shape.table
    properties = table._tbl.tblPr
    style_ref = properties.find(qn("a:tableStyleId")) if properties is not None else None
    style_id = style_ref.text.strip() if style_ref is not None and style_ref.text else ""
    style = table_styles.get(style_id, {})
    style_border = style.get("border")
    style_stroke: tuple[str, float] | None = None
    if style_border is not None:
        colour = _fill_colour(style_border.find(qn("a:solidFill")), scheme, qn)
        if colour:
            style_stroke = (colour, _line_px(style_border))
    style_text = style.get("text")
    style_colour = None
    if style_text is not None:
        holder = style_text.find(qn("a:solidFill"))
        style_colour = _fill_colour(holder if holder is not None else style_text, scheme, qn)

    widths = [int(column.width or 0) for column in table.columns]
    heights = [int(row.height or 0) for row in table.rows]
    rows: list[list[str]] = []
    looks: dict[tuple[int, int], dict[str, Any]] = {}
    strokes: list[tuple[str, float]] = []
    for r in range(len(heights)):
        row: list[str] = []
        for c in range(len(widths)):
            cell = table.cell(r, c)
            if cell.is_spanned:
                row.append("")
                continue
            row.append(cell.text_frame.text or "")
            own: dict[str, Any] = {}
            if cell.is_merge_origin:
                if cell.span_width > 1:
                    own["colSpan"] = cell.span_width
                if cell.span_height > 1:
                    own["rowSpan"] = cell.span_height
            tc_pr = cell._tc.tcPr
            if tc_pr is not None:
                solid = tc_pr.find(qn("a:solidFill"))
                if solid is not None:
                    fill = _fill_colour(solid, scheme, qn)
                    if fill:
                        own["fill"] = fill
                for side in ("a:lnB", "a:lnT", "a:lnL", "a:lnR"):
                    line = tc_pr.find(qn(side))
                    if line is None or line.find(qn("a:solidFill")) is None:
                        continue
                    colour = _fill_colour(line.find(qn("a:solidFill")), scheme, qn)
                    if colour:
                        strokes.append((colour, _line_px(line)))
                        break
            text = _text(_CellText(cell), scheme)
            if text is not None:
                for key in ("fontSize", "bold", "color", "align", "fontFamily"):
                    if key in text:
                        own[key] = text[key]
            looks[(r, c)] = own
        rows.append(row)
    if not rows or not rows[0]:
        return []

    element: dict[str, Any] = {**frame, "id": f"t{index}", "type": "table", "rows": rows}
    total_w, total_h = sum(widths), sum(heights)
    # PowerPoint draws a table from its column widths and row heights; the
    # frame around it is advisory and some writers leave it stale. One deck
    # declared a 3.3in square frame around a 7.9 x 1.7in grid, and the copy
    # showed a square table over the WordArt below it.
    if total_w > 0 and width:
        element["w"] = round(100 * total_w / width, 3)
    if total_h > 0 and height:
        element["h"] = round(100 * total_h / height, 3)
    if total_w > 0:
        element["colWidths"] = [round(100 * w / total_w, 3) for w in widths]
    if total_h > 0:
        element["rowHeights"] = [round(100 * h / total_h, 3) for h in heights]
    if table.first_row:
        element["header"] = True

    # What most cells do is the table's; the rest is each cell's own.
    shared: dict[str, Any] = {}
    for key in ("fontSize", "color", "align", "fontFamily", "fill", "bold"):
        values = [own[key] for own in looks.values() if key in own]
        if not values:
            continue
        common = Counter(json.dumps(v) for v in values).most_common(1)[0][0]
        if key in ("fill", "bold") and len(values) < len(looks):
            continue  # a fill or weight some cells lack is not the table's
        shared[key] = json.loads(common)
    if "color" not in shared and style_colour:
        shared["color"] = style_colour
    if "fontFamily" not in shared and style.get("font"):
        shared["fontFamily"] = style["font"]
    stroke = Counter(strokes).most_common(1)[0][0] if strokes else style_stroke
    if stroke:
        element["stroke"], element["strokeWidth"] = stroke
    element.update(shared)
    cells: list[dict[str, Any]] = []
    for (r, c), own in looks.items():
        entry = {"r": r, "c": c}
        for key, value in own.items():
            if key in ("colSpan", "rowSpan") or shared.get(key) != value:
                entry[key] = value
        if len(entry) > 2:
            cells.append(entry)
    if cells:
        element["cells"] = cells
    return [element]


def _line_px(line: Any) -> float:
    """A drawing line's weight in the model's px (PowerPoint's default when unset)."""
    return round(int(line.get("w") or 12700) / 12700 * _PT_TO_PX, 2)


def _with_stroke(
    frame: dict[str, Any], shape: Any, width: int, height: int
) -> dict[str, Any]:
    """Give a connector its stroke as thickness.

    A horizontal connector is zero pixels tall in the file — PowerPoint draws
    it from the line weight, not the box. The deck model has no stroke field
    and paints every shape as a box, so a zero stays invisible. Fill the flat
    side with the weight the line actually declares (or PowerPoint's default
    when it declares none, which is the common case).
    """
    try:
        stroke = shape.line.width or _DEFAULT_STROKE_EMU
    except (AttributeError, ValueError, TypeError):
        stroke = _DEFAULT_STROKE_EMU
    out = dict(frame)
    if not out["h"]:
        out["h"] = round(100 * stroke / height, 3)
    if not out["w"]:
        out["w"] = round(100 * stroke / width, 3)
    return out


def _frame(shape: Any, index: int, width: int, height: int) -> dict[str, Any] | None:
    """Percent geometry, or ``None`` when the shape declares no box.

    A placeholder that inherits its position from the layout reports ``None``
    for left/top; the layout is the skin's business, so such a shape is left
    to the skin rather than pinned at a guessed spot.
    """
    left, top = getattr(shape, "left", None), getattr(shape, "top", None)
    box_w, box_h = getattr(shape, "width", None), getattr(shape, "height", None)
    if left is None or top is None or box_w is None or box_h is None:
        return None
    frame: dict[str, Any] = {
        "id": f"e{index}",
        "x": round(100 * left / width, 3),
        "y": round(100 * top / height, 3),
        "w": round(100 * box_w / width, 3),
        "h": round(100 * box_h / height, 3),
    }
    # Rotation is stored on the shape's transform in degrees clockwise; a box
    # left upright reports 0, which we leave off so an unrotated deck stays
    # byte-for-byte what it was.
    rotation = getattr(shape, "rotation", None)
    if rotation:
        frame["rotation"] = round(float(rotation), 3)
    return frame


def _is_group(shape: Any) -> bool:
    return "GROUP" in str(getattr(shape, "shape_type", "") or "")


def _picture(shape: Any) -> str | None:
    """A picture as a ``data:`` URI the exporter can place back."""
    image = getattr(shape, "image", None)
    if image is None:
        return None
    kind = _IMAGE_TYPES.get((getattr(image, "ext", "") or "").lower())
    if kind is None:
        return None
    try:
        blob = image.blob
    except Exception:  # noqa: BLE001 - an unreadable part drops the picture
        return None
    return f"data:image/{kind};base64,{base64.b64encode(blob).decode('ascii')}"


#: How much of a WordArt box the glyphs fill vertically (cap height over
#: box height for the common warps).
_WORD_ART_FILL = 0.8


def _text_outline(font: Any, scheme: dict[str, str]) -> tuple[str, float] | None:
    """``(colour, px)`` of the run's text outline, or ``None`` when it has none.

    WordArt is usually a light fill drawn legible by a dark outline; with
    only the fill it faded into the page. The outline rides the element's
    ``stroke`` fields, the same ones a shape uses for its border.
    """
    rpr = getattr(font, "_rPr", None)
    if rpr is None:
        return None
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    line = rpr.find(qn("a:ln"))
    if line is None or line.find(qn("a:solidFill")) is None:
        return None
    colour = _fill_colour(line.find(qn("a:solidFill")), scheme, qn)
    if not colour:
        return None
    return colour, round(int(line.get("w") or 12700) / 12700 * _PT_TO_PX, 2)


def _is_word_art(shape: Any) -> bool:
    """A text shape whose body is warped — WordArt, in PowerPoint's own name."""
    element = getattr(shape, "_element", None)
    if element is None:
        return False
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    return element.find(".//" + qn("a:prstTxWarp")) is not None


def _text(shape: Any, scheme: dict[str, str]) -> dict[str, Any] | None:
    """Text plus the first run's formatting, or ``None`` when there is none.

    One element carries one set of formatting, so a box whose runs disagree
    is represented by its first — the words all survive, the variation does
    not. The original keeps it either way.
    """
    if not getattr(shape, "has_text_frame", False):
        return None
    body = (shape.text_frame.text or "").strip()
    if not body:
        return None

    out: dict[str, Any] = {"text": shape.text_frame.text}
    fit = _autofit(shape.text_frame)
    if fit:
        out["autofit"] = fit
    # A box PowerPoint never wraps must not wrap on the canvas either — a
    # one-line label that suddenly folds is the most visible way an import
    # stops looking like its original.
    if getattr(shape.text_frame, "word_wrap", None) is False:
        out["wrap"] = False
    font_scale, spacing_reduction = _stored_autofit_scale(shape.text_frame)
    paragraphs = list(shape.text_frame.paragraphs)
    runs = [run for paragraph in paragraphs for run in paragraph.runs]
    if runs:
        font = runs[0].font
        outline = _text_outline(font, scheme)
        if outline is not None:
            out["stroke"], out["strokeWidth"] = outline
        if font.size is not None:
            out["fontSize"] = round(font.size.pt * _PT_TO_PX, 1)
        elif _is_word_art(shape):
            # WordArt warps its text to fill the shape, whatever the run says
            # (usually nothing). The box height is the size the person sees;
            # without this a headline-sized word came across at body size.
            height = getattr(shape, "height", None)
            if height:
                out["fontSize"] = round(int(height) / 12700 * _PT_TO_PX * _WORD_ART_FILL, 1)
        if font.bold is not None:
            out["bold"] = bool(font.bold)
        colour = _colour(font, scheme)
        if colour:
            out["color"] = colour
    align = _align(paragraphs)
    if align:
        out["align"] = align
    if runs and runs[0].font.name:
        out["fontFamily"] = runs[0].font.name
    spacing = _line_height(paragraphs)
    if spacing:
        out["lineHeight"] = spacing
    # Apply the shrink PowerPoint already computed (normAutofit): the size
    # and leading in the file are pre-shrink, but what the person *sees* is
    # the scaled type — and matching what they see is the import's job.
    if font_scale < 1.0 and out.get("fontSize"):
        out["fontSize"] = round(out["fontSize"] * font_scale, 1)
    if spacing_reduction > 0.0:
        out["lineHeight"] = round((out.get("lineHeight") or 1.2) * (1.0 - spacing_reduction), 3)
    anchor = _vertical_align(shape.text_frame)
    if anchor:
        out["verticalAlign"] = anchor
    if runs:
        band = _highlight(runs[0])
        if band:
            out["highlight"] = band
    before, after = _paragraph_spacing(paragraphs)
    if before is not None:
        out["spaceBefore"] = before
    if after is not None:
        out["spaceAfter"] = after
    return out


def _highlight(run: Any) -> str | None:
    """The colour band behind a run, if it is highlighted.

    python-pptx exposes no accessor for ``a:highlight``, so this reads the run
    properties directly. A highlighted heading is a coloured bar in the deck;
    without it the heading reads as plain text and the slide looks unlike the
    file it came from.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return None
    properties = getattr(run, "_r", None)
    properties = properties.find(qn("a:rPr")) if properties is not None else None
    if properties is None:
        return None
    band = properties.find(qn("a:highlight"))
    if band is None:
        return None
    srgb = band.find(qn("a:srgbClr"))
    value = srgb.get("val") if srgb is not None else None
    return f"#{value}" if value else None


def _paragraph_spacing(paragraphs: list[Any]) -> tuple[float | None, float | None]:
    """Space above and below the text, in points, when the paragraphs agree."""

    def agreed(attr: str) -> float | None:
        values = {
            round(getattr(p, attr).pt * _PT_TO_PX, 2)
            for p in paragraphs
            if getattr(p, attr, None) is not None
        }
        return values.pop() if len(values) == 1 else None

    return agreed("space_before"), agreed("space_after")


def _line_height(paragraphs: list[Any]) -> float | None:
    """Line spacing as a multiple, when the paragraphs agree on one.

    python-pptx reports a float for a multiple and a Length for an exact
    height; only the multiple maps onto the model's field, and an exact height
    without the font size to divide by would be a guess.
    """
    values = {p.line_spacing for p in paragraphs if isinstance(p.line_spacing, float)}
    return round(values.pop(), 3) if len(values) == 1 else None


def master_has_content(data: bytes) -> bool:
    """Whether any master or layout draws shapes of its own (logo, footer).

    Placeholders are the slide's to fill and don't count; what counts is the
    decoration PowerPoint keeps out of reach on the slide — exactly what the
    editable copy used to lose.
    """
    from pptx import Presentation  # type: ignore[import-untyped]

    try:
        deck = Presentation(BytesIO(data))
    except Exception:  # noqa: BLE001 — unreadable deck: nothing to show
        return False
    holders = list(deck.slide_masters) + [
        layout for master in deck.slide_masters for layout in master.slide_layouts
    ]
    for holder in holders:
        for shape in holder.shapes:
            if not getattr(shape, "is_placeholder", False):
                return True
    return False


def blank_slides_pptx(data: bytes) -> bytes | None:
    """The deck with every slide's own shapes removed, as bytes.

    Rendering the result shows what the master and layout draw on their own —
    the backdrop the editable copy puts behind its elements. ``None`` when the
    deck cannot be read.
    """
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    try:
        deck = Presentation(BytesIO(data))
    except Exception:  # noqa: BLE001
        return None
    drawn = (
        qn("p:sp"),
        qn("p:pic"),
        qn("p:graphicFrame"),
        qn("p:cxnSp"),
        qn("p:grpSp"),
    )
    for slide in deck.slides:
        tree = slide.shapes._spTree
        for child in list(tree):
            if child.tag in drawn:
                tree.remove(child)
    out = BytesIO()
    deck.save(out)
    return out.getvalue()


def _autofit(frame: Any) -> str | None:
    """How the box and its text negotiate when the words outgrow the box.

    A box that grows with its text comes across as ``shape``; type that
    shrinks to its box as ``text``. A frame that says neither stays silent
    rather than becoming ``none``, so a deck carries only what its file set.
    """
    mode = getattr(frame, "auto_size", None)
    name = str(getattr(mode, "name", "")).upper()
    return {"SHAPE_TO_FIT_TEXT": "shape", "TEXT_TO_FIT_SHAPE": "text"}.get(name)


def _stored_autofit_scale(frame: Any) -> tuple[float, float]:
    """The shrink PowerPoint already computed for this box, from its file.

    ``normAutofit`` carries ``fontScale`` (thousandths of a percent) and
    ``lnSpcReduction``. Ignoring them is how a box PowerPoint shows at 62%
    arrived at 100% and ran past its frame — the single biggest source of
    "the canvas clips what the original shows". ``(1.0, 0.0)`` when unset.
    """
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    body = getattr(frame, "_bodyPr", None)
    fitted = body.find(qn("a:normAutofit")) if body is not None else None
    if fitted is None:
        return 1.0, 0.0

    def _thousandths(name: str) -> float | None:
        raw = fitted.get(name)
        if raw is None:
            return None
        try:
            return float(raw) / 100000.0
        except ValueError:
            return None

    scale = _thousandths("fontScale")
    reduction = _thousandths("lnSpcReduction")
    return (
        scale if scale is not None and 0 < scale <= 1 else 1.0,
        reduction if reduction is not None and 0 <= reduction < 1 else 0.0,
    )


#: PowerPoint's default text insets (EMU): 0.1in sides, 0.05in top/bottom.
_DEFAULT_INSETS_EMU = (91440, 45720, 91440, 45720)


def _text_insets_emu(frame: Any) -> tuple[int, int, int, int]:
    """``(left, top, right, bottom)`` text insets, defaults where unsaid.

    The box's text lives inside these margins; measuring wraps against the
    full box width is how every line broke a word or two earlier than the
    original. python-pptx reports the effective value (default applied).
    """
    left = getattr(frame, "margin_left", None)
    top = getattr(frame, "margin_top", None)
    right = getattr(frame, "margin_right", None)
    bottom = getattr(frame, "margin_bottom", None)
    defaults = _DEFAULT_INSETS_EMU
    return (
        int(left) if left is not None else defaults[0],
        int(top) if top is not None else defaults[1],
        int(right) if right is not None else defaults[2],
        int(bottom) if bottom is not None else defaults[3],
    )


def _vertical_align(frame: Any) -> str | None:
    """Where text sits in its box, when the file says."""
    anchor = getattr(frame, "vertical_anchor", None)
    if anchor is None:
        return None
    name = str(getattr(anchor, "name", "")).lower()
    return {"top": "top", "middle": "middle", "bottom": "bottom"}.get(name)


def _colour(font: Any, scheme: dict[str, str]) -> str | None:
    """The run's colour as ``#rrggbb``, theme references resolved.

    A run either names a value or names a theme slot. Both end up here as one
    colour, because the deck model has one field and the renderer needs a
    value to paint. A slot the scheme does not carry stays ``None`` rather
    than becoming a guess.
    """
    colour = getattr(font, "color", None)
    if colour is None:
        return None
    try:
        rgb = colour.rgb
        if rgb is not None:
            return f"#{rgb}"
    except (AttributeError, ValueError):
        pass
    # The third way a run names a colour: a preset name, which python-pptx
    # exposes as a type with no value — so the value is read from the XML.
    rpr = getattr(font, "_rPr", None)
    if rpr is not None:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]

        fill = rpr.find(qn("a:solidFill"))
        preset = fill.find(qn("a:prstClr")) if fill is not None else None
        if preset is not None:
            value = preset_colour(preset.get("val"))
            if value:
                return f"#{_with_brightness(preset, value, qn)}"
    try:
        name = _THEME_KEYS.get(str(getattr(colour.theme_color, "name", "")))
    except (AttributeError, ValueError):
        return None
    if not name or name not in scheme:
        return None
    brightness = getattr(colour, "brightness", 0) or 0
    return f"#{_shade(scheme[name], float(brightness))}"


def _align(paragraphs: list[Any]) -> str | None:
    """The alignment the paragraphs agree on, if they agree."""
    names = {
        str(p.alignment).split()[0].lower()
        for p in paragraphs
        if p.alignment is not None and (p.text or "").strip()
    }
    if len(names) != 1:
        return None
    name = names.pop()
    return name if name in {"left", "center", "right"} else None


def _drawing(shape: Any, scheme: dict[str, str]) -> dict[str, Any] | None:
    """A rectangle, ellipse or line, with its fill and outline resolved.

    Colours are read from the shape's XML so the sources a real deck uses
    all land: explicit RGB, theme references, preset names, the style
    reference the shape inherits from, and a gradient's first stop. An
    explicit ``noFill`` becomes ``"none"`` — the shape *says* it is
    transparent, which is different from saying nothing. Measured across
    58 decks: explicit RGB is only 23% of shape fills; a bank template was
    70% theme references, and every one of those used to fall back to the
    renderer's text colour — a black box over the slide.
    """
    kind = str(getattr(shape, "shape_type", "") or "").upper()
    name = str(getattr(shape, "name", "") or "").upper()
    if not kind or "PICTURE" in kind:
        return None

    if any(word in kind or word in name for word in _LINE_NAMES):
        drawn = "line"
    elif any(word in kind or word in name for word in _ELLIPSE_NAMES):
        drawn = "ellipse"
    elif "AUTO_SHAPE" in kind or "TEXT_BOX" in kind or "FREEFORM" in kind:
        drawn = "rect"
    else:
        return None

    out: dict[str, Any] = {"shape": drawn}
    stroke, weight = _shape_outline(shape, scheme)
    fill = stroke if drawn == "line" else _shape_fill(shape, scheme)
    # Unfilled, unbordered and textless draws nothing in PowerPoint either —
    # a spacer. Importing it would only add invisible elements to the copy.
    if drawn != "line" and fill in (None, "none") and not stroke:
        return None
    if fill:
        out["fill"] = fill
    if stroke:
        out["stroke"] = stroke
    if weight:
        out["strokeWidth"] = weight
    return out


def _shape_fill(shape: Any, scheme: dict[str, str]) -> str | None:
    """``#rrggbb``, ``"none"`` for an explicit noFill, or ``None`` when unsaid."""
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        if spPr.find(qn("a:noFill")) is not None:
            return "none"
        solid = spPr.find(qn("a:solidFill"))
        if solid is not None:
            colour = _fill_colour(solid, scheme, qn)
            if colour:
                return colour
        gradient = spPr.find(qn("a:gradFill"))
        if gradient is not None:
            # One field, one colour: the first stop stands for the ramp.
            stop = gradient.find(qn("a:gsLst") + "/" + qn("a:gs"))
            if stop is not None:
                colour = _fill_colour(stop, scheme, qn)
                if colour:
                    return colour
    # No fill of its own: the shape style names the theme fill it inherits.
    style = shape._element.find(qn("p:style"))
    if style is not None:
        ref = style.find(qn("a:fillRef"))
        if ref is not None and (ref.get("idx") or "0") != "0":
            colour = _fill_colour(ref, scheme, qn)
            if colour:
                return colour
    return None


def _shape_outline(shape: Any, scheme: dict[str, str]) -> tuple[str | None, float | None]:
    """The outline colour and weight in px, from the XML (theme resolved).

    Boxes drawn by their border alone — an empty rectangle around content —
    are the common annotation in a real deck; a border whose colour is a
    theme reference used to vanish, and the box went black with it.
    """
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    colour: str | None = None
    weight: float | None = None
    spPr = shape._element.find(qn("p:spPr"))
    line = spPr.find(qn("a:ln")) if spPr is not None else None
    if line is not None:
        if line.find(qn("a:noFill")) is not None:
            return None, None
        solid = line.find(qn("a:solidFill"))
        if solid is not None:
            colour = _fill_colour(solid, scheme, qn)
        raw_width = line.get("w")
        if raw_width:
            weight = round(int(raw_width) / 12700 * _PT_TO_PX, 2)
    if colour is None:
        style = shape._element.find(qn("p:style"))
        if style is not None:
            ref = style.find(qn("a:lnRef"))
            if ref is not None and (ref.get("idx") or "0") != "0":
                colour = _fill_colour(ref, scheme, qn)
    return colour, weight


