"""Low-level ``.pptx`` shape-reading helpers shared by :mod:`deck.extract`.

Moved here from the now-deleted ``pptx_import.py`` (the legacy percent-geometry
reader) when that module was removed — :mod:`deck.extract` is now the only
caller of geometry, text formatting, and drawing detection, so the helpers
live next to it rather than in a module that no longer exists.
"""

from __future__ import annotations

from typing import Any

# Picture formats the pptx exporter can write back out; anything else is
# dropped rather than embedded as bytes no exporter will read.
_IMAGE_TYPES = {"png": "png", "jpg": "jpeg", "jpeg": "jpeg", "gif": "gif"}

# PowerPoint's default line weight, used when a connector declares none.
_DEFAULT_STROKE_EMU = 12700  # 1pt

# The deck model measures type and spacing in px on a 1280px-wide slide;
# PowerPoint states them in points. Storing a point value in a pixel field
# renders every size a quarter too small and moves every line break.
_PT_TO_PX = 4 / 3

# Autoshape names that map onto the three shapes the deck model draws.
_ELLIPSE_NAMES = ("OVAL", "ELLIPSE", "CIRCLE")
_LINE_NAMES = ("LINE", "STRAIGHT_CONNECTOR", "CONNECTOR")

# Theme-colour names as the scheme spells them, keyed by the enum member name
# python-pptx reports. ``bg1``/``tx1`` are indirections the master's colour map
# resolves (usually bg1 -> lt1), so they are looked up through it.
_THEME_KEYS = {
    "BACKGROUND_1": "bg1",
    "BACKGROUND_2": "bg2",
    "TEXT_1": "tx1",
    "TEXT_2": "tx2",
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
    **{f"ACCENT_{n}": f"accent{n}" for n in range(1, 7)},
}


class PptxImportError(ValueError):
    """The bytes are not a presentation this reader can open."""


def _notes(slide: Any) -> str:
    """The speaker notes, which the exporter writes back out."""
    if not getattr(slide, "has_notes_slide", False):
        return ""
    frame = getattr(slide.notes_slide, "notes_text_frame", None)
    return (frame.text or "").strip() if frame is not None else ""


def _frame(shape: Any, index: int, width: int, height: int) -> dict[str, Any] | None:
    """Percent geometry, or ``None`` when the shape declares no box.

    A placeholder that inherits its position from the layout reports ``None``
    for left/top; the layout is the skin's business, so such a shape is left
    to the skin rather than pinned at a guessed spot.
    """
    left, top = getattr(shape, "left", None), getattr(shape, "top", None)
    box_w, box_h = getattr(shape, "width", None), getattr(shape, "height", None)
    if left is None or top is None or box_w is None or box_h is None:
        return None
    return {
        "id": f"e{index}",
        "x": round(100 * left / width, 3),
        "y": round(100 * top / height, 3),
        "w": round(100 * box_w / width, 3),
        "h": round(100 * box_h / height, 3),
    }


def _is_group(shape: Any) -> bool:
    return "GROUP" in str(getattr(shape, "shape_type", "") or "")


def _text(shape: Any, scheme: dict[str, str]) -> dict[str, Any] | None:
    """Text plus the first run's formatting, or ``None`` when there is none.

    One element carries one set of formatting, so a box whose runs disagree
    is represented by its first — the words all survive, the variation does
    not. The original keeps it either way.
    """
    if not getattr(shape, "has_text_frame", False):
        return None
    body = (shape.text_frame.text or "").strip()
    if not body:
        return None

    out: dict[str, Any] = {"text": shape.text_frame.text}
    paragraphs = list(shape.text_frame.paragraphs)
    runs = [run for paragraph in paragraphs for run in paragraph.runs]
    if runs:
        font = runs[0].font
        if font.size is not None:
            out["fontSize"] = round(font.size.pt * _PT_TO_PX, 1)
        if font.bold is not None:
            out["bold"] = bool(font.bold)
        colour = _colour(font, scheme)
        if colour:
            out["color"] = colour
    align = _align(paragraphs)
    if align:
        out["align"] = align
    if runs and runs[0].font.name:
        out["fontFamily"] = runs[0].font.name
    spacing = _line_height(paragraphs)
    if spacing:
        out["lineHeight"] = spacing
    anchor = _vertical_align(shape.text_frame)
    if anchor:
        out["verticalAlign"] = anchor
    if runs:
        band = _highlight(runs[0])
        if band:
            out["highlight"] = band
    before, after = _paragraph_spacing(paragraphs)
    if before is not None:
        out["spaceBefore"] = before
    if after is not None:
        out["spaceAfter"] = after
    return out


def _highlight(run: Any) -> str | None:
    """The colour band behind a run, if it is highlighted.

    python-pptx exposes no accessor for ``a:highlight``, so this reads the run
    properties directly. A highlighted heading is a coloured bar in the deck;
    without it the heading reads as plain text and the slide looks unlike the
    file it came from.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return None
    properties = getattr(run, "_r", None)
    properties = properties.find(qn("a:rPr")) if properties is not None else None
    if properties is None:
        return None
    band = properties.find(qn("a:highlight"))
    if band is None:
        return None
    srgb = band.find(qn("a:srgbClr"))
    value = srgb.get("val") if srgb is not None else None
    return f"#{value}" if value else None


def _paragraph_spacing(paragraphs: list[Any]) -> tuple[float | None, float | None]:
    """Space above and below the text, in points, when the paragraphs agree."""

    def agreed(attr: str) -> float | None:
        values = {
            round(getattr(p, attr).pt * _PT_TO_PX, 2)
            for p in paragraphs
            if getattr(p, attr, None) is not None
        }
        return values.pop() if len(values) == 1 else None

    return agreed("space_before"), agreed("space_after")


def _line_height(paragraphs: list[Any]) -> float | None:
    """Line spacing as a multiple, when the paragraphs agree on one.

    python-pptx reports a float for a multiple and a Length for an exact
    height; only the multiple maps onto the model's field, and an exact height
    without the font size to divide by would be a guess.
    """
    values = {p.line_spacing for p in paragraphs if isinstance(p.line_spacing, float)}
    return round(values.pop(), 3) if len(values) == 1 else None


def _vertical_align(frame: Any) -> str | None:
    """Where text sits in its box, when the file says."""
    anchor = getattr(frame, "vertical_anchor", None)
    if anchor is None:
        return None
    name = str(getattr(anchor, "name", "")).lower()
    return {"top": "top", "middle": "middle", "bottom": "bottom"}.get(name)


def _colour(font: Any, scheme: dict[str, str]) -> str | None:
    """The run's colour as ``#rrggbb``, theme references resolved.

    A run either names a value or names a theme slot. Both end up here as one
    colour, because the deck model has one field and the renderer needs a
    value to paint. A slot the scheme does not carry stays ``None`` rather
    than becoming a guess.
    """
    colour = getattr(font, "color", None)
    if colour is None:
        return None
    try:
        rgb = colour.rgb
        if rgb is not None:
            return f"#{rgb}"
    except (AttributeError, ValueError):
        pass
    try:
        name = _THEME_KEYS.get(str(getattr(colour.theme_color, "name", "")))
    except (AttributeError, ValueError):
        return None
    if not name or name not in scheme:
        return None
    brightness = getattr(colour, "brightness", 0) or 0
    return f"#{_shade(scheme[name], float(brightness))}"


def _align(paragraphs: list[Any]) -> str | None:
    """The alignment the paragraphs agree on, if they agree."""
    names = {
        str(p.alignment).split()[0].lower()
        for p in paragraphs
        if p.alignment is not None and (p.text or "").strip()
    }
    if len(names) != 1:
        return None
    name = names.pop()
    return name if name in {"left", "center", "right"} else None


def _shade(hex_rgb: str, brightness: float) -> str:
    """Apply PowerPoint's luminance change to a colour.

    Negative brightness is ``lumMod`` (toward black), positive is ``lumOff``
    (toward white) — white at -0.25 is #BFBFBF, which is what the deck shows.
    """
    try:
        red, green, blue = (int(hex_rgb[i : i + 2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return hex_rgb
    if brightness < 0:
        factor = 1 + brightness
        channels = (round(c * factor) for c in (red, green, blue))
    else:
        channels = (round(c + (255 - c) * brightness) for c in (red, green, blue))
    return "".join(f"{max(0, min(255, c)):02X}" for c in channels)


def _scheme(slide: Any) -> dict[str, str]:
    """The slide master's colour scheme, as ``name -> rrggbb``.

    Half the runs in a real deck name a theme colour instead of a value, and
    a reader that skips those drops half the deck's colour on the floor. The
    scheme is where the values live, and the master's colour map is what
    ``bg1`` and ``tx1`` mean for this deck.
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ModuleNotFoundError:  # pragma: no cover - install-time path
        return {}
    try:
        master = slide.slide_layout.slide_master
        theme = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
    except (AttributeError, KeyError):
        return {}
    try:
        from lxml import etree  # type: ignore[import-untyped]

        root = etree.fromstring(theme.blob)
    except Exception:  # noqa: BLE001 - an unreadable theme leaves colours alone
        return {}

    values: dict[str, str] = {}
    scheme = root.find(f".//{qn('a:clrScheme')}")
    for child in scheme if scheme is not None else []:
        name = child.tag.split("}")[-1]
        srgb = child.find(qn("a:srgbClr"))
        system = child.find(qn("a:sysClr"))
        if srgb is not None and srgb.get("val"):
            values[name] = srgb.get("val", "")
        elif system is not None and system.get("lastClr"):
            values[name] = system.get("lastClr", "")

    mapping = master._element.find(qn("p:clrMap"))
    if mapping is not None:
        for key, target in mapping.attrib.items():
            if target in values:
                values[key] = values[target]
    return values


def _drawing(shape: Any) -> dict[str, Any] | None:
    """A rectangle, ellipse or line, with its solid fill if it has one."""
    kind = str(getattr(shape, "shape_type", "") or "").upper()
    name = str(getattr(shape, "name", "") or "").upper()
    if not kind or "PICTURE" in kind:
        return None

    if any(word in kind or word in name for word in _LINE_NAMES):
        drawn = "line"
    elif any(word in kind or word in name for word in _ELLIPSE_NAMES):
        drawn = "ellipse"
    elif "AUTO_SHAPE" in kind or "TEXT_BOX" in kind or "FREEFORM" in kind:
        drawn = "rect"
    else:
        return None

    out: dict[str, Any] = {"shape": drawn}
    fill = _fill(shape, drawn)
    if fill:
        out["fill"] = fill
    stroke, weight = _outline(shape)
    if stroke:
        out["stroke"] = stroke
    if weight:
        out["strokeWidth"] = weight
    return out


def _outline(shape: Any) -> tuple[str | None, float | None]:
    """A shape's outline colour and weight in points.

    Boxes drawn by their border alone — an empty rectangle around content — are
    the common annotation in a real deck, and they carry no fill at all. Read
    separately from ``fill`` so both can be present, or just one.
    """
    line = getattr(shape, "line", None)
    if line is None:
        return None, None
    colour = None
    try:
        rgb = line.color.rgb
        colour = f"#{rgb}" if rgb is not None else None
    except (AttributeError, ValueError, TypeError):
        colour = None
    weight = None
    try:
        if line.width:
            weight = round(line.width / 12700 * _PT_TO_PX, 2)  # EMU -> px
    except (AttributeError, ValueError, TypeError):
        weight = None
    return colour, weight


def _fill(shape: Any, drawn: str) -> str | None:
    """The solid fill (or, for a line, its stroke colour)."""
    holder = shape.line if drawn == "line" else getattr(shape, "fill", None)
    if holder is None:
        return None
    try:
        colour = holder.color if drawn == "line" else holder.fore_color
        rgb = colour.rgb
    except (AttributeError, ValueError, TypeError):
        return None
    return f"#{rgb}" if rgb is not None else None


def _with_stroke(frame: dict[str, Any], shape: Any, width: int, height: int) -> dict[str, Any]:
    """Give a connector its stroke as thickness.

    A horizontal connector is zero pixels tall in the file — PowerPoint draws
    it from the line weight, not the box. The deck model has no stroke field
    and paints every shape as a box, so a zero stays invisible. Fill the flat
    side with the weight the line actually declares (or PowerPoint's default
    when it declares none, which is the common case).
    """
    try:
        stroke = shape.line.width or _DEFAULT_STROKE_EMU
    except (AttributeError, ValueError, TypeError):
        stroke = _DEFAULT_STROKE_EMU
    out = dict(frame)
    if not out["h"]:
        out["h"] = round(100 * stroke / height, 3)
    if not out["w"]:
        out["w"] = round(100 * stroke / width, 3)
    return out
