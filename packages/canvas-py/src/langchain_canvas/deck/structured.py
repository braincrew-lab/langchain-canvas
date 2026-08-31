"""Structured native table and categorical-chart HTML; never slide screenshots."""

from __future__ import annotations

import hashlib
import html
import json
import math
import re
from dataclasses import dataclass
from types import SimpleNamespace
from typing import Any, cast

from ._shapes import _colour, _fill, _text


@dataclass(frozen=True)
class StructuredShape:
    element_id: str
    x: float
    y: float
    w: float
    h: float
    kind: str
    data: dict[str, Any]


def native_chart_data(chart: Any, scheme: dict[str, str]) -> dict[str, Any] | None:
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_FILL

    kinds = {XL_CHART_TYPE.COLUMN_CLUSTERED: "column", XL_CHART_TYPE.BAR_CLUSTERED: "bar"}
    if chart.chart_type not in kinds or len(chart.plots) != 1:
        return None
    categories = [str(c.label) for c in chart.plots[0].categories]
    series = []
    for index, series_item in enumerate(chart.series):
        fill = series_item.format.fill
        if fill.type not in {None, MSO_FILL.SOLID}:
            return None
        color = (
            _colour(SimpleNamespace(color=series_item.format.fill.fore_color), scheme)
            if fill.type == MSO_FILL.SOLID
            else None
        )
        series.append(
            {
                "name": str(series_item.name),
                "values": list(series_item.values),
                "color": color or "#" + scheme.get(f"accent{index % 6 + 1}", "663399"),
            }
        )
    data = {
        "type": kinds[chart.chart_type],
        "categories": categories,
        "series": series,
        "title": chart.chart_title.text_frame.text
        if chart.has_title and chart.chart_title.has_text_frame
        else "",
        "font_family": chart.font.name or "Arial",
        "font_size": chart.font.size.pt * 4 / 3 if chart.font.size else 16,
        "color": _colour(chart.font, scheme) or "#222222",
        "legend": bool(chart.has_legend),
    }
    try:
        return validate_chart_data(data)
    except ValueError:
        return None


def _table_defaults(slide: Any) -> dict:
    """Resolve presentation defaults; leave unavailable built-in table styles explicit."""
    import lxml.etree as etree
    from pptx.oxml.ns import qn

    if slide is None:
        return {}
    presentation = slide.part.package.presentation_part.presentation
    master = slide.slide_layout.slide_master
    try:
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        theme = etree.fromstring(theme_part.blob)
    except (KeyError, AttributeError, etree.XMLSyntaxError):
        theme = etree.Element("unavailable-theme")
    fonts = {}
    for key, name in (("+mn-lt", "minorFont"), ("+mj-lt", "majorFont")):
        node = theme.find(f".//{qn('a:fontScheme')}/{qn('a:' + name)}/{qn('a:latin')}")
        if node is not None and node.get("typeface"):
            fonts[key] = node.get("typeface")
    inherited = presentation._element.find(
        f"{qn('p:defaultTextStyle')}/{qn('a:lvl1pPr')}/{qn('a:defRPr')}"
    )
    size = (
        float(inherited.get("sz")) / 75 if inherited is not None and inherited.get("sz") else None
    )
    latin = inherited.find(qn("a:latin")) if inherited is not None else None
    name = latin.get("typeface") if latin is not None else "+mn-lt"
    return {
        "fontFamily": fonts.get(name, name if not name.startswith("+") else None),
        "fontSize": size,
        "themeFonts": fonts,
    }


def _table_text(cell: Any, scheme: dict[str, str], defaults: dict) -> dict:
    from lxml import html as parser
    from pptx.oxml.ns import qn

    text = _text(
        SimpleNamespace(has_text_frame=True, text_frame=cell.text_frame), scheme, preserve_runs=True
    ) or {"text": "", "richHtml": ""}
    for key in ("fontFamily", "fontSize"):
        if defaults.get(key) is not None:
            text.setdefault(key, defaults[key])
    fonts = defaults.get("themeFonts", {})
    text["fontFamily"] = fonts.get(text.get("fontFamily"), text.get("fontFamily"))
    root = parser.fragment_fromstring(text["richHtml"], create_parent=True)
    spans = iter(root.xpath("./span"))
    if not text["richHtml"]:
        return text
    for paragraph in cell.text_frame.paragraphs:
        properties = paragraph._p.find(f"{qn('a:pPr')}/{qn('a:defRPr')}")
        for child in paragraph._p:
            if child.tag not in {qn("a:r"), qn("a:fld")}:
                continue
            node = next(spans)
            styles = dict(re.findall(r"([\w-]+):([^;]+)", node.get("style", "")))
            if properties is not None:
                if properties.get("sz"):
                    styles.setdefault("font-size", f"{float(properties.get('sz')) / 75:g}px")
                latin = properties.find(qn("a:latin"))
                if latin is not None:
                    styles.setdefault("font-family", latin.get("typeface"))
                for key, prop in (("b", "font-weight"), ("i", "font-style")):
                    if properties.get(key) is not None:
                        styles.setdefault(
                            prop,
                            ("700" if properties.get(key) == "1" else "400")
                            if key == "b"
                            else ("italic" if properties.get(key) == "1" else "normal"),
                        )
            if defaults.get("fontFamily"):
                styles.setdefault("font-family", defaults["fontFamily"])
            if defaults.get("fontSize"):
                styles.setdefault("font-size", f"{defaults['fontSize']:g}px")
            if "font-family" in styles:
                styles["font-family"] = fonts.get(styles["font-family"], styles["font-family"])
            node.set("style", ";".join(f"{k}:{v}" for k, v in styles.items()))
    text["richHtml"] = "".join(
        cast(str, parser.tostring(child, encoding="unicode")) for child in root
    )
    return text


def extract_structured(
    shape: Any, frame: dict, scheme: dict[str, str], *, slide: Any = None
) -> StructuredShape | None:
    if getattr(shape, "has_table", False):
        table = shape.table
        defaults = _table_defaults(slide)
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = _table_text(cell, scheme, defaults)
                cells.append(
                    {
                        "text": text,
                        "fill": _fill(cell, "rect"),
                        "span_width": cell.span_width,
                        "span_height": cell.span_height,
                        "spanned": cell.is_spanned,
                        "padding": [
                            float(getattr(cell, "margin_" + side) or 0) / 9525
                            for side in ("top", "right", "bottom", "left")
                        ],
                    }
                )
            rows.append(cells)
        data = {
            "rows": rows,
            "widths": [c.width / shape.width for c in table.columns],
            "heights": [r.height / shape.height for r in table.rows],
        }
        kind = "table"
    elif getattr(shape, "has_chart", False):
        data = native_chart_data(shape.chart, scheme)
        if data is None:
            return None
        kind = "chart"
    else:
        return None
    return StructuredShape(frame["id"], frame["x"], frame["y"], frame["w"], frame["h"], kind, data)


def validate_chart_data(data: dict) -> dict:
    """Reject malformed/unbounded chart metadata before rendering or exporting it."""
    if not isinstance(data, dict) or data.get("type") not in {"column", "bar"}:
        raise ValueError("Supported charts are clustered column and bar charts.")
    if not isinstance(data.get("font_family", "Arial"), str) or not re.fullmatch(
        r"[\w ,'-]{1,200}", data.get("font_family", "Arial")
    ):
        raise ValueError("Chart font family must be a safe CSS font name.")
    size = data.get("font_size", 16)
    if (
        isinstance(size, bool)
        or not isinstance(size, (float, int))
        or not math.isfinite(size)
        or not 6 <= size <= 96
    ):
        raise ValueError("Chart font size must be between 6 and 96 pixels.")
    if not isinstance(data.get("color", "#222222"), str) or not re.fullmatch(
        r"#[\da-fA-F]{6}", data.get("color", "#222222")
    ):
        raise ValueError("Chart text color must be six-digit hex.")
    if (
        not isinstance(data.get("title", ""), str)
        or len(data.get("title", "")) > 2000
        or not isinstance(data.get("legend", False), bool)
    ):
        raise ValueError("Chart title/legend metadata is invalid.")
    categories, series = data.get("categories"), data.get("series")
    if (
        not isinstance(categories, list)
        or not 1 <= len(categories) <= 50
        or any(not isinstance(c, str) or len(c) > 200 for c in categories)
    ):
        raise ValueError("Charts require 1..50 category labels of at most 200 characters.")
    if not isinstance(series, list) or not 1 <= len(series) <= 8:
        raise ValueError("Charts require 1..8 series.")
    for item in series:
        if (
            not isinstance(item, dict)
            or not isinstance(item.get("name"), str)
            or len(item["name"]) > 200
        ):
            raise ValueError("Each series requires a name of at most 200 characters.")
        values = item.get("values")
        if (
            not isinstance(values, list)
            or len(values) != len(categories)
            or any(
                isinstance(v, bool)
                or not isinstance(v, (int, float))
                or not math.isfinite(v)
                or abs(v) > 1e12
                for v in values
            )
        ):
            raise ValueError(
                "Series values must match categories and be finite numbers within +/-1e12."
            )
        if not isinstance(item.get("color", "#663399"), str) or not re.fullmatch(
            r"#[\da-fA-F]{6}", item.get("color", "#663399")
        ):
            raise ValueError("Chart series colors must be six-digit hex colors.")
    return data


def _label(text: str, x: float, y: float, w: float, h: float, extra: str = "") -> str:
    return (
        f'<p data-text-block="true" style="position:absolute;left:{x:g}px;top:{y:g}px;'
        f'width:{w:g}px;height:{h:g}px;margin:0;{extra}">{html.escape(text)}</p>'
    )


def chart_inner_html(data: dict, width: float, height: float, node_id: str | None = None) -> str:
    """Editable CSS bars/labels reflect the same data exported to the native chart."""
    validate_chart_data(data)
    if (
        not math.isfinite(width)
        or not math.isfinite(height)
        or not 160 <= width <= 20000
        or not 120 <= height <= 20000
    ):
        raise ValueError("Chart box must be at least 160 by 120 pixels.")
    categories, series = data["categories"], data["series"]
    values = [v for s in series for v in s["values"]]
    low, high = min(0, min(values)), max(0, max(values))
    span = high - low or 1
    x0, y0, plotw, ploth = 80.0, 42.0, width - 100, height - 100
    elements = [
        _label(
            str(data.get("title", "")), 8, 4, width - 16, 30, "font-weight:bold;text-align:center;"
        )
    ]
    if data["type"] == "column":
        baseline = y0 + high / span * ploth
        elements.append(
            f'<div style="position:absolute;left:{x0:g}px;top:{baseline:g}px;'
            f'width:{plotw:g}px;border-top:1px solid #888888"></div>'
        )
        step = plotw / len(categories)
        barw = step * 0.7 / len(series)
        for i, category in enumerate(categories):
            elements.append(
                _label(category, x0 + i * step, y0 + ploth + 6, step, 30, "text-align:center;")
            )
            for j, item in enumerate(series):
                value = item["values"][i]
                top = baseline - max(value, 0) / span * ploth
                h = abs(value) / span * ploth
                elements.append(
                    f'<div style="position:absolute;'
                    f"left:{x0 + i * step + step * 0.15 + j * barw:g}px;top:{top:g}px;"
                    f"width:{barw:g}px;height:{h:g}px;"
                    f'background:{item.get("color", "#663399")}"></div>'
                )
    else:
        baseline = x0 - low / span * plotw
        elements.append(
            f'<div style="position:absolute;left:{baseline:g}px;top:{y0:g}px;'
            f'height:{ploth:g}px;border-left:1px solid #888888"></div>'
        )
        step = ploth / len(categories)
        barh = step * 0.7 / len(series)
        for i, category in enumerate(categories):
            elements.append(_label(category, 0, y0 + i * step, 74, step, "text-align:right;"))
            for j, item in enumerate(series):
                value = item["values"][i]
                left = baseline + min(value, 0) / span * plotw
                w = abs(value) / span * plotw
                elements.append(
                    f'<div style="position:absolute;left:{left:g}px;'
                    f"top:{y0 + i * step + step * 0.15 + j * barh:g}px;"
                    f"width:{w:g}px;height:{barh:g}px;"
                    f'background:{item.get("color", "#663399")}"></div>'
                )
    if data.get("legend"):
        elements.append(
            _label(
                " · ".join(s["name"] for s in series),
                10,
                height - 24,
                width - 20,
                22,
                "text-align:center;",
            )
        )
    if node_id is not None:
        from lxml import html as parser

        prefix = hashlib.sha256(node_id.encode()).hexdigest()[:16]
        nodes = parser.fragments_fromstring("".join(elements))
        for index, node in enumerate(nodes):
            node.set("data-node-id", f"chart-{prefix}-{index}")
        return "".join(cast(str, parser.tostring(node, encoding="unicode")) for node in nodes)
    return "".join(elements)


def structured_html(element: StructuredShape, node_id: str, canvas: tuple[int, int]) -> str:
    x, y, w, h = (
        element.x * canvas[0] / 100,
        element.y * canvas[1] / 100,
        element.w * canvas[0] / 100,
        element.h * canvas[1] / 100,
    )
    geometry = f"position:absolute;left:{x:g}px;top:{y:g}px;width:{w:g}px;height:{h:g}px;"
    provenance = (
        f'data-node-id="{html.escape(node_id, quote=True)}" '
        f'data-pptx-shape-id="{element.element_id}"'
    )
    if element.kind == "chart":
        data = element.data
        metadata = html.escape(
            json.dumps(data, ensure_ascii=False, separators=(",", ":")), quote=True
        )
        font = html.escape(str(data.get("font_family", "Arial")), quote=True)
        return (
            f'<div {provenance} data-chart-type="{data["type"]}" data-chart-data="{metadata}" '
            f'style="{geometry}font-family:{font};font-size:{data.get("font_size", 16):g}px;'
            f'color:{data.get("color", "#222222")}">{chart_inner_html(data, w, h, node_id)}</div>'
        )
    rows = []
    for r, row in enumerate(element.data["rows"]):
        cells = []
        for c, cell in enumerate(row):
            if cell["spanned"]:
                continue
            text = cell["text"]
            cell_w = sum(element.data["widths"][c : c + cell["span_width"]]) * w
            cell_h = sum(element.data["heights"][r : r + cell["span_height"]]) * h
            style = (
                "box-sizing:border-box;text-align:left;font-weight:400;"
                f"width:{cell_w:g}px;height:{cell_h:g}px;"
            )
            style += "padding:" + " ".join(f"{v:g}px" for v in cell["padding"]) + ";"
            if cell["fill"]:
                style += f"background:{cell['fill']};"
            for prop, key in [
                ("font-family", "fontFamily"),
                ("color", "color"),
                ("text-align", "align"),
            ]:
                if text.get(key):
                    style += f"{prop}:{text[key]};"
            if text.get("fontSize"):
                style += f"font-size:{text['fontSize']:g}px;"
            tag = "th" if r == 0 else "td"
            content = text.get("richHtml", html.escape(text["text"]))
            cells.append(
                f'<{tag} data-node-id="{node_id}-r{r}-c{c}" data-text-block="true" '
                f'rowspan="{cell["span_height"]}" colspan="{cell["span_width"]}" '
                f'style="{html.escape(style, quote=True)}">{content}</{tag}>'
            )
        rows.append("<tr>" + "".join(cells) + "</tr>")
    return (
        f'<table {provenance} style="{geometry}border-collapse:collapse;table-layout:fixed">'
        + "".join(rows)
        + "</table>"
    )


def replace_chart_html(
    body_html: str, node_id: str, categories: list[str], series: list[dict]
) -> str:
    """Replace a registered chart's values, keeping topology and original styles."""
    import re

    from lxml import html as parser

    root = parser.fragment_fromstring(body_html, create_parent=True)
    matches = root.xpath(".//*[@data-node-id=$ident]", ident=node_id)
    if len(matches) != 1 or not matches[0].get("data-chart-data"):
        raise ValueError("Choose exactly one chart with data-chart-data metadata.")
    node = matches[0]
    old = validate_chart_data(json.loads(node.get("data-chart-data")))
    if node.get("data-chart-type") != old["type"]:
        raise ValueError("Chart type metadata does not match its root.")
    if len(categories) != len(old["categories"]) or len(series) != len(old["series"]):
        raise ValueError(
            "Preserve chart category and series counts; topology changes are unsupported."
        )
    updated = {
        **old,
        "categories": categories,
        "series": [
            {**original, "name": new["name"], "values": new["values"]}
            for original, new in zip(old["series"], series, strict=True)
        ],
    }
    validate_chart_data(updated)
    styles = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", node.get("style", "")))
    try:
        if not all(styles.get(k, "").endswith("px") for k in ("width", "height")):
            raise ValueError()
        width, height = (float(styles[k][:-2]) for k in ("width", "height"))
    except ValueError as exc:
        raise ValueError("Chart needs explicit pixel width and height.") from exc
    content = chart_inner_html(updated, width, height, node_id)
    if updated == old:
        from .sanitize import sanitize_slide_html

        inner = (node.text or "") + "".join(
            cast(str, parser.tostring(child, encoding="unicode")) for child in node
        )
        actual = sanitize_slide_html(inner)
        if not actual.removed and actual.html == sanitize_slide_html(content).html:
            return body_html
    for child in list(node):
        node.remove(child)
    node.text = None
    for child in parser.fragments_fromstring(content):
        node.append(child)
    node.set("data-chart-data", json.dumps(updated, ensure_ascii=False, separators=(",", ":")))
    return "".join(cast(str, parser.tostring(el, encoding="unicode")) for el in root)


def validate_chart_markup(attrs: dict, inner_html: str) -> dict:
    """Require rendered chart parts and structured data to describe the same chart."""
    from .sanitize import sanitize_slide_html

    data = validate_chart_data(json.loads(attrs["data-chart-data"]))
    if attrs.get("data-chart-type") != data["type"]:
        raise ValueError("Chart type metadata does not match its root.")
    styles = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", attrs.get("style", "")))
    allowed = {"position", "left", "top", "width", "height", "font-family", "font-size", "color"}
    if set(styles) - allowed:
        raise ValueError(
            "Chart CSS style changes are unsupported; preserve the native chart style."
        )
    expected_font = str(data.get("font_family", "Arial")).strip("\"' ").lower()
    if "font-family" in styles and styles["font-family"].strip("\"' ").lower() != expected_font:
        raise ValueError("Chart font family diverges from its metadata.")
    if "color" in styles and styles["color"].lower() != data.get("color", "#222222").lower():
        raise ValueError("Chart text color diverges from its metadata.")
    if "font-size" in styles and styles["font-size"] != f"{data.get('font_size', 16):g}px":
        raise ValueError("Chart font size diverges from its metadata.")
    if not all(styles.get(k, "").endswith("px") for k in ("width", "height")):
        raise ValueError("Chart needs explicit pixel width and height.")
    width, height = (float(styles[k][:-2]) for k in ("width", "height"))
    canonical = chart_inner_html(data, width, height, attrs.get("data-node-id"))
    actual = sanitize_slide_html(inner_html)
    expected = sanitize_slide_html(canonical)
    if actual.removed or actual.html != expected.html:
        raise ValueError(
            "Chart HTML diverges from its data/style metadata; "
            "use replace_chart_data to regenerate the preview."
        )
    return data
