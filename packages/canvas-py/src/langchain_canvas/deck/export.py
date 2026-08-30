"""``*.slides.html`` decks compiled to editable PowerPoint files.

The exporter for the canonical HTML deck dialect
(:mod:`langchain_canvas.deck.model`). Every supported node lands as a real
shape — text box, picture, or a rect/ellipse/line drawing — never a rendered
bitmap; a node carrying ``data-pptx-shape-id`` patches that shape on the
deck's own source presentation (its ``lcx:source`` skin) in place, so an
edited caption or moved box round-trips onto the exact shape it came from
instead of stacking a duplicate on top. A node the dialect cannot express as
a native shape (tagged ``data-lcx-fallback="raster"`` by
:func:`~langchain_canvas.deck.baseline.baseline_slide_html`, the same
attribute a generation pass may add for an effect it invents) rasterizes
through an injected :class:`RenderSlideAdapter`; without one it is skipped
and named in :attr:`DeckPptxExporter.degraded_nodes` rather than silently
dropped. Every export reopens what it just wrote and checks the slide count
before reporting success — a corrupt file must fail loudly here, not at the
user's PowerPoint.

Skin bytes: this exporter's ``export()`` stays content-only, matching the
:class:`~langchain_canvas.exporters.Exporter` contract (``export(content,
*, path, title)`` — no store access). The caller inlines the deck's
``lcx:source`` reference into a ``data:`` URI first
(:func:`langchain_canvas.tools.inline_deck_skin`, the deck twin of
``inline_slides_assets``'s template inlining), and :func:`skin_presentation`
decodes that URI internally.
"""

from __future__ import annotations

import base64
import binascii
import io
import re
from dataclasses import dataclass
from html.parser import HTMLParser
from typing import Any, Protocol

from ..converters import ensure_archive_within_limits
from ..exporters import (
    DEFAULT_SLIDE_PAGE_IN,
    PPTX_MIME,
    ExportedFile,
    MissingExporterDependencyError,
    pptx_page_size_inches,
)
from .model import SLIDES_HTML_SUFFIX, DeckParseError, SlideTemplate, parse_deck

__all__ = ["DeckPptxExporter", "RenderSlideAdapter", "percent_box_to_inches", "skin_presentation"]

# Deck node geometry is absolute px on a fixed canvas — the same convention
# `deck.baseline` uses to turn extracted percent geometry into markup.
_CANVAS_PX: dict[str, tuple[int, int]] = {"16:9": (1280, 720), "4:3": (1280, 960)}
_DEFAULT_CANVAS_PX = (1280, 720)
_RATIO_PAGE_IN: dict[str, tuple[float, float]] = {"16:9": DEFAULT_SLIDE_PAGE_IN, "4:3": (10.0, 7.5)}
_EMU_PER_INCH = 914400
_PX_TO_PT = 0.75
_DEFAULT_SHAPE_FILL = "5B5BD6"
_SHAPE_CLASSES = frozenset({"rect", "ellipse", "line"})
_ALIGN_NAMES = frozenset({"left", "center", "right"})


class RenderSlideAdapter(Protocol):
    """Renders one slide's HTML into layout metrics plus a PNG screenshot.

    Matches ``apps/server/app/agent/render.py::render_slide`` (the reference
    implementation) and ``verify.py::_render``'s return shape. This exporter
    consumes only the PNG half, to rasterize a node whose CSS effects have no
    native pptx equivalent — the metrics half belongs to layout verification,
    not export.
    """

    def __call__(self, html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]: ...


@dataclass(frozen=True)
class _SlideNode:
    """One parsed ``lcx-block`` node from a slide's body markup."""

    node_id: str
    pptx_shape_id: str | None
    tag: str
    classes: frozenset[str]
    style: dict[str, str]
    src: str | None
    text: str
    fallback_raster: bool


def _px(value: str | None) -> float:
    """The leading number in a ``"123.4px"`` style value, or ``0.0``."""
    if not value:
        return 0.0
    match = re.match(r"-?\d+(?:\.\d+)?", value.strip())
    return float(match.group(0)) if match else 0.0


_STYLE_DECL_RE = re.compile(r"([a-zA-Z-]+)\s*:\s*([^;]+)")


def _parse_style(style_attr: str) -> dict[str, str]:
    return {key.strip().lower(): value.strip() for key, value in _STYLE_DECL_RE.findall(style_attr)}


def percent_box_to_inches(
    style: dict[str, str],
    page_in: tuple[float, float],
    canvas_px: tuple[int, int] = _DEFAULT_CANVAS_PX,
) -> tuple[float, float, float, float]:
    """A node's ``(left, top, width, height)`` in inches on ``page_in``.

    Deck node geometry is absolute px on a fixed canvas (`deck.baseline`'s
    1280x720 / 1280x960 convention), not the legacy 0-100 percent geometry
    the pre-dialect exporter projected — so this is a fresh implementation
    rather than a reuse of that (now-removed) coordinate system.
    """
    canvas_w, canvas_h = canvas_px
    page_w_in, page_h_in = page_in
    left_in = (_px(style.get("left")) / canvas_w) * page_w_in
    top_in = (_px(style.get("top")) / canvas_h) * page_h_in
    width_in = (_px(style.get("width")) / canvas_w) * page_w_in
    height_in = (_px(style.get("height")) / canvas_h) * page_h_in
    return left_in, top_in, width_in, height_in


_PPTX_DATA_URI_RE = re.compile(rf"^data:{re.escape(PPTX_MIME)};base64,(.+)$", re.IGNORECASE)


def _decode_pptx_data_uri(value: str | None) -> bytes | None:
    if not value:
        return None
    match = _PPTX_DATA_URI_RE.match(value.strip())
    if not match:
        return None
    try:
        return base64.b64decode(match.group(1), validate=True)
    except (binascii.Error, ValueError):
        return None


def skin_presentation(skin_data_uri: str | None) -> Any | None:
    """The deck's source pptx opened as a live presentation, or ``None``.

    Keeps every slide of the skin — unlike a masters/layouts-only style
    template, a canonical-dialect deck's ``lcx:source`` names the very
    presentation its slides were extracted from, and provenance patching
    (:attr:`_SlideNode.pptx_shape_id`) needs each original slide's shapes
    intact to update in place. A missing or unreadable skin degrades to
    ``None`` — the caller then builds a blank presentation. An oversized
    skin raises ``UnsafeArchiveError`` via ``ensure_archive_within_limits`` —
    a decompression bomb is refused loudly, not absorbed.
    """
    data = _decode_pptx_data_uri(skin_data_uri)
    if data is None:
        return None
    ensure_archive_within_limits(data, path="the deck's source skin")
    from pptx import Presentation  # type: ignore[import-untyped]

    try:
        return Presentation(io.BytesIO(data))
    except Exception:  # noqa: BLE001 — any parse failure means "not a usable skin"
        return None


class _NodeParser(HTMLParser):
    """Flat ``lcx-block`` nodes (``data-node-id`` present) in a slide body.

    Nodes are leaf elements in the dialect today (`deck.baseline` emits
    plain ``<div>``/``<img>`` boxes with no nesting), so this tracks one
    open node at a time and a nesting depth for any incidental inner markup
    a generation pass might add (e.g. a ``<b>`` span), collapsing it to its
    plain text rather than trying to preserve rich runs.
    """

    _VOID_TAGS = frozenset({"img", "br", "hr", "input"})

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.nodes: list[_SlideNode] = []
        self._current: tuple[str, dict[str, str | None]] | None = None
        self._depth = 0
        self._text_chars: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attr_map = dict(attrs)
        if self._current is None:
            if "data-node-id" not in attr_map:
                return
            self._current = (tag, attr_map)
            self._text_chars = []
            self._depth = 0
            if tag in self._VOID_TAGS:
                self._flush()
            return
        if tag not in self._VOID_TAGS:
            self._depth += 1

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attr_map = dict(attrs)
        if self._current is None and "data-node-id" in attr_map:
            self._current = (tag, attr_map)
            self._text_chars = []
            self._flush()

    def handle_endtag(self, tag: str) -> None:
        if self._current is None:
            return
        if self._depth > 0:
            self._depth -= 1
            return
        self._flush()

    def handle_data(self, data: str) -> None:
        if self._current is not None:
            self._text_chars.append(data)

    def _flush(self) -> None:
        tag, attrs = self._current  # type: ignore[misc]
        self._current = None
        text = "".join(self._text_chars).strip()
        self._text_chars = []
        classes = frozenset((attrs.get("class") or "").split())
        self.nodes.append(
            _SlideNode(
                node_id=attrs.get("data-node-id") or "",
                pptx_shape_id=attrs.get("data-pptx-shape-id"),
                tag=tag,
                classes=classes,
                style=_parse_style(attrs.get("style") or ""),
                src=attrs.get("src"),
                text=text,
                fallback_raster=attrs.get("data-lcx-fallback") == "raster",
            )
        )


def _parse_slide_nodes(body_html: str) -> list[_SlideNode]:
    parser = _NodeParser()
    parser.feed(body_html)
    parser.close()
    return parser.nodes


_HEX_COLOR_RE = re.compile(r"^#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})$")


def _hex_color(value: str | None) -> str | None:
    if not value:
        return None
    match = _HEX_COLOR_RE.match(value.strip())
    if not match:
        return None
    digits = match.group(1)
    if len(digits) == 3:
        digits = "".join(ch * 2 for ch in digits)
    return digits.upper()


_BORDER_RE = re.compile(r"(\d+(?:\.\d+)?)px\s+\S+\s+(#[0-9a-fA-F]{3,6})")


def _border(value: str | None) -> tuple[str | None, float | None]:
    """``(color, width_px)`` from a CSS ``border`` shorthand, or ``(None, None)``."""
    if not value:
        return None, None
    match = _BORDER_RE.search(value)
    if not match:
        return None, None
    return _hex_color(match.group(2)), float(match.group(1))


_IMAGE_DATA_URI_RE = re.compile(r"^data:image/(?:png|jpe?g|gif);base64,(.+)$", re.IGNORECASE)


def _image_bytes(src: str | None) -> bytes | None:
    if not src:
        return None
    match = _IMAGE_DATA_URI_RE.match(src.strip())
    if not match:
        return None
    try:
        return base64.b64decode(match.group(1), validate=True)
    except (binascii.Error, ValueError):
        return None


def _stem(path: str) -> str:
    name = path.rstrip("/").rsplit("/", 1)[-1]
    if name.lower().endswith(SLIDES_HTML_SUFFIX):
        name = name[: -len(SLIDES_HTML_SUFFIX)]
    return name or "export"


def _safe_name(title: str | None) -> str | None:
    if not title or not title.strip():
        return None
    return re.sub(r"[\s/\\]+", "-", title.strip())


def _blank_layout(presentation: Any) -> Any:
    """The least-furnished layout on ``presentation`` — closest to blank."""
    layouts = [layout for master in presentation.slide_masters for layout in master.slide_layouts]
    for layout in layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    return min(layouts, key=lambda layout: len(layout.placeholders))


def _drop_slides_from(presentation: Any, start_index: int) -> None:
    """Remove every slide at or after ``start_index`` — the deck's own slides win."""
    id_list = presentation.slides._sldIdLst
    extra = list(id_list)[start_index:]
    for slide_id_elem in extra:
        presentation.part.drop_rel(slide_id_elem.rId)
        id_list.remove(slide_id_elem)


def _add_textbox(
    slide: Any, node: _SlideNode, page_in: tuple[float, float], canvas_px: tuple[int, int]
) -> None:
    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    left_in, top_in, width_in, height_in = percent_box_to_inches(node.style, page_in, canvas_px)
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(max(width_in, 0.01)), Inches(max(height_in, 0.01))
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = node.text
    paragraph = frame.paragraphs[0]
    font_size = _px(node.style.get("font-size"))
    if font_size:
        paragraph.font.size = Pt(font_size * _PX_TO_PT)
    if node.style.get("font-weight") == "bold":
        paragraph.font.bold = True
    color = _hex_color(node.style.get("color"))
    if color:
        from pptx.dml.color import RGBColor  # type: ignore[import-untyped]

        paragraph.font.color.rgb = RGBColor.from_string(color)
    align = node.style.get("text-align")
    if align in _ALIGN_NAMES:
        alignments = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        paragraph.alignment = alignments[align]


def _add_picture(
    slide: Any, node: _SlideNode, page_in: tuple[float, float], canvas_px: tuple[int, int]
) -> None:
    from pptx.util import Emu, Inches  # type: ignore[import-untyped]

    data = _image_bytes(node.src)
    if data is None:
        return  # not inlined / not an embeddable type — skip honestly
    left_in, top_in, width_in, height_in = percent_box_to_inches(node.style, page_in, canvas_px)
    left, top = Inches(left_in), Inches(top_in)
    width, height = Inches(max(width_in, 0.01)), Inches(max(height_in, 0.01))
    try:
        picture = slide.shapes.add_picture(io.BytesIO(data), left, top)
    except Exception:  # noqa: BLE001 — corrupt image data; keep the deck
        return
    native_w, native_h = picture.image.size
    if native_w and native_h:
        fit = min(int(width) / native_w, int(height) / native_h)
        picture.width = Emu(int(native_w * fit))
        picture.height = Emu(int(native_h * fit))
        picture.left = Emu(int(left) + (int(width) - int(picture.width)) // 2)
        picture.top = Emu(int(top) + (int(height) - int(picture.height)) // 2)


def _shape_kind(node: _SlideNode) -> str:
    matches = node.classes & _SHAPE_CLASSES
    return next(iter(matches)) if matches else "rect"


def _add_shape(
    slide: Any, node: _SlideNode, page_in: tuple[float, float], canvas_px: tuple[int, int]
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # type: ignore[import-untyped]
    from pptx.util import Emu, Inches, Pt  # type: ignore[import-untyped]

    left_in, top_in, width_in, height_in = percent_box_to_inches(node.style, page_in, canvas_px)
    left, top = Inches(left_in), Inches(top_in)
    width, height = Inches(max(width_in, 0.01)), Inches(max(height_in, 0.01))
    fill_color = _hex_color(node.style.get("background") or node.style.get("background-color"))
    stroke_color, stroke_width = _border(node.style.get("border"))
    kind = _shape_kind(node)

    if kind == "line":
        end_x, end_y = Emu(int(left) + int(width)), Emu(int(top) + int(height))
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, end_x, end_y)
        line_color = stroke_color or fill_color or _DEFAULT_SHAPE_FILL
        connector.line.color.rgb = RGBColor.from_string(line_color)
        connector.line.width = Pt(stroke_width or 2)
        return

    shape_type = MSO_SHAPE.OVAL if kind == "ellipse" else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    else:
        shape.fill.background()
    if stroke_color:
        shape.line.color.rgb = RGBColor.from_string(stroke_color)
        shape.line.width = Pt(stroke_width or 1)
    else:
        shape.line.fill.background()


def _patch_shape(
    shape: Any, node: _SlideNode, page_in: tuple[float, float], canvas_px: tuple[int, int]
) -> None:
    """Update an existing skin shape's text and geometry from a provenance node."""
    from pptx.util import Inches  # type: ignore[import-untyped]

    left_in, top_in, width_in, height_in = percent_box_to_inches(node.style, page_in, canvas_px)
    if width_in > 0 and height_in > 0:
        shape.left = Inches(left_in)
        shape.top = Inches(top_in)
        shape.width = Inches(width_in)
        shape.height = Inches(height_in)
    if node.text and getattr(shape, "has_text_frame", False):
        shape.text_frame.text = node.text


def _verify_reopen(data: bytes, *, expected_slide_count: int) -> None:
    """Reopen ``data`` and check its slide count — export's hard success gate."""
    from pptx import Presentation  # type: ignore[import-untyped]

    try:
        reopened = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — any reopen failure is the same finding
        raise ValueError(f"exported pptx failed to reopen: {exc}") from exc
    actual = len(reopened.slides)
    if actual != expected_slide_count:
        raise ValueError(f"exported pptx has {actual} slides, expected {expected_slide_count}")


class DeckPptxExporter:
    """``.slides.html`` decks as editable PowerPoint files.

    Every supported node lands as a real shape: a text-bearing ``lcx-block``
    becomes a text box, an ``<img>`` becomes a picture (contained in its box,
    centered, never stretched), and a ``rect``/``ellipse``/``line`` class
    becomes the matching drawn shape. A node carrying ``data-pptx-shape-id``
    patches that shape (text and geometry only) on the deck's own source
    presentation in place instead of adding a duplicate; deck slides beyond
    the source's own slide count are added fresh on its emptiest layout,
    and source slides beyond the deck's own slide count are dropped — the
    deck is the source of truth for what ships. A node tagged
    ``data-lcx-fallback="raster"`` (an effect with no native pptx
    equivalent) rasterizes through ``render_adapter`` when one is injected;
    without one it is named in :attr:`degraded_nodes` (``"slide_id:node_id"``)
    rather than silently dropped. Every export reopens its own output and
    checks the slide count before returning, raising on a mismatch.

    Honest limits: without a source skin there is no master or theme
    (elements sit on a blank page sized to the deck's ratio); nested inline
    markup inside a text node (e.g. a generation pass emitting ``<b>``)
    collapses to plain text, not a rich run; the raster fallback crops a
    full-slide screenshot to the node's own box rather than isolating the
    node's paint layer. Requires ``python-pptx`` — installed by the
    ``office`` extra.
    """

    suffixes: tuple[str, ...] = (SLIDES_HTML_SUFFIX,)
    target: str = "pptx"

    def __init__(self, *, render_adapter: RenderSlideAdapter | None = None) -> None:
        self._render_adapter = render_adapter
        self.degraded_nodes: list[str] = []

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingExporterDependencyError(
                "exporting .slides.html to pptx needs python-pptx — install "
                "langchain-canvas[office] or register your own exporter"
            ) from exc

        try:
            deck = parse_deck(content)
        except DeckParseError as exc:
            raise ValueError(f"{path} does not contain a valid deck: {exc}") from exc

        self.degraded_nodes = []
        canvas_px = _CANVAS_PX.get(deck.ratio, _DEFAULT_CANVAS_PX)
        presentation = skin_presentation(deck.source)
        page_in = _page_size_in(presentation, deck.source, deck.ratio)
        if presentation is None:
            from pptx.util import Inches  # type: ignore[import-untyped]

            presentation = Presentation()
            presentation.slide_width = Inches(page_in[0])
            presentation.slide_height = Inches(page_in[1])

        original_count = len(presentation.slides)
        blank_layout = _blank_layout(presentation)
        raster_cache: dict[str, bytes | None] = {}

        for index, slide_tpl in enumerate(deck.slides):
            if index < original_count:
                slide = presentation.slides[index]
                provenance = {f"e{i}": shape for i, shape in enumerate(list(slide.shapes))}
            else:
                slide = presentation.slides.add_slide(blank_layout)
                provenance = {}
            self._export_slide(
                slide, slide_tpl, provenance, deck.ratio, canvas_px, page_in, raster_cache
            )

        if original_count > len(deck.slides):
            _drop_slides_from(presentation, len(deck.slides))

        out = io.BytesIO()
        presentation.save(out)
        data = out.getvalue()
        _verify_reopen(data, expected_slide_count=len(deck.slides))

        name = _safe_name(title) or _safe_name(deck.title) or _stem(path)
        return ExportedFile(data, f"{name}.pptx", PPTX_MIME)

    def _export_slide(
        self,
        slide: Any,
        slide_tpl: SlideTemplate,
        provenance: dict[str, Any],
        ratio: str,
        canvas_px: tuple[int, int],
        page_in: tuple[float, float],
        raster_cache: dict[str, bytes | None],
    ) -> None:
        for node in _parse_slide_nodes(slide_tpl.body_html):
            if node.fallback_raster:
                self._apply_raster(slide, node, slide_tpl, ratio, canvas_px, page_in, raster_cache)
                continue
            shape = provenance.get(node.pptx_shape_id or "")
            if shape is not None:
                _patch_shape(shape, node, page_in, canvas_px)
            elif node.tag == "img":
                _add_picture(slide, node, page_in, canvas_px)
            elif node.classes & _SHAPE_CLASSES:
                _add_shape(slide, node, page_in, canvas_px)
            else:
                _add_textbox(slide, node, page_in, canvas_px)

    def _apply_raster(
        self,
        slide: Any,
        node: _SlideNode,
        slide_tpl: SlideTemplate,
        ratio: str,
        canvas_px: tuple[int, int],
        page_in: tuple[float, float],
        cache: dict[str, bytes | None],
    ) -> None:
        if self._render_adapter is None:
            self.degraded_nodes.append(f"{slide_tpl.slide_id}:{node.node_id}")
            return
        if slide_tpl.slide_id not in cache:
            cache[slide_tpl.slide_id] = self._render_slide_png(slide_tpl, ratio, canvas_px)
        png = cache[slide_tpl.slide_id]
        cropped = _crop_node(png, node) if png is not None else None
        if cropped is None:
            self.degraded_nodes.append(f"{slide_tpl.slide_id}:{node.node_id}")
            return
        from pptx.util import Inches  # type: ignore[import-untyped]

        left_in, top_in, width_in, height_in = percent_box_to_inches(node.style, page_in, canvas_px)
        slide.shapes.add_picture(
            io.BytesIO(cropped),
            Inches(left_in),
            Inches(top_in),
            Inches(max(width_in, 0.01)),
            Inches(max(height_in, 0.01)),
        )

    def _render_slide_png(
        self, slide_tpl: SlideTemplate, ratio: str, canvas_px: tuple[int, int]
    ) -> bytes | None:
        assert self._render_adapter is not None
        width, height = canvas_px
        doc = (
            "<!DOCTYPE html><html><head><meta charset='utf-8'>"
            f"<style>*{{margin:0;padding:0;box-sizing:border-box}} "
            f"html,body{{width:{width}px;height:{height}px;overflow:hidden;position:relative}} "
            f"{slide_tpl.style_css}</style></head><body>{slide_tpl.body_html}</body></html>"
        )
        try:
            _, png = self._render_adapter(doc, ratio=ratio)
        except Exception:  # noqa: BLE001 — a broken adapter degrades, not crashes, the export
            return None
        return png


def _crop_node(png: bytes, node: _SlideNode) -> bytes | None:
    try:
        from PIL import Image  # type: ignore[import-untyped]
    except ImportError:
        return None
    left = round(_px(node.style.get("left")))
    top = round(_px(node.style.get("top")))
    width = max(1, round(_px(node.style.get("width"))))
    height = max(1, round(_px(node.style.get("height"))))
    try:
        image = Image.open(io.BytesIO(png))
        cropped = image.crop((left, top, left + width, top + height))
        out = io.BytesIO()
        cropped.save(out, format="PNG")
        return out.getvalue()
    except Exception:  # noqa: BLE001 — a bad screenshot degrades, not crashes, the export
        return None


def _page_size_in(presentation: Any | None, source: str | None, ratio: str) -> tuple[float, float]:
    """The export page size in inches: the live presentation, then raw skin
    bytes, then the ratio's default (a stub's dimensions are typed Optional;
    a real file always carries them)."""
    if presentation is not None and presentation.slide_width and presentation.slide_height:
        return presentation.slide_width / _EMU_PER_INCH, presentation.slide_height / _EMU_PER_INCH
    skin_bytes = _decode_pptx_data_uri(source)
    if skin_bytes is not None:
        size = pptx_page_size_inches(skin_bytes)
        if size is not None:
            return size
    return _RATIO_PAGE_IN.get(ratio, DEFAULT_SLIDE_PAGE_IN)
