"""Page-limited, render-free census of original ``.pdf``/``.pptx`` pages.

Feeds :mod:`langchain_canvas.deck.patterns` (repeated-layout grouping) without
ever rendering a page, decoding an image, or calling a model — the plan's U1
requirement that observation precede transformation. Two safety layers run
before either parser touches a byte:

1. A hard byte-size cap on the whole source (:attr:`SourceLimits.max_source_bytes`).
2. For ``.pptx`` (a ZIP container), a central-directory-only bomb check
   (:func:`_ensure_pptx_archive_within_limits`) — same technique as
   :func:`~langchain_canvas.converters.ensure_archive_within_limits`, but with
   this module's own (smaller) v1 budget, since a census call is cheaper than
   a full export and can afford to be stricter.

Per-page census then runs bounded object/text counts before any dimension or
render decision, so an oversized single page is rejected
(:class:`OversizedPageError`) before a render budget is even spent on it.

PDF pages are read through ``pypdfium2``'s ``get_objects(textpage=...)``,
which returns text/image/path objects with their bounds without rasterizing
(the same technique :mod:`app.agent.pdf_source` uses for the writer path) —
only ``.get_bounds()``/``.extract()`` are called here, never ``.render()`` or
``.get_bitmap()``. PPTX pages reuse :mod:`langchain_canvas.deck._shapes`'
metadata readers, skipping the one call that decodes bytes
(``shape.image.blob``).
"""

from __future__ import annotations

import hashlib
import io
import zipfile
from dataclasses import dataclass, field
from typing import Any

from ._shapes import _fill, _frame, _is_group, _outline, _text

__all__ = [
    "OversizedPageError",
    "PageInventory",
    "SourceFingerprint",
    "SourceInventoryError",
    "SourceInventoryResult",
    "SourceLimits",
    "TextBoxCensus",
    "inspect_source_pages",
]


class SourceInventoryError(ValueError):
    """The source cannot be safely or meaningfully census'd."""


class OversizedPageError(SourceInventoryError):
    """A single page exceeds a census budget — refused before any render."""


@dataclass(frozen=True)
class SourceLimits:
    """v1 fixed numeric budgets — pre-parse/pre-render, before any decode.

    Cited by the plan as the shared ceiling for both the census in this
    module and, later, ``deck_template_models.py::TemplateBudget`` (task 3's
    per-stage runtime budget). Values are the plan's fixed v1 defaults.
    """

    max_source_bytes: int = 32 * 1024 * 1024
    max_archive_uncompressed_bytes: int = 128 * 1024 * 1024
    max_archive_entries: int = 2000
    max_pdf_page_side_pt: float = 2000.0
    max_objects_per_page: int = 500
    max_text_chars_per_page: int = 50_000
    max_render_pixels: int = 4_000_000
    max_crops_per_page: int = 32
    max_crop_pixels: int = 1_000_000
    max_decoded_image_pixels_per_page: int = 16_000_000


@dataclass(frozen=True)
class SourceFingerprint:
    """Identity for cursor pinning — the hash a follow-up call must match."""

    sha256: str
    page_count: int


@dataclass(frozen=True)
class TextBoxCensus:
    """One text-bearing box, normalized to the page's own width/height."""

    x: float
    y: float
    w: float
    h: float
    text: str
    role: str  # "title" | "body" | "unknown"


@dataclass(frozen=True)
class PageInventory:
    """One page's native census: text boxes, object counts, capability flags."""

    page_number: int  # 1-based
    text_boxes: tuple[TextBoxCensus, ...] = ()
    object_kind_counts: dict[str, int] = field(default_factory=dict)
    char_total: int = 0
    has_text: bool = False
    needs_visual_inspection: bool = False
    capability_issues: tuple[str, ...] = ()
    language: str | None = None


@dataclass(frozen=True)
class SourceInventoryResult:
    """One :func:`inspect_source_pages` call's bounded observation window."""

    fingerprint: SourceFingerprint
    pages: tuple[PageInventory, ...]
    scope_complete: bool
    next_start_page: int | None


def inspect_source_pages(
    data: bytes,
    *,
    path: str,
    start_page: int = 1,
    limit: int = 50,
    limits: SourceLimits | None = None,
) -> SourceInventoryResult:
    """Census pages ``[start_page, start_page + limit)`` without rendering.

    Raises :class:`SourceInventoryError` for an unreadable/unsafe source and
    :class:`OversizedPageError` for a single page whose declared dimensions
    or bounded object/text count exceed the budget — both checked before any
    render or image decode would happen.
    """
    limits = limits or SourceLimits()
    if len(data) > limits.max_source_bytes:
        raise SourceInventoryError(
            f"{path} is {len(data)} bytes — over the "
            f"{limits.max_source_bytes}-byte source limit"
        )
    if start_page < 1:
        raise SourceInventoryError("start_page must be 1 or greater")
    if limit < 1:
        raise SourceInventoryError("limit must be 1 or greater")

    lowered = path.lower()
    if lowered.endswith(".pdf"):
        return _inspect_pdf(data, path=path, start_page=start_page, limit=limit, limits=limits)
    if lowered.endswith(".pptx"):
        return _inspect_pptx(data, path=path, start_page=start_page, limit=limit, limits=limits)
    raise SourceInventoryError(f"{path} is neither a .pdf nor a .pptx source")


# --- PDF ---------------------------------------------------------------------------


def _inspect_pdf(
    data: bytes, *, path: str, start_page: int, limit: int, limits: SourceLimits
) -> SourceInventoryResult:
    try:
        import pypdfium2 as pdfium  # type: ignore[import-untyped]
    except ModuleNotFoundError as exc:  # pragma: no cover - install-time path
        raise SourceInventoryError(
            "reading .pdf needs pypdfium2 — install langchain-canvas[pdf-images]"
        ) from exc

    try:
        document = pdfium.PdfDocument(data)
    except Exception as exc:  # noqa: BLE001 — any malformed file is one failure
        raise SourceInventoryError(f"could not read {path}: {exc}") from exc

    try:
        page_count = len(document)
        sha256 = hashlib.sha256(data).hexdigest()
        end_page = min(start_page + limit - 1, page_count)
        pages: list[PageInventory] = []
        for number in range(start_page, end_page + 1):
            pages.append(_census_pdf_page(document, number, limits, path=path))
        return SourceInventoryResult(
            fingerprint=SourceFingerprint(sha256=sha256, page_count=page_count),
            pages=tuple(pages),
            scope_complete=end_page >= page_count,
            next_start_page=end_page + 1 if end_page < page_count else None,
        )
    finally:
        document.close()


def _census_pdf_page(
    document: Any, number: int, limits: SourceLimits, *, path: str
) -> PageInventory:
    import pypdfium2 as pdfium  # type: ignore[import-untyped]

    page = document[number - 1]
    try:
        width, height = page.get_size()
        if max(width, height) > limits.max_pdf_page_side_pt:
            raise OversizedPageError(
                f"{path} page {number} is {max(width, height):.0f}pt on its long "
                f"side — over the {limits.max_pdf_page_side_pt:.0f}pt limit"
            )
        textpage = page.get_textpage()
        try:
            objects = list(page.get_objects(textpage=textpage))

            if len(objects) > limits.max_objects_per_page:
                raise OversizedPageError(
                    f"{path} page {number} has {len(objects)} objects — over the "
                    f"{limits.max_objects_per_page}-object budget"
                )

            boxes: list[TextBoxCensus] = []
            object_kind_counts: dict[str, int] = {}
            char_total = 0
            for obj in objects:
                if isinstance(obj, pdfium.PdfTextObj):
                    object_kind_counts["text"] = object_kind_counts.get("text", 0) + 1
                    # .extract() needs the (still open) textpage, so text is
                    # pulled here rather than after the loop closes it.
                    text = obj.extract().strip()
                    if not text:
                        continue
                    char_total += len(text)
                    left, bottom, right, top = obj.get_bounds()
                    boxes.append(
                        TextBoxCensus(
                            x=round(max(0.0, left) / width, 4),
                            y=round(max(0.0, height - top) / height, 4),
                            w=round(max(0.0, right - left) / width, 4),
                            h=round(max(0.0, top - bottom) / height, 4),
                            text=text[:160],
                            role="unknown",
                        )
                    )
                elif isinstance(obj, pdfium.PdfImage):
                    object_kind_counts["image"] = object_kind_counts.get("image", 0) + 1
                else:
                    object_kind_counts["shape"] = object_kind_counts.get("shape", 0) + 1
        finally:
            textpage.close()

        if char_total > limits.max_text_chars_per_page:
            raise OversizedPageError(
                f"{path} page {number} has {char_total} text characters — over the "
                f"{limits.max_text_chars_per_page}-character budget"
            )

        has_text = bool(boxes)
        _assign_pdf_roles(boxes, height)
        return PageInventory(
            page_number=number,
            text_boxes=tuple(boxes),
            object_kind_counts=object_kind_counts,
            char_total=char_total,
            has_text=has_text,
            needs_visual_inspection=not has_text,
            capability_issues=(),
            language=None,
        )
    finally:
        page.close()


def _assign_pdf_roles(boxes: list[TextBoxCensus], page_height: float) -> None:
    """Mark the topmost, largest box on the page ``title``; leave the rest.

    A page's declared height sets what "topmost" means; boxes with no
    distinguishing position stay ``unknown`` per the plan's rule to keep
    ``unknown`` when title/body is not clearly resolvable.
    """
    if not boxes:
        return
    top_index = min(range(len(boxes)), key=lambda i: boxes[i].y)
    if boxes[top_index].y > 0.35:
        return  # nothing sits in the top third — no confident title
    boxes[top_index] = TextBoxCensus(
        x=boxes[top_index].x,
        y=boxes[top_index].y,
        w=boxes[top_index].w,
        h=boxes[top_index].h,
        text=boxes[top_index].text,
        role="title",
    )
    for index in range(len(boxes)):
        if index != top_index:
            boxes[index] = TextBoxCensus(
                x=boxes[index].x,
                y=boxes[index].y,
                w=boxes[index].w,
                h=boxes[index].h,
                text=boxes[index].text,
                role="body",
            )


# --- PPTX --------------------------------------------------------------------------

# This module's own (stricter) v1 census budget — see :class:`SourceLimits`.
# Deliberately separate from
# :func:`~langchain_canvas.converters.ensure_archive_within_limits`'s general
# export-path limits, which allow a larger, less time-sensitive export.
_RATIO_ENFORCE_FLOOR_BYTES = 1024 * 1024
_MAX_ARCHIVE_COMPRESSION_RATIO = 200


def _ensure_pptx_archive_within_limits(data: bytes, *, path: str, limits: SourceLimits) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            infos = archive.infolist()
    except zipfile.BadZipFile as exc:
        raise SourceInventoryError(f"{path} is not a readable .pptx container") from exc
    if len(infos) > limits.max_archive_entries:
        raise SourceInventoryError(
            f"{path} holds {len(infos)} parts — over the "
            f"{limits.max_archive_entries}-part census limit"
        )
    total_uncompressed = sum(info.file_size for info in infos)
    if total_uncompressed > limits.max_archive_uncompressed_bytes:
        raise SourceInventoryError(
            f"{path} unpacks to {total_uncompressed} bytes — over the "
            f"{limits.max_archive_uncompressed_bytes}-byte census limit"
        )
    total_compressed = sum(info.compress_size for info in infos)
    if (
        total_uncompressed > _RATIO_ENFORCE_FLOOR_BYTES
        and total_uncompressed > _MAX_ARCHIVE_COMPRESSION_RATIO * max(total_compressed, 1)
    ):
        raise SourceInventoryError(
            f"{path} compresses {total_uncompressed // max(total_compressed, 1)}:1 — "
            f"over the {_MAX_ARCHIVE_COMPRESSION_RATIO}:1 census limit"
        )


_TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
_BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "SUBTITLE"}


def _inspect_pptx(
    data: bytes, *, path: str, start_page: int, limit: int, limits: SourceLimits
) -> SourceInventoryResult:
    _ensure_pptx_archive_within_limits(data, path=path, limits=limits)
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ModuleNotFoundError as exc:  # pragma: no cover - install-time path
        raise SourceInventoryError(
            "reading .pptx needs python-pptx — install langchain-canvas[office]"
        ) from exc

    try:
        deck = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — any malformed file is one failure
        raise SourceInventoryError(f"could not read {path}: {exc}") from exc

    width = deck.slide_width or 0
    height = deck.slide_height or 0
    if not width or not height:
        raise SourceInventoryError(f"{path} declares no slide size")

    slides = list(deck.slides)
    page_count = len(slides)
    sha256 = hashlib.sha256(data).hexdigest()
    end_page = min(start_page + limit - 1, page_count)
    pages = [
        _census_pptx_slide(slides[number - 1], number, width, height, limits, path=path)
        for number in range(start_page, end_page + 1)
    ]
    return SourceInventoryResult(
        fingerprint=SourceFingerprint(sha256=sha256, page_count=page_count),
        pages=tuple(pages),
        scope_complete=end_page >= page_count,
        next_start_page=end_page + 1 if end_page < page_count else None,
    )


def _census_pptx_slide(
    slide: Any, number: int, width: int, height: int, limits: SourceLimits, *, path: str
) -> PageInventory:
    """One slide's native metadata census — groups/tables/charts flagged, not decoded.

    Runs independently of and before
    :func:`langchain_canvas.deck.extract._extract_slide` — the plan requires
    this census to see what extraction skips (groups, in particular) rather
    than counting only what a later conversion kept.
    """
    shapes = list(slide.shapes)
    if len(shapes) > limits.max_objects_per_page:
        raise OversizedPageError(
            f"{path} slide {number} has {len(shapes)} shapes — over the "
            f"{limits.max_objects_per_page}-object budget"
        )

    boxes: list[TextBoxCensus] = []
    object_kind_counts: dict[str, int] = {}
    char_total = 0
    capability_issues: list[str] = []

    for shape_index, shape in enumerate(shapes):
        if _is_group(shape):
            object_kind_counts["group"] = object_kind_counts.get("group", 0) + 1
            capability_issues.append("group")
            continue
        if getattr(shape, "has_table", False):
            object_kind_counts["table"] = object_kind_counts.get("table", 0) + 1
            capability_issues.append("native_table")
            continue
        if getattr(shape, "has_chart", False):
            object_kind_counts["chart"] = object_kind_counts.get("chart", 0) + 1
            capability_issues.append("chart")
            continue
        if _is_smart_art(shape):
            object_kind_counts["smartart"] = object_kind_counts.get("smartart", 0) + 1
            capability_issues.append("smartart")
            continue

        frame = _frame(shape, shape_index, width, height)
        kind = str(getattr(shape, "shape_type", "") or "").upper()
        if "PICTURE" in kind:
            object_kind_counts["picture"] = object_kind_counts.get("picture", 0) + 1
            continue

        text = _text(shape, {})
        if text is not None and frame is not None:
            object_kind_counts["text"] = object_kind_counts.get("text", 0) + 1
            char_total += len(text["text"])
            role = _pptx_placeholder_role(shape)
            boxes.append(
                TextBoxCensus(
                    x=round(frame["x"] / 100, 4),
                    y=round(frame["y"] / 100, 4),
                    w=round(frame["w"] / 100, 4),
                    h=round(frame["h"] / 100, 4),
                    text=text["text"].strip()[:160],
                    role=role,
                )
            )
            # The v1 baseline renderer never re-emits a text box's own fill or
            # outline (only a font colour and, for drawn shapes, their own
            # fill/stroke — see baseline.py::_text_style/_text_box), and it
            # never re-emits a non-default vertical anchor or explicit
            # paragraph spacing (space_before/space_after; line_spacing
            # *is* preserved via CSS `line-height`, so it is deliberately
            # excluded here). Flag each as an unresolved capability issue so
            # `deck_templates.py` can fail closed at finalize.
            if _fill(shape, "rect"):
                capability_issues.append("text_fill")
            stroke, _stroke_weight = _outline(shape)
            if stroke:
                capability_issues.append("text_outline")
            if text.get("verticalAlign") not in (None, "top"):
                capability_issues.append("text_vertical_anchor")
            if text.get("spaceBefore") is not None or text.get("spaceAfter") is not None:
                capability_issues.append("text_paragraph_spacing")
        elif frame is not None:
            object_kind_counts["shape"] = object_kind_counts.get("shape", 0) + 1

    if char_total > limits.max_text_chars_per_page:
        raise OversizedPageError(
            f"{path} slide {number} has {char_total} text characters — over the "
            f"{limits.max_text_chars_per_page}-character budget"
        )

    if _has_master_background_dependency(slide):
        capability_issues.append("master_background")

    has_text = bool(boxes)
    return PageInventory(
        page_number=number,
        text_boxes=tuple(boxes),
        object_kind_counts=object_kind_counts,
        char_total=char_total,
        has_text=has_text,
        needs_visual_inspection=not has_text and not shapes,
        capability_issues=tuple(dict.fromkeys(capability_issues)),
        language=None,
    )


def _pptx_placeholder_role(shape: Any) -> str:
    # python-pptx's `placeholder_format` is a property that *raises*
    # ValueError for a non-placeholder shape rather than returning None, so
    # `getattr(shape, "placeholder_format", None)` does not catch it — an
    # ordinary (non-placeholder) text box, the common case for a raw
    # `add_textbox` shape, must fall through to "unknown" instead of
    # propagating past this census.
    try:
        placeholder_format = shape.placeholder_format
    except ValueError:
        placeholder_format = None
    placeholder_type = getattr(placeholder_format, "type", None)
    name = str(getattr(placeholder_type, "name", "") or "").upper()
    if name in _TITLE_PLACEHOLDER_TYPES:
        return "title"
    if name in _BODY_PLACEHOLDER_TYPES:
        return "body"
    return "unknown"


def _is_smart_art(shape: Any) -> bool:
    """A graphicFrame holding a SmartArt diagram, not a table/chart.

    python-pptx exposes no dedicated accessor; the diagram's data relationship
    is the one native signal available without opening (let alone rendering)
    the diagram itself.
    """
    element = getattr(shape, "_element", None)
    if element is None:
        return False
    try:
        xml = element.xml
    except Exception:  # noqa: BLE001 - an unreadable element carries no flag
        return False
    return "dgm:relIds" in xml or "diagramData" in xml


def _has_master_background_dependency(slide: Any) -> bool:
    """The slide leaves its background to the layout/master rather than stating it.

    An unstated background is a dependency this v1 pipeline cannot resolve
    without opening (and interpreting) the master/layout part — reported so a
    template compiler downstream can refuse rather than silently drop it.
    """
    element = getattr(slide, "_element", None)
    if element is None:
        return False
    try:
        xml = element.xml
    except Exception:  # noqa: BLE001 - an unreadable element carries no flag
        return False
    return "<p:bg>" not in xml
