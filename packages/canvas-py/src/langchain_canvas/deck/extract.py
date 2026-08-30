"""Slide-level structural extraction from uploaded ``.pptx`` bytes.

Returns typed :class:`SlideExtraction` records that
:func:`langchain_canvas.deck.baseline.baseline_slide_html` renders into
dialect-compliant slide markup. The low-level shape-reading helpers this
module builds on — geometry, text formatting, and drawing detection — live
in :mod:`langchain_canvas.deck._shapes`.

Extraction is deterministic and re-runs cheaply from the stored original, so
there is no separate extraction cache: :func:`extract_slides` re-derives a
slide's structure from ``sources/`` bytes on every call rather than caching
it, keeping the store's idempotency story simple (see the module docstring
on ``convert_slide`` for the caller side of this decision).
"""

from __future__ import annotations

import hashlib
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any

from ..converters import ensure_archive_within_limits
from ._shapes import (
    _IMAGE_TYPES,
    PptxImportError,
    _drawing,
    _frame,
    _is_group,
    _notes,
    _scheme,
    _text,
    _with_stroke,
)

__all__ = [
    "ImageAsset",
    "PptxImportError",
    "ShapeGeom",
    "SlideExtraction",
    "TextRun",
    "extract_slides",
    "extracted_text",
]


@dataclass(frozen=True)
class TextRun:
    """One text-bearing shape's words, geometry, and first-run formatting."""

    element_id: str
    x: float
    y: float
    w: float
    h: float
    text: str
    font_size: float | None = None
    bold: bool | None = None
    color: str | None = None
    align: str | None = None
    font_family: str | None = None
    line_height: float | None = None
    vertical_align: str | None = None
    highlight: str | None = None
    space_before: float | None = None
    space_after: float | None = None


@dataclass(frozen=True)
class ShapeGeom:
    """A rectangle, ellipse, or line, with its geometry and drawn style."""

    element_id: str
    x: float
    y: float
    w: float
    h: float
    kind: str
    fill: str | None = None
    stroke: str | None = None
    stroke_width: float | None = None


@dataclass(frozen=True)
class ImageAsset:
    """A picture's geometry plus the bytes the caller stores under ``assets/``.

    ``sha``/``ext`` name the file the caller writes (``assets/{sha}.{ext}``);
    ``data`` is the raw image bytes to write there. Carrying geometry
    alongside the asset identity — rather than only ``(sha, ext, data)`` —
    is what lets :func:`~langchain_canvas.deck.baseline.baseline_slide_html`
    place the picture without a second pass over the presentation.
    """

    element_id: str
    x: float
    y: float
    w: float
    h: float
    sha: str
    ext: str
    data: bytes


@dataclass(frozen=True)
class SlideExtraction:
    """One slide's extracted text, shapes, images, and speaker notes."""

    index: int
    texts: list[TextRun] = field(default_factory=list)
    shapes: list[ShapeGeom] = field(default_factory=list)
    images: list[ImageAsset] = field(default_factory=list)
    # Reserved for relationship-id provenance (e.g. rId -> resolved target);
    # nothing in this change unit populates it yet.
    relationships: dict[str, str] = field(default_factory=dict)
    notes: str = ""


def extract_slides(data: bytes, *, path: str) -> list[SlideExtraction]:
    """Parse presentation bytes into one :class:`SlideExtraction` per slide.

    Runs :func:`~langchain_canvas.converters.ensure_archive_within_limits`
    first — a zip-bomb refusal happens before any parser touches the bytes.
    Raises :class:`PptxImportError` when the bytes are not a readable
    presentation.
    """
    ensure_archive_within_limits(data, path=path)
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ModuleNotFoundError as exc:  # pragma: no cover - install-time path
        raise PptxImportError(
            "reading .pptx needs python-pptx — install langchain-canvas[office]"
        ) from exc

    try:
        deck = Presentation(BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — any malformed file is one failure
        raise PptxImportError(f"could not read the presentation: {exc}") from exc

    width = deck.slide_width or 0
    height = deck.slide_height or 0
    if not width or not height:
        raise PptxImportError("the presentation declares no slide size")

    return [
        _extract_slide(slide, index, width, height) for index, slide in enumerate(deck.slides)
    ]


def _image_bytes(shape: Any) -> tuple[bytes, str] | None:
    """A picture's raw bytes and extension, or ``None`` when unreadable."""
    image = getattr(shape, "image", None)
    if image is None:
        return None
    kind = _IMAGE_TYPES.get((getattr(image, "ext", "") or "").lower())
    if kind is None:
        return None
    try:
        blob = image.blob
    except Exception:  # noqa: BLE001 - an unreadable part drops the picture
        return None
    return blob, kind


def _extract_slide(slide: Any, index: int, width: int, height: int) -> SlideExtraction:
    """One slide's shapes, sorted into text runs, drawn shapes, and images."""
    scheme = _scheme(slide)
    texts: list[TextRun] = []
    shapes: list[ShapeGeom] = []
    images: list[ImageAsset] = []

    for shape_index, shape in enumerate(slide.shapes):
        frame = _frame(shape, shape_index, width, height)
        if frame is None:
            continue
        if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
            continue
        if _is_group(shape):
            continue

        image = _image_bytes(shape)
        if image is not None:
            blob, ext = image
            images.append(
                ImageAsset(
                    element_id=frame["id"],
                    x=frame["x"],
                    y=frame["y"],
                    w=frame["w"],
                    h=frame["h"],
                    sha=hashlib.sha256(blob).hexdigest(),
                    ext=ext,
                    data=blob,
                )
            )
            continue

        text = _text(shape, scheme)
        if text is not None:
            texts.append(
                TextRun(
                    element_id=frame["id"],
                    x=frame["x"],
                    y=frame["y"],
                    w=frame["w"],
                    h=frame["h"],
                    text=text["text"],
                    font_size=text.get("fontSize"),
                    bold=text.get("bold"),
                    color=text.get("color"),
                    align=text.get("align"),
                    font_family=text.get("fontFamily"),
                    line_height=text.get("lineHeight"),
                    vertical_align=text.get("verticalAlign"),
                    highlight=text.get("highlight"),
                    space_before=text.get("spaceBefore"),
                    space_after=text.get("spaceAfter"),
                )
            )
            continue

        drawing = _drawing(shape)
        if drawing is not None:
            if drawing.get("shape") == "line":
                frame = _with_stroke(frame, shape, width, height)
            shapes.append(
                ShapeGeom(
                    element_id=frame["id"],
                    x=frame["x"],
                    y=frame["y"],
                    w=frame["w"],
                    h=frame["h"],
                    kind=drawing["shape"],
                    fill=drawing.get("fill"),
                    stroke=drawing.get("stroke"),
                    stroke_width=drawing.get("strokeWidth"),
                )
            )

    return SlideExtraction(
        index=index, texts=texts, shapes=shapes, images=images, notes=_notes(slide)
    )


def extracted_text(extraction: SlideExtraction) -> list[str]:
    """Non-blank source text this slide's HTML must preserve verbatim.

    Feeds :func:`langchain_canvas.deck.validate.ensure_text_equality` — the
    hard gate that fails a save when generated or hand-edited HTML drops or
    rewords text the presentation actually said.
    """
    return [run.text for run in extraction.texts if run.text.strip()]
