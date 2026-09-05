"""Exporters — canvas files converted into deliverable office formats.

The mirror of :mod:`converters`: converters absorb what users bring in
(``sources/`` bytes -> content blocks), exporters produce what users take
away (canvas files -> office formats). Collaboration stays on the canvas
source files; an exported ``.docx`` or ``.xlsx`` is a snapshot at the door.

An :class:`Exporter` turns one canvas file's text into an
:class:`ExportedFile` (bytes + filename + media type). Exporters are
pluggable: pass your own list to ``create_export_tool`` to replace or extend
the defaults — an in-house rendering pipeline implements the same one-method
contract. Built-in exporters that need a writer library import it lazily and
raise :class:`MissingExporterDependencyError` with an install hint when
absent, so the core package stays dependency-free.
"""

from __future__ import annotations

import base64
import binascii
import io
import json
import re
from collections import Counter
from dataclasses import dataclass
from html.parser import HTMLParser
from typing import Any, Protocol, runtime_checkable

from .converters import ensure_archive_within_limits
from .protocol.artifacts import Slide, SlideElement, SlidePage, SlidesData
from .slide_layout import BULLET_PREFIX, resolve_elements
from .slide_table import table_grid
from .slide_text import fit_scale, grown_height_pct
from .table_merge import merge_rows_into_sheet

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@dataclass(frozen=True)
class ExportedFile:
    """One export result: raw bytes plus how to save or serve them."""

    data: bytes
    filename: str
    media_type: str


@runtime_checkable
class Exporter(Protocol):
    """The pluggable contract: one canvas file's text in, a file out."""

    suffixes: tuple[str, ...]
    """Canvas-path suffixes this exporter reads (lowercase, with the dot)."""

    target: str
    """Output format key (``"docx"``, ``"xlsx"``, ...) — what users ask for."""

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        """Convert one file. ``path`` names the source for filenames/messages."""
        ...


class MissingExporterDependencyError(RuntimeError):
    """A built-in exporter's writer library is not installed.

    The message names the missing package and the extra that installs it, so
    the export tool can relay an honest, actionable error to the agent.
    """


def exporter_for(path: str, target: str, exporters: list[Exporter]) -> Exporter | None:
    """The first exporter matching ``path``'s suffix and the ``target`` format."""
    lowered = path.lower()
    wanted = target.lower().lstrip(".")
    for exporter in exporters:
        if exporter.target == wanted and lowered.endswith(exporter.suffixes):
            return exporter
    return None


def _stem(path: str) -> str:
    """The output filename stem for a canvas path (or directory prefix)."""
    name = path.rstrip("/").rsplit("/", 1)[-1]
    for suffix in (".table.json", ".slides.json", ".html", ".htm", ".md", ".markdown"):
        if name.lower().endswith(suffix):
            name = name[: -len(suffix)]
            break
    return name or "export"


class TableXlsxExporter:
    """``.table.json`` tables as Excel workbooks.

    Values first: the ``columns``/``rows`` shape becomes one sheet with a
    bold header row; an edited Fortune-sheet state (``data.sheet``) exports
    every sheet's cell values and merged ranges. Formula cells stay live
    formulas in the workbook (an ``=``-prefixed row value, or a typed
    formula's ``f`` field in the sheet state) — spreadsheet apps recalculate
    them on open. Rich per-cell styling stays with the browser export menu
    or an adopter's pipeline. Requires ``openpyxl`` — installed by the
    ``xlsx`` extra.
    """

    suffixes: tuple[str, ...] = (".table.json",)
    target: str = "xlsx"

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        try:
            from openpyxl import Workbook  # type: ignore[import-untyped]
            from openpyxl.styles import Font  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingExporterDependencyError(
                "exporting .table.json to xlsx needs openpyxl — install "
                "langchain-canvas[xlsx] or register your own exporter"
            ) from exc

        try:
            artifact = json.loads(content)
        except json.JSONDecodeError as exc:
            raise ValueError(f"{path} does not contain valid table JSON") from exc
        data = artifact.get("data") or {}

        workbook = Workbook()
        default_sheet = workbook.active
        if default_sheet is not None:
            workbook.remove(default_sheet)

        # Rows the agent wrote after the person's last edit are merged into the
        # sheet state first, so the export never drops an agent change.
        sheets = data.get("sheet") or []
        if sheets and data.get("columns"):
            sheets = (
                merge_rows_into_sheet(data["columns"], data.get("rows") or [], sheets) or sheets
            )
        for sheet in sheets:
            ws = workbook.create_sheet(str(sheet.get("name") or "Sheet1"))
            for cell in sheet.get("celldata") or []:
                v = cell.get("v")
                # A typed formula lives in ``f`` — keep it a formula (openpyxl
                # stores no cached value; spreadsheet apps recalculate on open).
                formula = v.get("f") if isinstance(v, dict) else None
                value: Any
                if isinstance(formula, str) and formula:
                    value = formula if formula.startswith("=") else "=" + formula
                elif isinstance(v, dict):
                    value = v.get("v") if v.get("v") is not None else v.get("m")
                else:
                    value = v
                ws.cell(row=int(cell["r"]) + 1, column=int(cell["c"]) + 1, value=value)
            for merge in ((sheet.get("config") or {}).get("merge") or {}).values():
                try:
                    ws.merge_cells(
                        start_row=merge["r"] + 1,
                        start_column=merge["c"] + 1,
                        end_row=merge["r"] + merge["rs"],
                        end_column=merge["c"] + merge["cs"],
                    )
                except (KeyError, ValueError):
                    continue  # overlapping/invalid merge ranges — keep the values

        if not workbook.worksheets:
            ws = workbook.create_sheet("Sheet1")
            columns = data.get("columns") or []
            if columns:
                keys = [column.get("key") for column in columns]
                ws.append([column.get("label") or column.get("key") for column in columns])
                for header_cell in ws[1]:
                    header_cell.font = Font(bold=True)
                for row in data.get("rows") or []:
                    ws.append([row.get(key) for key in keys])

        out = io.BytesIO()
        workbook.save(out)
        return ExportedFile(out.getvalue(), f"{_safe_name(title) or _stem(path)}.xlsx", XLSX_MIME)


# --- html -> docx ----------------------------------------------------------------

_SKIP_TAGS = {"style", "script", "title"}
_PARA_TAGS = {"p", "div", "section", "article", "header", "footer", "blockquote"}
_HEADING_TAGS = {"h1", "h2", "h3", "h4"}


@dataclass
class _Run:
    text: str
    bold: bool
    italic: bool
    strike: bool = False
    code: bool = False


class _HtmlOutline(HTMLParser):
    """Linearize an HTML document into export blocks.

    Blocks are tuples: ``("heading", level, runs)``, ``("para", runs)``,
    ``("bullet", runs)``, ``("table", rows, has_header)``, ``("rule",)``
    and ``("image", bytes)``. Inline bold/italic survive as run flags;
    style/script content is dropped.
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: list[tuple[Any, ...]] = []
        self._runs: list[_Run] = []
        self._kind: tuple[Any, ...] = ("para",)
        self._bold = 0
        self._italic = 0
        self._skip = 0
        self._table: list[list[str]] | None = None
        self._row: list[str] | None = None
        self._cell: list[str] | None = None
        self._has_header = False

    def finish(self) -> None:
        self.close()
        self._flush()

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in _SKIP_TAGS:
            self._skip += 1
            return
        if self._skip:
            return
        if self._table is not None:
            if tag == "tr":
                self._row = []
            elif tag in ("td", "th"):
                self._cell = []
                if tag == "th":
                    self._has_header = True
            return
        if tag == "table":
            self._flush()
            self._table = []
            self._has_header = False
        elif tag in _HEADING_TAGS:
            self._flush()
            self._kind = ("heading", int(tag[1]))
        elif tag == "li":
            self._flush()
            self._kind = ("bullet",)
        elif tag in _PARA_TAGS:
            self._flush()
        elif tag in ("b", "strong"):
            self._bold += 1
        elif tag in ("i", "em"):
            self._italic += 1
        elif tag == "br":
            self._runs.append(_Run("\n", False, False))
        elif tag == "hr":
            self._flush()
            self.blocks.append(("rule",))
        elif tag == "img":
            data = _data_uri_bytes(dict(attrs).get("src") or "")
            if data:
                self._flush()
                self.blocks.append(("image", data))

    def handle_endtag(self, tag: str) -> None:
        if tag in _SKIP_TAGS:
            self._skip = max(0, self._skip - 1)
            return
        if self._skip:
            return
        if self._table is not None:
            if tag in ("td", "th") and self._cell is not None:
                if self._row is None:
                    self._row = []
                self._row.append(_squash("".join(self._cell)).strip())
                self._cell = None
            elif tag == "tr" and self._row is not None:
                self._table.append(self._row)
                self._row = None
            elif tag == "table":
                if self._table:
                    self.blocks.append(("table", self._table, self._has_header))
                self._table = None
            return
        if tag in _HEADING_TAGS or tag == "li" or tag in _PARA_TAGS:
            self._flush()
        elif tag in ("b", "strong"):
            self._bold = max(0, self._bold - 1)
        elif tag in ("i", "em"):
            self._italic = max(0, self._italic - 1)

    def handle_data(self, data: str) -> None:
        if self._skip:
            return
        if self._cell is not None:
            self._cell.append(data)
            return
        if self._table is not None:
            return
        squashed = _squash(data)
        if squashed.strip() or self._runs:
            self._runs.append(_Run(squashed, self._bold > 0, self._italic > 0))

    def _flush(self) -> None:
        runs, self._runs = self._runs, []
        kind, self._kind = self._kind, ("para",)
        trimmed = _trim_runs(runs)
        if trimmed:
            self.blocks.append((*kind, trimmed))


def _squash(text: str) -> str:
    return re.sub(r"[ \t\r\f\v]*\n[ \t\r\f\v]*|[ \t\r\f\v]+", lambda m: "\n" if "\n" in m.group(0) else " ", text)


def _trim_runs(runs: list[_Run]) -> list[_Run]:
    """Strip outer whitespace from a run list; drop it entirely when blank."""
    if not "".join(run.text for run in runs).strip():
        return []
    trimmed = list(runs)
    trimmed[0] = _Run(trimmed[0].text.lstrip(), trimmed[0].bold, trimmed[0].italic)
    trimmed[-1] = _Run(trimmed[-1].text.rstrip(), trimmed[-1].bold, trimmed[-1].italic)
    return [run for run in trimmed if run.text]


def _data_uri_bytes(src: str) -> bytes | None:
    match = re.match(r"^data:image/(?:png|jpe?g|gif);base64,(.+)$", src.strip(), re.IGNORECASE)
    if not match:
        return None
    try:
        return base64.b64decode(match.group(1), validate=True)
    except (binascii.Error, ValueError):
        return None


def _pptx_data_uri_bytes(src: str | None) -> bytes | None:
    """The pptx bytes of an inlined template skin, or ``None``.

    A non-data-URI value here means the skin reference could not be inlined
    (file missing, wrong type) — the caller degrades to the blank export.
    """
    if not src:
        return None
    match = re.match(
        rf"^data:{re.escape(PPTX_MIME)};base64,(.+)$", src.strip(), re.IGNORECASE
    )
    if not match:
        return None
    try:
        return base64.b64decode(match.group(1), validate=True)
    except (binascii.Error, ValueError):
        return None


def pptx_page_size_inches(data: bytes) -> tuple[float, float] | None:
    """The (width, height) a pptx declares, in inches, or ``None``.

    Reads only ``ppt/presentation.xml``'s ``sldSz`` attributes out of the
    ZIP — no python-pptx needed — so tools can learn a skin's page size
    cheaply. Callers gate the bytes with ``ensure_archive_within_limits``
    first when they come from an untrusted upload.
    """
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            xml = archive.read("ppt/presentation.xml").decode("utf-8", "replace")
    except (zipfile.BadZipFile, KeyError):
        return None
    match = re.search(r'<p:sldSz[^>]*\bcx="(\d+)"[^>]*\bcy="(\d+)"', xml)
    if match is None:
        match = re.search(r'<p:sldSz[^>]*\bcy="(\d+)"[^>]*\bcx="(\d+)"', xml)
        if match is None:
            return None
        cy, cx = match.group(1), match.group(2)
    else:
        cx, cy = match.group(1), match.group(2)
    return int(cx) / _EMU_PER_INCH, int(cy) / _EMU_PER_INCH


def _skin_presentation(template: str | None, presentation_cls: Any) -> Any | None:
    """The template skin opened as the base presentation, or ``None``.

    Drops the skin's own slides — the deck's content replaces them — while
    its masters and layouts (logos, backgrounds, headers) stay and style
    every slide added on top. Unreadable skin bytes degrade to ``None`` so
    the export never dies on a bad template; a skin whose ZIP container
    exceeds the safety limits raises ``UnsafeArchiveError`` instead — a
    decompression bomb is an attack to refuse loudly, not a formatting
    mishap to absorb.
    """
    data = _pptx_data_uri_bytes(template)
    if data is None:
        return None
    ensure_archive_within_limits(data, path="the template skin")
    try:
        base = presentation_cls(io.BytesIO(data))
        id_list = base.slides._sldIdLst
        for slide_id in list(id_list):
            base.part.drop_rel(slide_id.rId)
            id_list.remove(slide_id)
    except Exception:  # noqa: BLE001 — any parse failure means "not a usable skin"
        return None
    return base


def _skin_typeface(template: str | None) -> str | None:
    """The face a template skin actually uses most, or ``None``.

    Counts the literal ``typeface`` values in the skin's own slides — the
    faces its author picked run by run — and falls back to its layouts when
    the file carries no slides, as a true template does. A value starting
    with ``+`` is a theme reference, not a face, so it is not counted: the
    theme is where a missing east-asian entry hides. Only the three script
    elements are read — a bullet's or a symbol's font is a dingbat picked
    for one glyph, never the face the deck is set in.

    A family and its weight variants can tie (``Pretendard`` /
    ``Pretendard Light`` / ``Pretendard SemiBold``, three uses each). The
    plainest name wins — shortest, then alphabetical — which lands on the
    family itself rather than on one of its weights, and keeps the answer
    the same whatever order the archive lists its parts in.
    """
    data = _pptx_data_uri_bytes(template)
    if data is None:
        return None
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = archive.namelist()
            for folder in ("ppt/slides/slide", "ppt/slideLayouts/slideLayout"):
                faces: Counter[str] = Counter()
                for name in names:
                    if not (name.startswith(folder) and name.endswith(".xml")):
                        continue
                    xml = archive.read(name).decode("utf-8", "replace")
                    faces.update(
                        face
                        for face in re.findall(
                            r'<a:(?:latin|ea|cs)\b[^>]*typeface="([^"]+)"', xml
                        )
                        if not face.startswith("+")
                    )
                if faces:
                    most = max(faces.values())
                    return min(
                        (face for face, uses in faces.items() if uses == most),
                        key=lambda face: (len(face), face),
                    )
    except (zipfile.BadZipFile, KeyError):
        return None
    return None


def _content_layout(presentation: Any) -> Any:
    """The least-furnished layout — the closest thing to a blank canvas.

    Prefers a layout literally named "Blank" (the stock template's index 6),
    else the one with the fewest placeholders, so skinned exports draw on
    the emptiest surface the template offers while inheriting its master.
    """
    layouts = [
        layout
        for master in presentation.slide_masters
        for layout in master.slide_layouts
    ]
    for layout in layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    return min(layouts, key=lambda layout: len(layout.placeholders))


def _safe_name(title: str | None) -> str | None:
    if not title or not title.strip():
        return None
    return re.sub(r"[\s/\\]+", "-", title.strip())


class HtmlDocxExporter:
    """``.html`` pages as Word documents.

    A deliberate subset survives the trip: headings (h1-h4), paragraph text
    with inline bold/italic, bullet items, flat tables, rules (``hr``)
    and ``data:``-URI images. CSS layout does not — the canvas file stays the
    source of truth; the ``.docx`` is a snapshot at the door. Requires
    ``python-docx`` — installed by the ``office`` extra.
    """

    suffixes: tuple[str, ...] = (".html", ".htm")
    target: str = "docx"

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        outline = _HtmlOutline()
        outline.feed(content)
        outline.finish()
        return _blocks_to_docx(outline.blocks, path=path, title=title, source=".html")


def _blocks_to_docx(
    blocks: list[tuple[Any, ...]], *, path: str, title: str | None, source: str
) -> ExportedFile:
    """Assemble export blocks into a ``.docx`` — one door for every text format."""
    try:
        from docx import Document  # type: ignore[import-untyped]
        from docx.enum.text import WD_BREAK  # type: ignore[import-untyped]
        from docx.shared import Inches  # type: ignore[import-untyped]
    except ImportError as exc:
        raise MissingExporterDependencyError(
            f"exporting {source} to docx needs python-docx — install "
            "langchain-canvas[office] or register your own exporter"
        ) from exc

    document = Document()
    for block in blocks:
        kind = block[0]
        if kind == "heading":
            paragraph = document.add_heading("", level=min(int(block[1]), 4))
            _add_runs(paragraph, block[2])
        elif kind == "bullet":
            _add_runs(document.add_paragraph(style=_list_style("List Bullet", block)), block[1])
        elif kind == "numbered":
            _add_runs(document.add_paragraph(style=_list_style("List Number", block)), block[1])
        elif kind == "para":
            _add_runs(document.add_paragraph(), block[1])
        elif kind == "quote":
            _add_runs(document.add_paragraph(style="Intense Quote"), block[1])
        elif kind == "code":
            _add_code(document, block[1])
        elif kind == "table":
            _add_table(document, block[1], block[2])
        elif kind == "rule":
            _add_rule(document)
        elif kind == "page_break":
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        elif kind == "image":
            try:
                document.add_picture(io.BytesIO(block[1]), width=Inches(6))
            except Exception:  # noqa: BLE001 — corrupt image data; keep the document
                continue

    out = io.BytesIO()
    document.save(out)
    return ExportedFile(out.getvalue(), f"{_safe_name(title) or _stem(path)}.docx", DOCX_MIME)


_MD_INLINE = re.compile(
    r"(\*\*\*[^*]+\*\*\*|\*\*[^*]+\*\*|\*[^*]+\*|___[^_]+___|__[^_]+__|_[^_]+_|~~[^~]+~~|`[^`]+`)"
)
_MD_IMAGE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")
_MD_NUMBERED = re.compile(r"^\s*\d+[.)]\s+")
_MD_BULLET = re.compile(r"^\s*[-*+]\s+")


def _md_runs(text: str) -> list[_Run]:
    """Inline markdown as runs: bold, italic, bold-italic, strikethrough and
    code (a monospace run without its ticks)."""
    runs: list[_Run] = []
    for part in _MD_INLINE.split(text):
        if not part:
            continue
        if part.startswith("***") and part.endswith("***") and len(part) > 6:
            runs.append(_Run(part[3:-3], True, True))
        elif part.startswith("___") and part.endswith("___") and len(part) > 6:
            runs.append(_Run(part[3:-3], True, True))
        elif part.startswith("~~") and part.endswith("~~") and len(part) > 4:
            runs.append(_Run(part[2:-2], False, False, strike=True))
        elif part.startswith("**") and part.endswith("**") and len(part) > 4:
            runs.append(_Run(part[2:-2], True, False))
        elif part.startswith("__") and part.endswith("__") and len(part) > 4:
            runs.append(_Run(part[2:-2], True, False))
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            runs.append(_Run(part[1:-1], False, True))
        elif part.startswith("_") and part.endswith("_") and len(part) > 2:
            runs.append(_Run(part[1:-1], False, True))
        elif part.startswith("`") and part.endswith("`") and len(part) > 2:
            runs.append(_Run(part[1:-1], False, False, code=True))
        else:
            runs.append(_Run(part, False, False))
    return runs


def _md_depth(line: str) -> int:
    """List nesting from leading spaces: 0, 1 or 2 (two spaces per level)."""
    return min(2, (len(line) - len(line.lstrip(" "))) // 2)


def _join_soft_lines(lines: list[str]) -> str:
    """One paragraph from its source lines: a line that ends in two spaces
    keeps a line break; the rest wrap into one line."""
    out = ""
    for i, raw in enumerate(lines):
        piece = raw.strip()
        if i:
            out += "\n" if lines[i - 1].endswith("  ") else " "
        out += piece
    return out


def _md_cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _md_is_separator(line: str) -> bool:
    body = line.strip().strip("|").replace("|", "").replace(" ", "")
    return bool(body) and set(body) <= {"-", ":"} and "-" in body


def _decode_data_uri(url: str) -> bytes | None:
    if not url.startswith("data:") or ";base64," not in url:
        return None
    try:
        return base64.b64decode(url.split(";base64,", 1)[1], validate=True)
    except (binascii.Error, ValueError):
        return None


def _markdown_blocks(text: str) -> list[tuple[Any, ...]]:
    """A deliberate markdown subset as export blocks.

    Headings, paragraphs with inline bold/italic/strike/code (a line ending
    in two spaces breaks the line), bullet and numbered items up to three
    levels deep, block quotes, pipe tables with the same inline marks in
    their cells, fenced code as monospace, thematic breaks as rules, and
    ``data:`` images. Everything else lands as plain text — the canvas file
    stays the source of truth; the ``.docx`` is a snapshot at the door,
    exactly like the HTML one.
    """
    blocks: list[tuple[Any, ...]] = []
    paragraph: list[str] = []
    lines = text.split("\n")

    def flush() -> None:
        if paragraph:
            blocks.append(("para", _md_runs(_join_soft_lines(paragraph))))
            paragraph.clear()

    index = 0
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()
        if not stripped:
            flush()
            index += 1
            continue
        if stripped.startswith("```"):
            flush()
            index += 1
            code: list[str] = []
            while index < len(lines) and not lines[index].strip().startswith("```"):
                code.append(lines[index])
                index += 1
            index += 1
            blocks.append(("code", "\n".join(code)))
            continue
        if stripped.startswith("#"):
            flush()
            level = len(stripped) - len(stripped.lstrip("#"))
            blocks.append(("heading", level, _md_runs(stripped[level:].strip())))
            index += 1
            continue
        if stripped in ("---", "***", "___"):
            flush()
            blocks.append(("rule",))
            index += 1
            continue
        image = _MD_IMAGE.match(stripped)
        if image:
            flush()
            data = _decode_data_uri(image.group(2))
            if data is not None:
                blocks.append(("image", data))
            elif image.group(1):
                blocks.append(("para", [_Run(image.group(1), False, False)]))
            index += 1
            continue
        if (
            stripped.startswith("|")
            and index + 1 < len(lines)
            and _md_is_separator(lines[index + 1])
        ):
            flush()
            rows = [_md_cells(stripped)]
            index += 2  # past the separator line
            while index < len(lines) and lines[index].strip().startswith("|"):
                rows.append(_md_cells(lines[index]))
                index += 1
            blocks.append(("table", [[_md_runs(cell) for cell in row] for row in rows], True))
            continue
        if _MD_BULLET.match(line):
            flush()
            blocks.append(("bullet", _md_runs(_MD_BULLET.sub("", line, count=1)), _md_depth(line)))
            index += 1
            continue
        if _MD_NUMBERED.match(line):
            flush()
            blocks.append(("numbered", _md_runs(_MD_NUMBERED.sub("", line, count=1)), _md_depth(line)))
            index += 1
            continue
        if stripped.startswith(">"):
            flush()
            quote: list[str] = []
            while index < len(lines) and lines[index].strip().startswith(">"):
                quote.append(lines[index].strip().lstrip(">").strip())
                index += 1
            blocks.append(("quote", _md_runs(_join_soft_lines(quote))))
            continue
        paragraph.append(line.rstrip("\n"))
        index += 1
    flush()
    return blocks


class MarkdownDocxExporter:
    """``.md`` documents as Word files.

    The same deliberate subset as the HTML door (headings, inline bold and
    italic, bullet and numbered lists, pipe tables, page breaks, ``data:``
    images), read from markdown instead. Requires ``python-docx`` —
    installed by the ``office`` extra.
    """

    suffixes: tuple[str, ...] = (".md", ".markdown")
    target: str = "docx"

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        return _blocks_to_docx(_markdown_blocks(content), path=path, title=title, source=".md")


_CODE_FACE = "Consolas"


def _add_runs(paragraph: Any, runs: list[_Run]) -> None:
    for run in runs:
        added = paragraph.add_run(run.text)
        if run.bold:
            added.bold = True
        if run.italic:
            added.italic = True
        if run.strike:
            added.font.strike = True
        if run.code:
            added.font.name = _CODE_FACE


def _list_style(base: str, block: tuple[Any, ...]) -> str:
    """``List Bullet`` / ``List Bullet 2`` / ``List Bullet 3`` by nesting depth."""
    depth = int(block[2]) if len(block) > 2 else 0
    return base if depth == 0 else f"{base} {depth + 1}"


def _add_code(document: Any, text: str) -> None:
    """A fenced block as one monospace paragraph, indentation and line breaks kept."""
    from docx.shared import Pt  # type: ignore[import-untyped]

    paragraph = document.add_paragraph()
    for i, line in enumerate(text.split("\n")):
        if i:
            paragraph.add_run().add_break()
        run = paragraph.add_run(line)
        run.font.name = _CODE_FACE
        run.font.size = Pt(9)


def _add_rule(document: Any) -> None:
    """A thematic break as an empty paragraph with a bottom border — a line,
    which is what ``---`` means in a document, not a new page."""
    from docx.oxml import OxmlElement  # type: ignore[import-untyped]
    from docx.oxml.ns import qn  # type: ignore[import-untyped]

    paragraph = document.add_paragraph()
    props = paragraph._p.get_or_add_pPr()
    border = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    for key, value in (("w:val", "single"), ("w:sz", "6"), ("w:space", "1"), ("w:color", "999999")):
        bottom.set(qn(key), value)
    border.append(bottom)
    props.append(border)


def _add_table(document: Any, rows: list[list[Any]], has_header: bool) -> None:
    """Rows of cells; a cell is a string or a list of runs (inline marks kept)."""
    width = max(len(row) for row in rows)
    if not width:
        return
    table = document.add_table(rows=len(rows), cols=width)
    table.style = "Table Grid"
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            if isinstance(value, str):
                cell.text = value
            else:
                _add_runs(cell.paragraphs[0], value)
            if has_header and i == 0:
                for run in cell.paragraphs[0].runs:
                    run.bold = True


# --- slides -> pptx --------------------------------------------------------

# The editor's slide canvas is 1280x720 px; the exported deck is 10 x 5.625
# inches (16:9) — the same page the browser-side pptx export uses, so the two
# doors produce the same-looking deck. Element geometry is percent-based.
_SLIDE_WIDTH_IN = 10.0
_SLIDE_HEIGHT_IN = 5.625
# The classic canvas page — what percent geometry means when a deck carries
# no `page` of its own. Tools use it to re-fit decks onto another page.
DEFAULT_SLIDE_PAGE_IN = (_SLIDE_WIDTH_IN, _SLIDE_HEIGHT_IN)
# python-pptx page dimensions are Emu integers (914400 per inch).
_EMU_PER_INCH = 914400
# Element font sizes are px on the 1280px-wide slide; PowerPoint wants points.
_PX_TO_PT = 0.75
_EMU_PER_POINT = 12700
# The hanging indent a bulleted line gets, as a multiple of its own type size.
_BULLET_HANG = 1.2
_DEFAULT_FONT_PX = 24.0
_DEFAULT_SHAPE_FILL = "5B5BD6"

_HEX_COLOR_PATTERN = re.compile(r"^#?([0-9a-fA-F]{3}|[0-9a-fA-F]{6})$")


def _hex_rgb(value: str | None) -> str | None:
    """A 6-digit RGB hex string for a ``#rgb`` / ``#rrggbb`` color, else None."""
    if not value:
        return None
    match = _HEX_COLOR_PATTERN.match(value.strip())
    if not match:
        return None
    digits = match.group(1)
    if len(digits) == 3:
        digits = "".join(ch * 2 for ch in digits)
    return digits.upper()


def _resolved_slide_elements(slide: Slide, page: SlidePage | None = None) -> list[SlideElement]:
    """What is actually on the slide: explicit edits win, else derive."""
    return resolve_elements(slide, page)


# PowerPoint's "No Style, No Grid" table style: the file then draws only the
# fills and lines the element states, which is what the canvas draws.
_TABLE_STYLE_NONE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _name_face(run: Any, face: str | None, qn: Any) -> None:
    """Name the run's face for every script.

    ``a:latin`` covers Latin script only — Korean and the rest of CJK read
    ``a:ea``, and complex scripts read ``a:cs``. A theme's east-asian entry
    is often empty, so a run naming only its Latin face leaves Hangul with
    nowhere to go. python-pptx stops at ``a:latin``; the schema fixes the
    sibling order.
    """
    if not face:
        return
    latin = run.font._rPr.get_or_add_latin()
    latin.set("typeface", face)
    latin.addnext(latin.makeelement(qn("a:cs"), {"typeface": face}))
    latin.addnext(latin.makeelement(qn("a:ea"), {"typeface": face}))


def _add_slide_table(
    slide: Any,
    element: SlideElement,
    box: tuple[Any, Any, Any, Any],
    *,
    scale: float,
    text_color: str | None,
    skin_face: str | None,
    alignments: dict[str, Any],
    anchors: dict[str, Any],
) -> None:
    """A table element as a real PowerPoint table — columns at their widths,
    merged cells merged, each cell with the fill, grid line and text the
    element states — so the received file's table is still a table."""
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    from pptx.util import Emu, Pt  # type: ignore[import-untyped]

    grid = table_grid(element.model_dump(by_alias=True, exclude_none=True))
    if grid is None:
        return
    left, top, width, height = box
    frame = slide.shapes.add_table(grid.n_rows, grid.n_cols, left, top, width, height)
    table = frame.table
    properties = table._tbl.tblPr
    style_id = properties.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = properties.makeelement(qn("a:tableStyleId"), {})
        properties.append(style_id)
    style_id.text = _TABLE_STYLE_NONE
    table.first_row = bool(element.header)
    table.horz_banding = False
    for index, share in enumerate(grid.col_widths):
        table.columns[index].width = Emu(int(int(width) * share / 100.0))
    for index, share in enumerate(grid.row_heights):
        table.rows[index].height = Emu(int(int(height) * share / 100.0))
    for (r, c), (row_span, col_span) in grid.spans.items():
        table.cell(r, c).merge(table.cell(r + row_span - 1, c + col_span - 1))

    stroke = _hex_rgb(element.stroke)
    line_width = str(int((element.stroke_width or 1) * _PX_TO_PT * _EMU_PER_POINT))
    for r in range(grid.n_rows):
        for c in range(grid.n_cols):
            if (r, c) in grid.covered:
                continue
            cell = table.cell(r, c)
            own = grid.styles.get((r, c), {})
            fill = _hex_rgb(own.get("fill")) or _hex_rgb(element.fill)
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(fill)
            else:
                cell.fill.background()
            if element.vertical_align:
                cell.vertical_anchor = anchors[element.vertical_align]
            # The grid: every edge of every cell, drawn or declared absent, so
            # neither the style nor the viewer invents a line. The schema
            # wants the four edges ahead of the fill.
            cell_properties = cell._tc.get_or_add_tcPr()
            for position, side in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
                line = cell_properties.makeelement(qn(side), {"w": line_width} if stroke else {})
                if stroke:
                    solid = line.makeelement(qn("a:solidFill"), {})
                    solid.append(solid.makeelement(qn("a:srgbClr"), {"val": stroke}))
                    line.append(solid)
                else:
                    line.append(line.makeelement(qn("a:noFill"), {}))
                cell_properties.insert(position, line)

            text_frame = cell.text_frame
            text_frame.word_wrap = True
            size_px = own.get("fontSize") or element.font_size or _DEFAULT_FONT_PX
            size_pt = float(size_px) * _PX_TO_PT * scale
            bold = own.get("bold")
            if bold is None:
                bold = (
                    element.bold if element.bold is not None else (bool(element.header) and r == 0)
                )
            colour = _hex_rgb(own.get("color")) or _hex_rgb(element.color) or text_color
            alignment = alignments.get(own.get("align") or element.align or "left")
            for index, line_text in enumerate(grid.rows[r][c].split("\n")):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.alignment = alignment
                run = paragraph.add_run()
                run.text = line_text
                run.font.size = Pt(size_pt)
                run.font.bold = bool(bold)
                if colour is not None:
                    run.font.color.rgb = RGBColor.from_string(colour)
                _name_face(run, element.font_family or skin_face, qn)


class SlidesPptxExporter:
    """``.slides.json`` decks as editable PowerPoint files.

    Every element lands as a real shape — text boxes with runs, pictures,
    drawing shapes — never a rendered bitmap, so the received file reopens
    fully editable. What survives the trip: percent geometry mapped onto a
    16:9 page (padding insets included), text with size / bold / color /
    alignment, ``data:``-URI images (png / jpeg / gif) placed contained in
    their box, rect / ellipse / line shapes, solid ``#hex`` slide
    backgrounds, and speaker notes. Structured slides (title / bullets /
    layout) derive the same elements the canvas renders.

    Template skins: when the deck's ``template`` field references a pptx
    (inlined to a data URI before export), the export opens that file as its
    base — the skin's masters and layouts style every slide, so the
    original's logos, backgrounds, and headers survive; the skin's own
    slides are dropped and its native page size is kept (percent geometry
    projects onto any page). Text takes the face the skin uses most, named
    for Latin, east-asian, and complex scripts alike, so the deck reads in
    the template's own type. A missing or unreadable skin degrades to the
    blank default below.

    Honest limits: without a skin there is no master or theme (elements sit
    on a blank 16:9 layout); the canvas preview never renders the skin
    (export-time only); no animations, transitions, or SmartArt; image/url
    backgrounds are skipped; an explicit slide ``background`` paints over
    the skin's; non-data-URI image references are skipped (inline assets
    before exporting); without a skin there is no face to name, so fonts
    fall back to whatever the viewer has installed.
    Requires ``python-pptx`` — installed by the ``office`` extra.
    """

    suffixes: tuple[str, ...] = (".slides.json",)
    target: str = "pptx"

    def export(self, content: str, *, path: str, title: str | None = None) -> ExportedFile:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
            from pptx.enum.shapes import (  # type: ignore[import-untyped]
                MSO_CONNECTOR,
                MSO_SHAPE,
            )
            from pptx.enum.text import (  # type: ignore[import-untyped]
                MSO_ANCHOR,
                MSO_AUTO_SIZE,
                PP_ALIGN,
            )
            from pptx.oxml.ns import qn  # type: ignore[import-untyped]
            from pptx.util import Emu, Inches, Pt  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingExporterDependencyError(
                "exporting .slides.json to pptx needs python-pptx — install "
                "langchain-canvas[office] or register your own exporter"
            ) from exc

        try:
            envelope = json.loads(content)
        except json.JSONDecodeError as exc:
            raise ValueError(f"{path} does not contain valid slides JSON") from exc
        if not isinstance(envelope, dict):
            raise ValueError(f"{path} does not contain a slides envelope")
        if not isinstance(envelope.get("data"), dict):
            # A deck written without the envelope used to export as one blank
            # slide — silence that hides the real mistake. Name the shape so
            # a tool-calling model can correct itself.
            raise ValueError(
                f'{path} has no "data" envelope — write slides files as '
                '{"type": "slides", "data": {"slides": [...]}} '
                "(element x/y/w/h are percent of the slide, 0-100)"
            )
        try:
            deck = SlidesData.model_validate(envelope.get("data") or {})
        except Exception as exc:  # noqa: BLE001 — pydantic detail relayed honestly
            raise ValueError(f"{path} does not contain a valid slide deck: {exc}") from exc
        envelope_title = envelope.get("title")
        if not isinstance(envelope_title, str):
            envelope_title = None

        alignments = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }

        anchors = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        # The deck's own page is the coordinate space its percent geometry
        # refers to; absent means the classic 16:9 canvas.
        canvas_w = deck.page.width_in if deck.page else _SLIDE_WIDTH_IN
        canvas_h = deck.page.height_in if deck.page else _SLIDE_HEIGHT_IN
        presentation = _skin_presentation(deck.template, Presentation)
        skin_face = _skin_typeface(deck.template) if presentation is not None else None
        if presentation is None:
            presentation = Presentation()
            presentation.slide_width = Inches(canvas_w)
            presentation.slide_height = Inches(canvas_h)
        # A skin keeps its native page size. (The stub types the dimensions
        # Optional; a real file always carries them.)
        page_w_in = (presentation.slide_width or Inches(canvas_w)) / _EMU_PER_INCH
        page_h_in = (presentation.slide_height or Inches(canvas_h)) / _EMU_PER_INCH
        # When the page differs from the deck's canvas (a skin with another
        # aspect ratio), project with ONE uniform scale and center the
        # content — width and height stretched separately would distort
        # every shape (a circle approved on the 16:9 preview must stay a
        # circle on a 4:3 page). The leftover margins show the skin's own
        # background, which is exactly what a branded page is for.
        scale = min(page_w_in / canvas_w, page_h_in / canvas_h)
        offset_x = (page_w_in - canvas_w * scale) / 2.0
        offset_y = (page_h_in - canvas_h * scale) / 2.0
        blank_layout = _content_layout(presentation)

        for slide_model in deck.slides or [Slide()]:
            slide = presentation.slides.add_slide(blank_layout)
            background = _hex_rgb(slide_model.background)
            if background is not None:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(background)

            # A slide `padding` (percent) insets the content area, exactly as
            # the editor and the browser-side exports apply it.
            pad = (slide_model.padding or 0.0) / 100.0
            span = 1.0 - 2.0 * pad

            def inch_box(element: SlideElement) -> tuple[Any, Any, Any, Any]:
                left = offset_x + (pad + (element.x / 100.0) * span) * canvas_w * scale
                top = offset_y + (pad + (element.y / 100.0) * span) * canvas_h * scale
                width = (element.w / 100.0) * span * canvas_w * scale
                height = (element.h / 100.0) * span * canvas_h * scale
                return Inches(left), Inches(top), Inches(width), Inches(height)

            for element in _resolved_slide_elements(slide_model, deck.page):
                left, top, width, height = inch_box(element)
                if element.type == "text":
                    fit = element.autofit or "none"
                    words = element.text or ""
                    font_px = element.font_size or _DEFAULT_FONT_PX
                    if fit == "shape":
                        # The box grows with its text, as it did in the file
                        # it came from. The height written is the grown one,
                        # so a viewer that does not re-fit on open still
                        # shows every line.
                        grown = grown_height_pct(
                            words, font_px, element.w, element.h, element.line_height,
                            (canvas_w, canvas_h),
                        )
                        if grown > element.h:
                            _, _, _, height = inch_box(element.model_copy(update={"h": grown}))
                    box = slide.shapes.add_textbox(left, top, width, height)
                    if element.rotation:
                        box.rotation = element.rotation
                    frame = box.text_frame
                    # An element marked wrap: false stays the one line its
                    # original was; everything else wraps like the canvas.
                    frame.word_wrap = element.wrap is not False
                    if fit == "shape":
                        frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    elif fit == "text":
                        # The type shrinks to its box. PowerPoint re-fits
                        # when the text is next edited; the scale written
                        # here is what shows until then.
                        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        shrink = fit_scale(
                            words, font_px, element.w, element.h, element.line_height,
                            (canvas_w, canvas_h),
                        )
                        if shrink < 1.0:
                            frame._bodyPr.find(qn("a:normAutofit")).set(
                                "fontScale", str(int(round(shrink * 100000)))
                            )
                    else:
                        # Boxes are measured to hold their own text, so
                        # nothing needs shrinking — and shrink-to-fit is what
                        # made two bullets of the same size render at
                        # different sizes.
                        frame.auto_size = MSO_AUTO_SIZE.NONE
                    frame.margin_left = frame.margin_right = Emu(0)
                    frame.margin_top = frame.margin_bottom = Emu(0)
                    # Where the text sits inside its box. Left unset, a box
                    # measured for centred text draws it against the top edge.
                    if element.vertical_align:
                        frame.vertical_anchor = anchors[element.vertical_align]
                    color = _hex_rgb(element.color) or _hex_rgb(slide_model.text_color)
                    alignment = alignments.get(element.align or "left")
                    # Text rides the same scale as the geometry, so type and
                    # shapes keep their relative proportions on any page size.
                    size_pt = (element.font_size or _DEFAULT_FONT_PX) * _PX_TO_PT * scale
                    for index, line in enumerate((element.text or "").split("\n")):
                        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                        paragraph.alignment = alignment
                        if element.line_height:
                            paragraph.line_spacing = element.line_height
                        if element.space_before is not None:
                            paragraph.space_before = Pt(element.space_before * _PX_TO_PT)
                        if element.space_after is not None:
                            paragraph.space_after = Pt(element.space_after * _PX_TO_PT)
                        if line.startswith(BULLET_PREFIX):
                            line = line[len(BULLET_PREFIX):]
                            # A literal bullet inside the run is drawn by
                            # whichever font covers the text after it, so a
                            # deck of mixed scripts gets mixed bullet glyphs
                            # and ragged left edges. A paragraph bullet is
                            # drawn once by the list, and its hanging indent
                            # puts wrapped lines under the text. python-pptx
                            # has no public API for either.
                            hang = str(int(size_pt * _BULLET_HANG * _EMU_PER_POINT))
                            properties = paragraph._p.get_or_add_pPr()
                            properties.set("marL", hang)
                            properties.set("indent", "-" + hang)
                            properties.append(
                                properties.makeelement(
                                    qn("a:buChar"), {"char": BULLET_PREFIX.strip()}
                                )
                            )
                        run = paragraph.add_run()
                        run.text = line
                        run.font.size = Pt(size_pt)
                        run.font.bold = bool(element.bold)
                        edge = _hex_rgb(element.stroke)
                        if edge:
                            # A text outline (WordArt): python-pptx has no
                            # accessor, and the schema wants `a:ln` first
                            # among the run properties' children.
                            properties = run.font._rPr
                            line = properties.makeelement(
                                qn("a:ln"),
                                {"w": str(int((element.stroke_width or 1) * _PX_TO_PT * 12700))},
                            )
                            fill = line.makeelement(qn("a:solidFill"), {})
                            fill.append(fill.makeelement(qn("a:srgbClr"), {"val": edge}))
                            line.append(fill)
                            properties.insert(0, line)
                        band = _hex_rgb(element.highlight)
                        if band:
                            # python-pptx exposes no highlight accessor; the
                            # element goes straight into the run properties.
                            properties = run.font._rPr
                            mark = properties.makeelement(qn("a:highlight"), {})
                            colour = mark.makeelement(qn("a:srgbClr"), {"val": band})
                            mark.append(colour)
                            properties.append(mark)
                        if color is not None:
                            run.font.color.rgb = RGBColor.from_string(color)
                        _name_face(run, element.font_family or skin_face, qn)
                elif element.type == "table":
                    _add_slide_table(
                        slide,
                        element,
                        (left, top, width, height),
                        scale=scale,
                        text_color=_hex_rgb(slide_model.text_color),
                        skin_face=skin_face,
                        alignments=alignments,
                        anchors=anchors,
                    )
                elif element.type == "shape":
                    outline = _hex_rgb(element.stroke)
                    # A shape may be drawn by its outline alone, and "none"
                    # says the shape is explicitly unfilled. Defaulting the
                    # fill in either case would paint over what the border
                    # frames. The default colour stays for a shape that says
                    # nothing and has no border — an authored box should show.
                    fill_color = (
                        None
                        if element.fill == "none"
                        else _hex_rgb(element.fill) or (None if outline else _DEFAULT_SHAPE_FILL)
                    )
                    if element.shape == "line":
                        connector = slide.shapes.add_connector(
                            MSO_CONNECTOR.STRAIGHT, left, top, Emu(int(left) + int(width)), Emu(int(top) + int(height))
                        )
                        connector.line.color.rgb = RGBColor.from_string(
                            outline or fill_color or _DEFAULT_SHAPE_FILL
                        )
                        connector.line.width = Pt((element.stroke_width or 2) * _PX_TO_PT)
                    else:
                        shape_type = MSO_SHAPE.OVAL if element.shape == "ellipse" else MSO_SHAPE.RECTANGLE
                        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                        if element.rotation:
                            shape.rotation = element.rotation
                        if fill_color:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                        else:
                            shape.fill.background()
                        if outline:
                            shape.line.color.rgb = RGBColor.from_string(outline)
                            shape.line.width = Pt((element.stroke_width or 1) * _PX_TO_PT)
                        else:
                            shape.line.fill.background()
                elif element.src:
                    data = _data_uri_bytes(element.src)
                    if data is None:
                        continue  # not inlined / not an embeddable type — skip honestly
                    try:
                        picture = slide.shapes.add_picture(io.BytesIO(data), left, top)
                    except Exception:  # noqa: BLE001 — corrupt image data; keep the deck
                        continue
                    # Contain the picture in its box (never stretch), centered —
                    # the same object-fit the canvas and browser exports use.
                    native_w, native_h = picture.image.size
                    if native_w and native_h:
                        # `fit`, not `scale` — this is the per-picture
                        # contain factor; rebinding `scale` here silently
                        # corrupted the page projection for every element
                        # after an image (fonts exploded past pptx limits).
                        fit = min(int(width) / native_w, int(height) / native_h)
                        picture.width = Emu(int(native_w * fit))
                        picture.height = Emu(int(native_h * fit))
                        picture.left = Emu(int(left) + (int(width) - int(picture.width)) // 2)
                        picture.top = Emu(int(top) + (int(height) - int(picture.height)) // 2)
                    if element.rotation:
                        picture.rotation = element.rotation

            if slide_model.notes:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    notes_frame.text = slide_model.notes

        out = io.BytesIO()
        presentation.save(out)
        # The deck's own title names the file when the caller has none — a
        # slides envelope carries one, unlike the html/table sources.
        name = _safe_name(title) or _safe_name(envelope_title) or _stem(path)
        return ExportedFile(out.getvalue(), f"{name}.pptx", PPTX_MIME)


def default_exporters() -> list[Exporter]:
    """The built-in exporters, in routing order."""
    return [TableXlsxExporter(), HtmlDocxExporter(), MarkdownDocxExporter(), SlidesPptxExporter()]
