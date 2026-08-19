"""Source converters — stored source files rendered as model-usable content.

Uploaded files land in the store under ``sources/`` (the user's original
material, read-only for agents). A :class:`SourceConverter` turns one such
file's raw bytes into a :class:`ConvertedSource`: a list of LangChain-standard
content blocks (``{"type": "text", ...}`` and, for vision formats,
``{"type": "image", ...}``) plus metadata. The standard ``read_canvas`` tool
routes binary files through the converter matching their suffix, so the agent
never parses bytes.

Converters are pluggable: pass your own list to ``create_canvas_tools`` to
replace or extend the defaults — an in-house OCR or document-AI pipeline
implements the same one-method contract. Built-in converters that need a
parser library import it lazily and raise
:class:`MissingConverterDependencyError` with an install hint when absent, so
the core package stays dependency-free.
"""

from __future__ import annotations

import base64
import csv
import io
from dataclasses import dataclass, field
from typing import Any, Protocol, runtime_checkable


@dataclass(frozen=True)
class ConvertedSource:
    """A source file in model-usable form.

    ``blocks`` are LangChain-standard content blocks — text blocks always,
    image blocks for vision formats. ``metadata`` carries format facts a tool
    may surface to the agent (sheet names, page count, detected encoding, ...).
    """

    blocks: list[dict[str, Any]]
    metadata: dict[str, Any] = field(default_factory=dict)


@runtime_checkable
class SourceConverter(Protocol):
    """The pluggable contract: bytes of one source file in, blocks out."""

    suffixes: tuple[str, ...]
    """Path suffixes this converter handles (lowercase, with the dot)."""

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        """Convert one file. ``path`` is the store path, for messages/metadata."""
        ...


class MissingConverterDependencyError(RuntimeError):
    """A built-in converter's parser library is not installed.

    The message names the missing package and the extra that installs it, so
    the standard tools can relay an honest, actionable error to the agent.
    """


@runtime_checkable
class PageRenderable(Protocol):
    """Optional converter extension: a paged source rendered as images.

    A separate protocol (not new methods on :class:`SourceConverter`) so
    existing custom converters keep working unchanged. Formats whose pages
    are worth *seeing* — scans, charts, layout — implement it; the standard
    ``read_canvas`` tool routes its ``pages`` parameter here. The PDF
    converter is the first implementation; office-document renderers plug
    into the same slot.
    """

    def render_pages(self, data: bytes, *, path: str, pages: list[int]) -> ConvertedSource:
        """The requested 1-based pages as labeled image blocks.

        Raises ``ValueError`` with an honest message for an out-of-range
        request, so the tool can relay it verbatim.
        """
        ...

    def render_grid(self, data: bytes, *, path: str) -> ConvertedSource:
        """Every page as a small labeled thumbnail, tiled into grid images.

        The cheap overview between "read the text" and "render these pages":
        one call shows the whole document's shape so the agent can pick the
        pages worth a full render.
        """
        ...


def converter_for(path: str, converters: list[SourceConverter]) -> SourceConverter | None:
    """The first converter whose suffix matches ``path`` (case-insensitive)."""
    lowered = path.lower()
    for converter in converters:
        if lowered.endswith(converter.suffixes):
            return converter
    return None


class TextSourceConverter:
    """Plain-text formats, decoded with UTF-8 (BOM-aware) then CP949.

    The trivial reference implementation of the contract. Files that decode
    with neither encoding fall back to lossy UTF-8 and say so in metadata —
    the agent gets honest content instead of an opaque failure.
    """

    suffixes: tuple[str, ...] = (".txt", ".md", ".markdown", ".csv", ".json", ".html", ".htm")

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        for encoding in ("utf-8-sig", "cp949"):
            try:
                return ConvertedSource(
                    blocks=[{"type": "text", "text": data.decode(encoding)}],
                    metadata={"encoding": encoding},
                )
            except UnicodeDecodeError:
                continue
        return ConvertedSource(
            blocks=[{"type": "text", "text": data.decode("utf-8", errors="replace")}],
            metadata={"encoding": "unknown (lossy utf-8)"},
        )


class XlsxSourceConverter:
    """Excel workbooks, one CSV-shaped section per sheet.

    Reads cached values (what Excel last computed), so formula cells show
    their results. Requires ``openpyxl`` — installed by the ``xlsx`` extra.
    """

    suffixes: tuple[str, ...] = (".xlsx",)

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingConverterDependencyError(
                "reading .xlsx needs openpyxl — install langchain-canvas[xlsx] "
                "or register your own converter for .xlsx"
            ) from exc

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sections: list[str] = []
        sheet_names: list[str] = []
        for sheet in workbook.worksheets:
            sheet_names.append(sheet.title)
            out = io.StringIO()
            writer = csv.writer(out)
            for row in sheet.iter_rows(values_only=True):
                writer.writerow(["" if cell is None else str(cell) for cell in row])
            body = out.getvalue().rstrip("\n") or "(empty)"
            sections.append(f"### sheet: {sheet.title}\n{body}")
        workbook.close()
        return ConvertedSource(
            blocks=[{"type": "text", "text": "\n\n".join(sections)}],
            metadata={"sheets": ", ".join(sheet_names), "values": "cached results"},
        )


#: Anthropic-family request limit is 5 MB per image *after* base64 (+33%).
MAX_INLINE_IMAGE_BYTES = 3_750_000

#: At most this many image blocks per tool call — pages or grid sheets alike.
MAX_IMAGES_PER_CALL = 8


class ImageSourceConverter:
    """Images as vision blocks — the model sees the picture, not a description.

    No parser dependency: the bytes go straight into a base64 image block.
    Delivery to the model requires a vision-capable model and a provider whose
    tool messages accept image blocks (Anthropic and Bedrock Converse do; some
    chat-completions providers accept text only). Oversized images degrade to
    an honest note instead of an oversized request.
    """

    suffixes: tuple[str, ...] = (".png", ".jpg", ".jpeg", ".gif", ".webp")

    max_bytes: int = MAX_INLINE_IMAGE_BYTES

    _MIME = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
    }

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        suffix = "." + path.rsplit(".", 1)[-1].lower()
        mime = self._MIME.get(suffix, "application/octet-stream")
        if len(data) > self.max_bytes:
            return ConvertedSource(
                blocks=[
                    {
                        "type": "text",
                        "text": (
                            f"(image is {len(data):,} bytes — too large to inline; "
                            f"the limit is {self.max_bytes:,} bytes. Ask the user "
                            "for a smaller version.)"
                        ),
                    }
                ],
                metadata={"mime": mime, "bytes": len(data), "inlined": False},
            )
        return ConvertedSource(
            blocks=[
                {"type": "text", "text": f"(image {path}, {mime}, {len(data):,} bytes)"},
                {
                    "type": "image",
                    "source_type": "base64",
                    "mime_type": mime,
                    "data": base64.b64encode(data).decode(),
                },
            ],
            metadata={"mime": mime, "bytes": len(data), "inlined": True},
        )


class PdfSourceConverter:
    """PDFs as per-page extracted text, and — on request — as page images.

    Text extraction requires ``pypdf`` (the ``pdf`` extra). Pages without an
    extractable text layer (scans) say so honestly instead of silently
    contributing nothing.

    Page rendering (the :class:`PageRenderable` side) requires ``pypdfium2``
    and ``pillow`` (the ``pdf-images`` extra) — both plain wheels, no external
    binary. Rendering is recomputed on every call by design: a persistent
    preview belongs to the file-artifact track, not here.
    """

    suffixes: tuple[str, ...] = (".pdf",)

    #: Render scales tried in order until the PNG fits the inline byte cap.
    render_scales: tuple[float, ...] = (2.0, 1.5, 1.0)
    max_bytes: int = MAX_INLINE_IMAGE_BYTES
    #: Grid overview tiling: 4 columns x 5 rows of thumbnails per sheet.
    grid_columns: int = 4
    grid_rows: int = 5
    grid_cell_width: int = 280

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        try:
            from pypdf import PdfReader  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingConverterDependencyError(
                "reading .pdf needs pypdf — install langchain-canvas[pdf] "
                "or register your own converter for .pdf"
            ) from exc

        reader = PdfReader(io.BytesIO(data))
        sections: list[str] = []
        empty_pages = 0
        for number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                empty_pages += 1
                text = "(no extractable text on this page — it may be a scan or an image)"
            sections.append(f"### page {number}\n{text}")
        metadata: dict[str, Any] = {"pages": len(reader.pages)}
        if empty_pages:
            metadata["pages without text"] = empty_pages
        return ConvertedSource(
            blocks=[{"type": "text", "text": "\n\n".join(sections)}], metadata=metadata
        )

    def _document(self, data: bytes) -> Any:
        try:
            import pypdfium2 as pdfium  # type: ignore[import-untyped]
            from PIL import Image  # noqa: F401  — needed by bitmap.to_pil()
        except ImportError as exc:
            raise MissingConverterDependencyError(
                "rendering .pdf pages needs pypdfium2 and pillow — install "
                "langchain-canvas[pdf-images] or register your own renderer for .pdf"
            ) from exc
        return pdfium.PdfDocument(data)

    def _render_png(self, document: Any, index: int, scale: float) -> bytes:
        bitmap = document[index].render(scale=scale)
        out = io.BytesIO()
        bitmap.to_pil().save(out, format="PNG")
        return out.getvalue()

    @staticmethod
    def _image_block(png: bytes) -> dict[str, Any]:
        return {
            "type": "image",
            "source_type": "base64",
            "mime_type": "image/png",
            "data": base64.b64encode(png).decode(),
        }

    def render_pages(self, data: bytes, *, path: str, pages: list[int]) -> ConvertedSource:
        """The requested 1-based pages as labeled PNG image blocks.

        Each page gets a ``### page N`` text label directly before its image,
        so the model never has to infer which image is which. A page is
        rendered at the largest scale whose PNG fits the inline byte cap;
        one that fits at no scale degrades to an honest text note.
        """
        document = self._document(data)
        total = len(document)
        bad = [p for p in pages if p < 1 or p > total]
        if bad:
            raise ValueError(f"{path} has {total} page(s) — valid range is 1-{total}")

        blocks: list[dict[str, Any]] = []
        scales_used: set[float] = set()
        for page in pages:
            png = b""
            for scale in self.render_scales:
                png = self._render_png(document, page - 1, scale)
                if len(png) <= self.max_bytes:
                    scales_used.add(scale)
                    break
            if len(png) > self.max_bytes:
                blocks.append(
                    {
                        "type": "text",
                        "text": (
                            f"### page {page} — could not be inlined: the rendered image is "
                            f"{len(png):,} bytes even at the smallest scale (limit "
                            f"{self.max_bytes:,}). Read this page as text instead."
                        ),
                    }
                )
                continue
            blocks.append({"type": "text", "text": f"### page {page} (image follows)"})
            blocks.append(self._image_block(png))
        return ConvertedSource(
            blocks=blocks,
            metadata={
                "pages": total,
                "rendered": ", ".join(str(p) for p in pages),
                "scale": ", ".join(str(s) for s in sorted(scales_used, reverse=True)) or "none",
            },
        )

    def render_grid(self, data: bytes, *, path: str) -> ConvertedSource:
        """Every page as a numbered thumbnail, tiled into grid sheets.

        The overview layer of the observation ladder: 20 thumbnails per
        sheet, each labeled with its page number so the agent can follow up
        with a full render of just the pages that matter. Raises
        ``ValueError`` when the document needs more sheets than one call may
        carry.
        """
        document = self._document(data)
        from PIL import Image, ImageDraw, ImageFont

        total = len(document)
        per_sheet = self.grid_columns * self.grid_rows
        sheets = (total + per_sheet - 1) // per_sheet
        if sheets > MAX_IMAGES_PER_CALL:
            raise ValueError(
                f"{path} has {total} pages — a grid overview would need {sheets} sheets, "
                f"more than the {MAX_IMAGES_PER_CALL}-image limit per call. Read the text "
                "layer first and request specific pages instead."
            )

        font = ImageFont.load_default(size=18)
        blocks: list[dict[str, Any]] = []
        for sheet_index in range(sheets):
            first = sheet_index * per_sheet
            last = min(first + per_sheet, total)
            thumbs: list[Any] = []
            for index in range(first, last):
                width = document[index].get_size()[0]
                scale = self.grid_cell_width / width if width else 0.4
                bitmap = document[index].render(scale=scale)
                thumb = bitmap.to_pil()
                draw = ImageDraw.Draw(thumb)
                label = str(index + 1)
                box = draw.textbbox((0, 0), label, font=font)
                draw.rectangle((0, 0, box[2] + 10, box[3] + 8), fill="#111111")
                draw.text((5, 3), label, fill="#ffffff", font=font)
                thumbs.append(thumb)

            cell_h = max(t.height for t in thumbs)
            columns = min(self.grid_columns, len(thumbs))
            rows = (len(thumbs) + columns - 1) // columns
            width = columns * self.grid_cell_width + 8 * (columns + 1)
            height = rows * cell_h + 8 * (rows + 1)
            sheet = Image.new("RGB", (width, height), "#e5e5e5")
            for i, thumb in enumerate(thumbs):
                r, c = divmod(i, columns)
                sheet.paste(thumb, (8 + c * (self.grid_cell_width + 8), 8 + r * (cell_h + 8)))
            out = io.BytesIO()
            sheet.save(out, format="PNG")
            encoded = out.getvalue()
            if len(encoded) > self.max_bytes:
                # A photo-heavy sheet can blow the PNG cap — JPEG is far smaller.
                out = io.BytesIO()
                sheet.save(out, format="JPEG", quality=80)
                encoded = out.getvalue()
            label = f"### pages {first + 1}-{last} (grid overview, numbered thumbnails)"
            if len(encoded) > self.max_bytes:
                blocks.append(
                    {
                        "type": "text",
                        "text": label
                        + " — could not be inlined (too large); request page ranges instead.",
                    }
                )
                continue
            blocks.append({"type": "text", "text": label})
            block = self._image_block(encoded)
            block["mime_type"] = "image/png" if encoded[:4] == b"\x89PNG" else "image/jpeg"
            blocks.append(block)
        return ConvertedSource(
            blocks=blocks,
            metadata={
                "pages": total,
                "grid": f"{sheets} sheet(s), up to {per_sheet} thumbnails each",
            },
        )


def _table_lines(rows: list[list[str]]) -> str:
    out = io.StringIO()
    writer = csv.writer(out)
    for row in rows:
        writer.writerow(row)
    return out.getvalue().rstrip("\n")


class DocxSourceConverter:
    """Word documents as text — paragraphs and tables, in document order.

    Requires ``python-docx`` — installed by the ``office`` extra.
    """

    suffixes: tuple[str, ...] = (".docx",)

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        try:
            from docx import Document  # type: ignore[import-untyped]
            from docx.table import Table  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingConverterDependencyError(
                "reading .docx needs python-docx — install langchain-canvas[office] "
                "or register your own converter for .docx"
            ) from exc

        document = Document(io.BytesIO(data))
        parts: list[str] = []
        paragraphs = 0
        tables = 0
        for item in document.iter_inner_content():
            if isinstance(item, Table):
                tables += 1
                rows = [[cell.text.strip() for cell in row.cells] for row in item.rows]
                parts.append(_table_lines(rows))
            else:
                text = item.text.strip()
                if text:
                    paragraphs += 1
                    parts.append(text)
        return ConvertedSource(
            blocks=[{"type": "text", "text": "\n\n".join(parts)}],
            metadata={"paragraphs": paragraphs, "tables": tables},
        )


class PptxSourceConverter:
    """PowerPoint decks as per-slide text — titles, body text, and tables.

    Requires ``python-pptx`` — installed by the ``office`` extra.
    """

    suffixes: tuple[str, ...] = (".pptx",)

    def convert(self, data: bytes, *, path: str) -> ConvertedSource:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError as exc:
            raise MissingConverterDependencyError(
                "reading .pptx needs python-pptx — install langchain-canvas[office] "
                "or register your own converter for .pptx"
            ) from exc

        deck = Presentation(io.BytesIO(data))
        sections: list[str] = []
        for number, slide in enumerate(deck.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells] for row in shape.table.rows
                    ]
                    parts.append(_table_lines(rows))
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
            body = "\n".join(parts) or "(no text on this slide)"
            sections.append(f"### slide {number}\n{body}")
        return ConvertedSource(
            blocks=[{"type": "text", "text": "\n\n".join(sections)}],
            metadata={"slides": len(deck.slides)},
        )


def default_converters() -> list[SourceConverter]:
    """The built-in converter set. Grows as format tiers land."""
    return [
        TextSourceConverter(),
        XlsxSourceConverter(),
        ImageSourceConverter(),
        PdfSourceConverter(),
        DocxSourceConverter(),
        PptxSourceConverter(),
    ]
