"""Source-grounded slide templates: task 1's minimal U3/U4 foundation, plus
task 6's intent-routing and instrumentation workflow tests.

Proves the frame-clone + slot-fill path (`instantiate_archetype`) against a
supported rich-text+shape exemplar instantiated 10 times, that the result
reopens as an editable rich PPTX through the app's native blank export path,
and that the SDK exporter fails closed on the same template metadata.
"""

from __future__ import annotations

import hashlib
import io
from types import SimpleNamespace

import pytest
from app.agent.deck_editing import create_deck_editing_tools
from app.agent.deck_template_models import (
    ArtifactRef,
    CompiledTemplateManifest,
    FinalizeRequest,
    PrepareRequest,
    SlideContentRequest,
    SourceRef,
    VariableBinding,
)
from app.agent.deck_template_tools import create_deck_template_tools
from app.agent.deck_template_writer import (
    ArchetypeFrame,
    TemplateInstantiationError,
    instantiate_archetype,
    write_deck_from_template,
)
from app.agent.deck_templates import finalize_template, prepare_template
from app.agent.exports import EditableDeckPptxExporter
from langchain_canvas.deck import Deck, DeckPptxExporter, SlideTemplate, serialize_deck
from langchain_canvas.store import InMemoryCanvasStore
from pptx import Presentation

from template_source_fixtures import (
    injection_pdf_source,
    korean_pptx_source,
    text_pdf_source,
)


def _exemplar_frame() -> ArchetypeFrame:
    """A frame with one rich-run headline slot, one plain body slot, a shape."""
    body_html = (
        '<section class="slide">'
        '<div data-node-id="node-headline">'
        "<b>OLD_BOLD</b><span>OLD_PLAIN</span>"
        "</div>"
        '<p data-node-id="node-body">Old body text</p>'
        '<div class="rect" data-pptx-shape-id="e0" '
        'style="position:absolute;left:20px;top:200px;width:100px;height:60px;'
        'background:#334155"></div>'
        "</section>"
    )
    return ArchetypeFrame(
        archetype_id="body",
        style_css=".slide { color: #111827; }",
        body_html=body_html,
        slot_node_ids={"headline": "node-headline", "body": "node-body"},
    )


def test_supported_exemplar_to_ten_slide_export():
    frame = _exemplar_frame()
    source_business_text = frozenset({"OLD_BOLD", "OLD_PLAIN", "Old body text"})

    slides = [
        instantiate_archetype(
            frame,
            index,
            {
                "headline": [f"New Bold {index}", f"New Plain {index}"],
                "body": f"New body text {index}",
            },
            source_business_text=source_business_text,
        )
        for index in range(1, 11)
    ]

    assert [slide.slide_id for slide in slides] == [
        f"slide-{i:03d}" for i in range(1, 11)
    ]
    for slide in slides:
        assert "data-pptx-shape-id" not in slide.body_html
        assert "OLD_BOLD" not in slide.body_html
        assert "Old body text" not in slide.body_html
        assert '<div class="rect"' in slide.body_html  # the locked shape survives

    deck = Deck(title="Generated Deck", ratio="16:9", source=None, slides=slides)
    content = serialize_deck(deck)

    # App rich PPTX reopen: native blank path (source=None) with rich runs.
    exported = EditableDeckPptxExporter().export(content, path="deck.slides.html")
    reopened = Presentation(io.BytesIO(exported.data))
    assert len(reopened.slides) == 10
    first_texts = {
        run.text
        for shape in reopened.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    }
    assert "New Bold 1" in first_texts
    assert "New Plain 1" in first_texts

    # SDK explicit rejection: this SDK exporter never exports template decks.
    template_deck = Deck(
        title="Generated Deck",
        ratio="16:9",
        source=None,
        slides=slides,
        template={
            "schema_version": 1,
            "template": {"path": "templates/h.template.json", "revision": "r1", "sha256": "abc"},
            "instances": {},
        },
    )
    with pytest.raises(ValueError, match="unsupported_template_export"):
        DeckPptxExporter().export(serialize_deck(template_deck), path="deck.slides.html")


def test_instantiate_archetype_rejects_missing_and_unknown_slots():
    frame = _exemplar_frame()

    with pytest.raises(TemplateInstantiationError):
        instantiate_archetype(frame, 1, {"headline": ["Only headline"]})

    with pytest.raises(TemplateInstantiationError):
        instantiate_archetype(
            frame,
            1,
            {"headline": ["A", "B"], "body": "C", "unknown_slot": "D"},
        )


def test_instantiate_archetype_rejects_leftover_source_business_text():
    frame = _exemplar_frame()

    with pytest.raises(TemplateInstantiationError):
        instantiate_archetype(
            frame,
            1,
            {"headline": ["New Bold", "New Plain"], "body": "New body text"},
            source_business_text=frozenset({"New Bold"}),
        )


# --- task 6: intent routing, tool-boundary rejection, compile/write counters --------


def _runtime(canvas_id: str, events: list | None = None) -> SimpleNamespace:
    return SimpleNamespace(
        config={"configurable": {"thread_id": canvas_id}},
        context=None,
        stream_writer=events.append if events is not None else None,
    )


def _patch_measure_slide(monkeypatch: pytest.MonkeyPatch) -> None:
    """A geometry backend stub so verification never depends on a real renderer."""

    def _fake_measure(_document: str, *, ratio: str) -> dict:
        return {"width": 960, "height": 540, "elements": [], "unsupported": []}

    monkeypatch.setattr("app.agent.deck_template_verification.measure_slide", _fake_measure)


def _patch_reference_review(monkeypatch: pytest.MonkeyPatch) -> None:
    """A clean reference-comparison stub: no renderer, no reviewer model call."""
    monkeypatch.setattr(
        "app.agent.deck_template_verification.render_slide",
        lambda _document, *, ratio: ({}, b"rendered-png"),
    )
    monkeypatch.setattr(
        "app.agent.deck_template_verification.review_rendered_against_reference",
        lambda _reference, _rendered: [],
    )


def _finalize_bindings(archetype: dict) -> list[dict]:
    return [
        {
            "archetype_id": archetype["id"],
            "node_id": slot["node_id"],
            "disposition": "variable",
            "slot_key": slot["key"],
            "role": slot["role"],
            "required": True,
        }
        for slot in archetype["slots"]
    ]


def test_style_request_avoids_full_conversion_and_scratch_writer(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A source-reuse-for-a-new-topic request runs only the four template
    tools — never the scratch full-conversion or batch-writer pipelines."""
    canvas_id = "thread-style"
    store = InMemoryCanvasStore()
    data = korean_pptx_source(
        [("연간 매출 요약", "지난 한 해 동안의 매출 실적을 정리한 설명입니다.")]
    )
    store.write_bytes(canvas_id, "sources/deck.pptx", data, "Upload")
    runtime = _runtime(canvas_id)
    tools = {t.name: t for t in create_deck_template_tools(store)}

    def _boom(*_args, **_kwargs):
        raise AssertionError(
            "full conversion / scratch batch writer must not run for a template-reuse request"
        )

    import app.agent.deck_batch as deck_batch_module
    import app.agent.pdf_deck as pdf_deck_module

    monkeypatch.setattr(pdf_deck_module, "open_pdf_as_slides", _boom)
    monkeypatch.setattr(deck_batch_module, "generate_slide_bodies", _boom)
    _patch_measure_slide(monkeypatch)

    census = tools["inspect_deck_patterns"].func(source="sources/deck.pptx", runtime=runtime)
    assert census["groups"]

    define = tools["define_deck_template"]
    candidate = define.func(
        mode="prepare",
        runtime=runtime,
        source="sources/deck.pptx",
        source_sha256=hashlib.sha256(data).hexdigest(),
        pages=[1],
    )
    assert candidate["status"] == "candidate"
    archetype = candidate["archetypes"][0]

    ready = define.func(
        mode="finalize",
        runtime=runtime,
        candidate_ref=candidate["candidate_ref"],
        bindings=_finalize_bindings(archetype),
    )
    assert ready["status"] == "ready"

    slot_keys = [slot["key"] for slot in archetype["slots"]]
    write_result = tools["write_deck_from_template"].func(
        template_ref=ready["template_ref"],
        destination="new-deck.slides.html",
        title="New Topic Deck",
        slides=[
            {
                "archetype_id": archetype["id"],
                "mode": "verbatim",
                "slots": {
                    slot_keys[0]: "새로운 주제 제목",
                    slot_keys[1]: "새로운 주제에 대한 설명입니다.",
                },
            }
        ],
        runtime=runtime,
    )
    assert write_result["status"] == "ok"

    verified = tools["verify_template_deck"].func(
        path="new-deck.slides.html", revision=write_result["revision"], runtime=runtime
    )
    # A PPTX archetype carries no pinned reference render, so the U4 original
    # comparison cannot run for it — that degrades visual fidelity rather than
    # passing an unchecked dimension as verified.
    assert verified["visual_fidelity"]["status"] == "degraded", verified
    assert verified["content"]["status"] == "verified", verified
    # A verbatim slide has no new authored voice to judge — writing_style is
    # not_checked (never vacuously verified), so `complete` is correctly
    # False even though every actually-checked dimension passed.
    assert verified["writing_style"]["status"] == "not_checked", verified
    assert verified["complete"] is False


def test_reproduction_and_edit_routes_remain_legacy(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Exact-reproduction (clone) and edit-existing-content tools keep working
    unchanged, and never call into the new template compile/write path."""
    canvas_id = "thread-legacy"
    store = InMemoryCanvasStore()
    body = (
        '<section class="slide"><h1 data-node-id="title">Original Title</h1>'
        '<p data-node-id="body">Original body</p></section>'
    )
    source_content = serialize_deck(
        Deck(
            "Template",
            "16:9",
            "sources/original.pptx",
            [SlideTemplate("one", "One", "", body)],
        )
    )
    revision = store.write(canvas_id, "deck.slides.html", source_content, "Seed").revision
    runtime = _runtime(canvas_id)
    editing_tools = {t.name: t for t in create_deck_editing_tools(store)}

    def _boom(*_args, **_kwargs):
        raise AssertionError("legacy reproduction/edit routes must not call the template path")

    monkeypatch.setattr("app.agent.deck_templates.prepare_template", _boom)
    monkeypatch.setattr("app.agent.deck_templates.finalize_template", _boom)
    monkeypatch.setattr("app.agent.deck_template_writer.write_deck_from_template", _boom)

    cloned = editing_tools["clone_deck_template"].func(
        source="deck.slides.html", destination="clone.slides.html", revision=revision, runtime=runtime
    )
    assert cloned["status"] == "committed"

    edited = editing_tools["replace_slide_text"].func(
        path="clone.slides.html",
        slide_id="one",
        revision=cloned["revision"],
        node_id="body",
        runtime=runtime,
        text="Edited body",
    )
    assert edited["status"] == "committed"


def test_forged_cross_canvas_refs_and_source_instructions_rejected(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A forged candidate commit, a cross-canvas template ref, and an
    injected-instruction source page are all rejected at the tool boundary."""
    store = InMemoryCanvasStore()
    canvas_a, canvas_b = "thread-a", "thread-b"
    tools = {t.name: t for t in create_deck_template_tools(store)}

    # --- forged candidate: right shape, wrong (non-compiler) actor ---------------
    forged_manifest = CompiledTemplateManifest(
        status="candidate",
        source=SourceRef(path="sources/deck.pptx", revision="head", sha256="0" * 64),
        selected_pages=[1],
        ratio="16:9",
        archetypes=[],
    )
    payload = forged_manifest.model_dump_json()
    forged_commit = store.write(canvas_a, "templates/forged.candidate.json", payload, "human edit", actor="human")
    forged_ref = {
        "path": "templates/forged.candidate.json",
        "revision": forged_commit.revision,
        "sha256": hashlib.sha256(payload.encode("utf-8")).hexdigest(),
    }
    result = tools["define_deck_template"].func(
        mode="finalize",
        runtime=_runtime(canvas_a),
        candidate_ref=forged_ref,
        bindings=[
            {
                "archetype_id": "archetype-1",
                "node_id": "node-0",
                "disposition": "omit",
            }
        ],
    )
    assert result["status"] == "error"
    assert result["code"] == "verification_failed"

    # --- cross-canvas ref: a real ready template from canvas A, read as canvas B --
    data = korean_pptx_source([("보고서 제목", "보고서 본문 설명입니다.")])
    store.write_bytes(canvas_a, "sources/deck.pptx", data, "Upload")
    prepared = prepare_template(
        PrepareRequest(
            mode="prepare",
            source="sources/deck.pptx",
            source_sha256=hashlib.sha256(data).hexdigest(),
            pages=[1],
        ),
        store=store,
        canvas_id=canvas_a,
    )
    archetype = prepared["archetypes"][0]
    finalized = finalize_template(
        FinalizeRequest(
            mode="finalize",
            candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=[VariableBinding(**binding) for binding in _finalize_bindings(archetype)],
        ),
        store=store,
        canvas_id=canvas_a,
        current_source_sha256=hashlib.sha256(data).hexdigest(),
    )
    assert finalized["status"] == "ready"

    cross_canvas_result = tools["write_deck_from_template"].func(
        template_ref=finalized["template_ref"],
        destination="deck.slides.html",
        title="Cross-canvas attempt",
        slides=[
            {
                "archetype_id": archetype["id"],
                "mode": "verbatim",
                "slots": {slot["key"]: "text" for slot in archetype["slots"]},
            }
        ],
        runtime=_runtime(canvas_b),
    )
    assert cross_canvas_result["status"] == "error"
    assert cross_canvas_result["code"] == "verification_failed"

    # --- source prompt injection: page text is returned as inert data, never
    # interpreted or specially stripped ------------------------------------------
    store.write_bytes(canvas_a, "sources/injected.pdf", injection_pdf_source(), "Upload")
    census = tools["inspect_deck_patterns"].func(
        source="sources/injected.pdf", runtime=_runtime(canvas_a)
    )
    assert census["groups"]
    all_examples = [
        text
        for group in census["groups"]
        for texts in group["examples"].values()
        for text in texts
    ]
    assert any("Ignore all previous instructions" in text for text in all_examples)


def test_html_generation_counters_distinguish_compile_and_write(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """U2's per-selected-page PDF reconstruction and U3's verbatim write are
    counted separately, and neither one calls the scratch full-conversion or
    batch-writer HTML-generation path."""
    canvas_id = "thread-counters"
    store = InMemoryCanvasStore()
    data = text_pdf_source([[("Page one", 100, 700)], [("Page two", 100, 700)]])
    store.write_bytes(canvas_id, "sources/deck.pdf", data, "Upload")

    reconstruct_calls = 0

    def _counting_reconstruct(source):
        nonlocal reconstruct_calls
        reconstruct_calls += 1
        return f'<section class="slide"><p data-node-id="text-0">{source.texts[0]["text"]}</p></section>', []

    def _boom(*_args, **_kwargs):
        raise AssertionError("scratch full-conversion/batch-writer HTML generation must not run")

    import app.agent.deck_batch as deck_batch_module
    import app.agent.pdf_deck as pdf_deck_module

    monkeypatch.setattr(pdf_deck_module, "write_pdf_html", _boom)
    monkeypatch.setattr(deck_batch_module, "generate_slide_bodies", _boom)

    prepared = prepare_template(
        PrepareRequest(
            mode="prepare",
            source="sources/deck.pdf",
            source_sha256=hashlib.sha256(data).hexdigest(),
            pages=[1, 2],
        ),
        store=store,
        canvas_id=canvas_id,
        reconstruct_pdf_page_fn=_counting_reconstruct,
    )
    assert prepared["status"] == "candidate"
    assert reconstruct_calls == 2  # exactly the k=2 selected pages, no more

    finalized = finalize_template(
        FinalizeRequest(
            mode="finalize",
            candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=[
                VariableBinding(**binding)
                for archetype in prepared["archetypes"]
                for binding in _finalize_bindings(archetype)
            ],
        ),
        store=store,
        canvas_id=canvas_id,
        current_source_sha256=hashlib.sha256(data).hexdigest(),
    )
    assert finalized["status"] == "ready"

    result = write_deck_from_template(
        ArtifactRef(**finalized["template_ref"]),
        "deck.slides.html",
        "Generated Deck",
        [
            SlideContentRequest(
                archetype_id=archetype["id"],
                mode="verbatim",
                slots={archetype["slots"][0]["key"]: f"New text {index}"},
            )
            for index, archetype in enumerate(prepared["archetypes"])
        ],
        _runtime(canvas_id),
        store=store,
        canvas_id=canvas_id,
        writer_model="unused-model",
    )
    assert result["status"] == "ok"
    assert result["slide_count"] == 2
    assert reconstruct_calls == 2  # unchanged: writing never re-compiles the source


# --- CV-P1-002: PDF node ids are derived from the compiled frame -------------------


def _finalize_bindings_all(archetype: dict) -> list[dict]:
    """Bind every candidate slot/static node — text and image — for finalize."""
    bindings = [
        {
            "archetype_id": archetype["id"],
            "node_id": slot["node_id"],
            "disposition": "variable",
            "slot_key": slot["key"],
            "role": slot["role"],
            "required": True,
        }
        for slot in archetype["slots"]
    ]
    bindings += [
        {
            "archetype_id": archetype["id"],
            "node_id": node["node_id"],
            "disposition": "retain",
            "retain_reason": "static",
        }
        for node in archetype["static_nodes"]
    ]
    return bindings


def test_pdf_end_to_end_prepare_finalize_write_verify_with_model_chosen_node_ids(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """prepare(pdf) -> finalize -> write -> verify succeeds when the writer
    model assigns its own ``data-node-id`` values, distinct from the
    ``text-{i}``/``image-{i}`` ids a prior version of this compiler assumed."""
    from app.agent.deck_template_verification import verify_template_deck_snapshot

    _patch_measure_slide(monkeypatch)
    _patch_reference_review(monkeypatch)
    canvas_id = "thread-pdf-node-ids"
    store = InMemoryCanvasStore()
    data = text_pdf_source([[("Original title text", 100, 700)]])
    store.write_bytes(canvas_id, "sources/deck.pdf", data, "Upload")

    def _model_reconstruct(source):
        # Node ids a real writer model would freely choose — deliberately
        # not the text-{i}/image-{i} shape this compiler used to assume.
        nodes = "".join(
            f'<p data-node-id="model-chosen-node-{i}" data-text-block="true">{t["text"]}</p>'
            for i, t in enumerate(source.texts)
        )
        return f"<section class='slide'>{nodes}</section>", []

    prepared = prepare_template(
        PrepareRequest(
            mode="prepare",
            source="sources/deck.pdf",
            source_sha256=hashlib.sha256(data).hexdigest(),
            pages=[1],
        ),
        store=store,
        canvas_id=canvas_id,
        reconstruct_pdf_page_fn=_model_reconstruct,
    )
    assert prepared["status"] == "candidate"
    archetype = prepared["archetypes"][0]
    assert archetype["slots"][0]["node_id"] == "model-chosen-node-0"

    finalized = finalize_template(
        FinalizeRequest(
            mode="finalize",
            candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_finalize_bindings_all(archetype),
        ),
        store=store,
        canvas_id=canvas_id,
        current_source_sha256=hashlib.sha256(data).hexdigest(),
    )
    assert finalized["status"] == "ready"

    write_result = write_deck_from_template(
        ArtifactRef(**finalized["template_ref"]),
        "deck.slides.html",
        "PDF Node Id Deck",
        [
            SlideContentRequest(
                archetype_id=archetype["id"],
                mode="verbatim",
                slots={archetype["slots"][0]["key"]: "New title text"},
            )
        ],
        _runtime(canvas_id),
        store=store,
        canvas_id=canvas_id,
        writer_model="unused-model",
    )
    assert write_result["status"] == "ok"

    verified = verify_template_deck_snapshot(
        "deck.slides.html",
        write_result["revision"],
        store=store,
        canvas_id=canvas_id,
        judge_model="unused-judge-model",
    )
    assert verified["visual_fidelity"]["status"] == "verified", verified
    assert verified["content"]["status"] == "verified", verified
