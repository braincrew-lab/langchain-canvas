"""``POST /api/canvas/{thread_id}/export`` — the export menu's office door.

Decks and tables convert server-side (the browser has no pptx/xlsx writer);
these tests exercise the route end-to-end against the real store, with
distinct thread ids per test and cleanup after each (module-level ``STORE``
is shared process-wide).
"""

from __future__ import annotations

import io
import json
import shutil
from typing import Any

import pytest
from app.agent.store import DATA_DIR, STORE
from app.main import app
from fastapi.testclient import TestClient
from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
from langchain_canvas.exporters import PPTX_MIME, XLSX_MIME
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_TEST_THREAD_IDS = (
    "t-export-deck",
    "t-export-content",
    "t-export-errors",
    "t-export-table",
)


@pytest.fixture(autouse=True)
def _cleanup_canvas_data() -> Any:
    yield
    for thread_id in _TEST_THREAD_IDS:
        shutil.rmtree(DATA_DIR / thread_id, ignore_errors=True)


client = TestClient(app)


def _deck_html(text: str = "Hello export") -> str:
    slide = SlideTemplate(
        slide_id="slide-001",
        title=None,
        style_css="",
        body_html=(
            '<section class="slide"><div class="lcx-block" data-node-id="node-1" '
            'style="position: absolute; left: 100.00px; top: 50.00px; '
            f'width: 400.00px; height: 80.00px;">{text}</div></section>'
        ),
    )
    return serialize_deck(Deck(title="Deck", ratio="16:9", source=None, slides=[slide]))


def test_export_stored_deck_to_pptx():
    STORE.write(
        "t-export-deck", "deck.slides.html", _deck_html(), "Create deck", actor="agent"
    )
    res = client.post(
        "/api/canvas/t-export-deck/export",
        json={"path": "deck.slides.html", "target": "pptx", "title": "My Deck"},
    )
    assert res.status_code == 200
    assert res.headers["content-type"].startswith(PPTX_MIME)
    assert "attachment" in res.headers["content-disposition"]
    deck = Presentation(io.BytesIO(res.content))
    assert len(deck.slides) == 1
    assert "Hello export" in [s.text for s in deck.slides[0].shapes if s.has_text_frame]


def test_export_client_content_wins_over_store():
    # The browser posts its current copy — no stored file is needed at all.
    res = client.post(
        "/api/canvas/t-export-content/export",
        json={
            "path": "deck.slides.html",
            "target": "pptx",
            "content": _deck_html("Unsaved"),
        },
    )
    assert res.status_code == 200
    assert res.headers["content-type"].startswith(PPTX_MIME)


def test_export_table_to_xlsx():
    envelope = json.dumps(
        {
            "type": "table",
            "title": "Sales",
            "data": {
                "columns": [
                    {"key": "q", "label": "Quarter"},
                    {"key": "rev", "label": "Revenue"},
                ],
                "rows": [{"q": "Q1", "rev": 12}, {"q": "Q2", "rev": 18}],
            },
        }
    )
    res = client.post(
        "/api/canvas/t-export-table/export",
        json={"path": "sales.table.json", "target": "xlsx", "content": envelope},
    )
    assert res.status_code == 200
    assert res.headers["content-type"].startswith(XLSX_MIME)


def test_export_unroutable_target_is_422():
    res = client.post(
        "/api/canvas/t-export-errors/export",
        json={"path": "deck.slides.html", "target": "gif"},
    )
    assert res.status_code == 422


def test_export_missing_file_is_404():
    res = client.post(
        "/api/canvas/t-export-errors/export",
        json={"path": "missing.slides.html", "target": "pptx"},
    )
    assert res.status_code == 404


def test_export_rejects_unresolved_external_image():
    content = serialize_deck(
        Deck(
            title="Unsafe image",
            ratio="16:9",
            source=None,
            slides=[
                SlideTemplate(
                    slide_id="s1",
                    title=None,
                    style_css="",
                    body_html='<img src="http://127.0.0.1:9/private.png">',
                )
            ],
        )
    )
    res = client.post(
        "/api/canvas/t-export-content/export",
        json={"path": "image.slides.html", "target": "pptx", "content": content},
    )
    assert res.status_code == 422
    assert "image could not be loaded" in res.json()["detail"]


def test_export_arbitrary_html_paints_every_slide_in_order():
    # Generated decks use tables/flex/CSS, with no native lcx node ids.
    # The old exporter returned 200 with zero shapes on all these slides.
    content = serialize_deck(
        Deck(
            title="Paint",
            ratio="4:3",
            source=None,
            slides=[
                SlideTemplate(
                    slide_id=f"s{i}",
                    title=None,
                    style_css=f"body{{background:{color}}}",
                    body_html="<h1>PDF reproduction</h1><table><tr><td>Table</td></tr></table>",
                )
                for i, color in enumerate(("#ff0000", "#0000ff"))
            ],
        )
    )
    res = client.post(
        "/api/canvas/t-export-content/export",
        json={"path": "paint.slides.html", "target": "pptx", "content": content},
    )
    assert res.status_code == 200
    deck = Presentation(io.BytesIO(res.content))
    assert len(deck.slides) == 2
    assert deck.slide_width / deck.slide_height == pytest.approx(4 / 3)
    for slide in deck.slides:
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        assert "PDF reproduction" in texts
        assert "Table" in texts
        assert not any(
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
        )


def test_export_resolves_stored_images_without_browser_inlining():
    image = Image.new("RGB", (1280, 720), (0, 255, 0))
    output = io.BytesIO()
    image.save(output, format="PNG")
    STORE.write_bytes(
        "t-export-content", "assets/page.png", output.getvalue(), "Asset", actor="human"
    )
    content = serialize_deck(
        Deck(
            title="Image",
            ratio="16:9",
            source=None,
            slides=[
                SlideTemplate(
                    slide_id="s1",
                    title=None,
                    style_css="",
                    body_html='<img src="assets/page.png" style="width:1280px;height:720px">',
                )
            ],
        )
    )
    res = client.post(
        "/api/canvas/t-export-content/export",
        json={"path": "image.slides.html", "target": "pptx", "content": content},
    )
    assert res.status_code == 200
    picture = next(
        s
        for s in Presentation(io.BytesIO(res.content)).slides[0].shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert Image.open(io.BytesIO(picture.image.blob)).convert("RGB").getpixel(
        (600, 400)
    ) == (0, 255, 0)


def test_css_arrowhead_exports_as_native_triangle_not_a_rectangle():
    content = _deck_html().replace(
        "</section>",
        '<div style="position:absolute;left:500px;top:200px;width:0;height:0;'
        "border-top:4px solid transparent;border-bottom:4px solid transparent;"
        'border-left:8px solid red"></div></section>',
    )
    res = client.post(
        "/api/canvas/t-export-content/export",
        json={"path": "arrow.slides.html", "target": "pptx", "content": content},
    )
    assert res.status_code == 200
    triangle = next(
        shape
        for shape in Presentation(io.BytesIO(res.content)).slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    )
    assert str(triangle.fill.fore_color.rgb) == "FF0000"
    assert len(triangle._element.xpath(".//a:lnTo")) == 2
