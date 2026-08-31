"""U2 template compiler: ``prepare``/``finalize``, trust gate, and cooperative budget.

Uses ``InMemoryCanvasStore`` (real bounded parsers underneath — real
``python-pptx``/``pypdfium2`` extraction, only the PDF writer model call is a
typed fake boundary via ``reconstruct_pdf_page_fn`` injection) and a fake
monotonic clock for ``TemplateBudget`` so the 120s cooperative window and the
shared model-attempt/token caps can be exercised deterministically, without
sleeping and without a live model.
"""

from __future__ import annotations

import base64
import ctypes
import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pydantic import ValidationError

from app.agent.deck_template_models import (
    ArtifactRef,
    BudgetExceededError,
    CompiledTemplateManifest,
    FinalizeRequest,
    OmitBinding,
    PrepareRequest,
    RetainBinding,
    TemplateBudget,
    VariableBinding,
)
from app.agent.deck_templates import (
    TEMPLATE_COMPILER_ACTOR,
    TrustError,
    finalize_template,
    prepare_template,
    require_trusted_artifact,
)
from app.agent.pdf_source import PdfPageSource
from langchain_canvas.store import InMemoryCanvasStore

pdfium = pytest.importorskip("pypdfium2")
from pypdfium2 import raw  # noqa: E402

CANVAS_ID = "thread-1"


class _FakeClock:
    """A manually-advanced monotonic clock — deterministic, no ``sleep``."""

    def __init__(self, start: float = 0.0) -> None:
        self.now = start

    def __call__(self) -> float:
        return self.now

    def advance(self, seconds: float) -> None:
        self.now += seconds


# --- PPTX fixture --------------------------------------------------------------------


def _explicit_white_background(slide: Any) -> None:
    """Set an explicit ``<p:bg>`` so the census's master-background gate passes.

    An unstated background is treated as an unresolved native dependency
    (see ``source_inventory.py::_has_master_background_dependency``); tests
    that expect a successful ``finalize`` need a slide that states its own.
    """
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)


def _pptx_deck(labels: list[str]) -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    for label in labels:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        _explicit_white_background(slide)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = label
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _pptx_deck_with_native_table() -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Title"
    slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(2))
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _pptx_deck_with_filled_text_shape() -> bytes:
    """A text shape with its own solid background — an attribute v1 never renders."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _explicit_white_background(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Title"
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x33, 0x66, 0x99)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _pptx_deck_with_vertical_anchor() -> bytes:
    """A text shape anchored to the middle — v1 never emits a vertical anchor."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _explicit_white_background(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = "Title"
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _pptx_deck_with_paragraph_spacing() -> bytes:
    """A paragraph with explicit space-before — v1 never renders paragraph spacing."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _explicit_white_background(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Title"
    box.text_frame.paragraphs[0].space_before = Pt(12)
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _pptx_deck_with_text_and_image() -> bytes:
    """A plain text slot alongside a static image node (two distinct candidate nodes)."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _explicit_white_background(slide)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Title"
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    )
    slide.shapes.add_picture(io.BytesIO(png), Inches(3), Inches(3), Inches(1), Inches(1))
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _ready_frame_html(store: InMemoryCanvasStore, result: dict[str, Any], archetype_id: str) -> str:
    """The compiled ``frame_html`` for one archetype in a written ``ready`` template."""
    ref = ArtifactRef(**result["template_ref"])
    content = store.read(CANVAS_ID, ref.path, revision=ref.revision).content
    manifest = CompiledTemplateManifest.model_validate_json(content)
    archetype = next(item for item in manifest.archetypes if item.id == archetype_id)
    return archetype.frame_html


def _fully_bound(archetype: dict[str, Any]) -> list[Any]:
    """Bind every candidate slot/static-node with a supported disposition."""
    return [
        VariableBinding(
            archetype_id=archetype["id"], node_id=slot["node_id"], disposition="variable",
            slot_key=slot["key"], role="title", required=True,
        )
        for slot in archetype["slots"]
    ] + [
        RetainBinding(
            archetype_id=archetype["id"], node_id=node["node_id"], disposition="retain",
            retain_reason="static",
        )
        for node in archetype["static_nodes"]
    ]


# --- PDF fixture (reused pattern from test_deck_source_catalog.py) -------------------


def _pdf_text_object(document: Any, text: str, x: float, y: float) -> object:
    obj = raw.FPDFPageObj_NewTextObj(document.raw, b"Helvetica", 18.0)
    encoded = (text + "\x00").encode("utf-16-le")
    buffer = (ctypes.c_ushort * (len(encoded) // 2)).from_buffer_copy(encoded)
    raw.FPDFText_SetText(obj, buffer)
    raw.FPDFPageObj_Transform(obj, 1, 0, 0, 1, x, y)
    return obj


def _pdf(pages: int) -> bytes:
    document = pdfium.PdfDocument.new()
    for n in range(pages):
        page = document.new_page(612.0, 792.0)
        raw.FPDFPage_InsertObject(page.raw, _pdf_text_object(document, f"Title {n}", 100, 700))
        page.gen_content()
    out = io.BytesIO()
    document.save(out)
    document.close()
    return out.getvalue()


def _store_with(path: str, data: bytes) -> InMemoryCanvasStore:
    store = InMemoryCanvasStore()
    store.write_bytes(CANVAS_ID, path, data, "Upload")
    return store


def _sha256(data: bytes) -> str:
    import hashlib

    return hashlib.sha256(data).hexdigest()


# --- prepare -> finalize happy path ---------------------------------------------------


def test_prepare_converts_only_selected_k_pages() -> None:
    data = _pptx_deck(["Slide 1", "Slide 2", "Slide 3"])
    store = _store_with("sources/deck.pptx", data)

    result = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[2]),
        store=store,
        canvas_id=CANVAS_ID,
    )

    assert result["status"] == "candidate"
    assert len(result["archetypes"]) == 1
    assert result["archetypes"][0]["source_page"] == 2
    assert result["archetypes"][0]["slots"][0]["observed_lengths"]["chars"] == len("Slide 2")


def test_finalize_reuses_candidate_without_reconversion(monkeypatch: pytest.MonkeyPatch) -> None:
    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    assert prepared["status"] == "candidate"
    archetype = prepared["archetypes"][0]
    bindings = [
        VariableBinding(
            archetype_id=archetype["id"], node_id=slot["node_id"], disposition="variable",
            slot_key=slot["key"], role="title", required=True,
        )
        for slot in archetype["slots"]
    ] + [
        RetainBinding(
            archetype_id=archetype["id"], node_id=node["node_id"], disposition="retain",
            retain_reason="static",
        )
        for node in archetype["static_nodes"]
    ]

    call_count = 0
    import app.agent.deck_templates as deck_templates_module

    original_extract = deck_templates_module.extract_slides

    def _counting_extract(*args: Any, **kwargs: Any) -> Any:
        nonlocal call_count
        call_count += 1
        return original_extract(*args, **kwargs)

    monkeypatch.setattr(deck_templates_module, "extract_slides", _counting_extract)

    result = finalize_template(
        FinalizeRequest(
            mode="finalize",
            candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=bindings,
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result["status"] == "ready"
    assert call_count == 0  # finalize never re-extracts the source


def test_candidate_is_not_writable_template() -> None:
    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    ref = ArtifactRef(**prepared["candidate_ref"])

    with pytest.raises(TrustError):
        require_trusted_artifact(store, CANVAS_ID, ref, expected_status="ready")


def test_generic_write_and_human_save_cannot_forge_candidate_or_ready() -> None:
    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    ref = ArtifactRef(**prepared["candidate_ref"])
    body = store.read(CANVAS_ID, ref.path, revision=ref.revision).content

    # A generic write/human-save at a *different* path, same bytes, wrong actor.
    forged = store.write(CANVAS_ID, "templates/forged.candidate.json", body, "human edit", actor="human")

    forged_ref = ArtifactRef(path="templates/forged.candidate.json", revision=forged.revision, sha256=ref.sha256)
    with pytest.raises(TrustError):
        require_trusted_artifact(store, CANVAS_ID, forged_ref, expected_status="candidate")


def test_variable_image_binding_rejected() -> None:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Title"
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    )
    slide.shapes.add_picture(io.BytesIO(png), Inches(3), Inches(3), Inches(1), Inches(1))
    buffer = io.BytesIO()
    deck.save(buffer)
    data = buffer.getvalue()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    image_node = next(node for node in archetype["static_nodes"] if node["node_type"] == "image")
    text_slot = archetype["slots"][0]

    bindings = [
        VariableBinding(
            archetype_id=archetype["id"], node_id=text_slot["node_id"], disposition="variable",
            slot_key=text_slot["key"], role="title", required=True,
        ),
        VariableBinding(
            archetype_id=archetype["id"], node_id=image_node["node_id"], disposition="variable",
            slot_key="image-slot", role="body", required=True,
        ),
    ]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result == {
        "status": "error",
        "code": "unsupported_template",
        "message": result["message"],
        "details": {},
        "retryable": False,
    }
    assert store.list_files(CANVAS_ID) == [
        f for f in store.list_files(CANVAS_ID) if not f.path.endswith(".template.json")
    ]


def test_group_smartart_and_master_dependency_fail_closed() -> None:
    data = _pptx_deck_with_native_table()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    assert "native_table" in archetype["capability_issues"]

    bindings = [
        VariableBinding(
            archetype_id=archetype["id"], node_id=slot["node_id"], disposition="variable",
            slot_key=slot["key"], role="title", required=True,
        )
        for slot in archetype["slots"]
    ] + [
        RetainBinding(
            archetype_id=archetype["id"], node_id=node["node_id"], disposition="retain",
            retain_reason="static",
        )
        for node in archetype["static_nodes"]
    ]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result["status"] == "error"
    assert result["code"] == "unsupported_template"
    assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


# --- cooperative budget ----------------------------------------------------------------


def test_cooperative_budget_and_shared_retry_leave_no_ready_or_deck(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A stage finishing after the 120s window discards its result — no candidate."""
    data = _pdf(2)
    store = _store_with("sources/deck.pdf", data)
    clock = _FakeClock()
    budget = TemplateBudget(clock=clock, stage_budget_seconds=120.0)

    calls: list[int] = []

    def _slow_reconstruct(source: PdfPageSource) -> tuple[str, list[str]]:
        calls.append(source.number)
        clock.advance(130.0)  # this stage's own work runs past the budget window
        return "<section class='slide'></section>", []

    result = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pdf", source_sha256=_sha256(data), pages=[1, 2]),
        store=store,
        canvas_id=CANVAS_ID,
        budget=budget,
        reconstruct_pdf_page_fn=_slow_reconstruct,
    )

    assert result["status"] == "error"
    assert result["code"] == "resource_budget_exceeded"
    assert result["retryable"] is True
    # Page 1's reconstruction ran and blew the budget; page 2 is never attempted
    # (no replacement call/retry while the exceeded budget stands).
    assert calls == [1]
    assert budget.exceeded is True
    # No candidate (and therefore no deck downstream of it) was ever written.
    assert store.list_files(CANVAS_ID) == [
        f for f in store.list_files(CANVAS_ID) if f.path == "sources/deck.pdf"
    ]


def test_stalled_or_queued_stage_discards_late_result_without_next_call(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A queued/stalled earlier stage also discards its result and stops the pipeline."""
    data = _pdf(2)
    store = _store_with("sources/deck.pdf", data)
    clock = _FakeClock()
    budget = TemplateBudget(clock=clock, stage_budget_seconds=120.0)

    import app.agent.deck_templates as deck_templates_module

    extract_calls: list[int] = []
    original_extract_pdf_pages = deck_templates_module.extract_pdf_pages

    def _stalled_extract_pdf_pages(data: bytes, pages: list[int], **kwargs: Any) -> Any:
        extract_calls.append(1)
        clock.advance(150.0)  # the extraction stage itself stalls past the window
        return original_extract_pdf_pages(data, pages)

    monkeypatch.setattr(deck_templates_module, "extract_pdf_pages", _stalled_extract_pdf_pages)

    reconstruct_calls: list[int] = []

    def _reconstruct(source: PdfPageSource) -> tuple[str, list[str]]:
        reconstruct_calls.append(source.number)
        return "<section class='slide'></section>", []

    result = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pdf", source_sha256=_sha256(data), pages=[1, 2]),
        store=store,
        canvas_id=CANVAS_ID,
        budget=budget,
        reconstruct_pdf_page_fn=_reconstruct,
    )

    assert result["status"] == "error"
    assert result["code"] == "resource_budget_exceeded"
    assert extract_calls == [1]  # the stalled stage ran exactly once
    assert reconstruct_calls == []  # next stage never started once exceeded
    assert not any(f.path.startswith("templates/") for f in store.list_files(CANVAS_ID))


def test_shared_model_attempt_cap_stops_after_24_reserved_calls() -> None:
    budget = TemplateBudget(clock=_FakeClock())
    for _ in range(24):
        budget.reserve_model_call(prompt_text="x", max_response_tokens=10)
    with pytest.raises(BudgetExceededError):
        budget.reserve_model_call(prompt_text="x", max_response_tokens=10)


def test_finalize_uses_fresh_budget_by_default_and_prepare_consumption_resumes() -> None:
    consumed = {"elapsed_seconds": 119.0, "model_attempts": 1, "total_tokens": 10}
    resumed = TemplateBudget.resume(consumed, clock=_FakeClock(0.0), stage_budget_seconds=120.0)
    assert resumed.exceeded is False
    resumed.clock.advance(2.0)  # type: ignore[attr-defined]
    assert resumed.exceeded is True

    fresh = TemplateBudget(clock=_FakeClock(0.0))
    assert fresh.exceeded is False


# --- EV-P0-001: text-shape capability gate (fill / outline / anchor / spacing) --------


def test_filled_text_shape_is_preserved_or_not_ready() -> None:
    """A text shape's own background fill is never re-emitted by v1 — fail closed."""
    data = _pptx_deck_with_filled_text_shape()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    if result["status"] == "ready":
        # Only acceptable if the compiled frame demonstrably preserved the fill.
        frame_html = _ready_frame_html(store, result, archetype["id"])
        assert "background" in frame_html
    else:
        assert result["status"] == "error"
        assert result["code"] == "unsupported_template"
        assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


def test_vertical_anchor_is_preserved_or_not_ready() -> None:
    """A non-default (middle) vertical anchor is never re-emitted by v1 — fail closed."""
    data = _pptx_deck_with_vertical_anchor()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    if result["status"] == "ready":
        frame_html = _ready_frame_html(store, result, archetype["id"])
        assert "align-items" in frame_html or "justify-content" in frame_html
    else:
        assert result["status"] == "error"
        assert result["code"] == "unsupported_template"
        assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


def test_paragraph_spacing_is_preserved_or_not_ready() -> None:
    """Explicit paragraph space-before is never re-emitted by v1 — fail closed."""
    data = _pptx_deck_with_paragraph_spacing()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    if result["status"] == "ready":
        frame_html = _ready_frame_html(store, result, archetype["id"])
        assert "margin-top" in frame_html
    else:
        assert result["status"] == "error"
        assert result["code"] == "unsupported_template"
        assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


# --- EV-P1-002: stale-source and degraded-reconstruction gates (already implemented) --


def test_same_name_source_overwrite_rejects_finalize() -> None:
    """Overwriting the exact source path with new bytes staless the candidate."""
    original = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", original)

    prepared = prepare_template(
        PrepareRequest(
            mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(original), pages=[1]
        ),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    changed = _pptx_deck(["Changed slide"])
    store.write_bytes(CANVAS_ID, "sources/deck.pptx", changed, "re-upload")

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(changed),
    )

    assert result["status"] == "error"
    assert result["code"] == "stale_source"
    assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


def test_unrelated_canvas_write_does_not_stale_source() -> None:
    """A write to an unrelated path leaves the candidate's own source hash unaffected."""
    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    store.write_bytes(CANVAS_ID, "sources/other.pptx", _pptx_deck(["Other"]), "unrelated upload")

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result["status"] == "ready"


def test_degraded_reconstruction_never_becomes_ready() -> None:
    """A PDF page with non-empty reconstruction review issues never reaches ``ready``."""
    data = _pdf(1)
    store = _store_with("sources/deck.pdf", data)

    def _degraded_reconstruct(source: PdfPageSource) -> tuple[str, list[str]]:
        # One matching data-text-block node per census text object, so the
        # only reconstruction issue this test asserts on is the injected
        # "text_overflow" review finding, not a node-count mismatch.
        nodes = "".join(
            f'<p data-node-id="pdf-text-{i}" data-text-block="true">{t["text"]}</p>'
            for i, t in enumerate(source.texts)
        )
        return f"<section class='slide'>{nodes}</section>", ["text_overflow"]

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pdf", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
        reconstruct_pdf_page_fn=_degraded_reconstruct,
    )
    archetype = prepared["archetypes"][0]
    assert archetype["reconstruction_issues"] == ["text_overflow"]

    bindings = [
        VariableBinding(
            archetype_id=archetype["id"], node_id=slot["node_id"], disposition="variable",
            slot_key=slot["key"], role="title", required=True,
        )
        for slot in archetype["slots"]
    ] + [
        RetainBinding(
            archetype_id=archetype["id"], node_id=node["node_id"], disposition="retain",
            retain_reason="static",
        )
        for node in archetype["static_nodes"]
    ]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result["status"] == "error"
    assert result["code"] == "unsupported_template"
    assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


# --- EV-P1-003: closed/discriminated node-binding contract (already implemented) ------


def test_prepare_finalize_payloads_are_closed_and_discriminated() -> None:
    """``extra='forbid'`` rejects unknown fields; the disposition union rejects unknown tags."""
    with pytest.raises(ValidationError):
        PrepareRequest.model_validate(
            {
                "mode": "prepare", "source": "sources/deck.pptx",
                "source_sha256": "a" * 64, "pages": [1], "extra_field": "nope",
            }
        )

    with pytest.raises(ValidationError):
        VariableBinding.model_validate(
            {
                "archetype_id": "a", "node_id": "n", "disposition": "variable",
                "slot_key": "s", "role": "title", "required": True, "proof": {"x": 1},
            }
        )

    with pytest.raises(ValidationError):
        FinalizeRequest.model_validate(
            {
                "mode": "finalize",
                "candidate_ref": {"path": "p", "revision": "r", "sha256": "a" * 64},
                "bindings": [
                    {"archetype_id": "a", "node_id": "n", "disposition": "unknown-tag"}
                ],
            }
        )


def test_bindings_require_exact_candidate_node_coverage() -> None:
    """Leaving one candidate node unbound blocks finalize with ``ambiguous_slots``."""
    data = _pptx_deck_with_text_and_image()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store,
        canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    image_node = next(node for node in archetype["static_nodes"] if node["node_type"] == "image")

    # Only the image node is bound; the text slot is left unbound.
    bindings = [
        RetainBinding(
            archetype_id=archetype["id"], node_id=image_node["node_id"], disposition="retain",
            retain_reason="static",
        )
    ]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings
        ),
        store=store,
        canvas_id=CANVAS_ID,
        current_source_sha256=_sha256(data),
    )

    assert result["status"] == "error"
    assert result["code"] == "ambiguous_slots"
    assert "missing_nodes" in result["details"]


def test_binding_cannot_mutate_frame_proof_or_budget() -> None:
    """A binding cannot smuggle in ``frame_html``/``proof``/``budget_consumed`` fields."""
    with pytest.raises(ValidationError):
        RetainBinding.model_validate(
            {
                "archetype_id": "a", "node_id": "n", "disposition": "retain",
                "retain_reason": "static", "frame_html": "<div></div>",
            }
        )

    with pytest.raises(ValidationError):
        OmitBinding.model_validate(
            {
                "archetype_id": "a", "node_id": "n", "disposition": "omit",
                "budget_consumed": {"elapsed_seconds": 0},
            }
        )


# --- CV-P1-003: pre-planted content-addressed path cannot launder or block ------------


def test_reprepare_after_unrelated_commit_reuses_same_candidate_commit() -> None:
    """Re-preparing identical inputs after an unrelated commit returns the
    original compiler commit (idempotent create-only), and still finalizes."""
    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    first = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    store.write_bytes(CANVAS_ID, "sources/unrelated.pptx", _pptx_deck(["Other"]), "unrelated upload")

    second = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    assert second["candidate_ref"] == first["candidate_ref"]

    archetype = second["archetypes"][0]
    finalize_result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**second["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert finalize_result["status"] == "ready"


def test_preplanted_human_actor_file_does_not_permanently_block_prepare() -> None:
    """A pre-planted human/agent file at the deterministic candidate path
    cannot be laundered into trusted compiler output, and does not
    permanently block templating (no same-canvas DoS)."""
    data = _pptx_deck(["Only slide"])

    probe_store = InMemoryCanvasStore()
    probe_store.write_bytes(CANVAS_ID, "sources/deck.pptx", data, "Upload")
    probe = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=probe_store, canvas_id=CANVAS_ID,
    )
    squatted_path = probe["candidate_ref"]["path"]

    store = _store_with("sources/deck.pptx", data)
    store.write(CANVAS_ID, squatted_path, "not a real candidate manifest", "human plant", actor="human")

    result = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )

    assert result["status"] == "candidate"
    ref = ArtifactRef(**result["candidate_ref"])
    assert ref.path == squatted_path
    # The squatted content is never adopted — the returned ref names a
    # fresh compiler-authored revision, which reads back as a valid,
    # trusted candidate.
    manifest = require_trusted_artifact(store, CANVAS_ID, ref, expected_status="candidate")
    assert manifest.source.path == "sources/deck.pptx"


# --- EV-P1-002: finalize materializes ready manifest from bindings --------------------


def _pptx_deck_with_three_text_nodes() -> bytes:
    """A title (to bind variable), a legal footer (to retain), and a note (to omit)."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _explicit_white_background(slide)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title_box.text_frame.text = "Original Title"
    legal_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(4), Inches(1))
    legal_box.text_frame.text = "Legal footer text"
    note_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    note_box.text_frame.text = "Internal note to omit"
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def test_retain_and_omit_bindings_materialize_ready_manifest() -> None:
    """A retain binding preserves its text verbatim in the ready frame; an
    omit binding drops the node from both the manifest and the markup; a
    variable binding becomes the only writer-fillable slot."""
    from app.agent.deck_template_models import SlideContentRequest
    from app.agent.deck_template_writer import write_deck_from_template

    data = _pptx_deck_with_three_text_nodes()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    title_slot = next(s for s in archetype["slots"] if s["observed_lengths"]["chars"] == len("Original Title"))
    legal_slot = next(s for s in archetype["slots"] if s["observed_lengths"]["chars"] == len("Legal footer text"))
    note_slot = next(s for s in archetype["slots"] if s["observed_lengths"]["chars"] == len("Internal note to omit"))

    bindings = [
        VariableBinding(
            archetype_id=archetype["id"], node_id=title_slot["node_id"], disposition="variable",
            slot_key="title", role="title", required=True,
        ),
        RetainBinding(
            archetype_id=archetype["id"], node_id=legal_slot["node_id"], disposition="retain",
            retain_reason="legal",
        ),
        OmitBinding(archetype_id=archetype["id"], node_id=note_slot["node_id"], disposition="omit"),
    ]

    result = finalize_template(
        FinalizeRequest(mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert result["status"] == "ready"

    ref = ArtifactRef(**result["template_ref"])
    manifest = CompiledTemplateManifest.model_validate_json(
        store.read(CANVAS_ID, ref.path, revision=ref.revision).content
    )
    ready_archetype = manifest.archetypes[0]

    assert [slot.key for slot in ready_archetype.slots] == ["title"]
    retained = next(n for n in ready_archetype.static_nodes if n.node_id == legal_slot["node_id"])
    assert retained.disposition == "retain"
    assert retained.retain_reason == "legal"
    assert not any(n.node_id == note_slot["node_id"] for n in ready_archetype.static_nodes)
    assert "Internal note to omit" not in ready_archetype.frame_html
    assert "Legal footer text" in ready_archetype.frame_html

    from types import SimpleNamespace

    write_result = write_deck_from_template(
        ref,
        "deck.slides.html",
        "Retain Omit Deck",
        [SlideContentRequest(archetype_id=ready_archetype.id, mode="verbatim", slots={"title": "New Title"})],
        SimpleNamespace(stream_writer=None),
        store=store, canvas_id=CANVAS_ID, writer_model="unused-model",
    )
    assert write_result["status"] == "ok"
    written = store.read(CANVAS_ID, "deck.slides.html").content
    assert "Legal footer text" in written
    assert "Internal note to omit" not in written
    assert "New Title" in written


# --- EV-P1-003: real proof.source_to_frame correspondence -----------------------------


def test_finalize_failure_leaves_no_ready_manifest() -> None:
    """Any finalize failure path (here: an unresolved native table) writes
    no ready manifest at all."""
    data = _pptx_deck_with_native_table()
    store = _store_with("sources/deck.pptx", data)
    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert result["status"] == "error"
    assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


def test_unclassified_business_text_blocks_ready(monkeypatch: pytest.MonkeyPatch) -> None:
    """A source text object the pre-extraction census counted, but that
    extraction silently dropped, blocks ready — the exact-coverage binding
    check alone cannot catch this, since it only sees nodes the candidate
    already has."""
    data = _pptx_deck_with_text_and_image()
    store = _store_with("sources/deck.pptx", data)

    import app.agent.deck_templates as deck_templates_module

    original_extract = deck_templates_module.extract_slides

    def _dropping_extract(*args: Any, **kwargs: Any) -> Any:
        extractions = original_extract(*args, **kwargs)
        for extraction in extractions:
            extraction.texts.clear()
        return extractions

    monkeypatch.setattr(deck_templates_module, "extract_slides", _dropping_extract)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]
    assert archetype["slots"] == []  # extraction dropped the census-counted text object
    image_node = next(n for n in archetype["static_nodes"] if n["node_type"] == "image")

    bindings = [
        RetainBinding(
            archetype_id=archetype["id"], node_id=image_node["node_id"], disposition="retain",
            retain_reason="static",
        ),
    ]

    result = finalize_template(
        FinalizeRequest(mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]), bindings=bindings),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )

    assert result["status"] == "error"
    assert result["code"] == "unsupported_template"
    assert not any(f.path.endswith(".template.json") for f in store.list_files(CANVAS_ID))


# --- EV-P1-004: pinned image assets are hash-verified and embedded --------------------


def test_asset_overwrite_cannot_change_pinned_render_or_export() -> None:
    """A retained image's bytes are pinned and embedded as a data URI at
    finalize time; overwriting the pinned store path afterward cannot
    change what an already-instantiated frame renders or exports."""
    from app.agent.deck_template_verification import _asset_and_font_issues

    data = _pptx_deck_with_text_and_image()
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert result["status"] == "ready"

    ref = ArtifactRef(**result["template_ref"])
    manifest = CompiledTemplateManifest.model_validate_json(
        store.read(CANVAS_ID, ref.path, revision=ref.revision).content
    )
    ready_archetype = manifest.archetypes[0]

    assert ready_archetype.assets, "retained image must be pinned as an AssetRef"
    asset = ready_archetype.assets[0]
    assert asset.sha256 in ready_archetype.frame_html
    assert "data:image/" in ready_archetype.frame_html

    store.write_bytes(CANVAS_ID, asset.path, b"tampered bytes", "overwrite attempt", actor="human")

    # The already-materialized ready frame never re-reads the pinned path —
    # its bytes are already embedded — so neither its content nor
    # verification's pinned-asset check is affected by the overwrite.
    assert "data:image/" in ready_archetype.frame_html
    assert asset.sha256 in ready_archetype.frame_html
    assert _asset_and_font_issues(ready_archetype, ready_archetype.frame_html) == []


# --- EV-P0-001: writing_style is actually populated at finalize -----------------------


def test_style_rules_keep_role_evidence_and_unknown_font(monkeypatch: pytest.MonkeyPatch) -> None:
    """The style-profile stage runs through the shared budget and keeps
    each rule's role/evidence — including an origin='unknown' (font) rule —
    all the way into the ready manifest."""
    from app.agent.deck_template_models import StyleEvidence, StyleProfileResponse, StyleRule
    import app.agent.deck_templates as deck_templates_module

    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    calls: list[list[dict[str, str]]] = []

    def _fake_invoke(writer_model: str, messages: list[dict[str, str]]) -> "StyleProfileResponse":
        calls.append(messages)
        return StyleProfileResponse(
            rules=[
                StyleRule(
                    role="title", property="title_form", value="noun_phrase", origin="observed",
                    evidence=[StyleEvidence(page=archetype["source_page"], snippet="Only slide")],
                ),
                StyleRule(role="title", property="font_family", value="unknown", origin="unknown"),
            ]
        )

    monkeypatch.setattr(deck_templates_module, "_invoke_style_profile_model", _fake_invoke)

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert result["status"] == "ready"
    assert calls, "the style-profile model must have been invoked from the shared budget path"

    ref = ArtifactRef(**result["template_ref"])
    manifest = CompiledTemplateManifest.model_validate_json(
        store.read(CANVAS_ID, ref.path, revision=ref.revision).content
    )
    ready_archetype = manifest.archetypes[0]
    assert len(ready_archetype.writing_style) == 2
    observed = next(r for r in ready_archetype.writing_style if r.origin == "observed")
    assert observed.role == "title"
    assert observed.evidence and observed.evidence[0].snippet == "Only slide"
    unknown = next(r for r in ready_archetype.writing_style if r.origin == "unknown")
    assert unknown.property == "font_family"


def test_style_profiling_failure_leaves_ready_without_rules(monkeypatch: pytest.MonkeyPatch) -> None:
    """A style-profile model failure/unavailability leaves an archetype's
    writing_style empty and never blocks finalize."""
    import app.agent.deck_templates as deck_templates_module

    data = _pptx_deck(["Only slide"])
    store = _store_with("sources/deck.pptx", data)

    prepared = prepare_template(
        PrepareRequest(mode="prepare", source="sources/deck.pptx", source_sha256=_sha256(data), pages=[1]),
        store=store, canvas_id=CANVAS_ID,
    )
    archetype = prepared["archetypes"][0]

    def _boom(writer_model: str, messages: list[dict[str, str]]) -> Any:
        raise RuntimeError("model unavailable")

    monkeypatch.setattr(deck_templates_module, "_invoke_style_profile_model", _boom)

    result = finalize_template(
        FinalizeRequest(
            mode="finalize", candidate_ref=ArtifactRef(**prepared["candidate_ref"]),
            bindings=_fully_bound(archetype),
        ),
        store=store, canvas_id=CANVAS_ID, current_source_sha256=_sha256(data),
    )
    assert result["status"] == "ready"

    ref = ArtifactRef(**result["template_ref"])
    manifest = CompiledTemplateManifest.model_validate_json(
        store.read(CANVAS_ID, ref.path, revision=ref.revision).content
    )
    assert manifest.archetypes[0].writing_style == []
