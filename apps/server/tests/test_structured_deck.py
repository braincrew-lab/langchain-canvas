"""Native tables/charts keep structure when their content changes."""

import io

from langchain_canvas.deck import baseline_slide_html, extract_slides
from lxml import html
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def source_presentation():
    p = Presentation()
    p.slide_width = Inches(13.333333)
    p.slide_height = Inches(7.5)
    slide = p.slides.add_slide(p.slide_layouts[6])
    table = slide.shapes.add_table(
        2, 2, Inches(0.5), Inches(0.5), Inches(5), Inches(2)
    ).table
    for row, values in enumerate([["Region", "Revenue"], ["East", "20"]]):
        for col, value in enumerate(values):
            cell = table.cell(row, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                "663399" if row == 0 else "EEEEEE"
            )
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(16)
            run.font.bold = row == 0
    data = CategoryChartData()
    data.categories = ["Q1", "Q2"]
    data.add_series("Sales", [10, 20])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6),
        Inches(0.5),
        Inches(6),
        Inches(4),
        data,
    ).chart
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor.from_string("663399")
    chart.has_title = True
    chart.chart_title.text_frame.text = "Revenue"
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def test_import_native_table_and_chart_as_structured_html():
    data = source_presentation()
    extracted = extract_slides(data, path="source.pptx")[0]
    body = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    root = html.fromstring(body)
    tables = root.xpath('.//table[@data-pptx-shape-id="e0"]')
    assert len(tables) == 1
    assert [cell.text_content() for cell in tables[0].xpath(".//td|.//th")] == [
        "Region",
        "Revenue",
        "East",
        "20",
    ]
    assert len(root.xpath('.//*[@data-chart-type="column"]')) == 1
    assert not root.xpath(".//img")
    from langchain_canvas.deck import sanitize_slide_html

    assert not sanitize_slide_html(body).removed


def test_chart_update_changes_preview_and_native_data_without_changing_style():
    import base64
    import json

    from app.agent.exports import EditableDeckPptxExporter
    from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
    from langchain_canvas.deck.structured import replace_chart_html
    from langchain_canvas.exporters import PPTX_MIME

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    node = html.fromstring(body).xpath(".//*[@data-chart-data]")[0]
    updated = replace_chart_html(
        body,
        node.get("data-node-id"),
        ["H1", "H2"],
        [{"name": "Profit", "values": [25, 15]}],
    )
    chart_node = html.fromstring(updated).xpath(".//*[@data-chart-data]")[0]
    assert json.loads(chart_node.get("data-chart-data"))["series"][0]["values"] == [
        25,
        15,
    ]
    assert "H1" in chart_node.text_content()
    original_node = html.fromstring(body).xpath(".//*[@data-chart-data]")[0]
    assert chart_node.get("style") == original_node.get("style")
    content = serialize_deck(
        Deck(
            "Native",
            "16:9",
            f"data:{PPTX_MIME};base64,{base64.b64encode(source).decode()}",
            [SlideTemplate("s1", None, "", updated)],
        )
    )
    result = EditableDeckPptxExporter().export(content, path="native.slides.html")
    slide = Presentation(io.BytesIO(result.data)).slides[0]
    assert len(slide.shapes) == 2
    assert slide.shapes[0].has_table
    chart = slide.shapes[1].chart
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert list(chart.series[0].values) == [25, 15]
    assert chart.series[0].name == "Profit"
    assert str(chart.series[0].format.fill.fore_color.rgb) == "663399"
    assert chart.chart_title.text_frame.text == "Revenue"


def native_export(source, body, style_css=""):
    import base64

    from app.agent.exports import EditableDeckPptxExporter
    from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
    from langchain_canvas.exporters import PPTX_MIME

    deck = serialize_deck(
        Deck(
            "Native",
            "16:9",
            f"data:{PPTX_MIME};base64,{base64.b64encode(source).decode()}",
            [SlideTemplate("s1", None, style_css, body)],
        )
    )
    return Presentation(
        io.BytesIO(
            EditableDeckPptxExporter().export(deck, path="native.slides.html").data
        )
    )


def test_table_cell_updates_keep_native_structure_and_rich_font():
    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    cell = root.xpath(".//table//td")[-1]
    cell.xpath(".//span")[0].text = "35"
    p = native_export(source, html.tostring(root, encoding="unicode"))
    table = p.slides[0].shapes[0].table
    assert table.cell(1, 1).text == "35"
    assert str(table.cell(0, 0).fill.fore_color.rgb) == "663399"
    run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert run.font.bold and run.font.name == "Arial"
    assert len(table.rows) == 2 and len(table.columns) == 2
    assert p.slides[0].shapes[1].has_chart


def test_native_table_export_rejects_css_bypasses_and_resizing():
    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    with pytest.raises(ValueError, match="custom stylesheets"):
        native_export(source, body, "td{color:red!important}")
    with pytest.raises(ValueError, match="custom stylesheets"):
        native_export(
            source, body.replace("</section>", "<style>td{color:red}</style></section>")
        )
    root = html.fromstring(body)
    root.set("style", root.get("style", "") + ";color:red")
    with pytest.raises(ValueError, match="inherited styling"):
        native_export(source, html.tostring(root, encoding="unicode"))
    root = html.fromstring(body)
    table = root.xpath(".//table")[0]
    table.set("style", table.get("style") + ";width:600px")
    with pytest.raises(ValueError, match="sizing and styling"):
        native_export(source, html.tostring(root, encoding="unicode"))
    root = html.fromstring(body)
    table = root.xpath(".//table")[0]
    table.set("style", table.get("style") + ";left:80px;top:80px")
    moved = native_export(source, html.tostring(root, encoding="unicode"))
    assert moved.slides[0].shapes[0].left == 80 * 9525


def test_native_table_export_rejects_row_and_rowgroup_styling():
    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    for row_style in ("color:#ff0000", "background:#ff0000", "height:300px"):
        root = html.fromstring(body)
        root.xpath(".//tr")[1].set("style", row_style)
        with pytest.raises(ValueError, match="Native table"):
            native_export(source, html.tostring(root, encoding="unicode"))
    root = html.fromstring(body)
    table = root.xpath(".//table")[0]
    group = html.Element("tbody")
    for row in list(table):
        group.append(row)
    table.append(group)
    assert (
        native_export(source, html.tostring(root, encoding="unicode"))
        .slides[0]
        .shapes[0]
        .has_table
    )
    group.set("style", "color:#ff0000")
    with pytest.raises(ValueError, match="Native table"):
        native_export(source, html.tostring(root, encoding="unicode"))


def test_native_table_export_rejects_unmapped_outer_content():
    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    root.xpath(".//table")[0].insert(
        0, html.fromstring("<caption>Unmapped caption</caption>")
    )
    with pytest.raises(ValueError, match="Native table"):
        native_export(source, html.tostring(root, encoding="unicode"))
    root = html.fromstring(body)
    root.xpath(".//td")[0].tail = "Unmapped text"
    with pytest.raises(ValueError, match="Native table"):
        native_export(source, html.tostring(root, encoding="unicode"))


def test_merged_native_table_roundtrips_without_duplicate_cells():
    p = Presentation(io.BytesIO(source_presentation()))
    p.slides[0].shapes[0].table.cell(0, 0).merge(p.slides[0].shapes[0].table.cell(0, 1))
    out = io.BytesIO()
    p.save(out)
    source = out.getvalue()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    assert html.fromstring(body).xpath('.//th[@colspan="2"]')
    result = native_export(source, body)
    cell = result.slides[0].shapes[0].table.cell(0, 0)
    assert cell.is_merge_origin and cell.span_width == 2
    assert "Region" in cell.text and "Revenue" in cell.text


def test_bar_chart_negative_zero_values_and_unchanged_update():
    from langchain_canvas.deck.structured import (
        StructuredShape,
        replace_chart_html,
        structured_html,
    )

    data = {
        "type": "bar",
        "categories": ["A", "B", "C"],
        "series": [{"name": "Results", "values": [-2, 0, 4], "color": "#663399"}],
        "legend": True,
    }
    body = (
        '<section class="slide">'
        + structured_html(
            StructuredShape("e0", 10, 10, 70, 60, "chart", data), "chart", (1280, 720)
        )
        + "</section>"
    )
    assert "Results" in body and "<img" not in body
    assert (
        replace_chart_html(
            body,
            "chart",
            data["categories"],
            [{"name": "Results", "values": [-2, 0, 4]}],
        )
        == body
    )
    updated = replace_chart_html(
        body, "chart", ["X", "Y", "Z"], [{"name": "Next", "values": [0, 0, 0]}]
    )
    assert "Next" in html.fromstring(updated).text_content()
    from app.agent.render import measure_slide

    layout = measure_slide(updated, ratio="16:9")
    assert not layout["unsupported"]
    assert all(i["w"] >= 0 for i in layout["items"] if "w" in i)


def test_chart_metadata_rejects_css_injection_and_nonfinite_dimensions():
    import pytest
    from langchain_canvas.deck.structured import chart_inner_html, validate_chart_data

    base = {
        "type": "bar",
        "categories": ["A"],
        "series": [{"name": "S", "values": [1], "color": "#123456"}],
    }
    for changes in (
        {"font_family": "Arial;color:red"},
        {"font_size": float("nan")},
        {"color": "red;opacity:0"},
        {"title": "X" * 2001},
    ):
        with pytest.raises(ValueError):
            validate_chart_data({**base, **changes})
    with pytest.raises(ValueError):
        chart_inner_html(base, float("inf"), 400)


def test_chart_metadata_guards_malformed_values_and_unsupported_edits():
    import copy
    import json

    import pytest
    from langchain_canvas.deck.structured import replace_chart_html, validate_chart_data

    base = {
        "type": "column",
        "categories": ["A"],
        "series": [{"name": "S", "values": [1], "color": "#123456"}],
    }
    cases = [
        {"type": "pie"},
        {"categories": []},
        {"categories": [1]},
        {"series": []},
        {"series": [{"name": 3, "values": [1]}]},
        {"series": [{"name": "S", "values": [float("nan")]}]},
        {"series": [{"name": "S", "values": [1], "color": "red"}]},
    ]
    for delta in cases:
        with pytest.raises(ValueError):
            validate_chart_data({**base, **delta})
    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    node = root.xpath(".//*[@data-chart-data]")[0]
    ident = node.get("data-node-id")
    for target, categories, series in [
        ("missing", ["A", "B"], [{"name": "S", "values": [1, 2]}]),
        (ident, ["A"], [{"name": "S", "values": [1]}]),
    ]:
        with pytest.raises(ValueError):
            replace_chart_html(body, target, categories, series)
    for attr, value in [
        ("data-chart-type", "bar"),
        ("style", "width:50%;height:400px"),
        ("style", "width:oops-px;height:400px"),
    ]:
        altered = copy.deepcopy(root)
        chart = altered.xpath(".//*[@data-chart-data]")[0]
        chart.set(attr, value)
        with pytest.raises(ValueError):
            replace_chart_html(
                html.tostring(altered, encoding="unicode"),
                ident,
                ["X", "Y"],
                [{"name": "New", "values": [20, 30]}],
            )
    assert json.loads(node.get("data-chart-data"))["series"][0]["values"] == [10, 20]


def test_native_export_rejects_chart_or_table_topology_mismatch():
    import copy
    import json

    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    original = html.fromstring(body)
    for change in (
        "chart_type",
        "chart_count",
        "unknown_chart",
        "table_shape",
        "table_rows",
        "table_cols",
        "table_span",
        "table_semantic_id",
    ):
        root = copy.deepcopy(original)
        chart = root.xpath(".//*[@data-chart-data]")[0]
        table = root.xpath(".//table")[0]
        if change.startswith("chart_"):
            data = json.loads(chart.get("data-chart-data"))
            if change == "chart_type":
                data["type"] = "bar"
            else:
                data["categories"].append("Extra")
                data["series"][0]["values"].append(40)
            chart.set("data-chart-data", json.dumps(data))
        elif change == "unknown_chart":
            chart.set("data-pptx-shape-id", "e99")
        elif change == "table_shape":
            table.set("data-pptx-shape-id", "e1")
            chart.getparent().remove(chart)
        elif change == "table_rows":
            table.remove(table.xpath("./tr")[-1])
        elif change == "table_cols":
            table.xpath("./tr")[-1].remove(table.xpath("./tr")[-1][-1])
        elif change == "table_span":
            table.xpath(".//th")[0].set("colspan", "2")
        else:
            table.xpath(".//td")[-1].attrib.pop("data-node-id")
        with pytest.raises(ValueError):
            native_export(source, html.tostring(root, encoding="unicode"))


def test_unsupported_native_chart_is_reported_by_extraction():
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("S", [1, 2])
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(5), Inches(3), data
    )
    output = io.BytesIO()
    p.save(output)
    extracted = extract_slides(output.getvalue(), path="line.pptx")[0]
    assert not extracted.structured
    assert any("chart" in message.lower() for message in extracted.warnings)


def test_chart_parts_have_stable_ids_across_data_edits():
    from langchain_canvas.deck.structured import replace_chart_html

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    chart = root.xpath(".//*[@data-chart-data]")[0]
    ids = chart.xpath("./*/@data-node-id")
    assert len(ids) >= 5 and len(ids) == len(set(ids))
    updated = replace_chart_html(
        body, chart.get("data-node-id"), ["X", "Y"], [{"name": "S", "values": [30, 10]}]
    )
    assert (
        html.fromstring(updated).xpath(".//*[@data-chart-data]/*/@data-node-id") == ids
    )


def test_native_table_geometry_is_independent_of_host_box_sizing():
    from app.agent.render import measure_slide

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    layout = measure_slide("<style>body{margin:0}</style>" + body, ratio="16:9")
    table = next(e for e in layout["elements"] if e["tag"] == "table")
    assert abs(table["w"] - 480) < 1
    cells = [e for e in layout["elements"] if e["tag"] in {"td", "th"}]
    assert all(e["style"]["text-align"] == "left" for e in cells)


def test_legacy_text_edit_cannot_desynchronize_chart_preview_from_data():
    from types import SimpleNamespace

    from app.agent.deck_editing import create_deck_editing_tools
    from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
    from langchain_canvas.store import InMemoryCanvasStore

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    label = root.xpath(".//*[@data-chart-data]/p")[-1].get("data-node-id")
    store = InMemoryCanvasStore()
    revision = store.write(
        "test",
        "deck.slides.html",
        serialize_deck(
            Deck("Test", "16:9", None, [SlideTemplate("s1", None, "", body)])
        ),
        "Seed",
    ).revision
    runtime = SimpleNamespace(
        context=None, config={"configurable": {"thread_id": "test"}}, stream_writer=None
    )
    tool = next(
        t for t in create_deck_editing_tools(store) if t.name == "replace_slide_text"
    )
    result = tool.func(
        path="deck.slides.html",
        slide_id="s1",
        revision=revision,
        node_id=label,
        text="Changed label",
        runtime=runtime,
    )
    assert result["status"] == "error"
    assert store.read("test", "deck.slides.html").revision == revision


def test_chart_data_operation_repairs_corrupt_preview_even_with_same_values():
    import json

    from langchain_canvas.deck.structured import replace_chart_html

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    chart = root.xpath(".//*[@data-chart-data]")[0]
    chart.xpath("./p")[-1].text = "Wrong label"
    corrupt = html.tostring(root, encoding="unicode")
    data = json.loads(chart.get("data-chart-data"))
    repaired = replace_chart_html(
        corrupt,
        chart.get("data-node-id"),
        data["categories"],
        [{"name": s["name"], "values": s["values"]} for s in data["series"]],
    )
    assert "Wrong label" not in html.fromstring(repaired).text_content()
    assert "Q2" in html.fromstring(repaired).text_content()


def test_native_chart_style_changes_cannot_be_silently_dropped_on_export():
    import json

    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    chart = root.xpath(".//*[@data-chart-data]")[0]
    chart.set("style", chart.get("style") + ";background:red")
    with pytest.raises(ValueError):
        native_export(source, html.tostring(root, encoding="unicode"))
    root = html.fromstring(body)
    chart = root.xpath(".//*[@data-chart-data]")[0]
    data = json.loads(chart.get("data-chart-data"))
    data["title"] = "Changed title"
    chart.set("data-chart-data", json.dumps(data))
    from langchain_canvas.deck.structured import chart_inner_html

    for child in list(chart):
        chart.remove(child)
    for child in html.fragments_fromstring(
        chart_inner_html(data, 576, 383.984, chart.get("data-node-id"))
    ):
        chart.append(child)
    with pytest.raises(ValueError):
        native_export(source, html.tostring(root, encoding="unicode"))


def inherited_table_source():
    from pptx.oxml.xmlchemy import OxmlElement

    presentation = Presentation(io.BytesIO(source_presentation()))
    table = presentation.slides[0].shapes[0].table
    for row in table.rows:
        for cell in row.cells:
            cell.text = "Inherited"
    paragraph = table.cell(1, 0).text_frame.paragraphs[0]
    paragraph.clear()
    for text in ("Adjacent ", "unstyled "):
        paragraph.add_run().text = text
    run = paragraph.add_run()
    run.text = "linked"
    run.font.bold = True
    run.font.color.theme_color = __import__(
        "pptx.enum.dml", fromlist=["MSO_THEME_COLOR"]
    ).MSO_THEME_COLOR.ACCENT_1
    run.hyperlink.address = "https://example.com/source"
    paragraph.space_before = Pt(7)
    paragraph.space_after = Pt(9)
    table.cell(1, 0).margin_left = Inches(0.2)
    node = OxmlElement("a:latin")
    node.set("typeface", "+mn-lt")
    paragraph._p.get_or_add_pPr().get_or_add_defRPr().append(node)
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_gradient_chart_does_not_abort_table_import_or_drop_native_chart():
    presentation = Presentation(io.BytesIO(source_presentation()))
    chart = presentation.slides[0].shapes[1].chart
    chart.series[0].format.fill.gradient()
    output = io.BytesIO()
    presentation.save(output)
    source = output.getvalue()
    extracted = extract_slides(source, path="gradient.pptx")[0]
    assert [shape.kind for shape in extracted.structured] == ["table"]
    assert any("chart" in warning.lower() for warning in extracted.warnings)
    body = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    result = native_export(source, body)
    assert (
        result.slides[0].shapes[1].chart.series[0].format.fill.type
        == chart.series[0].format.fill.type
    )
    assert result.slides[0].shapes[0].has_table


def test_inherited_table_preview_resolves_theme_font_and_neutralizes_header_bold():
    from app.agent.render import measure_slide

    source = inherited_table_source()
    extracted = extract_slides(source, path="table.pptx")[0]
    body = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    layout = measure_slide(body, ratio="16:9")
    cell = next(e for e in layout["elements"] if e["tag"] == "th")
    assert cell["style"]["font-family"] == "Calibri"
    assert cell["style"]["font-size"] == "24px"
    assert cell["style"]["font-weight"] == "400"
    assert any("table style" in warning.lower() for warning in extracted.warnings)


def test_unchanged_native_table_text_xml_is_identical_after_export():
    source = inherited_table_source()
    original = Presentation(io.BytesIO(source)).slides[0].shapes[0].table
    body = baseline_slide_html(
        extract_slides(source, path="t.pptx")[0], slide_id="s1", ratio="16:9"
    )
    result = native_export(source, body).slides[0].shapes[0].table
    for r in range(2):
        for c in range(2):
            assert (
                result.cell(r, c).text_frame._txBody.xml
                == original.cell(r, c).text_frame._txBody.xml
            )


def test_native_table_content_edit_only_changes_xml_text_nodes():
    from lxml import etree

    source = inherited_table_source()
    original = Presentation(io.BytesIO(source)).slides[0].shapes[0].table.cell(1, 0)
    body = baseline_slide_html(
        extract_slides(source, path="t.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    cell = root.xpath(".//table/tr[2]/td")[0]
    spans = cell.xpath("./span")
    assert len(spans) == 3
    for node, text in zip(spans, ["New ", "coherent ", "content"], strict=True):
        node.text = text
    exported = (
        native_export(source, html.tostring(root, encoding="unicode"))
        .slides[0]
        .shapes[0]
        .table.cell(1, 0)
    )
    import copy

    expected = copy.deepcopy(original.text_frame._txBody)
    for node, text in zip(
        expected.xpath(".//a:t"), ["New ", "coherent ", "content"], strict=True
    ):
        node.text = text
    assert etree.tostring(
        exported.text_frame._txBody, method="c14n", exclusive=True
    ) == etree.tostring(expected, method="c14n", exclusive=True)
    assert exported.text == "New coherent content"
    assert (
        exported.text_frame.paragraphs[0].runs[-1].hyperlink.address
        == "https://example.com/source"
    )


def test_native_table_style_edit_fails_closed_instead_of_flattening_inheritance():
    import pytest

    source = inherited_table_source()
    body = baseline_slide_html(
        extract_slides(source, path="t.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    cell = root.xpath(".//table/tr[2]/td")[0]
    cell.set("style", cell.get("style") + ";font-family:Courier New")
    with pytest.raises(ValueError, match="Native table.*style"):
        native_export(source, html.tostring(root, encoding="unicode"))


def test_canonical_chart_changes_reach_specific_native_safety_guards():
    import json
    import re

    import pytest
    from langchain_canvas.deck.structured import chart_inner_html, validate_chart_markup

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    for change, message in [
        ("type", "Chart metadata does not match the original native chart type"),
        ("count", "Native chart category/series counts must be preserved"),
        ("title", "Native chart style/title changes are unsupported"),
        ("color", "Native chart style/title changes are unsupported"),
    ]:
        root = html.fromstring(body)
        chart = root.xpath(".//*[@data-chart-data]")[0]
        data = json.loads(chart.get("data-chart-data"))
        if change == "type":
            data["type"] = "bar"
            chart.set("data-chart-type", "bar")
        elif change == "count":
            data["categories"].append("Q3")
            data["series"][0]["values"].append(30)
        elif change == "title":
            data["title"] = "New title"
        else:
            data["series"][0]["color"] = "#00FF00"
        chart.set("data-chart-data", json.dumps(data))
        styles = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", chart.get("style")))
        width, height = (float(styles[k][:-2]) for k in ("width", "height"))
        for child in list(chart):
            chart.remove(child)
        for child in html.fragments_fromstring(
            chart_inner_html(data, width, height, chart.get("data-node-id"))
        ):
            chart.append(child)
        candidate = html.tostring(root, encoding="unicode")
        assert (
            validate_chart_markup(
                dict(chart.attrib),
                "".join(html.tostring(child, encoding="unicode") for child in chart),
            )
            == data
        )
        with pytest.raises(ValueError, match=message):
            native_export(source, candidate)


def test_legacy_unstyled_run_markup_remains_exportable_but_ambiguous_edits_fail():
    import pytest

    source = inherited_table_source()
    body = baseline_slide_html(
        extract_slides(source, path="t.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    for cell in root.xpath(".//td|.//th"):
        cell.set(
            "style",
            cell.get("style")
            .replace("font-weight:400;", "")
            .replace("font-family:Calibri;", "")
            .replace("font-size:24px;", ""),
        )
        for span in list(cell.xpath("./span")):
            import re

            style = re.sub(
                r"(?:font-family:Calibri|font-size:24px);?", "", span.get("style", "")
            ).strip(";")
            span.set("style", style)
            if not style:
                span.drop_tag()
    legacy = html.tostring(root, encoding="unicode")
    result = native_export(source, legacy)
    original = Presentation(io.BytesIO(source)).slides[0].shapes[0].table
    assert (
        result.slides[0].shapes[0].table.cell(1, 0).text_frame._txBody.xml
        == original.cell(1, 0).text_frame._txBody.xml
    )
    cell = root.xpath(".//table/tr[2]/td")[0]
    cell.text = "Ambiguous replacement "
    with pytest.raises(ValueError, match="rich slots are ambiguous"):
        native_export(source, html.tostring(root, encoding="unicode"))


def test_empty_native_table_cell_can_be_filled_without_rewriting_body_properties():
    presentation = Presentation(io.BytesIO(source_presentation()))
    cell = presentation.slides[0].shapes[0].table.cell(1, 1)
    cell.text = ""
    cell.margin_top = Inches(0.15)
    original_body = cell.text_frame._txBody.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
    )
    from lxml import etree

    before = etree.tostring(original_body)
    output = io.BytesIO()
    presentation.save(output)
    source = output.getvalue()
    body = baseline_slide_html(
        extract_slides(source, path="empty.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    root.xpath(".//table/tr[2]/td")[1].text = "Filled"
    result = (
        native_export(source, html.tostring(root, encoding="unicode"))
        .slides[0]
        .shapes[0]
        .table.cell(1, 1)
    )
    assert result.text == "Filled"
    assert (
        etree.tostring(
            result.text_frame._txBody.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
            )
        )
        == before
    )


def test_non_solid_chart_fills_are_retained_without_html_conversion():
    from langchain_canvas.deck.structured import native_chart_data

    for fill_kind in ("patterned", "background"):
        presentation = Presentation(io.BytesIO(source_presentation()))
        chart = presentation.slides[0].shapes[1].chart
        getattr(chart.series[0].format.fill, fill_kind)()
        assert native_chart_data(chart, {}) is None


def test_table_mixed_explicit_and_inherited_fonts_keep_each_native_run_style():
    from app.agent.render import measure_slide

    presentation = Presentation(io.BytesIO(source_presentation()))
    cell = presentation.slides[0].shapes[0].table.cell(1, 0)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.clear()
    explicit = paragraph.add_run()
    explicit.text = "Explicit "
    explicit.font.name = "Arial"
    explicit.font.size = Pt(24)
    paragraph.add_run().text = "Inherited"
    output = io.BytesIO()
    presentation.save(output)
    source = output.getvalue()
    body = baseline_slide_html(
        extract_slides(source, path="fonts.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    ident = root.xpath(".//table/tr[2]/td")[0].get("data-node-id")
    block = next(
        b for b in measure_slide(body, ratio="16:9")["textBlocks"] if b["id"] == ident
    )
    runs = block["paragraphs"][0]["runs"]
    assert runs[0]["font"] == "Arial" and runs[0]["size"] == 32
    assert runs[1]["font"] == "Calibri" and runs[1]["size"] == 24
    exported = native_export(source, body).slides[0].shapes[0].table.cell(1, 0)
    original = Presentation(io.BytesIO(source)).slides[0].shapes[0].table.cell(1, 0)
    assert exported.text_frame._txBody.xml == original.text_frame._txBody.xml


def test_extra_table_text_cannot_be_silently_dropped_when_rich_slots_match():
    import pytest

    source = inherited_table_source()
    body = baseline_slide_html(
        extract_slides(source, path="t.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    root.xpath(".//table/tr[2]/td")[0].text = "Unmapped content "
    with pytest.raises(ValueError, match="rich slots are ambiguous"):
        native_export(source, html.tostring(root, encoding="unicode"))


def test_table_paragraph_defaults_survive_missing_theme_and_presentation_defaults():
    from app.agent.render import measure_slide
    from langchain_canvas.deck._shapes import _frame
    from langchain_canvas.deck.structured import extract_structured
    from pptx.oxml.ns import qn

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(5), Inches(2))
    cell = shape.table.cell(0, 0)
    cell.text = "Paragraph inheritance"
    font = cell.text_frame.paragraphs[0].font
    font.name = "Arial"
    font.size = Pt(30)
    font.bold = False
    font.italic = True
    master = slide.slide_layout.slide_master.part
    relation = next(r for r in master.rels.values() if r.reltype.endswith("/theme"))
    master.drop_rel(relation.rId)
    defaults = presentation._element.find(qn("p:defaultTextStyle"))
    presentation._element.remove(defaults)
    output = io.BytesIO()
    presentation.save(output)
    source = output.getvalue()
    extracted = extract_slides(source, path="missing-theme.pptx")[0]
    body = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    run = measure_slide(body, ratio="16:9")["textBlocks"][0]["paragraphs"][0]["runs"][0]
    assert (
        run["font"] == "Arial"
        and run["size"] == 40
        and run["weight"] == "400"
        and run["italic"]
    )
    assert (
        native_export(source, body).slides[0].shapes[0].table.cell(0, 0).text
        == cell.text
    )
    assert (
        extract_structured(
            shape,
            _frame(shape, 0, presentation.slide_width, presentation.slide_height),
            {},
        )
        is not None
    )


def test_native_table_outer_containers_allow_only_harmless_metadata_and_whitespace():
    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    root = html.fromstring(body)
    table = root.xpath(".//table")[0]
    table.text = "\n  "
    group = html.Element("thead", {"data-node-id": "header-group", "style": "  "})
    group.text = "\n "
    row = table[0]
    row.set("title", "Source header")
    row.text = " "
    row[0].tail = "\n "
    group.append(row)
    table.insert(0, group)
    group.tail = "\n "
    assert (
        native_export(source, html.tostring(root, encoding="unicode"))
        .slides[0]
        .shapes[0]
        .has_table
    )
    for selector in (".//table", ".//tr", ".//thead"):
        for attribute in (
            "bgcolor",
            "align",
            "height",
            "hidden",
            "border",
            "cellspacing",
            "width",
        ):
            altered = html.fromstring(html.tostring(root, encoding="unicode"))
            altered.xpath(selector)[0].set(attribute, "100")
            with pytest.raises(ValueError, match="Native table containers"):
                native_export(source, html.tostring(altered, encoding="unicode"))


def test_native_table_rejects_unknown_containers_text_and_invalid_nesting():
    import pytest

    source = source_presentation()
    body = baseline_slide_html(
        extract_slides(source, path="s.pptx")[0], slide_id="s1", ratio="16:9"
    )
    for extra in (
        "<div>Other</div>",
        "<colgroup><col></colgroup>",
        "<tfoot><tr><td>Footer</td></tr></tfoot>",
        "<tbody><tbody></tbody></tbody>",
        "<td>Orphan</td>",
    ):
        root = html.fromstring(body)
        root.xpath(".//table")[0].append(html.fromstring(extra))
        with pytest.raises(ValueError, match="Native table structure"):
            native_export(source, html.tostring(root, encoding="unicode"))
    for selector in (".//table", ".//tr"):
        root = html.fromstring(body)
        root.xpath(selector)[0].text = "Unmapped content"
        with pytest.raises(ValueError, match="Native table content"):
            native_export(source, html.tostring(root, encoding="unicode"))
