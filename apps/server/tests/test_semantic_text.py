"""Editing units must survive rich text and paragraph export."""

import io

from app.agent.exports import EditableDeckPptxExporter
from langchain_canvas.deck import Deck, SlideTemplate, serialize_deck
from pptx import Presentation


def test_rich_multiline_paragraph_is_one_powerpoint_textbox():
    body = '<section style="width:1280px;height:720px"><p data-node-id="message" data-text-block="true" style="position:absolute;left:100px;top:50px;width:500px;height:120px;margin:0;font:24px Arial">A sentence with <strong>bold words</strong> and <span style="color:red">red words</span>.<br>Second line remains the same paragraph.</p></section>'
    deck = serialize_deck(
        Deck("Paragraph", "16:9", None, [SlideTemplate("s1", None, "", body)])
    )
    result = EditableDeckPptxExporter().export(deck, path="paragraph.slides.html")
    slide = Presentation(io.BytesIO(result.data)).slides[0]
    texts = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    assert len(texts) == 1
    assert "A sentence with bold words and red words." in texts[0].text
    assert "Second line remains the same paragraph." in texts[0].text
    runs = [r for p in texts[0].text_frame.paragraphs for r in p.runs]
    assert any(r.text == "bold words" and r.font.bold for r in runs)
    assert any(
        r.text == "red words" and str(r.font.color.rgb) == "FF0000" for r in runs
    )


def test_table_cells_are_separate_editing_units():
    body = "<table><tr><td>Left <b>cell</b></td><td>Right cell</td></tr></table>"
    deck = serialize_deck(
        Deck("Cells", "16:9", None, [SlideTemplate("s1", None, "", body)])
    )
    result = EditableDeckPptxExporter().export(deck, path="cells.slides.html")
    slide = Presentation(io.BytesIO(result.data)).slides[0]
    assert [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()] == [
        "Left cell",
        "Right cell",
    ]


def test_legacy_words_merge_without_merging_columns():
    from app.agent.semantic_text import consolidate_slide_html

    body = '<section style="width:1280px;height:720px">'
    for ident, x, text in [
        ("a", 40, "Hello "),
        ("b", 94, "world"),
        ("c", 700, "Other column"),
    ]:
        body += f'<div data-node-id="{ident}" style="position:absolute;left:{x}px;top:50px;font:20px Arial;white-space:pre">{text}</div>'
    html, report = consolidate_slide_html(body + "</section>", ratio="16:9")
    assert report["merged"] == 1
    from lxml import html as parser

    root = parser.fromstring(html)
    blocks = root.xpath('//*[@data-text-block="true"]')
    assert [b.text_content() for b in blocks] == ["Hello world", "Other column"]


def test_source_pptx_keeps_native_table_while_replacing_text():
    import base64

    from langchain_canvas.exporters import PPTX_MIME
    from pptx.util import Inches

    original = Presentation()
    slide = original.slides.add_slide(original.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text = "Original"
    slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(2)).table.cell(
        0, 0
    ).text = "Native table"
    source = io.BytesIO()
    original.save(source)
    body = '<section><p data-node-id="title" data-pptx-shape-id="e0" data-text-block="true" style="position:absolute;left:100px;top:100px;width:500px;height:100px;margin:0;font:24px Arial">Updated <b>content</b></p></section>'
    deck = serialize_deck(
        Deck(
            "Template",
            "16:9",
            f"data:{PPTX_MIME};base64,{base64.b64encode(source.getvalue()).decode()}",
            [SlideTemplate("s1", None, "", body)],
        )
    )
    result = EditableDeckPptxExporter().export(deck, path="template.slides.html")
    result_slide = Presentation(io.BytesIO(result.data)).slides[0]
    assert len(result_slide.shapes) == 2
    assert result_slide.shapes[0].text == "Updated content"
    assert result_slide.shapes[0].text_frame.paragraphs[0].runs[-1].font.bold
    assert result_slide.shapes[1].table.cell(0, 0).text == "Native table"


def test_pptx_import_keeps_mixed_styles_inside_one_semantic_block():
    from langchain_canvas.deck.baseline import baseline_slide_html
    from langchain_canvas.deck.extract import extract_slides
    from lxml import html as parser
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    frame = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(2)
    ).text_frame
    first = frame.paragraphs[0].add_run()
    first.text = "Original "
    first.font.size = Pt(20)
    second = frame.paragraphs[0].add_run()
    second.text = "emphasis"
    second.font.bold = True
    second.font.color.rgb = RGBColor(255, 0, 0)
    frame.paragraphs[0].add_line_break()
    frame.paragraphs[0].add_run().text = "Same paragraph"
    data = io.BytesIO()
    presentation.save(data)
    extracted = extract_slides(data.getvalue(), path="source.pptx")[0]
    html = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    root = parser.fromstring(html)
    blocks = root.xpath('.//*[@data-text-block="true"]')
    assert len(blocks) == 1
    assert "Original emphasisSame paragraph" == blocks[0].text_content()
    assert "font-weight:700" in html and "color:#FF0000" in html
    assert "<br>" in html


def test_neighboring_words_across_subpixel_baselines_keep_reading_order():
    from app.agent.semantic_text import consolidate_slide_html
    from lxml import html as parser

    body = '<section><div data-node-id="first" style="position:absolute;left:40px;top:409.6px;font:20px Arial;white-space:pre">Hello </div><div data-node-id="second" style="position:absolute;left:94px;top:409.3px;font:20px Arial;white-space:pre">world</div></section>'
    clean, report = consolidate_slide_html(body)
    assert report["merged"] == 1
    assert [
        n.text_content()
        for n in parser.fromstring(clean).xpath('.//*[@data-text-block="true"]')
    ] == ["Hello world"]


def test_empty_source_placeholder_clears_native_text_and_remains_editable():
    import base64

    from langchain_canvas.exporters import PPTX_MIME
    from pptx.util import Inches

    original = Presentation()
    slide = original.slides.add_slide(original.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text = "Remove me"
    source = io.BytesIO()
    original.save(source)
    body = '<section><p data-node-id="title" data-pptx-shape-id="e0" data-text-block="true" style="position:absolute;left:100px;top:100px;width:500px;height:100px;margin:0"></p></section>'
    deck = serialize_deck(
        Deck(
            "Empty",
            "16:9",
            f"data:{PPTX_MIME};base64,{base64.b64encode(source.getvalue()).decode()}",
            [SlideTemplate("s1", None, "", body)],
        )
    )
    result = EditableDeckPptxExporter().export(deck, path="empty.slides.html")
    native = Presentation(io.BytesIO(result.data)).slides[0]
    assert len(native.shapes) == 1
    assert native.shapes[0].has_text_frame and native.shapes[0].text == ""


def test_normalization_preserves_nested_containing_block_coordinates():
    from app.agent.render import measure_slide
    from app.agent.semantic_text import consolidate_slide_html

    body = '<section style="position:absolute;left:100px;top:100px;width:1000px;height:600px"><div data-node-id="first" style="position:absolute;left:40px;top:50px;font:20px Arial;white-space:pre">Hello </div><div data-node-id="second" style="position:absolute;left:94px;top:50px;font:20px Arial;white-space:pre">world</div></section>'
    clean, report = consolidate_slide_html(body)
    assert report["merged"] == 1
    after = measure_slide("<style>body{margin:0}</style>" + clean, ratio="16:9")
    block = next(b for b in after["textBlocks"] if b["id"] == "first")
    assert abs(block["x"] - 140) < 1 and abs(block["y"] - 150) < 1


def test_pptx_import_retains_fields_mixed_with_rich_runs():
    from langchain_canvas.deck.baseline import baseline_slide_html
    from langchain_canvas.deck.extract import extract_slides
    from lxml import html as parser
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    frame = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame
    field = OxmlElement("a:fld")
    field.set("id", "{32D9264F-8F48-4CF0-B17E-07247BAD2111}")
    field.set("type", "slidenum")
    text = OxmlElement("a:t")
    text.text = "5"
    field.append(text)
    frame.paragraphs[0]._p.append(field)
    frame.paragraphs[0].add_run().text = " - Progress"
    data = io.BytesIO()
    presentation.save(data)
    extracted = extract_slides(data.getvalue(), path="field.pptx")[0]
    html = baseline_slide_html(extracted, slide_id="s1", ratio="16:9")
    assert parser.fromstring(html).text_content().strip() == "5 - Progress"


def test_explicit_block_with_nested_paragraphs_exports_no_duplicate_boxes():
    body = '<div data-node-id="unit" data-text-block="true" style="width:600px;height:200px;padding:12px;font:20px Arial;letter-spacing:1px"><p data-node-id="first" style="text-align:center;margin:4px 0">First <i>paragraph</i></p><p data-node-id="second">Second paragraph</p></div>'
    deck = serialize_deck(
        Deck("Rich", "16:9", None, [SlideTemplate("s1", None, "", body)])
    )
    exported = EditableDeckPptxExporter().export(deck, path="rich.slides.html")
    shapes = list(Presentation(io.BytesIO(exported.data)).slides[0].shapes)
    assert len(shapes) == 1
    frame = shapes[0].text_frame
    assert len(frame.paragraphs) == 2
    assert [p.text for p in frame.paragraphs] == ["First paragraph", "Second paragraph"]
    assert frame.margin_left > 0
    assert frame.paragraphs[0].runs[-1].font.italic
    assert frame.paragraphs[0].runs[0]._r.get_or_add_rPr().get("spc")


def test_duplicate_ids_fail_and_generated_ids_do_not_collide():
    import pytest
    from app.agent.semantic_text import consolidate_slide_html
    from lxml import html as parser

    with pytest.raises(ValueError, match="Duplicate"):
        consolidate_slide_html(
            '<p data-node-id="same">A</p><p data-node-id="same">B</p>'
        )
    clean, _ = consolidate_slide_html(
        '<section data-node-id="text-unit-1"><p>Hello</p><style>p{color:red}</style><p>World<br>Again</p></section>'
    )
    ids = parser.fromstring(clean).xpath("//@data-node-id")
    assert len(ids) == len(set(ids)) == 3
    assert "text-unit-2" in ids
    again, report = consolidate_slide_html(clean)
    assert report["merged"] == 0
    assert again == clean


def _fragment(ident, x, y, text, extra=""):
    return f'<div data-node-id="{ident}" style="position:absolute;left:{x}px;top:{y}px;font:20px Arial;line-height:24px;white-space:pre;{extra}">{text}</div>'


def test_grid_dividers_keep_neighboring_cells_and_rows_separate():
    from app.agent.semantic_text import consolidate_slide_html
    from lxml import html as parser

    body = (
        "<section>" + _fragment("a", 40, 40, "Hello") + _fragment("b", 93, 40, "world")
    )
    body += '<div style="position:absolute;left:90px;top:30px;height:50px;border-left:1px solid black"></div>'
    body += _fragment("c", 300, 40, "Long first sentence") + _fragment(
        "d", 300, 68, "Long second sentence"
    )
    body += '<div style="position:absolute;left:290px;top:66px;width:250px;border-top:1px solid black"></div></section>'
    clean, report = consolidate_slide_html(body)
    assert report["merged"] == 0
    assert len(parser.fromstring(clean).xpath('.//*[@data-text-block="true"]')) == 4


def test_continuation_lines_merge_but_bullets_stay_separate():
    from app.agent.semantic_text import consolidate_slide_html
    from lxml import html as parser

    body = (
        "<section>"
        + _fragment("a", 40, 40, "A complete first line")
        + _fragment("b", 40, 66, "A continuing second line")
        + _fragment("c", 40, 92, "• A separate bullet")
        + "</section>"
    )
    clean, report = consolidate_slide_html(body)
    assert report["merged"] == 1
    root = parser.fromstring(clean)
    assert root.xpath('.//*[@data-node-id="a"]/br')
    assert root.xpath('.//*[@data-node-id="c"]')


def test_join_preserves_inline_emphasis_punctuation_and_column_width():
    from app.agent.render import measure_slide
    from app.agent.semantic_text import consolidate_slide_html
    from lxml import html as parser

    body = (
        "<section>"
        + _fragment("a", 40, 40, "Hello")
        + _fragment("b", 94, 40, '<b data-node-id="bold">world</b>', "color:red;")
        + _fragment("c", 150, 40, "!")
    )
    body += '<div style="position:absolute;left:180px;top:30px;height:60px;border-left:1px solid black"></div></section>'
    clean, report = consolidate_slide_html(body)
    assert report["merged"] == 2
    root = parser.fromstring(clean)
    block = root.xpath('.//*[@data-node-id="a"]')[0]
    assert block.text_content() == "Hello world!"
    assert block.xpath(".//b") and not block.xpath(".//b/@data-node-id")
    assert "color: rgb(255, 0, 0)" in clean
    measured = measure_slide("<style>body{margin:0}</style>" + clean, ratio="16:9")
    box = next(b for b in measured["textBlocks"] if b["id"] == "a")
    assert box["x"] + box["w"] <= 178
