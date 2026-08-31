"""Native objects and source templates must survive export without rasterization."""

import base64
import io

import pytest
from app.agent.exports import EditableDeckPptxExporter, _add_item, _fill
from app.agent.style_tokens import BackgroundToken, StyleTokens
from langchain_canvas.deck import (
    Deck,
    SlideTemplate,
    sanitize_slide_html,
    serialize_deck,
)
from langchain_canvas.exporters import PPTX_MIME
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches


def _export(body, source=None):
    content = serialize_deck(
        Deck("Test", "16:9", source, [SlideTemplate("s1", None, "", body)])
    )
    return Presentation(
        io.BytesIO(
            EditableDeckPptxExporter().export(content, path="test.slides.html").data
        )
    )


def _png(size, color):
    data = io.BytesIO()
    Image.new("RGB", size, color).save(data, format="PNG")
    return data.getvalue()


def _url(data, mime="image/png"):
    return f"data:{mime};base64,{base64.b64encode(data).decode()}"


def test_native_gradients_rounded_shapes_ovals_and_border_alpha():
    deck = _export(
        '<section><div style="position:absolute;left:20px;top:20px;width:240px;height:100px;border-radius:12px;background:linear-gradient(45deg,rgba(255,0,0,.5) 20%,blue 90%);border-bottom:2px solid rgba(0,0,0,.4)"></div><div style="position:absolute;left:300px;top:20px;width:100px;height:100px;border-radius:50%;background:radial-gradient(red,rgba(0,0,255,.7))"></div></section>'
    )
    shapes = list(deck.slides[0].shapes)
    rounded = next(
        s
        for s in shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    )
    oval = next(
        s
        for s in shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.auto_shape_type == MSO_SHAPE.OVAL
    )
    assert rounded.fill.gradient_angle == 45
    assert rounded._element.xpath(".//a:gs/@pos") == ["20000", "90000"]
    assert rounded._element.xpath(".//a:gs/a:srgbClr/a:alpha/@val") == [
        "50000",
        "100000",
    ]
    assert oval._element.xpath(".//a:path/@path") == ["circle"]
    line = next(s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.LINE)
    assert line._element.xpath(".//a:alpha/@val") == ["40000"]
    assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes)


@pytest.mark.parametrize(
    ("size", "fit", "cropped"),
    [
        ((200, 100), "contain", None),
        ((200, 100), "cover", "horizontal"),
        ((100, 200), "cover", "vertical"),
        ((200, 100), "fill", None),
    ],
)
def test_original_images_keep_bytes_and_expected_fit(size, fit, cropped):
    original = _png(size, "green")
    deck = _export(
        f'<img src="{_url(original)}" style="position:absolute;left:100px;top:100px;width:100px;height:100px;object-fit:{fit}">'
    )
    picture = deck.slides[0].shapes[0]
    assert picture.image.blob == original
    if fit == "contain":
        assert picture.width == 2 * picture.height
    else:
        assert picture.width == picture.height
    assert (picture.crop_left > 0) == (cropped == "horizontal")
    assert (picture.crop_top > 0) == (cropped == "vertical")


def test_source_template_patches_shape_image_and_keeps_native_table():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    slide.shapes.add_picture(io.BytesIO(_png((20, 20), "red")), Inches(4), Inches(1))
    slide.shapes.add_table(1, 1, Inches(1), Inches(4), Inches(2), Inches(1)).table.cell(
        0, 0
    ).text = "Keep table"
    buffer = io.BytesIO()
    source.save(buffer)
    replacement = _png((40, 20), "blue")
    deck = _export(
        '<div data-node-id="box" data-pptx-shape-id="e0" style="position:absolute;left:40px;top:40px;width:200px;height:100px;background:green"></div>'
        + f'<img data-node-id="picture" data-pptx-shape-id="e1" src="{_url(replacement)}" style="position:absolute;left:300px;top:40px;width:200px;height:100px">',
        _url(buffer.getvalue(), PPTX_MIME),
    )
    result = deck.slides[0]
    assert len(result.shapes) == 3
    assert str(result.shapes[0].fill.fore_color.rgb) == "008000"
    assert result.shapes[1].image.blob == replacement
    assert result.shapes[2].table.cell(0, 0).text == "Keep table"


def test_source_decks_fail_closed_for_unresolved_invalid_count_and_provenance():
    for source, message in [
        ("source.pptx", "must be loaded"),
        (_url(b"bad", PPTX_MIME), "Invalid source"),
    ]:
        with pytest.raises(ValueError, match=message):
            _export("<p>Body</p>", source)
    original = Presentation()
    original.slides.add_slide(original.slide_layouts[6])
    buf = io.BytesIO()
    original.save(buf)
    source = _url(buf.getvalue(), PPTX_MIME)
    with pytest.raises(ValueError, match="Unknown native"):
        _export('<p data-node-id="bad" data-pptx-shape-id="e99">Body</p>', source)
    original.slides.add_slide(original.slide_layouts[6])
    buf = io.BytesIO()
    original.save(buf)
    with pytest.raises(ValueError, match="order and count"):
        _export("<p>Body</p>", _url(buf.getvalue(), PPTX_MIME))
    content = serialize_deck(Deck("Empty", "16:9", None, []))
    with pytest.raises(ValueError, match="empty deck"):
        EditableDeckPptxExporter().export(content, path="empty.slides.html")


def test_ten_template_instances_export_editable_without_old_source_objects():
    from app.agent.deck_template_writer import ArchetypeFrame, instantiate_archetype

    frame = ArchetypeFrame(
        archetype_id="body",
        style_css=".slide { color: #111827; }",
        body_html=(
            '<section class="slide">'
            '<p data-node-id="node-headline">OLD_HEADLINE</p>'
            '<p data-node-id="node-body">OLD_BODY</p>'
            '<div class="rect" data-pptx-shape-id="e0" '
            'style="position:absolute;left:20px;top:200px;width:100px;height:60px;'
            'background:#334155"></div>'
            "</section>"
        ),
        slot_node_ids={"headline": "node-headline", "body": "node-body"},
    )
    source_business_text = frozenset({"OLD_HEADLINE", "OLD_BODY"})
    slides = [
        instantiate_archetype(
            frame,
            index,
            {"headline": f"Headline {index}", "body": f"Body {index}"},
            source_business_text=source_business_text,
        )
        for index in range(1, 11)
    ]
    content = serialize_deck(Deck("Generated", "16:9", None, slides))

    result = EditableDeckPptxExporter().export(content, path="deck.slides.html")

    reopened = Presentation(io.BytesIO(result.data))
    assert len(reopened.slides) == 10
    all_text = " ".join(
        run.text
        for slide in reopened.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert "OLD_HEADLINE" not in all_text
    assert "OLD_BODY" not in all_text
    assert "Headline 1" in all_text
    assert "Body 10" in all_text


def _text_block(key, x, y, w, h, text):
    return {
        "key": key,
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "paragraphs": [
            {
                "align": "left",
                "lineHeight": 20,
                "runs": [
                    {
                        "text": text,
                        "font": "Arial",
                        "size": 16,
                        "weight": "400",
                        "italic": False,
                        "underline": False,
                        "color": "rgb(0, 0, 0)",
                    }
                ],
            }
        ],
    }


def _stub_layout(monkeypatch, *, unsupported, items, blocks):
    """Drive the exporter with a fixed layout and a fixed raster replacement."""
    calls = {}

    def fake_measure_slide(document, *, ratio):
        calls["measure_ratio"] = ratio
        return {
            "width": 1280,
            "height": 720,
            "items": items,
            "textBlocks": blocks,
            "unsupported": unsupported,
        }

    def fake_fallback(document, entries, ratio):
        calls["fallback_entries"] = entries
        calls["fallback_ratio"] = ratio
        return [
            {
                "kind": "image",
                "src": _url(_png((20, 10), "red")),
                "fit": "fill",
                "x": entry["x"],
                "y": entry["y"],
                "w": entry["w"],
                "h": entry["h"],
            }
            for entry in entries
        ]

    monkeypatch.setattr("app.agent.exports.measure_slide", fake_measure_slide)
    monkeypatch.setattr("app.agent.exports._raster_fallback_items", fake_fallback)
    return calls


def test_unsupported_filter_element_exports_as_positioned_picture():
    presentation = _export(
        '<section data-text-block="b1">'
        '<div style="position:absolute;left:100px;top:80px;width:200px;height:120px;'
        'filter:blur(2px);background:#c00"></div></section>'
    )
    pictures = [
        s
        for s in presentation.slides[0].shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 1
    unit = presentation.slide_width / 1280
    assert abs(pictures[0].left - round(100 * unit)) <= unit
    assert abs(pictures[0].top - round(80 * unit)) <= unit
    # A real crop of the rendered slide, not a placeholder.
    assert pictures[0].image.size == (200, 120)


def test_unsupported_element_does_not_double_draw_covered_items(monkeypatch):
    calls = _stub_layout(
        monkeypatch,
        unsupported=[{"reason": "filter/clip-path", "x": 100, "y": 80, "w": 200, "h": 120}],
        items=[],
        blocks=[
            _text_block("covered", 110, 90, 100, 40, "COVERED"),
            _text_block("outside", 600, 400, 100, 40, "OUTSIDE"),
        ],
    )
    presentation = _export('<section data-text-block="b1">x</section>')
    texts = [
        s.text_frame.text
        for s in presentation.slides[0].shapes
        if s.has_text_frame
    ]
    assert calls["fallback_ratio"] == "16:9"
    assert calls["fallback_entries"][0]["reason"] == "filter/clip-path"
    assert "COVERED" not in texts
    assert "OUTSIDE" in texts


def test_partially_overlapping_item_is_still_drawn(monkeypatch):
    _stub_layout(
        monkeypatch,
        unsupported=[{"reason": "filter/clip-path", "x": 100, "y": 80, "w": 200, "h": 120}],
        items=[],
        blocks=[_text_block("straddling", 250, 90, 200, 40, "STRADDLING")],
    )
    presentation = _export('<section data-text-block="b1">x</section>')
    texts = [
        s.text_frame.text
        for s in presentation.slides[0].shapes
        if s.has_text_frame
    ]
    assert "STRADDLING" in texts


def test_source_backed_deck_still_rejects_unsupported_css():
    source = Presentation()
    source.slides.add_slide(source.slide_layouts[6])
    buffer = io.BytesIO()
    source.save(buffer)
    content = serialize_deck(
        Deck(
            "Test",
            "16:9",
            _url(buffer.getvalue(), PPTX_MIME),
            [
                SlideTemplate(
                    "s1",
                    None,
                    "",
                    '<section><div style="position:absolute;left:10px;top:10px;'
                    'width:80px;height:80px;clip-path:circle(40px);'
                    'background:#0c0"></div></section>',
                )
            ],
        )
    )
    with pytest.raises(ValueError, match="clip-path"):
        EditableDeckPptxExporter().export(content, path="test.slides.html")


def test_unsupported_gradient_and_fragment_contract_fail_clearly():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
    with pytest.raises(ValueError, match="Unsupported CSS gradient"):
        _fill(shape, {"gradient": "linear-gradient(red,blue)"})
    with pytest.raises(ValueError, match="semantic editing block"):
        _add_item(slide, {"kind": "text", "x": 0, "y": 0, "w": 100, "h": 20}, 1)


def _image_item(x, y, w, h, size):
    return {
        "kind": "image",
        "src": _url(_png(size, "blue")),
        "fit": "fill",
        "x": x,
        "y": y,
        "w": w,
        "h": h,
    }


def test_full_bleed_background_image_is_added_before_other_shapes(monkeypatch):
    _stub_layout(
        monkeypatch,
        unsupported=[],
        items=[
            _image_item(400, 300, 100, 50, (100, 50)),
            _image_item(0, 0, 1280, 720, (160, 90)),
        ],
        blocks=[],
    )
    presentation = _export('<section data-text-block="b1">x</section>')
    shapes = list(presentation.slides[0].shapes)
    unit = presentation.slide_width / 1280
    assert shapes[0].width == round(1280 * unit)
    assert shapes[0].top == 0


def test_export_normalizes_measured_font_to_matched_open_source_name(monkeypatch):
    block = _text_block("b1", 10, 10, 400, 40, "분기 실적")
    block["paragraphs"][0]["runs"][0]["font"] = "Malgun Gothic, sans-serif"
    _stub_layout(monkeypatch, unsupported=[], items=[], blocks=[block])

    presentation = _export('<section data-text-block="b1">x</section>')

    shape = next(s for s in presentation.slides[0].shapes if s.has_text_frame)
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Noto Sans KR"
    assert run._r.xpath(".//a:ea/@typeface") == ["Noto Sans KR"]


def _deck_with_tokens_attribute(attribute_value):
    body = sanitize_slide_html(
        f'<section class="slide" data-style-tokens=\'{attribute_value}\'>'
        '<div style="position:absolute;left:10px;top:10px;width:80px;height:40px;'
        'background:#0c0"></div></section>'
    ).html
    return _export(body)


def test_theme_sets_master_background_from_tokens():
    tokens = StyleTokens(background=BackgroundToken(kind="solid", value="#112233"))
    presentation = _deck_with_tokens_attribute(tokens.model_dump_json())
    fill = presentation.slide_master.background.fill
    assert fill.type == MSO_FILL.SOLID
    assert fill.fore_color.rgb == RGBColor.from_string("112233")


@pytest.mark.parametrize(
    "attribute_value",
    [
        "{not json",
        '{"colors": 3}',
        # A named CSS color validates as a plain string but is not a color the
        # PPTX writer can parse: the export must still succeed, unthemed.
        '{"background": {"kind": "solid", "value": "red"}}',
    ],
)
def test_theme_skipped_when_style_tokens_attribute_is_invalid_json(attribute_value):
    presentation = _deck_with_tokens_attribute(attribute_value)
    assert presentation.slide_master.background.fill.type != MSO_FILL.SOLID
