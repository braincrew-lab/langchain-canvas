"""Tests for verify.py: `check_slide_layout` / `screenshot_slide`, including
the `slide_id` deck-slide path and the render-failure error boundary.

Playwright is faked throughout via monkeypatching `app.agent.verify.render_slide`
(the function verify.py imports from render.py) — these tests exercise the
tool pipeline (deck-slide extraction, error handling), not the headless-
Chromium adapter itself, which is render.py's own concern.
"""

from __future__ import annotations

import io
import shutil
from collections.abc import Callable
from dataclasses import dataclass, field
from typing import Any

import pytest
from app.agent import render as render_module
from app.agent.render import shutdown_browser
from app.agent.store import DATA_DIR, STORE
from app.agent.tools import _DECK_TOOLS_BY_NAME
from app.agent.verify import check_slide_layout, screenshot_slide
from langchain_canvas.deck import parse_deck
from pptx import Presentation
from pptx.util import Inches

_TEST_THREAD_IDS = ("t-slide", "t-error", "t-page")


@pytest.fixture(autouse=True)
def _cleanup_canvas_data() -> Any:
    yield
    for thread_id in _TEST_THREAD_IDS:
        shutil.rmtree(DATA_DIR / thread_id, ignore_errors=True)


@dataclass
class _Runtime:
    """The slice of ToolRuntime the server tools read (mirrors test_convert_slide.py)."""

    context: Any = None
    config: dict[str, Any] = field(default_factory=dict)
    events: list[dict] = field(default_factory=list)

    @property
    def stream_writer(self) -> Callable[[dict], None]:
        return self.events.append


def _runtime(thread_id: str) -> _Runtime:
    return _Runtime(config={"configurable": {"thread_id": thread_id}})


def _deck_bytes(text: str = "Hello deck") -> bytes:
    """One-slide presentation bytes with a single text box saying `text`."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _seed_deck(thread_id: str, *, text: str = "Hello deck") -> tuple[str, str, str]:
    """Upload a `.pptx` and copy it into an editable deck via the SDK's own
    `open_deck_for_editing` — the exact state check_slide_layout expects.

    Returns `(path, slide_id, ratio)` for the deck's first (and only) slide.
    """
    STORE.write_bytes(thread_id, "sources/deck.pptx", _deck_bytes(text), "Upload", actor="human")
    result = _DECK_TOOLS_BY_NAME["open_deck_for_editing"].func(
        source="sources/deck.pptx", runtime=_runtime(thread_id)
    )
    assert "Copied" in result, result
    path = next(
        info.path for info in STORE.list_files(thread_id) if info.path.endswith(".slides.html")
    )
    deck = parse_deck(STORE.read(thread_id, path).content)
    return path, deck.slides[0].slide_id, deck.ratio


_FAKE_METRICS = {
    "overflowX": 0,
    "overflowY": 0,
    "offCanvas": [],
    "maxBottom": 700,
    "textLength": 100,
    "brokenImages": 0,
}


def test_check_slide_layout_slide_id_renders_only_that_slide_body(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path, slide_id, ratio = _seed_deck("t-slide", text="Hello deck")
    calls: list[tuple[str, str]] = []

    def fake_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        calls.append((html, ratio))
        return dict(_FAKE_METRICS), b""

    monkeypatch.setattr("app.agent.verify.render_slide", fake_render)

    result = check_slide_layout.func(file=path, slide_id=slide_id, runtime=_runtime("t-slide"))

    assert "LAYOUT CHECK" in result
    assert "Error" not in result
    assert len(calls) == 1
    rendered_html, rendered_ratio = calls[0]
    assert rendered_ratio == ratio
    # Only the one slide's body was rendered — not the whole multi-<template> deck,
    # which would render empty (browsers never render <template> content).
    assert "<template" not in rendered_html
    assert "Hello deck" in rendered_html


def test_check_slide_layout_render_failure_returns_error_string(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    STORE.write(
        "t-error", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )

    def raising_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        raise RuntimeError("Chromium crashed")

    monkeypatch.setattr("app.agent.verify.render_slide", raising_render)

    result = check_slide_layout.func(file="index.html", runtime=_runtime("t-error"))

    assert result.startswith("Error:")


def test_screenshot_slide_render_failure_returns_error_content(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    STORE.write(
        "t-page", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )

    def raising_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        raise RuntimeError("Chromium crashed")

    monkeypatch.setattr("app.agent.verify.render_slide", raising_render)

    result = screenshot_slide.func(file="index.html", runtime=_runtime("t-page"))

    assert len(result) == 1
    assert result[0]["type"] == "text"
    assert result[0]["text"].startswith("Error:")


def test_shutdown_browser_with_nothing_launched_does_not_raise() -> None:
    shutdown_browser()


def test_check_slide_layout_missing_file_returns_hint() -> None:
    result = check_slide_layout.func(file="missing.html", runtime=_runtime("t-error"))

    assert result == "No file missing.html exists yet. Use plan_deck / write_slide first."


def test_check_slide_layout_slide_id_on_non_deck_file_returns_error() -> None:
    STORE.write(
        "t-error", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )

    result = check_slide_layout.func(
        file="index.html", slide_id="slide-001", runtime=_runtime("t-error")
    )

    assert result.startswith("Error:")
    assert "not a valid slide deck" in result


def test_check_slide_layout_reports_every_problem_kind(monkeypatch: pytest.MonkeyPatch) -> None:
    STORE.write(
        "t-page", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )
    metrics = {
        "overflowX": 5,
        "overflowY": 3,
        "offCanvas": ["<div> 10x10 at (0, 0)"],
        "maxBottom": 100,
        "textLength": 5,
        "brokenImages": 2,
    }

    def fake_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        return dict(metrics), b""

    monkeypatch.setattr("app.agent.verify.render_slide", fake_render)

    result = check_slide_layout.func(file="index.html", runtime=_runtime("t-page"))

    assert "2 error(s), 2 warning(s)" in result
    assert "content overflows the 1280x720 slide" in result
    assert "element extends outside the slide" in result
    assert "image(s) failed to load" in result
    assert "almost no text content rendered" in result


def test_check_slide_layout_reports_empty_bottom_warning(monkeypatch: pytest.MonkeyPatch) -> None:
    STORE.write(
        "t-page", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )
    metrics = dict(_FAKE_METRICS, maxBottom=100)  # textLength stays >= 20

    def fake_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        return dict(metrics), b""

    monkeypatch.setattr("app.agent.verify.render_slide", fake_render)

    result = check_slide_layout.func(file="index.html", runtime=_runtime("t-page"))

    assert "the bottom of the slide looks empty" in result


def test_screenshot_slide_missing_file_returns_hint() -> None:
    result = screenshot_slide.func(file="missing.html", runtime=_runtime("t-error"))

    assert result == [{"type": "text", "text": "No file missing.html exists yet."}]


def test_screenshot_slide_slide_id_on_non_deck_file_returns_error() -> None:
    STORE.write(
        "t-error", "index.html", "<html><body>hi</body></html>", "Create page", actor="agent"
    )

    result = screenshot_slide.func(
        file="index.html", slide_id="slide-001", runtime=_runtime("t-error")
    )

    assert len(result) == 1
    assert result[0]["type"] == "text"
    assert "not a valid slide deck" in result[0]["text"]


def test_screenshot_slide_success_saves_png_and_returns_image(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path, slide_id, _ratio = _seed_deck("t-slide", text="Hello deck")

    def fake_render(html: str, *, ratio: str) -> tuple[dict[str, Any], bytes]:
        return dict(_FAKE_METRICS), b"fake-png-bytes"

    monkeypatch.setattr("app.agent.verify.render_slide", fake_render)

    result = screenshot_slide.func(file=path, slide_id=slide_id, runtime=_runtime("t-slide"))

    assert len(result) == 2
    assert result[0]["type"] == "text"
    assert f"Screenshot of {path}#{slide_id}" in result[0]["text"]
    assert result[1]["type"] == "image"
    assert result[1]["source"]["media_type"] == "image/png"

    saved = DATA_DIR / "t-slide" / "screenshots" / f"{path}__{slide_id}.png"
    assert saved.read_bytes() == b"fake-png-bytes"


# --- render.py: exercised directly against a fake Playwright, so these run
# without a real Chromium install. `_browser`/`_playwright` are module
# globals cached across calls (the whole point of the shared-browser design),
# so every test resets them before and tears down with `shutdown_browser()`.


class _FakePage:
    def __init__(self, viewport: dict[str, int], *, fail: bool = False) -> None:
        self.viewport = viewport
        self._fail = fail
        self.closed = False

    def set_content(self, html: str, wait_until: str = "load") -> None:
        if self._fail:
            raise RuntimeError("render failed")
        self.html = html

    def evaluate(self, script: str) -> dict[str, Any]:
        return dict(_FAKE_METRICS)

    def screenshot(self, type: str = "png") -> bytes:
        return b"fake-png"

    def close(self) -> None:
        self.closed = True


class _FakeBrowser:
    def __init__(self, *, fail: bool = False) -> None:
        self._fail = fail
        self.closed = False

    def new_page(self, viewport: dict[str, int]) -> _FakePage:
        return _FakePage(viewport, fail=self._fail)

    def close(self) -> None:
        self.closed = True


class _FakeChromium:
    def __init__(self, *, fail: bool = False) -> None:
        self._fail = fail
        self.launch_count = 0
        self.launched: list[_FakeBrowser] = []

    def launch(self) -> _FakeBrowser:
        self.launch_count += 1
        browser = _FakeBrowser(fail=self._fail)
        self.launched.append(browser)
        return browser


class _FakePlaywrightContext:
    def __init__(self, *, fail: bool = False) -> None:
        self.chromium = _FakeChromium(fail=fail)
        self.stopped = False

    def stop(self) -> None:
        self.stopped = True


class _FakeSyncPlaywright:
    def __init__(self, *, fail: bool = False) -> None:
        self.context = _FakePlaywrightContext(fail=fail)

    def start(self) -> _FakePlaywrightContext:
        return self.context


def _reset_render_module_state() -> None:
    render_module._browser = None
    render_module._playwright = None


def test_viewport_for_ratio_known_and_unknown_ratios() -> None:
    assert render_module.viewport_for_ratio("16:9") == (1280, 720)
    assert render_module.viewport_for_ratio("4:3") == (1280, 960)
    assert render_module.viewport_for_ratio("unknown") == (1280, 720)


def test_render_slide_caches_browser_across_calls(monkeypatch: pytest.MonkeyPatch) -> None:
    fake = _FakeSyncPlaywright()
    monkeypatch.setattr("playwright.sync_api.sync_playwright", lambda: fake)
    _reset_render_module_state()
    try:
        metrics, png = render_module.render_slide("<html></html>", ratio="16:9")
        assert png == b"fake-png"
        assert metrics == _FAKE_METRICS
        render_module.render_slide("<html></html>", ratio="4:3")
        assert fake.context.chromium.launch_count == 1  # second call reused the cached browser
    finally:
        render_module.shutdown_browser()
    assert fake.context.stopped is True


def test_render_slide_drops_cached_browser_on_failure_then_relaunches(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    failing = _FakeSyncPlaywright(fail=True)
    working = _FakeSyncPlaywright(fail=False)
    factories = iter([failing, working])
    monkeypatch.setattr("playwright.sync_api.sync_playwright", lambda: next(factories))
    _reset_render_module_state()
    try:
        with pytest.raises(RuntimeError, match="render failed"):
            render_module.render_slide("<html></html>", ratio="16:9")
        assert render_module._browser is None  # dropped so the next call relaunches

        _metrics, png = render_module.render_slide("<html></html>", ratio="16:9")
        assert png == b"fake-png"
        assert working.context.chromium.launch_count == 1
    finally:
        render_module.shutdown_browser()
