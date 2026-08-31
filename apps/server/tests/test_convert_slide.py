"""Tests for `convert_slide`: the reference pipeline that turns one imported
deck slide's raw extracted layout into a polished slide.

The writer model is faked throughout — these tests never call a real LLM.
Each test builds its own store (module-level `STORE` is shared process-wide,
so tests use distinct thread ids to stay isolated) and seeds it with a
`.pptx` source plus a `deck.slides.html` copy via the SDK's own
`open_deck_for_editing`, exactly the state `open_deck_for_editing` leaves
behind in production.
"""

from __future__ import annotations

import io
import shutil
from collections.abc import Callable
from dataclasses import dataclass, field
from types import SimpleNamespace
from typing import Any

import pytest
from app.agent.store import DATA_DIR, STORE
from app.agent.tools import _DECK_TOOLS_BY_NAME, convert_slide
from langchain_canvas.deck import parse_deck, read_slide
from pptx import Presentation
from pptx.util import Inches

# Thread ids this module writes under the real (gitignored) canvas-data/ dir
# via the module-global filesystem-backed STORE — cleaned up after each test
# so the module leaves no artifacts behind across runs.
_TEST_THREAD_IDS = ("t-success", "t-degrade", "t-retry")


@pytest.fixture(autouse=True)
def _cleanup_canvas_data() -> Any:
    yield
    for thread_id in _TEST_THREAD_IDS:
        shutil.rmtree(DATA_DIR / thread_id, ignore_errors=True)


@dataclass
class _Runtime:
    """The slice of ToolRuntime the server tools (and the SDK's deck tools,
    via `_canvas_id`) read."""

    context: Any = None
    config: dict[str, Any] = field(default_factory=dict)
    events: list[dict] = field(default_factory=list)

    @property
    def stream_writer(self) -> Callable[[dict], None]:
        return self.events.append


def _runtime(thread_id: str) -> _Runtime:
    return _Runtime(config={"configurable": {"thread_id": thread_id}})


def _deck_bytes(text: str = "Hello deck") -> bytes:
    """One-slide presentation bytes with a single text box saying `text`."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _seed_deck(thread_id: str, *, text: str = "Hello deck") -> tuple[str, str]:
    """Upload a `.pptx` and copy it into an editable deck via the SDK's own
    `open_deck_for_editing` — the exact state convert_slide expects.

    Returns `(path, slide_id)` for the deck's first (and only) slide.
    """
    STORE.write_bytes(thread_id, "sources/deck.pptx", _deck_bytes(text), "Upload", actor="human")
    result = _DECK_TOOLS_BY_NAME["open_deck_for_editing"].func(
        source="sources/deck.pptx", runtime=_runtime(thread_id)
    )
    assert "Copied" in result, result
    path = STORE.list_files(thread_id)
    deck_path = next(info.path for info in path if info.path.endswith(".slides.html"))
    deck = parse_deck(STORE.read(thread_id, deck_path).content)
    return deck_path, deck.slides[0].slide_id


def _fake_model(respond: Callable[[str], str]) -> Any:
    """A stand-in for `init_chat_model(...)`'s return value: `.invoke(prompt)`
    returns an object with a `.content` string, exactly what `_text_of` reads.
    """
    return SimpleNamespace(invoke=lambda prompt: SimpleNamespace(content=respond(prompt)))


def _patch_model(monkeypatch: pytest.MonkeyPatch, respond: Callable[[str], str]) -> None:
    monkeypatch.setattr("app.agent.tools.init_chat_model", lambda *_a, **_k: _fake_model(respond))


_FAKE_METRICS = {
    "overflowX": 0,
    "overflowY": 0,
    "offCanvas": [],
    "maxBottom": 700,
    "textLength": 100,
    "brokenImages": 0,
}


def _patch_render(monkeypatch: pytest.MonkeyPatch) -> None:
    """Fake the Playwright render boundary — these tests exercise
    convert_slide's pipeline logic, not the headless-Chromium adapter
    itself (that is `render.py`'s own concern)."""
    monkeypatch.setattr(
        "app.agent.tools.render_slide", lambda html, *, ratio: (dict(_FAKE_METRICS), b"")
    )


def test_convert_slide_success_applies_correction_and_verifies(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path, slide_id = _seed_deck("t-success", text="Hello deck")

    def respond(prompt: str) -> str:
        # Extract the baseline markup the prompt embeds and add a class
        # marker — `class` survives sanitize_slide_html's allowlist, and the
        # text inside stays byte-identical, so ensure_text_equality passes,
        # but the saved slide differs from the baseline (verifying a
        # correction actually landed).
        baseline = prompt.split("Baseline markup:\n", 1)[1].split("\n\nReturn ONLY", 1)[0]
        return baseline.replace('<section class="slide">', '<section class="slide corrected">')

    _patch_model(monkeypatch, respond)
    _patch_render(monkeypatch)

    result = convert_slide.func(path=path, slide_id=slide_id, runtime=_runtime("t-success"))

    assert "Error" not in result
    assert "LAYOUT CHECK" in result  # _layout_report ran once

    saved = read_slide(STORE.read("t-success", path).content, slide_id)
    assert "corrected" in saved.body_html
    assert "Hello deck" in saved.body_html


def test_convert_slide_degrades_on_text_mutation(monkeypatch: pytest.MonkeyPatch) -> None:
    path, slide_id = _seed_deck("t-degrade", text="Hello deck")
    before = read_slide(STORE.read("t-degrade", path).content, slide_id)

    # The model drops the extracted text entirely — a text-mutating response.
    _patch_model(monkeypatch, lambda _prompt: '<section class="slide"><div>Bye</div></section>')

    runtime = _runtime("t-degrade")
    result = convert_slide.func(path=path, slide_id=slide_id, runtime=runtime)

    assert "degraded" in result.lower()
    assert "retry" in result.lower()

    after = read_slide(STORE.read("t-degrade", path).content, slide_id)
    assert after.body_html == before.body_html  # baseline kept, nothing written

    degraded_events = [e for e in runtime.events if e.get("type") == "canvas.slide_status"]
    assert len(degraded_events) == 1
    assert degraded_events[0]["stage"] == "degraded"
    assert degraded_events[0]["slideId"] == slide_id


def test_convert_slide_retry_is_idempotent(monkeypatch: pytest.MonkeyPatch) -> None:
    path, slide_id = _seed_deck("t-retry", text="Hello deck")

    def respond(prompt: str) -> str:
        baseline = prompt.split("Baseline markup:\n", 1)[1].split("\n\nReturn ONLY", 1)[0]
        return baseline.replace('<section class="slide">', '<section class="slide corrected">')

    _patch_model(monkeypatch, respond)
    _patch_render(monkeypatch)

    first_result = convert_slide.func(path=path, slide_id=slide_id, runtime=_runtime("t-retry"))
    first_body = read_slide(STORE.read("t-retry", path).content, slide_id).body_html

    second_result = convert_slide.func(path=path, slide_id=slide_id, runtime=_runtime("t-retry"))
    second_body = read_slide(STORE.read("t-retry", path).content, slide_id).body_html

    assert "Error" not in first_result
    assert "Error" not in second_result
    assert first_body == second_body  # re-corrects from the same immutable source, every time


def test_convert_slide_prompt_excludes_deck_style(monkeypatch: pytest.MonkeyPatch) -> None:
    path, slide_id = _seed_deck("t-degrade", text="Hello deck")
    captured_prompts: list[str] = []

    def respond(prompt: str) -> str:
        captured_prompts.append(prompt)
        baseline = prompt.split("Baseline markup:\n", 1)[1].split("\n\nReturn ONLY", 1)[0]
        return baseline

    _patch_model(monkeypatch, respond)
    _patch_render(monkeypatch)

    result = convert_slide.func(path=path, slide_id=slide_id, runtime=_runtime("t-degrade"))

    assert "Error" not in result
    assert len(captured_prompts) == 1
    prompt = captured_prompts[0]
    assert "Generous margins" not in prompt
    assert "display serif" not in prompt
    assert "Preserve the source design exactly" in prompt
