"""Tools never raise: every failure path returns an ``"Error: ..."`` string
instead of letting an exception escape and abort the LangGraph ``ToolNode``
run (see `.omb/plans/2026-08-30-fix-chat-agent-stream.md` RC2 / Change 3).

Covers the three failure sources shared by all six writer tools:
  1. missing `thread_id` in the run config (`thread_id()` raises `ValueError`)
  2. `STORE.write` raising
  3. the writer model's `.invoke`/`.stream` call raising
"""

from __future__ import annotations

import io
import shutil
from dataclasses import dataclass, field
from types import SimpleNamespace
from typing import Any

import pytest
from app.agent.store import DATA_DIR, STORE
from app.agent.tools import (
    _DECK_TOOLS_BY_NAME,
    build_page,
    build_table,
    convert_slide,
    plan_deck,
    write_report,
    write_slide,
)
from langchain_canvas.deck import parse_deck
from pptx import Presentation
from pptx.util import Inches

# Thread ids this module writes under the real (gitignored) canvas-data/ dir
# via the module-global filesystem-backed STORE — cleaned up after each test,
# same convention as test_convert_slide.py.
_RENDER_THREAD_IDS = ("t-write-slide-render-error", "t-convert-slide-render-error")


@pytest.fixture(autouse=True)
def _cleanup_canvas_data() -> Any:
    yield
    for thread_id in _RENDER_THREAD_IDS:
        shutil.rmtree(DATA_DIR / thread_id, ignore_errors=True)


@dataclass
class _Runtime:
    """The slice of ToolRuntime the server tools read — same shape as
    `test_convert_slide.py`'s fixture."""

    context: Any = None
    config: dict[str, Any] = field(default_factory=dict)
    events: list[dict] = field(default_factory=list)

    @property
    def stream_writer(self):
        return self.events.append


def _runtime(thread_id: str | None) -> _Runtime:
    if thread_id is None:
        return _Runtime(config={})
    return _Runtime(config={"configurable": {"thread_id": thread_id}})


def _deck_bytes(text: str = "Hello deck") -> bytes:
    """One-slide presentation bytes with a single text box saying `text`."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _seed_imported_deck(thread_id: str) -> tuple[str, str]:
    """Upload a `.pptx` and copy it into an editable deck via the SDK's own
    `open_deck_for_editing` — the exact state convert_slide expects."""
    STORE.write_bytes(thread_id, "sources/deck.pptx", _deck_bytes(), "Upload", actor="human")
    result = _DECK_TOOLS_BY_NAME["open_deck_for_editing"].func(
        source="sources/deck.pptx", runtime=_runtime(thread_id)
    )
    assert "Copied" in result, result
    path = next(
        info.path for info in STORE.list_files(thread_id) if info.path.endswith(".slides.html")
    )
    deck = parse_deck(STORE.read(thread_id, path).content)
    return path, deck.slides[0].slide_id


def test_build_page_missing_thread_id_returns_error_string(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")

    result = build_page.func(brief="A landing page", runtime=_runtime(None))

    assert isinstance(result, str)
    assert result.startswith("Error:")


def test_build_table_store_write_raises_returns_error_string(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")

    write_calls: list[tuple[str, str]] = []

    def _raising_write(thread_id: str, path: str, *_args: object, **_kwargs: object):
        write_calls.append((thread_id, path))
        raise RuntimeError("disk full")

    monkeypatch.setattr("app.agent.tools.STORE.write", _raising_write)

    result = build_table.func(
        title="Revenue",
        columns=["quarter", "amount"],
        rows=[{"quarter": "Q1", "amount": 100}],
        runtime=_runtime("t-store-error"),
    )

    assert isinstance(result, str)
    assert result.startswith("Error:")
    assert "disk full" in result
    assert len(write_calls) == 1
    assert write_calls[0][0] == "t-store-error"


def test_write_report_writer_model_raises_returns_error_string(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")

    invoked_with: list[str] = []

    def _raising_stream(prompt: str):
        invoked_with.append(prompt)
        raise ConnectionError("writer model unavailable")
        yield  # pragma: no cover - makes this a generator function

    fake_model = SimpleNamespace(stream=_raising_stream)
    monkeypatch.setattr("app.agent.tools.init_chat_model", lambda *_a, **_k: fake_model)

    result = write_report.func(topic="AI safety governance", runtime=_runtime("t-model-error"))

    assert isinstance(result, str)
    assert result.startswith("Error:")
    assert "writer model unavailable" in result
    assert len(invoked_with) == 1
    assert "AI safety governance" in invoked_with[0]


def test_write_slide_render_raises_degrades_layout_check(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """render_slide's Playwright boundary raises (per render.py's re-raise
    after dropping the cached browser) — write_slide already saved the
    slide, so it degrades the layout check instead of aborting the run."""
    thread_id = "t-write-slide-render-error"
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")

    plan_result = plan_deck.func(
        title="Deck", slide_titles=["Cover"], runtime=_runtime(thread_id)
    )
    assert "Error" not in plan_result, plan_result
    path = next(
        info.path for info in STORE.list_files(thread_id) if info.path.endswith(".slides.html")
    )
    slide_id = parse_deck(STORE.read(thread_id, path).content).slides[0].slide_id

    monkeypatch.setattr(
        "app.agent.tools.init_chat_model",
        lambda *_a, **_k: SimpleNamespace(
            invoke=lambda _prompt: SimpleNamespace(content='<section class="slide">Body</section>')
        ),
    )

    render_calls: list[str] = []

    def _raising_render(html: str, *, ratio: object):
        render_calls.append(html)
        raise RuntimeError("playwright renderer crashed")

    monkeypatch.setattr("app.agent.tools.render_slide", _raising_render)

    result = write_slide.func(
        path=path, slide_id=slide_id, title="Cover", brief="Say hello", runtime=_runtime(thread_id)
    )

    assert isinstance(result, str)
    assert "Error" not in result  # the slide save itself succeeded
    assert "layout check skipped" in result
    assert "playwright renderer crashed" in result
    assert len(render_calls) == 1

    saved = parse_deck(STORE.read(thread_id, path).content)
    assert "Body" in saved.slides[0].body_html


def test_convert_slide_render_raises_returns_error_string(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """render_slide raising after a successful correction+save must still
    surface as an Error: string, never an unhandled exception."""
    thread_id = "t-convert-slide-render-error"
    path, slide_id = _seed_imported_deck(thread_id)

    def respond(prompt: str) -> str:
        baseline = prompt.split("Baseline markup:\n", 1)[1].split("\n\nReturn ONLY", 1)[0]
        return baseline.replace('<section class="slide">', '<section class="slide corrected">')

    monkeypatch.setattr(
        "app.agent.tools.init_chat_model",
        lambda *_a, **_k: SimpleNamespace(invoke=lambda prompt: SimpleNamespace(content=respond(prompt))),
    )

    render_calls: list[str] = []

    def _raising_render(html: str, *, ratio: object):
        render_calls.append(html)
        raise RuntimeError("playwright renderer crashed")

    monkeypatch.setattr("app.agent.tools.render_slide", _raising_render)

    result = convert_slide.func(path=path, slide_id=slide_id, runtime=_runtime(thread_id))

    assert isinstance(result, str)
    assert result.startswith("Error:")
    assert "playwright renderer crashed" in result
    assert len(render_calls) == 1
