"""HTML layout -> editable PowerPoint text, vector shapes, and source images."""

from __future__ import annotations

import base64
import io
import re
from pathlib import PurePosixPath
from typing import cast

from langchain_canvas.deck import parse_deck
from langchain_canvas.exporters import (
    PPTX_MIME,
    ExportedFile,
    Exporter,
    MissingExporterDependencyError,
    default_exporters,
)

from .export_fallback import (
    _apply_theme_from_tokens,
    _background_first,
    _covered_by_raster,
    _raster_fallback_items,
    _style_tokens_from_body,
)
from .pdf_fonts import contains_hangul, match_open_source_family
from .render import measure_slide, viewport_for_ratio


def _color(value: str) -> tuple[str, float]:
    match = re.fullmatch(r"rgba?\(([^)]+)\)", value.strip())
    if not match:
        raise ValueError(f"Unsupported CSS color: {value}")
    values = [float(v.strip()) for v in match[1].split(",")]
    return "".join(f"{round(v):02X}" for v in values[:3]), values[3] if len(
        values
    ) == 4 else 1.0


def _set_color(color_format, value: str, opacity: float = 1):
    from pptx.dml.color import RGBColor
    from pptx.oxml.xmlchemy import OxmlElement

    rgb, alpha = _color(value)
    color_format.rgb = RGBColor.from_string(rgb)
    if alpha * opacity < 1:
        node = OxmlElement("a:alpha")
        node.set("val", str(round(alpha * opacity * 100000)))
        color_format._color._xClr.append(node)


def _fill(shape, item):
    from pptx.oxml.xmlchemy import OxmlElement

    gradient = item.get("gradient", "none")
    if gradient == "none":
        shape.fill.solid()
        _set_color(shape.fill.fore_color, item["fill"], item.get("alpha", 1))
        return
    colors = list(re.finditer(r"rgba?\([^)]+\)(?:\s+([\d.]+)%)?", gradient))
    if len(colors) < 2 or gradient.count("gradient(") != 1:
        raise ValueError("Unsupported CSS gradient: simplify it before export")
    shape.fill.gradient()
    grad = shape.fill._xPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill"
    )
    stops = grad.find("{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst")
    for child in list(stops):
        stops.remove(child)
    for i, match in enumerate(colors):
        gs = OxmlElement("a:gs")
        pos = float(match[1]) / 100 if match[1] else i / (len(colors) - 1)
        gs.set("pos", str(round(pos * 100000)))
        rgb, alpha = _color(match[0].split(")")[0] + ")")
        color = OxmlElement("a:srgbClr")
        color.set("val", rgb)
        a = OxmlElement("a:alpha")
        a.set("val", str(round(alpha * item.get("alpha", 1) * 100000)))
        color.append(a)
        gs.append(color)
        stops.append(gs)
    if gradient.startswith("radial-gradient"):
        for child in list(grad):
            if child.tag.endswith("}lin"):
                grad.remove(child)
        path = OxmlElement("a:path")
        path.set("path", "circle")
        rect = OxmlElement("a:fillToRect")
        for side in ("l", "t", "r", "b"):
            rect.set(side, "50000")
        path.append(rect)
        grad.append(path)
    else:
        angle = re.search(r"([-\d.]+)deg", gradient)
        shape.fill.gradient_angle = (450 - (float(angle[1]) if angle else 180)) % 360


def _add_item(slide, item, unit):
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    kind = item["kind"]
    x, y = Emu(round(item["x"] * unit)), Emu(round(item["y"] * unit))
    w, h = (
        Emu(max(1, round(item.get("w", 1) * unit))),
        Emu(max(1, round(item.get("h", 1) * unit))),
    )
    points = unit / 12700
    if kind == "text":
        raise ValueError("Text must belong to a semantic editing block before export")
    elif kind == "polygon":
        points = item["points"]
        builder = slide.shapes.build_freeform(*points[0], scale=unit)
        builder.add_line_segments(points[1:], close=True)
        shape = builder.convert_to_shape()
        shape.line.fill.background()
        _fill(shape, item)
    elif kind == "shape":
        radius = item.get("radius", 0)
        shape_type = (
            MSO_SHAPE.OVAL
            if radius >= min(item["w"], item["h"]) / 2
            else MSO_SHAPE.ROUNDED_RECTANGLE
            if radius > 0
            else MSO_SHAPE.RECTANGLE
        )
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            shape.adjustments[0] = min(0.5, radius / min(item["w"], item["h"]))
        shape.line.fill.background()
        _fill(shape, item)
    elif kind == "line":
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x,
            y,
            Emu(round(item["x2"] * unit)),
            Emu(round(item["y2"] * unit)),
        )
        line.line.width = Pt(item["thickness"] * points)
        _set_color(line.line.color, item["color"], item.get("alpha", 1))
    elif kind == "image":
        src = item["src"]
        if not src.startswith("data:image/") or ";base64," not in src:
            raise ValueError(
                "image could not be loaded: export requires inlined original assets"
            )
        data = base64.b64decode(src.split(",", 1)[1], validate=True)
        picture = slide.shapes.add_picture(io.BytesIO(data), x, y)
        native_w, native_h = picture.image.size
        if item["fit"] == "contain":
            scale = min(w / native_w, h / native_h)
            picture.width = Emu(round(native_w * scale))
            picture.height = Emu(round(native_h * scale))
            picture.left = Emu(x + (w - picture.width) // 2)
            picture.top = Emu(y + (h - picture.height) // 2)
        else:
            picture.width = w
            picture.height = h
            if item["fit"] == "cover":
                ratio = (w / h) / (native_w / native_h)
                if ratio < 1:
                    picture.crop_left = picture.crop_right = (1 - ratio) / 2
                else:
                    picture.crop_top = picture.crop_bottom = (1 - 1 / ratio) / 2


def _add_text_block(slide, item, unit):
    """Keep sentences and rich inline runs in one native editable text frame."""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Emu, Pt

    shape = slide.shapes.add_textbox(
        Emu(round(item["x"] * unit)),
        Emu(round(item["y"] * unit)),
        Emu(max(1, round(item["w"] * unit))),
        Emu(max(1, round(item["h"] * unit))),
    )
    if item.get("id"):
        shape.name = item["id"]
    frame = shape.text_frame
    top, right, bottom, left = item.get("padding", [0, 0, 0, 0])
    frame.margin_top, frame.margin_right, frame.margin_bottom, frame.margin_left = (
        Emu(round(v * unit)) for v in (top, right, bottom, left)
    )
    frame.word_wrap = item.get("whiteSpace") not in {"pre", "nowrap"}
    frame.vertical_anchor = MSO_ANCHOR.TOP
    alignments = {
        "left": PP_ALIGN.LEFT,
        "start": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "end": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    for index, source in enumerate(item["paragraphs"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = alignments.get(source["align"], PP_ALIGN.LEFT)
        paragraph.space_before = Pt(source.get("spaceBefore", 0) * unit / 12700)
        paragraph.space_after = Pt(source.get("spaceAfter", 0) * unit / 12700)
        paragraph.line_spacing = Pt(source["lineHeight"] * unit / 12700)
        for styled in source["runs"]:
            if styled.get("break"):
                paragraph.add_line_break()
                continue
            run = paragraph.add_run()
            run.text = styled["text"]
            family = match_open_source_family(
                styled["font"].split(",")[0].strip("\"' "),
                int(styled["weight"]),
                has_hangul=contains_hangul(styled["text"]),
            )
            run.font.name = family
            run.font.size = Pt(styled["size"] * unit / 12700)
            run.font.bold = int(styled["weight"]) >= 600
            run.font.italic = styled["italic"]
            run.font.underline = styled["underline"]
            _set_color(run.font.color, styled["color"], styled.get("alpha", 1))
            properties = run._r.get_or_add_rPr()
            ea = OxmlElement("a:ea")
            ea.set("typeface", family)
            properties.append(ea)
            if styled.get("letterSpacing"):
                properties.set("spc", str(round(styled["letterSpacing"] * unit / 127)))
    return shape


def _guard_native_table_css(body, style_css):
    """Reject CSS that cannot safely round-trip through native table formatting."""
    from lxml import html as parser

    root = parser.fragment_fromstring(body, create_parent=True)
    tables = root.xpath(".//table[@data-pptx-shape-id]")
    if not tables:
        return
    if style_css.strip() or any((n.text or "").strip() for n in root.xpath(".//style")):
        raise ValueError(
            "Native tables cannot export custom stylesheets; preserve source cell styles"
        )

    def check_container(node):
        allowed = {
            "table": {"tr", "thead", "tbody"},
            "thead": {"tr"},
            "tbody": {"tr"},
            "tr": {"td", "th"},
        }
        metadata = {"data-node-id", "id", "title", "aria-label"}
        permitted = metadata | (
            {"style", "data-pptx-shape-id"} if node.tag == "table" else {"style"}
        )
        if set(node.attrib) - permitted or (
            node.tag != "table" and node.get("style", "").strip()
        ):
            raise ValueError(
                "Native table containers cannot export presentation attributes or row styling"
            )
        if (node.text or "").strip():
            raise ValueError("Native table content must remain inside its source cells")
        for child in node:
            if child.tag not in allowed[node.tag] or (child.tail or "").strip():
                raise ValueError(
                    "Native table structure contains unsupported or unmapped content"
                )
            if child.tag not in {"td", "th"}:
                check_container(child)

    for table in tables:
        check_container(table)
        for ancestor in (table, *table.iterancestors()):
            properties = re.findall(r"([\w-]+)\s*:", ancestor.get("style", ""))
            if any(
                re.fullmatch(
                    r"font(?:-.*)?|color|line-height|letter-spacing|word-spacing|text-.*|white-space|direction|writing-mode|opacity|filter|transform",
                    name.lower(),
                )
                for name in properties
            ):
                raise ValueError(
                    "Native tables cannot export inherited styling changes; preserve source cell styles"
                )


def _native_table_cells(shape, slide, node_id, unit, canvas=None):
    """Regenerate comparison HTML on a clone; font accessors may create rPr XML."""
    import copy

    from langchain_canvas.deck._shapes import _frame, _scheme
    from langchain_canvas.deck.structured import extract_structured, structured_html
    from lxml import html as parser
    from pptx.shapes.graphfrm import GraphicFrame

    source = slide.part.package.presentation_part.presentation
    clone = GraphicFrame(copy.deepcopy(shape._element), shape._parent)
    frame = _frame(clone, 0, source.slide_width, source.slide_height)
    if frame is None:
        raise ValueError("Native table has no explicit source geometry")
    element = extract_structured(clone, frame, _scheme(slide), slide=slide)
    if element is None:
        raise ValueError(
            "Native table cannot be reconstructed for safe content editing"
        )
    canvas = canvas or (
        round(source.slide_width / unit),
        round(source.slide_height / unit),
    )
    root = parser.fromstring(structured_html(element, node_id, canvas))
    return {cell.get("data-node-id"): cell for cell in root.xpath(".//td|.//th")}


def _cell_signature(node, defaults):
    """Ignore content and legacy omitted defaults, never explicit style changes."""
    from langchain_canvas.deck import sanitize_slide_html
    from lxml import html as parser
    from lxml.etree import tostring

    clean = sanitize_slide_html(cast(str, parser.tostring(node, encoding="unicode")))
    if clean.removed:
        raise ValueError("Native table style contains unsupported markup")
    root = parser.fromstring(clean.html)
    styles = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", root.get("style", "")))
    for key in ("font-family", "font-size", "font-weight"):
        if key in defaults:
            styles.setdefault(key, defaults[key])
    root.set("style", ";".join(f"{k}:{v}" for k, v in sorted(styles.items())))
    for child in reversed(root.xpath(".//span")):
        inline = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", child.get("style", "")))
        inline = {k: v for k, v in inline.items() if styles.get(k) != v}
        if not inline:
            child.drop_tag()
        else:
            child.set("style", ";".join(f"{k}:{v}" for k, v in sorted(inline.items())))
    for child in root.iter():
        child.text = child.tail = None
        for key in list(child.attrib):
            if key.startswith("data-"):
                child.attrib.pop(key)
    return tostring(root, method="c14n")


def _patch_native_cell(cell, node, original):
    """Replace only original a:t values, retaining all native text formatting XML."""
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    styles = dict(re.findall(r"([\w-]+)\s*:\s*([^;]+)", original.get("style", "")))
    if _cell_signature(node, styles) != _cell_signature(original, styles):
        raise ValueError(
            "Native table cell style/structure changes are unsupported; replace content using its original rich slots"
        )
    if node.text_content() == original.text_content():
        return
    native = cell.text_frame._txBody.findall(".//" + qn("a:t"))
    spans = node.xpath("./span")
    values = (
        [span.text or "" for span in spans]
        if len(spans) == len(native)
        else list(node.itertext())
    )
    if not native and not len(node):
        paragraph = cell.text_frame._txBody.find(qn("a:p"))
        run, text = OxmlElement("a:r"), OxmlElement("a:t")
        text.text = node.text or ""
        run.append(text)
        end = paragraph.find(qn("a:endParaRPr"))
        paragraph.insert(
            list(paragraph).index(end) if end is not None else len(paragraph), run
        )
        return
    if len(native) != len(values) or "".join(values) != node.text_content():
        raise ValueError(
            "Native table rich slots are ambiguous; reopen the source before editing collapsed runs"
        )
    for target, value in zip(native, values, strict=True):
        target.text = value


def _patch_source_slide(slide, layout, unit, body_html):
    """Patch represented native shapes; retain unsupported original objects intact."""
    import copy

    from pptx.oxml.ns import qn
    from pptx.util import Emu

    provenance = {f"e{i}": shape for i, shape in enumerate(list(slide.shapes))}
    blocks = {
        b.get("pptxId"): b
        for b in layout["textBlocks"]
        if b.get("pptxId") and b.get("pptxRoot")
    }

    from langchain_canvas.deck._shapes import _scheme
    from langchain_canvas.deck.structured import (
        native_chart_data,
        validate_chart_markup,
    )
    from lxml import html as parser
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    root = parser.fragment_fromstring(body_html, create_parent=True)
    structured = {
        n.get("data-pptx-shape-id"): n
        for n in root.xpath(".//*[@data-pptx-shape-id]")
        if n.tag == "table" or n.get("data-chart-data")
    }
    by_id = {b["id"]: b for b in layout["textBlocks"] if b.get("id")}
    for ident, node in structured.items():
        if ident not in provenance:
            raise ValueError(f"Unknown native PowerPoint shape: {ident}")
        shape = provenance[ident]
        if node.get("data-chart-data"):
            inner = (node.text or "") + "".join(
                cast(str, parser.tostring(child, encoding="unicode")) for child in node
            )
            data = validate_chart_markup(dict(node.attrib), inner)
            kinds = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            }
            if (
                not shape.has_chart
                or shape.chart.chart_type != kinds[data["type"]]
                or len(shape.chart.plots) != 1
            ):
                raise ValueError(
                    "Chart metadata does not match the original native chart type"
                )
            if len(data["series"]) != len(shape.chart.series) or len(
                data["categories"]
            ) != len(shape.chart.plots[0].categories):
                raise ValueError(
                    "Native chart category/series counts must be preserved"
                )
            original = native_chart_data(shape.chart, _scheme(slide))

            def formatting(chart_data):
                return {
                    **{
                        k: v
                        for k, v in chart_data.items()
                        if k not in {"series", "categories"}
                    },
                    "series_colors": [s.get("color") for s in chart_data["series"]],
                }

            if original is None or formatting(original) != formatting(data):
                raise ValueError(
                    "Native chart style/title changes are unsupported; replace data only"
                )
            replacement = CategoryChartData()
            replacement.categories = data["categories"]
            for series in data["series"]:
                replacement.add_series(series["name"], series["values"])
            shape.chart.replace_data(replacement)
        elif node.tag == "table":
            if not shape.has_table:
                raise ValueError("Table provenance must reference a native table")
            table = shape.table
            originals = _native_table_cells(
                shape,
                slide,
                node.get("data-node-id"),
                unit,
                (layout["width"], layout["height"]),
            )
            original_table = next(iter(originals.values())).getparent().getparent()

            def table_style(element):
                return {
                    key.lower(): value.strip()
                    for key, value in re.findall(
                        r"([\w-]+)\s*:\s*([^;]+)", element.get("style", "")
                    )
                    if key.lower() not in {"position", "left", "top"}
                }

            if table_style(node) != table_style(original_table):
                raise ValueError(
                    "Native table sizing and styling must remain unchanged; move or replace content only"
                )
            rows = node.xpath("./tr|./thead/tr|./tbody/tr")
            if len(rows) != len(table.rows):
                raise ValueError("Native table row count must be preserved")
            occupied = set()
            for r, row in enumerate(rows):
                c = 0
                for cell_node in row.xpath("./td|./th"):
                    while (r, c) in occupied:
                        c += 1
                    cs, rs = (
                        int(cell_node.get("colspan", "1")),
                        int(cell_node.get("rowspan", "1")),
                    )
                    if (
                        cs < 1
                        or rs < 1
                        or c + cs > len(table.columns)
                        or r + rs > len(table.rows)
                    ):
                        raise ValueError("Native table topology must be preserved")
                    cell = table.cell(r, c)
                    if (cell.span_width, cell.span_height) != (cs, rs):
                        raise ValueError(
                            "Native table merged-cell topology must be preserved"
                        )
                    block = by_id.get(cell_node.get("data-node-id"))
                    if block is None:
                        raise ValueError("Native table cells require semantic node IDs")
                    original = originals.get(cell_node.get("data-node-id"))
                    if original is None:
                        raise ValueError(
                            "Native table cell IDs must retain their original row/column mapping"
                        )
                    _patch_native_cell(cell, cell_node, original)
                    occupied.update(
                        (rr, cc) for rr in range(r, r + rs) for cc in range(c, c + cs)
                    )
                    c += cs
                if any((r, cc) not in occupied for cc in range(len(table.columns))):
                    raise ValueError("Native table column count must be preserved")
    for node in layout["elements"]:
        ident = node.get("pptxId")
        if not ident or not node.get("pptxRoot", True):
            continue
        if ident not in provenance:
            raise ValueError(f"Unknown native PowerPoint shape: {ident}")
        shape = provenance[ident]
        shape.left, shape.top, shape.width, shape.height = (
            Emu(round(node[k] * unit)) for k in ("x", "y", "w", "h")
        )
        if ident in structured:
            continue
        if ident in blocks and shape.has_text_frame:
            replacement = _add_text_block(slide, blocks[ident], unit)
            old = shape.text_frame._txBody
            old.getparent().replace(old, copy.deepcopy(replacement.text_frame._txBody))
            replacement._element.getparent().remove(replacement._element)
        for item in layout["items"]:
            if item.get("pptxId") != ident:
                continue
            if item["kind"] == "shape":
                _fill(shape, item)
            elif item["kind"] == "image" and hasattr(shape, "image"):
                src = item["src"]
                if not src.startswith("data:image/") or ";base64," not in src:
                    raise ValueError("Source image must be inlined for export")
                data = base64.b64decode(src.split(",", 1)[1], validate=True)
                _, relationship = slide.part.get_or_add_image_part(io.BytesIO(data))
                shape._element.find(".//" + qn("a:blip")).set(
                    qn("r:embed"), relationship
                )
    return set(provenance)


class EditableDeckPptxExporter:
    """Measure CSS once, emit native PowerPoint elements; no page rasterization."""

    suffixes: tuple[str, ...] = (".slides.html",)
    target: str = "pptx"

    def export(
        self, content: str, *, path: str, title: str | None = None
    ) -> ExportedFile:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise MissingExporterDependencyError(
                "PPTX export requires python-pptx"
            ) from exc
        deck = parse_deck(content)
        if not deck.slides:
            raise ValueError("Cannot export an empty deck")
        width, height = viewport_for_ratio(deck.ratio)
        from langchain_canvas.deck.export import skin_presentation

        presentation = skin_presentation(deck.source)
        source_backed = presentation is not None
        if source_backed:
            if len(presentation.slides) != len(deck.slides):
                raise ValueError(
                    "Source-backed decks must preserve original slide order and count"
                )
        else:
            if deck.source and deck.source.startswith("data:" + PPTX_MIME):
                raise ValueError(
                    "Invalid source PowerPoint; refusing to discard original objects"
                )
            if deck.source and deck.source.lower().endswith(".pptx"):
                raise ValueError(
                    "Source PowerPoint must be loaded before export; refusing to discard original objects"
                )
            presentation = Presentation()
            presentation.slide_width = Inches(13.333333)
            presentation.slide_height = Inches(13.333333 * height / width)
        if presentation.slide_width is None:
            raise ValueError("Presentation has no slide width")
        unit = presentation.slide_width / width
        for slide_index, template in enumerate(deck.slides):
            body = template.body_html
            if source_backed:
                _guard_native_table_css(body, template.style_css)
            if not source_backed and "data-text-block=" not in body:
                from .semantic_text import consolidate_slide_html

                body, _ = consolidate_slide_html(
                    body, ratio=deck.ratio, style_css=template.style_css
                )
            doc = (
                '<!doctype html><html><head><meta charset="utf-8">'
                "<style>html{margin:0;overflow:hidden;background:transparent}"
                f"body{{margin:0;width:{width}px;height:{height}px}}*{{box-sizing:border-box}}</style>"
                f"<style>{template.style_css}</style></head><body>{body}</body></html>"
            )
            layout = measure_slide(doc, ratio=deck.ratio)
            if not source_backed:
                tokens = _style_tokens_from_body(template.body_html)
                if tokens is not None:
                    _apply_theme_from_tokens(presentation, tokens)
            replacements: list[dict] = []
            if layout["unsupported"]:
                if source_backed:
                    # Painting over original PPTX objects would destroy their
                    # editability, so keep rejecting here.
                    raise ValueError(
                        f"{template.slide_id}: "
                        + "; ".join(u["reason"] for u in layout["unsupported"])
                    )
                replacements = _raster_fallback_items(
                    doc, layout["unsupported"], deck.ratio
                )
            slide = (
                presentation.slides[slide_index]
                if source_backed
                else presentation.slides.add_slide(presentation.slide_layouts[6])
            )
            patched = (
                _patch_source_slide(slide, layout, unit, body)
                if source_backed
                else set()
            )
            blocks = {b["key"]: b for b in layout.get("textBlocks", [])}
            items = _background_first(layout["items"], width, height)
            last = {
                item["blockKey"]: i
                for i, item in enumerate(items)
                if item["kind"] == "text" and item.get("blockKey") in blocks
            }
            for index, item in enumerate(items):
                if item.get("pptxId") in patched:
                    continue
                key = item.get("blockKey")
                if item["kind"] == "text" and key in blocks:
                    if index == last[key] and not _covered_by_raster(
                        blocks[key], replacements
                    ):
                        _add_text_block(slide, blocks[key], unit)
                elif not _covered_by_raster(item, replacements):
                    _add_item(slide, item, unit)
            for key, block in blocks.items():
                if (
                    key not in last
                    and block.get("pptxId") not in patched
                    and not _covered_by_raster(block, replacements)
                ):
                    _add_text_block(slide, block, unit)
            # Added last so the replacement paints over whatever it replaced.
            for item in replacements:
                _add_item(slide, item, unit)
        output = io.BytesIO()
        presentation.save(output)
        data = output.getvalue()
        reopened = Presentation(io.BytesIO(data))
        if len(reopened.slides) != len(deck.slides):
            raise ValueError("PPTX verification failed: missing slides")
        name = (
            title or deck.title or PurePosixPath(path).name.removesuffix(".slides.html")
        )
        name = re.sub(r'[\x00-\x1f<>:"/\\|?*]', "-", name).strip(" .") or "slides"
        return ExportedFile(data, f"{name}.pptx", PPTX_MIME)


def app_exporters() -> list[Exporter]:
    return [
        EditableDeckPptxExporter(),
        *(e for e in default_exporters() if e.target != "pptx"),
    ]
